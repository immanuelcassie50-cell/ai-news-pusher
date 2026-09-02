# -*- coding: utf-8 -*-
"""
知行：学习落地工作坊 · 授课PPT生成脚本
- 杂志风：浅色背景 + 深色头部 + 信息密集
- 约130页
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# 配色
BG = RGBColor(0xFA, 0xF7, 0xF2)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
RED = RGBColor(0xC8, 0x44, 0x2C)
GREEN = RGBColor(0x2E, 0x5A, 0x4A)
GOLD = RGBColor(0xB8, 0x8A, 0x44)
TXT = RGBColor(0x22, 0x22, 0x22)
GRY = RGBColor(0x6B, 0x6B, 0x6B)
LGT = RGBColor(0xE8, 0xE2, 0xD6)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
LBG = RGBColor(0xF2, 0xEC, 0xE0)

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 130
NOTES = []

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def bg(s, c=BG):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.line.fill.background()
    r.fill.solid()
    r.fill.fore_color.rgb = c
    r.shadow.inherit = False
    sp = r._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)


def tx(s, x, y, w, h, t, sz=14, c=TXT, b=False, al=PP_ALIGN.LEFT, an=MSO_ANCHOR.TOP, fn="Microsoft YaHei"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = an
    p = tf.paragraphs[0]
    p.alignment = al
    if isinstance(t, str):
        run = p.add_run()
        run.text = t
        run.font.name = fn
        run.font.size = Pt(sz)
        run.font.color.rgb = c
        run.font.bold = b
        rPr = run._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', 'Microsoft YaHei')
    else:
        first = True
        for line, size, cc, bb in t:
            if first:
                first = False
                para = p
            else:
                para = tf.add_paragraph()
            para.alignment = al
            run = para.add_run()
            run.text = line
            run.font.name = fn
            run.font.size = Pt(size)
            run.font.color.rgb = cc
            run.font.bold = bb
            rPr = run._r.get_or_add_rPr()
            ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', 'Microsoft YaHei')


def rc(s, x, y, w, h, fill, line=None, lw=0.5):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(lw)
    r.shadow.inherit = False


def hdr(s, pn, sec, title):
    rc(s, 0, 0, W, Inches(0.08), DARK)
    rc(s, 0, Inches(0.08), Inches(3.6), Inches(0.95), DARK)
    tx(s, Inches(0.25), Inches(0.18), Inches(3.2), Inches(0.35),
       "知行 · 学习落地工作坊", sz=11, c=WHT, b=True)
    tx(s, Inches(0.25), Inches(0.55), Inches(3.2), Inches(0.4),
       sec, sz=13, c=WHT)
    tx(s, Inches(3.85), Inches(0.18), Inches(6.5), Inches(0.45),
       title, sz=18, c=DARK, b=True, an=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(10.4), Inches(0.5), Inches(2.6), Inches(0.4),
       "P.{:03d} / {:03d}".format(pn, TOTAL), sz=10, c=GRY, al=PP_ALIGN.RIGHT)
    rc(s, Inches(3.85), Inches(1.0), Inches(9.2), Pt(1.5), RED)


def ftr(s, t="知行：学习落地工作坊 · V1.0 · 2026"):
    rc(s, 0, Inches(7.18), W, Pt(0.75), DARK)
    tx(s, Inches(0.3), Inches(7.22), Inches(13.0), Inches(0.25),
       t, sz=9, c=WHT, an=MSO_ANCHOR.MIDDLE)


def note(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# 1. 封面
def P1():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    rc(s, 0, 0, Inches(5.5), H, DARK)
    rc(s, Inches(5.5), 0, Inches(0.12), H, RED)
    tx(s, Inches(0.6), Inches(1.2), Inches(4.7), Inches(0.5),
       "ZHIXING WORKSHOP", sz=14, c=RED, b=True)
    tx(s, Inches(0.6), Inches(1.8), Inches(4.7), Inches(1.5),
       "知行", sz=96, c=WHT, b=True)
    tx(s, Inches(0.6), Inches(3.4), Inches(4.7), Inches(0.7),
       "学习落地工作坊", sz=32, c=WHT, b=True)
    tx(s, Inches(0.6), Inches(4.3), Inches(4.7), Inches(0.5),
       "From Knowing to Doing", sz=14, c=RED)
    rc(s, Inches(0.6), Inches(5.1), Inches(1.0), Pt(2), RED)
    tx(s, Inches(0.6), Inches(5.3), Inches(4.7), Inches(0.4),
       "1天 · 7小时 · 5章 · 4R循环", sz=12, c=WHT)
    tx(s, Inches(0.6), Inches(5.7), Inches(4.7), Inches(0.4),
       "目标学员：内训师 + 区域主管", sz=11, c=LGT)
    tx(s, Inches(0.6), Inches(6.1), Inches(4.7), Inches(0.4),
       "贯穿案例：静远的销售技能落地工作坊", sz=11, c=LGT)
    rc(s, Inches(6.0), Inches(0.6), Inches(7.0), Inches(6.3), LBG)
    tx(s, Inches(6.2), Inches(0.8), Inches(6.5), Inches(1.5),
       "01", sz=140, c=DARK, b=True, al=PP_ALIGN.RIGHT)
    rc(s, Inches(6.2), Inches(2.4), Inches(2.5), Pt(4), RED)
    tx(s, Inches(6.2), Inches(2.7), Inches(6.5), Inches(0.5),
       "三核技能", sz=18, c=DARK, b=True)
    tx(s, Inches(6.2), Inches(3.25), Inches(6.5), Inches(0.5),
       "点 燃  ·  深 问  ·  共 创", sz=24, c=RED, b=True)
    tx(s, Inches(6.2), Inches(4.3), Inches(6.5), Inches(0.5),
       "核心框架", sz=18, c=DARK, b=True)
    tx(s, Inches(6.2), Inches(4.85), Inches(6.5), Inches(0.5),
       "4R 循 环", sz=24, c=GREEN, b=True)
    tx(s, Inches(6.2), Inches(5.7), Inches(6.5), Inches(0.4),
       "回顾 → 现实 → 共创 → 行动", sz=12, c=GRY)
    tx(s, Inches(6.2), Inches(6.3), Inches(6.5), Inches(0.4),
       "© 2026 知行工作坊教学团队", sz=9, c=GRY)
    note(s, "【讲师开场】\n\n各位学员早上好。欢迎来到《知行：学习落地工作坊》。\n\n在我们正式开始之前，我想先讲一个让我自己反复思考的事实：培训做得热闹，行为不一定改变。\n\n我曾跟踪过一个销售技能培训项目：完课率 92%，满意度 4.7/5。但是 90 天后回到工作岗位，业务结果几乎纹丝未动。学员不是不努力，他们是真的知道了——但知道和用起来之间，存在着一条看不见的鸿沟。\n\n这条鸿沟，我们叫它知行鸿沟。这门课，就是关于如何跨过这条鸿沟的。\n\n我们今天有 7 小时，分成 5 个章节，覆盖 4R 循环和三核技能。\n贯穿案例是静远——一位零售企业的区域培训主任。她遇到的问题和您可能遇到的非常类似。\n\n请翻开您的手册第一页。我们开始。")


# 2. 导览
def P2():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 2, "导览 · AGENDA", "今天的旅程")
    chapters = [
        ("第一章", "知行之间", "60 min", "让知行鸿沟从抽象变真实", RED),
        ("第二章", "流程规划", "120 min", "4R循环 + 三维分析", GREEN),
        ("第三章", "点燃参与", "60 min", "5种点燃开场方法", GOLD),
        ("第四章", "深问与共创", "120 min", "4类提问 + 演练", RED),
        ("第五章", "训后跟进", "60 min", "30天跟进路线图", GREEN),
    ]
    cw = Inches(2.3)
    ch = Inches(3.4)
    gap = Inches(0.18)
    sx = Inches(0.45)
    sy = Inches(2.0)
    for i, (n, nm, tm, ds, c) in enumerate(chapters):
        x = sx + (cw + gap) * i
        rc(s, x, sy, cw, ch, WHT, line=LGT)
        rc(s, x, sy, cw, Inches(0.65), c)
        tx(s, x, sy + Inches(0.12), cw, Inches(0.4),
           n, sz=14, c=WHT, b=True, al=PP_ALIGN.CENTER)
        tx(s, x, sy + Inches(0.85), cw, Inches(0.7),
           nm, sz=22, c=DARK, b=True, al=PP_ALIGN.CENTER)
        tx(s, x, sy + Inches(1.7), cw, Inches(0.5),
           tm, sz=24, c=c, b=True, al=PP_ALIGN.CENTER)
        tx(s, x + Inches(0.15), sy + Inches(2.4),
           cw - Inches(0.3), Inches(0.8),
           ds, sz=11, c=TXT, al=PP_ALIGN.CENTER)
    tx(s, Inches(0.45), Inches(1.2), Inches(12.5), Inches(0.4),
       "时间轴：09:00 → 18:00 · 含茶歇与午餐", sz=12, c=GRY)
    rc(s, Inches(0.45), Inches(1.6), Inches(12.4), Pt(2), DARK)
    for tm, x in [("09:00", 0.45), ("10:00", 2.85), ("12:15", 5.25),
                  ("13:30", 6.7), ("14:30", 8.6), ("16:45", 10.6), ("18:00", 12.5)]:
        tx(s, Inches(x - 0.3), Inches(1.65), Inches(0.7), Inches(0.3),
           tm, sz=9, c=DARK, b=True)
    rc(s, Inches(0.45), Inches(5.7), Inches(12.4), Inches(1.2), LBG)
    tx(s, Inches(0.6), Inches(5.8), Inches(12.0), Inches(0.4),
       "课程目标", sz=14, c=RED, b=True)
    tx(s, Inches(0.6), Inches(6.15), Inches(12.0), Inches(0.35),
       "1. 看见知行鸿沟，理解讲师→引导师角色转变", sz=11, c=TXT)
    tx(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.35),
       "2. 掌握 4R 循环 + 三维分析，能设计一场完整工作坊", sz=11, c=TXT)
    ftr(s)
    note(s, "【课程导览】\n\n五个章节，总共 7 小时。\n\n第一章 60 分钟，让我们看见鸿沟。\n第二章 120 分钟，是今天的核心：4R 循环。\n第三章 60 分钟，讲如何点燃开场。\n第四章 120 分钟，是技术含量最高的部分：深度提问。\n第五章 60 分钟，让转化真的发生。\n\n茶歇和午餐时间请大家准时回来。我们的课程极易超时间，每个章节会严格按时间控制。")


# 3. 学习目标
def P3():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 3, "导览", "您将带走的 4 份产出")
    items = [
        ("01", "知行鸿沟认知", "看见传统培训的盲区", "理解讲师→引导师角色转变", RED),
        ("02", "4R 循环设计草案", "针对您真实场景", "完成 4R 流程完整设计", GREEN),
        ("03", "5种点燃开场方法", "可立即使用的开场设计", "为下一个工作坊准备", GOLD),
        ("04", "30天跟进路线图", "包含实践任务", "主管跟进节点 + 反馈机制", RED),
    ]
    for i, (n, t, d1, d2, c) in enumerate(items):
        y = Inches(1.5 + i * 1.35)
        rc(s, Inches(0.6), y, Inches(12.1), Inches(1.2), WHT, line=LGT)
        rc(s, Inches(0.6), y, Inches(1.4), Inches(1.2), c)
        tx(s, Inches(0.6), y, Inches(1.4), Inches(1.2), n,
           sz=40, c=WHT, b=True, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
        tx(s, Inches(2.2), y + Inches(0.15), Inches(10.3), Inches(0.45),
           t, sz=20, c=DARK, b=True)
        tx(s, Inches(2.2), y + Inches(0.6), Inches(5.0), Inches(0.4),
           d1, sz=12, c=TXT)
        tx(s, Inches(7.5), y + Inches(0.6), Inches(5.0), Inches(0.4),
           d2, sz=12, c=TXT)
    ftr(s)
    note(s, "【学习目标】\n\n今天结束的时候，您会带走 4 份具体的产出：\n1. 知行鸿沟的认知\n2. 一份针对您真实场景的 4R 工作坊设计草案\n3. 5种开场方法中至少 2 种为下一个工作坊准备\n4. 一份 30 天跟进路线图\n\n请注意第二份和第四份——不是虚拟案例，是您自己的真实工作场景。")


# ============== 第一章 (12页) ==============
def P4():  # 章节封面
    s = prs.slides.add_slide(BLANK)
    bg(s)
    rc(s, 0, 0, W, H, DARK)
    tx(s, Inches(0.6), Inches(2.0), Inches(8), Inches(0.5),
       "CHAPTER 01", sz=18, c=RED, b=True)
    tx(s, Inches(0.6), Inches(2.6), Inches(10), Inches(1.4),
       "知行之间", sz=72, c=WHT, b=True)
    tx(s, Inches(0.6), Inches(4.2), Inches(10), Inches(0.6),
       "为什么知道了还是不做？", sz=24, c=LGT)
    rc(s, Inches(0.6), Inches(5.0), Inches(2), Pt(3), RED)
    tx(s, Inches(0.6), Inches(5.2), Inches(8), Inches(0.4),
       "60 min · 让知行鸿沟从抽象概念变成您的真实体感", sz=14, c=LGT)
    tx(s, Inches(0.6), Inches(5.7), Inches(8), Inches(0.4),
       "1.1 我的培训困境  ·  1.2 70-20-10 与知行鸿沟", sz=12, c=LGT)
    tx(s, Inches(0.6), Inches(6.05), Inches(8), Inches(0.4),
       "1.3 讲师 vs 引导师  ·  1.4 落地工作坊理想状态", sz=12, c=LGT)
    tx(s, Inches(9.5), Inches(4.5), Inches(3.5), Inches(2.5),
       "01", sz=200, c=RED, b=True, al=PP_ALIGN.RIGHT)
    note(s, "【第一章：知行之间】\n\n60 分钟。让知行鸿沟从抽象概念变成您的真实体感。\n\n这一章有 4 个小节：\n1.1 我的培训困境\n1.2 70-20-10 与知行鸿沟\n1.3 讲师 vs 引导师\n1.4 落地工作坊理想状态\n\n最后会收口到 4R 循环。")


def P5():  # 1.1 开场
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 5, "第一章 · 知行之间", "1.1 · 我的培训困境")
    rc(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.8), DARK)
    tx(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
       "个人书写 · 90 秒", sz=14, c=RED, b=True)
    tx(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.0),
       "你负责的培训，结束之后，\n学员有什么行为改变了？\n你有没有亲眼见过？", sz=28, c=WHT, b=True)
    tx(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.4),
       "—— 一句话回答，真实、未经过滤", sz=12, c=LGT)
    steps = [("个人书写", "90 秒", "写下来"),
             ("两人分享", "2 分钟", "旁边学员"),
             ("全体分享", "3-4 位", "30 秒/人")]
    for i, (t, time, sub) in enumerate(steps):
        x = Inches(0.6 + i * 4.2)
        rc(s, x, Inches(4.7), Inches(3.9), Inches(1.4), WHT, line=LGT)
        rc(s, x, Inches(4.7), Inches(3.9), Inches(0.12), RED)
        tx(s, x, Inches(4.95), Inches(3.9), Inches(0.4),
           t, sz=18, c=DARK, b=True, al=PP_ALIGN.CENTER)
        tx(s, x, Inches(5.4), Inches(3.9), Inches(0.3),
           time, sz=14, c=RED, b=True, al=PP_ALIGN.CENTER)
        tx(s, x, Inches(5.75), Inches(3.9), Inches(0.3),
           sub, sz=11, c=GRY, al=PP_ALIGN.CENTER)
    tx(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.4),
       "讲师观察清单", sz=12, c=RED, b=True)
    tx(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
       "· 学员分享时是否真诚？  · 是我(个人)还是我们(集体)？  · 是困扰还是无所谓？",
       sz=11, c=TXT)
    ftr(s)
    note(s, "【1.1 开场：我的培训困境】\n\n讲师话术 1.1.1：\n在开始之前，我想请大家做一件事——不需要打开手册，就在脑子里或者纸上，用一句话回答这个问题：\n你负责的培训，结束之后，学员有什么行为改变了？你有没有亲眼见过？\n\n我知道这个问题对很多人来说不那么舒服。我自己的版本是这个问题的反向——我负责过的一门课，满意度 4.7，完课率 95%，但 90 天后，学员没有一个人真正在工作中用。\n这是真的。我希望你花 90 秒，把你的答案写下来——真实的、未经过滤的答案。\n\n[给学员 90 秒]\n\n讲师话术 1.1.2：\n请和您旁边的学员分享。不需要分享您公司机密的内容，只需要分享那个行为有没有变的真相。\n\n[2 分钟后]\n\n讲师话术 1.1.3：\n有没有 2-3 位愿意和全体分享的？不需要说太多——30 秒以内，一句话。\n\n应变：\n- 没人愿意分享 → 讲师自己先讲一个真实例子\n- 防御性回答 → 追问还可以的具体表现是什么\n- 情绪激烈 → 是您自己的观察还是您的猜测")


def P6():  # 1.1 静远
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 6, "第一章 · 知行之间", "1.1 · 静远的故事")
    rc(s, Inches(0.6), Inches(1.4), Inches(4.5), Inches(5.4), DARK)
    tx(s, Inches(0.8), Inches(1.6), Inches(4.1), Inches(0.5),
       "案例人物", sz=12, c=RED, b=True)
    tx(s, Inches(0.8), Inches(2.1), Inches(4.1), Inches(1.2),
       "静 远", sz=54, c=WHT, b=True)
    rc(s, Inches(0.8), Inches(3.4), Inches(1.0), Pt(2), RED)
    tx(s, Inches(0.8), Inches(3.6), Inches(4.1), Inches(0.5),
       "零售企业 · 区域培训主任", sz=14, c=WHT)
    tx(s, Inches(0.8), Inches(4.0), Inches(4.1), Inches(0.4),
       "3 年培训经验 · 管 5 个区域 / 80 位员工", sz=11, c=LGT)
    tx(s, Inches(0.8), Inches(5.2), Inches(4.1), Inches(1.4),
       "我上课的时候，大家都很认真，满意度也高。可是回到门店呢？",
       sz=13, c=WHT)
    tx(s, Inches(5.4), Inches(1.4), Inches(7.6), Inches(0.4),
       "她遇到的问题", sz=16, c=RED, b=True)
    tx(s, Inches(5.4), Inches(1.8), Inches(7.6), Inches(0.4),
       "上个月，公司刚结束一轮销售技能线上培训：", sz=12, c=TXT)
    data = [("92%", "线上课完课率", RED),
            ("4.7/5", "学员满意度", GREEN),
            ("8%", "90天后行为改变", RED)]
    for i, (n, l, c) in enumerate(data):
        y = Inches(2.3 + i * 1.1)
        rc(s, Inches(5.4), y, Inches(2.0), Inches(0.95), WHT, line=LGT)
        rc(s, Inches(5.4), y, Inches(0.1), Inches(0.95), c)
        tx(s, Inches(5.5), y + Inches(0.05), Inches(1.9), Inches(0.5),
           n, sz=22, c=DARK, b=True, al=PP_ALIGN.CENTER)
        tx(s, Inches(5.5), y + Inches(0.55), Inches(1.9), Inches(0.4),
           l, sz=10, c=GRY, al=PP_ALIGN.CENTER)
    rc(s, Inches(7.6), Inches(2.3), Inches(5.4), Inches(3.4), LBG)
    tx(s, Inches(7.8), Inches(2.4), Inches(5.0), Inches(0.4),
       "业务方反馈", sz=12, c=RED, b=True)
    tx(s, Inches(7.8), Inches(2.75), Inches(5.0), Inches(0.4),
       "我们花了钱做培训，结果呢？", sz=13, c=DARK, b=True)
    tx(s, Inches(7.8), Inches(3.1), Inches(5.0), Inches(0.4),
       "门店该怎么做还是怎么做。", sz=13, c=DARK, b=True)
    tx(s, Inches(7.8), Inches(3.5), Inches(5.0), Inches(0.4),
       "你能给我们一个落地工作坊吗？", sz=13, c=DARK, b=True)
    rc(s, Inches(7.8), Inches(4.0), Inches(4.0), Pt(1), GRY)
    tx(s, Inches(7.8), Inches(4.2), Inches(5.0), Inches(0.4),
       "—— 静远面对的核心问题", sz=11, c=GRY, b=True)
    tx(s, Inches(7.8), Inches(4.6), Inches(5.0), Inches(0.7),
       "学员不是不努力，他们是知道了，但是知道和用起来之间有一条看不见的鸿沟。",
       sz=12, c=TXT)
    tx(s, Inches(5.4), Inches(5.85), Inches(7.6), Inches(0.4),
       "这不只是静远的问题。", sz=14, c=RED, b=True)
    tx(s, Inches(5.4), Inches(6.2), Inches(7.6), Inches(0.4),
       "这是 70% 的企业培训都面临的问题。", sz=14, c=DARK, b=True)
    ftr(s)
    note(s, "【1.1 案例：静远】\n\n贯穿案例。\n\n静远是某零售企业的区域培训主任，负责 5 个区域、80 位一线员工。\n上个月，公司刚结束一轮销售技能线上培训——完课率 92%，满意度 4.7/5。\n但 90 天后，业务方反馈：门店该怎么做还是怎么做。\n他们要求静远 2 周内为每个区域做一次落地工作坊。\n\n这是真实的、典型的知行鸿沟困境。\n静远会带着这个问题，和我们一起走完今天 5 个章节。\n您可以把她当成自己工作的影子——您遇到的，可能比她还难，也可能比她还简单。但本质是一样的。")


def P7():  # 1.1 知行鸿沟
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 7, "第一章 · 知行之间", "1.1 · 看见：知行鸿沟")
    rc(s, Inches(2.5), Inches(1.5), Inches(8.3), Inches(1.5), DARK)
    tx(s, Inches(2.5), Inches(1.6), Inches(8.3), Inches(0.4),
       "核心概念", sz=12, c=RED, b=True, al=PP_ALIGN.CENTER)
    tx(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(1.0),
       "知行鸿沟", sz=48, c=WHT, b=True, al=PP_ALIGN.CENTER)
    rc(s, Inches(0.6), Inches(3.4), Inches(5.9), Inches(3.4), WHT, line=LGT)
    rc(s, Inches(0.6), Inches(3.4), Inches(5.9), Inches(0.5), GREEN)
    tx(s, Inches(0.6), Inches(3.45), Inches(5.9), Inches(0.4),
       "知 (Knowing)", sz=16, c=WHT, b=True, al=PP_ALIGN.CENTER)
    tx(s, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
       "学员能复述知识、能通过测试", sz=12, c=TXT, b=True)
    tx(s, Inches(0.8), Inches(4.4), Inches(5.5), Inches(0.4),
       "· 上课认真听  · 知道应该怎么做", sz=11, c=TXT)
    tx(s, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.4),
       "· 课堂上能给出标准答案  · 满意度评分高", sz=11, c=TXT)
    rc(s, Inches(0.8), Inches(5.3), Inches(5.5), Pt(1), GRY)
    tx(s, Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.4),
       "传统培训评估的就是这一层", sz=11, c=GREEN, b=True)
    tx(s, Inches(0.8), Inches(5.85), Inches(5.5), Inches(0.4),
       "但这并不是学习的终点", sz=11, c=GRY)
    tx(s, Inches(0.8), Inches(6.3), Inches(5.5), Inches(0.4),
       "→ 满意度 4.7 ≠ 行为改变", sz=12, c=RED, b=True)
    rc(s, Inches(6.8), Inches(3.4), Inches(5.9), Inches(3.4), WHT, line=LGT)
    rc(s, Inches(6.8), Inches(3.4), Inches(5.9), Inches(0.5), RED)
    tx(s, Inches(6.8), Inches(3.45), Inches(5.9), Inches(0.4),
       "行 (Doing)", sz=16, c=WHT, b=True, al=PP_ALIGN.CENTER)
    tx(s, Inches(7.0), Inches(4.0), Inches(5.5), Inches(0.4),
       "学员在工作中真的用起来", sz=12, c=TXT, b=True)
    tx(s, Inches(7.0), Inches(4.4), Inches(5.5), Inches(0.4),
       "· 主动在场景中应用  · 遇到困难能调整", sz=11, c=TXT)
    tx(s, Inches(7.0), Inches(4.7), Inches(5.5), Inches(0.4),
       "· 形成新的工作习惯  · 业务结果可量化", sz=11, c=TXT)
    rc(s, Inches(7.0), Inches(5.3), Inches(5.5), Pt(1), GRY)
    tx(s, Inches(7.0), Inches(5.5), Inches(5.5), Inches(0.4),
       "这是企业真正买单的原因", sz=11, c=RED, b=True)
    tx(s, Inches(7.0), Inches(5.85), Inches(5.5), Inches(0.4),
       "但传统培训几乎不评估这一层", sz=11, c=GRY)
    tx(s, Inches(7.0), Inches(6.3), Inches(5.5), Inches(0.4),
       "→ 区域销售业绩才是结果", sz=12, c=RED, b=True)
    tx(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.3),
       "→ 引导师的工作，就是跨越这条鸿沟", sz=12, c=DARK, b=True)
    ftr(s)
    note(s, "【1.1 知行鸿沟定义】\n\n知行鸿沟是今天贯穿的概念。\n\n左边是知——能复述、能通过测试、知道应该怎么做。\n右边是行——在工作中真的用起来、形成新习惯。\n\n传统培训评估的是左边，满意度、考试、结业证书。\n但企业真正买单的，是右边。\n这两者之间，有一条看不见的鸿沟。\n\n引导师的工作，就是跨越这条鸿沟。这和讲得好不好是两件不同的事。")


def P8():  # 1.2 70-20-10
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 8, "第一章 · 知行之间", "1.2 · 70-20-10 学习法则")
    tx(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.4),
       "Lombardo & Eichinger, 1996", sz=11, c=GRY)
    blocks = [
        ("10%", "正式课堂", "听讲师讲、读教材、看视频", "传统培训 90% 资源投入", RED),
        ("20%", "同事/导师交流", "观察他人、被辅导、讨论案例", "几乎完全自然发生", GREEN),
        ("70%", "实际工作", "做项目、解决真问题、应用所学", "几乎完全自然发生", GOLD),
    ]
    for i, (pct, name, detail, status, color) in enumerate(blocks):
        x = Inches(0.6 + i * 4.2)
        rc(s, x, Inches(1.8), Inches(3.9), Inches(3.6), WHT, line=LGT)
        rc(s, x, Inches(1.8), Inches(3.9), Inches(0.6), color)
        tx(s, x, Inches(1.9), Inches(3.9), Inches(0.5),
           pct, sz=36, c=WHT, b=True, al=PP_ALIGN.CENTER)
        tx(s, x, Inches(2.55), Inches(3.9), Inches(0.4),
           name, sz=18, c=DARK, b=True, al=PP_ALIGN.CENTER)
        rc(s, x + Inches(1.0), Inches(3.05), Inches(1.9), Pt(1), GRY)
        tx(s, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(0.6),
           detail, sz=12, c=TXT, al=PP_ALIGN.CENTER)
        tx(s, x + Inches(0.2), Inches(4.3), Inches(3.5), Inches(0.4),
           status, sz=11, c=color, b=True, al=PP_ALIGN.CENTER)
    rc(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.4), DARK)
    tx(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
       "传统培训的资源分配", sz=14, c=RED, b=True)
    tx(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.4),
       "90% 预算  →  投入到了那 10% (正式课堂)", sz=14, c=WHT, b=True)
    tx(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
       "剩下 10%  →  投入到了那 90% (自然发生，几乎没设计)",
       sz=14, c=WHT)
    tx(s, Inches(8.5), Inches(6.85), Inches(4.2), Inches(0.3),
       "→ 这就是知行鸿沟的根因", sz=11, c=RED, b=True, al=PP_ALIGN.RIGHT)
    ftr(s)
    note(s, "【1.2 70-20-10】\n\n讲师话术 1.2.1：\n学习科学家 Lombardo 和 Eichinger 在 1996 年提出了一个学习法则——\n- 10% 的学习来自正式课堂\n- 20% 来自和同事/导师的交流\n- 70% 来自实际工作和挑战\n\n讲师话术 1.2.2（视角翻转）：\n传统培训行业的资源分配是这样的：\n- 90% 的预算投入到了那 10%\n- 剩下的 20% 和 70% 几乎是自然发生\n\n为什么？因为最容易量化——讲师好找、课程好开发、教室好订。\n所以传统培训的最大盲区，是它只设计了 10%，而把 90% 留给了自然发生。\n\n今天我们专门解决那 90% 怎么从自然发生变成被设计。\n\n应变：学员挑战我觉得培训还是有效的→\n我完全同意培训是有价值的。我想说的是，培训做完之后，行为是否真的改变了——这是另一个问题。")


def P9():  # 1.2 小组讨论
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 9, "第一章 · 知行之间", "1.2 · 小组讨论 · 5 分钟")
    tx(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.4),
       "在您过去主持/参与的培训中，下面三个 70% 的发生场景：", sz=12, c=TXT)
    tx(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.4),
       "哪一个被显式设计过？哪一个完全靠学员自己？", sz=12, c=DARK, b=True)
    rows = [
        ("70% 的发生场景", "是否被显式设计过", "实际效果"),
        ("同事交流", "□ 是  □ 否", "________________"),
        ("主管辅导", "□ 是  □ 否", "________________"),
        ("实际应用", "□ 是  □ 否", "________________"),
    ]
    tx_t = Inches(0.6)
    ty = Inches(2.4)
    cw = [Inches(4.0), Inches(4.0), Inches(4.1)]
    rh = Inches(0.85)
    for r, row in enumerate(rows):
        x = tx_t
        for c, cell in enumerate(row):
            is_h = (r == 0)
            rc(s, x, ty + rh * r, cw[c], rh,
               DARK if is_h else WHT, line=LGT)
            tx(s, x + Inches(0.15), ty + rh * r,
               cw[c] - Inches(0.3), rh, cell,
               sz=13 if is_h else 12,
               c=WHT if is_h else TXT, b=is_h, an=MSO_ANCHOR.MIDDLE)
            x += cw[c]
    rc(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0), LBG)
    tx(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
       "讨论指引", sz=12, c=RED, b=True)
    tx(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.4),
       "· 5 分钟小组讨论  · 1 位记录员  · 1 位汇报员  · 待会分享 1-2 个典型场景",
       sz=11, c=TXT)
    ftr(s)
    note(s, "【1.2 小组讨论】\n\n讲师话术：\n在我继续之前，请小组花 5 分钟讨论一下——\n在你过去主持或参与的培训中，下面这三个 70% 的发生场景：\n- 同事交流\n- 主管辅导\n- 实际应用\n\n哪一个被显式设计过？哪一个完全靠学员自己？实际效果如何？\n1 位记录员，1 位汇报员。待会我会请 1-2 个小组分享。\n\n巡视重点：\n- 学员能否区分被设计和自然发生\n- 学员是否意识到自己之前的设计盲点")


def P10():  # 1.3 角色对比
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 10, "第一章 · 知行之间", "1.3 · 讲师 vs 引导师")
    tx(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.4),
       "区分不是哪个更好，而是在什么场景下用哪个", sz=14, c=RED, b=True)
    rows = [
        ("维度", "讲师 Trainer", "引导师 Facilitator"),
        ("核心任务", "传递知识和信息", "激活学习和转化"),
        ("核心技能", "讲得清楚、讲得好", "问得好、引得好"),
        ("权威来源", "内容专家", "流程设计者"),
        ("成就感", "我讲得好", "他们用起来了"),
        ("学员状态", "接收者", "共同创造者"),
    ]
    tx_t = Inches(0.6)
    ty = Inches(2.0)
    cw = [Inches(2.5), Inches(4.8), Inches(4.8)]
    rh = Inches(0.65)
    for r, row in enumerate(rows):
        x = tx_t
        for c, cell in enumerate(row):
            is_h = (r == 0)
            fill = DARK if is_h else (LBG if r % 2 == 0 else WHT)
            color = WHT if is_h else (DARK if c == 0 else TXT)
            rc(s, x, ty + rh * r, cw[c], rh, fill, line=LGT)
            tx(s, x + Inches(0.15), ty + rh * r,
               cw[c] - Inches(0.3), rh, cell,
               sz=12, c=color, b=is_h or c == 0, an=MSO_ANCHOR.MIDDLE)
            x += cw[c]
    rc(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.95), DARK)
    tx(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.4),
       "关键认知", sz=12, c=RED, b=True)
    tx(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.4),
       "讲师 = 首次内容输入  ·  引导师 = 让知识转化为行为  ·  二者不是替代，是叠加",
       sz=14, c=WHT, b=True)
    ftr(s)
    note(s, "【1.3 讲师 vs 引导师】\n\n讲师话术 1.3.1（重要 · 防御管理）：\n在讲讲师和引导师的差异前，我想先说一句——区分不是哪个更好，而是在什么场景下用哪个。\n\n讲师的核心价值在首次内容输入——让学员从不知道到知道。\n引导师的核心价值在让知识转化为行为——让学员从知道到用起来。\n\n各位作为内训师和区域主管，您今天学的是引导师能力——这是您原有的讲师能力之上的第二层能力，不是替代。\n\n应变（学员防御性反应）：\n学员：我讲得好有什么问题？\n讲得好是非常重要的能力——我们这门课不教怎么不讲。我们教的是什么时候从讲切换到引。\n想象您的学员已经学完了录播课，现在您要做 90 分钟的落地工作坊。\n您需要的是用引导让学员把已经知道的用起来。这就是引导师的工作。")


def P11():  # 1.3 角色转变
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 11, "第一章 · 知行之间", "1.3 · 您的角色转变")
    tx(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.4),
       "个人反思 · 2 分钟", sz=14, c=RED, b=True)
    tx(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.4),
       "在您自己的角色中，您更多是讲师还是引导师？", sz=14, c=DARK)
    tx(s, Inches(0.6), Inches(2.05), Inches(12.1), Inches(0.4),
       "您的下一个工作坊，您希望偏向哪边？为什么？", sz=14, c=DARK)
    rc(s, Inches(0.6), Inches(4.95), Inches(8.5), Inches(0.1), RED)
    tx(s, Inches(0.6), Inches(5.1), Inches(8.5), Inches(0.4),
       "讲师  ←——————————→  引导师", sz=16, c=DARK, b=True, al=PP_ALIGN.CENTER)
    roles = [
        ("纯讲师", "内容专家\n单向传递\n知识测试", GRY),
        ("引导师型讲师", "流程设计者\n激发讨论\n行为改变", RED),
        ("纯引导师", "中立促进\n不参与内容\n深度倾听", GREEN),
    ]
    for i, (n, d, c) in enumerate(roles):
        x = Inches(9.5)
        y = Inches(1.5 + i * 1.7)
        rc(s, x, y, Inches(3.2), Inches(1.5), WHT, line=LGT)
        rc(s, x, y, Inches(0.12), Inches(1.5), c)
        tx(s, x + Inches(0.2), y + Inches(0.1), Inches(3.0), Inches(0.4),
           n, sz=14, c=DARK, b=True)
        tx(s, x + Inches(0.2), y + Inches(0.55), Inches(3.0), Inches(0.9),
           d, sz=10, c=TXT)
    tx(s, Inches(0.6), Inches(5.85), Inches(8.5), Inches(0.4),
       "各位作为内训师和区域主管，您今天学的是", sz=12, c=TXT)
    tx(s, Inches(0.6), Inches(6.25), Inches(8.5), Inches(0.4),
       "引导师型讲师——既懂内容，又会引导。", sz=18, c=RED, b=True)
    tx(s, Inches(0.6), Inches(6.7), Inches(8.5), Inches(0.4),
       "这是您原有的讲师能力之上的第二层能力。", sz=12, c=GRY)
    ftr(s)
    note(s, "【1.3 角色转变】\n\n讲师话术：\n请花 2 分钟想一想——在您自己的角色中，您更多是讲师还是引导师？您的下一个工作坊，您希望偏向哪边？为什么？\n\n巡视重点：\n- 学员是否能区分两种角色\n- 学员是否清楚自己的下一步调整方向\n\n核心收口：\n各位作为内训师和区域主管，您今天学的是引导师型讲师——既懂内容，又会引导。这是您原有的讲师能力之上的第二层能力。")


def P12():  # 1.4 理想状态
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 12, "第一章 · 知行之间", "1.4 · 落地工作坊理想状态")
    tx(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.4),
       "全体讨论 · 5 分钟", sz=14, c=RED, b=True)
    tx(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.4),
       "请用 3 个关键词描述您心目中理想中的落地工作坊应该是什么样子。",
       sz=14, c=DARK)
    rc(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(2.2), DARK)
    tx(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(0.4),
       "请写下您的 3 个关键词", sz=12, c=RED, b=True, al=PP_ALIGN.CENTER)
    for i, color in enumerate([RED, GREEN, GOLD]):
        x = Inches(1.5 + i * 3.6)
        rc(s, x, Inches(3.1), Inches(3.0), Inches(1.4), WHT)
        rc(s, x, Inches(3.1), Inches(3.0), Inches(0.12), color)
        tx(s, x, Inches(3.2), Inches(3.0), Inches(0.4),
           "关键词 {}".format(i + 1), sz=11, c=GRY, al=PP_ALIGN.CENTER)
        tx(s, x, Inches(3.7), Inches(3.0), Inches(0.6),
           "____________", sz=24, c=LGT, al=PP_ALIGN.CENTER)
    tx(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.4),
       "学员常见的高频词", sz=14, c=DARK, b=True)
    words = [("有参与", RED), ("有产出", GREEN), ("有改变", GOLD),
             ("有能量", RED), ("有跟进", GREEN), ("有温度", GOLD)]
    for i, (w, c) in enumerate(words):
        x = Inches(0.6 + (i % 3) * 4.0)
        y = Inches(5.5 + (i // 3) * 0.7)
        rc(s, x, y, Inches(3.8), Inches(0.55), WHT, line=LGT)
        rc(s, x, y, Inches(0.1), Inches(0.55), c)
        tx(s, x + Inches(0.2), y, Inches(3.6), Inches(0.55),
           w, sz=14, c=DARK, b=True, an=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.18),
       "这些关键词都有解法——就是接下来要讲的 4R 循环。", sz=11, c=RED, b=True)
    ftr(s)
    note(s, "【1.4 落地工作坊理想状态】\n\n讲师话术 1.4.1：\n我们刚才看到了培训没效果的真实样子。让我反过来问一下——\n您心目中一场理想中的落地工作坊应该是什么样子？\n请用 3 个关键词描述。3 分钟个人书写。\n\n[3 分钟后收集 4-5 位]\n\n讲师话术 1.4.2（引出 4R 循环）：\n我听到了很多有意思的关键词——有参与、有产出、有改变、有能量。\n接下来 5 分钟，我想用一个小小的心理模型来回应——这些关键词其实都有解法。\n这个模型叫 4R 循环——它不是一个理论，是一个工具。\n- R1 回顾 = 让学员记起来\n- R2 现实 = 让学员连起来（连接真实工作）\n- R3 共创 = 让学员深下去（深度讨论产生洞见）\n- R4 行动 = 让学员走出去（具体承诺）\n这个工具，就是我们今天接下来 6 小时要学的内容。它不复杂，但它需要练习。")


def P13():  # 1.X 章节小结
    s = prs.slides.add_slide(BLANK)
    bg(s)
    hdr(s, 13, "第一章 · 知行之间", "1.X · 第一章小结")
    quotes = [
        ("01", "培训做得热闹，行为不一定改变。", RED),
        ("02", "讲师的核心价值是输入；引导师的核心价值是转化。", GREEN),
        ("03", "4R 循环不是理论，是工具——它需要练习。", GOLD),
    ]
    for i, (n, q, c) in enumerate(quotes):
        y = Inches(1.5 + i * 1.5)
        rc(s, Inches(0.6), y, Inches(12.1), Inches(1.3), WHT, line=LGT)
        rc(s, Inches(0.6), y, Inches(1.3), Inches(1.3), c)
        tx(s, Inches(0.6), y, Inches(1.3), Inches(1.3), n,
           sz=40, c=WHT, b=True, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
        tx(s, Inches(2.2), y, Inches(10.3), Inches(1.3), q,
           sz=18, c=DARK, b=True, an=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.4),
       "下一章：流程规划——用 4R 循环设计一场完整工作坊", sz=14, c=RED, b=True)
    tx(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
       "茶歇 15 分钟后开始", sz=12, c=GRY)
    ftr(s)
    note(s, "【第一章小结】\n\n三个金句：\n1. 培训做得热闹，行为不一定改变。\n2. 讲师的核心价值是输入；引导师的核心价值是转化。\n3. 4R 循环不是理论，是工具——它需要练习。\n\n接下来 15 分钟茶歇。\n回来后我们进入第二章：流程规划。\n这是今天的核心技术章节——用 4R 循环设计一场完整工作坊。")


print("第一章代码就绪")