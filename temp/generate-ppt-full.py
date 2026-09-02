# -*- coding: utf-8 -*-
"""
《经营者讲党课》完整PPT生成脚本 (120-160页)
红灰配色，浅底背景
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
THEME = {
    'primary': RGBColor(0xC4, 0x12, 0x30),   # 党建红
    'secondary': RGBColor(0x4A, 0x4A, 0x4A),  # 深灰
    'accent': RGBColor(0x8B, 0x45, 0x13),     # 红褐
    'light': RGBColor(0xF5, 0xF5, 0xF5),      # 浅灰白
    'bg': RGBColor(0xFF, 0xFF, 0xFF),         # 白色背景
    'text': RGBColor(0x33, 0x33, 0x33),       # 深色文字
    'gold': RGBColor(0xD4, 0xAF, 0x37),       # 金色
    'light_red': RGBColor(0xFF, 0xF0, 0xF0),  # 浅红
}

SLIDES_DIR = 'D:/新课开发/党业融合/经营者讲党课/完整课程包/003-授课PPT'
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_cover_slide(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['primary']
    background.line.fill.background()

    # 装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME['secondary']
    top_bar.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "讲党课"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "业务干部的登台表达赋能工作坊"
    p.font.size = Pt(32)
    p.font.color.rgb = THEME['light']
    p.alignment = PP_ALIGN.CENTER

    # 底部信息
    bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    tf = bottom_box.text_frame
    p = tf.paragraphs[0]
    p.text = "一天学会 当众讲好一个故事"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, num, title, subtitle=""):
    """章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 左侧装饰条
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = THEME['primary']
    left_bar.line.fill.background()

    # 章节号
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(3), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = THEME['primary']

    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(3.5), Inches(2.5), Inches(6), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = THEME['secondary']

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(3.5), Inches(4), Inches(6), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = THEME['accent']

    return slide

def add_content_slide(prs, title, bullets, icon=""):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = THEME['primary']
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 左侧装饰
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), Inches(0.08), prs.slide_height - Inches(1.2))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = THEME['secondary']
    left_bar.line.fill.background()

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▶ " + item
        p.font.size = Pt(18)
        p.font.color.rgb = THEME['text']
        p.space_before = Pt(14)
        p.space_after = Pt(8)

    return slide

def add_two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    """双栏页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = THEME['primary']
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 左栏
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.4), Inches(4.5), Inches(5.6))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = THEME['light_red']
    left_box.line.color.rgb = THEME['primary']

    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = THEME['primary']

    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.1), Inches(4.6))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = THEME['text']
        p.space_before = Pt(10)

    # 右栏
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.4), Inches(4.5), Inches(5.6))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    right_box.line.color.rgb = THEME['secondary']

    right_title_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(4), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = THEME['secondary']

    right_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(4.1), Inches(4.6))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = THEME['text']
        p.space_before = Pt(10)

    return slide

def add_process_slide(prs, title, steps):
    """流程页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = THEME['primary']
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 流程步骤
    n = len(steps)
    step_w = Inches(2.2)
    gap = Inches(0.3)
    start_x = Inches(0.5)

    for i, (num, step_title, step_desc) in enumerate(steps):
        x = start_x + i * (step_w + gap)

        # 圆形编号
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.5), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = THEME['primary']
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(x + Inches(0.7), Inches(1.55), Inches(0.8), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # 步骤卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), step_w, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xF8)
        card.line.color.rgb = THEME['primary']

        # 标题
        t_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.7), step_w - Inches(0.2), Inches(0.6))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = THEME['primary']
        p.alignment = PP_ALIGN.CENTER

        # 描述
        d_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.3), step_w - Inches(0.2), Inches(3.5))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(11)
        p.font.color.rgb = THEME['text']
        p.alignment = PP_ALIGN.CENTER

        # 箭头
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_w + Inches(0.05), Inches(1.7), Inches(0.2), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = THEME['secondary']
            arrow.line.fill.background()

    return slide

def add_quote_slide(prs, quote, author=""):
    """金句页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['light_red']
    background.line.fill.background()

    # 引号装饰
    quote_mark = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(2), Inches(2))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = """
    p.font.size = Pt(120)
    p.font.color.rgb = THEME['primary']

    # 引用内容
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(28)
    p.font.color.rgb = THEME['secondary']
    p.alignment = PP_ALIGN.CENTER

    if author:
        author_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = "— " + author
        p.font.size = Pt(16)
        p.font.color.rgb = THEME['accent']
        p.alignment = PP_ALIGN.RIGHT

    return slide

def add_card_slide(prs, title, cards):
    """卡片页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = THEME['primary']
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 卡片
    n = len(cards)
    card_w = Inches(2.8) if n <= 3 else Inches(2.2)
    gap = Inches(0.2)
    start_x = Inches(0.5)

    for i, (card_title, card_items) in enumerate(cards):
        x = start_x + i * (card_w + gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), card_w, Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xF8)
        card.line.color.rgb = THEME['primary']

        # 卡片标题
        ct_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(1.7), card_w - Inches(0.2), Inches(0.5))
        tf = ct_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = THEME['primary']
        p.alignment = PP_ALIGN.CENTER

        # 卡片内容
        cc_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.3), card_w - Inches(0.2), Inches(4.5))
        tf = cc_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(card_items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(11)
            p.font.color.rgb = THEME['text']
            p.space_before = Pt(8)

    return slide

# ==================== 开始生成幻灯片 ====================

# ===== 封面 =====
add_cover_slide(prs)

# ===== 目录 =====
add_content_slide(prs, "课程目录", [
    "第一章：认知重塑 — 业务干部为什么能讲好党课",
    "第二章：内容设计 — 案例转化四步法",
    "第三章：登台表达 — 分层递进练习体系",
    "第四章：实战演练 — 全班展示与点评",
    "第五章：工具与转化 — 带走可迁移的方法论"
])

# ===== 第一章：认知重塑 =====
add_section_slide(prs, 1, "认知重塑", "业务干部为什么能讲好党课")

# 痛点导入
add_two_col_slide(prs, "业务干部讲党课的两大困境",
    "困境一：不会讲", [
        "完全不会讲党课",
        "只会念PPT或复述文件",
        "不知道如何开头",
        "不知道如何让学员愿意听"
    ],
    "困境二：不敢讲", [
        "怕在同事面前显得不专业",
        "怕被议论讲得差",
        "觉得自己不是党务出身",
        "表达焦虑比内容焦虑更严重"
    ])

# 课程定位
add_content_slide(prs, "课程核心定位", [
    "不代写讲稿，不代做PPT",
    "只做一件事：",
    "教业务口干部把自己真实经历的管理故事，",
    "转化成一堂能讲、敢讲、学员爱听的党课，",
    "并且现场练到能自然讲出来为止。"
])

# 业务干部的优势
add_two_col_slide(prs, "业务干部讲党课的独特优势",
    "优势一：有真实故事", [
        "带领团队克服困难的故事",
        "做出艰难抉择的故事",
        "团队成长的故事",
        "这些经历本身就是最好的党课素材"
    ],
    "优势二：有说服力", [
        "因为真实，所以可信",
        "因为亲身经历，所以有细节",
        "因为有感情，所以能打动人",
        "这恰恰是党课最需要的"
    ])

# 课程边界
add_content_slide(prs, "课程边界：五个不做", [
    "不代写讲稿、不代做PPT — 学员必须自己参与转化",
    "不介入党课内容的具体思想政治内核 — 学员自行对照组织部门要求",
    "不做党史知识讲授、不做党建理论培训 — 假设学员已确定主题",
    "不涉及具体政策文件的解读争议 — 引导对照本单位口径",
    "AI工具不作为课程卖点 — 仅用于辅助整理讲稿逻辑框架草稿"
])

# 课程承诺
add_quote_slide(prs,
    "一天学会 当众讲好一个故事\n离开时带着自己的党课框架",
    "本课程的核心承诺")

# ===== 第二章：案例转化四步法 =====
add_section_slide(prs, 2, "内容设计", "案例转化四步法")

# 四步法总览
add_process_slide(prs, "案例转化四步法", [
    (1, "故事盘点", "从真实经历中找出有戏剧张力的素材"),
    (2, "主题锚定", "找到故事与主题的真实连接点"),
    (3, "结构搭建", "用讲故事的方式重新组织内容"),
    (4, "语言转译", "把业务语言转化成讲台语言")
])

# 第一步详解
add_section_slide(prs, 2, "第一步：故事盘点", "从真实经历中找出有戏剧张力的素材")

add_content_slide(prs, "什么是好故事素材", [
    "好的党课素材必须有'冲突'和'抉择'",
    "平铺直叙的'我们完成了任务'没有戏剧张力",
    "真正的张力来自：当时面临的真实两难选择是什么？",
    "引导学员回忆那个'纠结到睡不着觉'的时刻"
])

add_content_slide(prs, "引导提问的技巧", [
    "\"你带团队这些年，有没有一次真的很难抉择的时刻？\"",
    "\"有没有一次你原本想放弃，但最后咬牙坚持下来的事？\"",
    "\"当时最让你睡不着觉的是什么？\"",
    "\"如果重新来过一次，你会做不同的选择吗？\""
])

add_content_slide(prs, "常见误区：苦劳型素材", [
    "学员第一反应往往会讲：'团队多么努力、多么辛苦'",
    "这类素材只有苦劳，没有冲突，没有抉择",
    "讲师需要引导其往'两难选择'方向深挖：",
    "\"那当时有没有想过放弃？是什么让你坚持下来的？\""
])

# 第二步详解
add_section_slide(prs, 2, "第二步：主题锚定", "找到故事与主题的真实连接点")

add_content_slide(prs, "主题锚定的原则", [
    "不是生硬地把故事'扣'到主题上",
    "而是问：'这个故事让你自己悟到了什么道理？'",
    "再问：这个道理和今天要讲的主题，有没有真实的呼应？",
    "如果有真实的呼应，故事自然会和主题产生共鸣"
])

add_quote_slide(prs,
    "如果学员选的故事和主题连接生硬，\n宁可换一个故事，也不要强行嫁接。\n牵强的连接是党课让人听着尴尬的主要原因。")

add_content_slide(prs, "主题锚定的练习", [
    "给出党课主题，让学员思考：",
    "这个主题让你想起自己经历中的哪个故事？",
    "为什么这个故事让你想起这个主题？",
    "它们之间最真实的连接点是什么？"
])

# 第三步详解
add_section_slide(prs, 3, "第三步：结构搭建", "用讲故事的方式重新组织内容")

add_content_slide(prs, "五段式叙事结构", [
    "开场：设置悬念/抛出真实困境（30秒内让听众进入情境）",
    "抉择：还原当时的纠结和几种可能的选择",
    "转折：揭示当时的真实选择和过程中的转折",
    "感悟：引导听众自己感受到那个道理（不是替听众总结）",
    "收尾：回扣主题，用一句简短有力的话收尾"
])

add_content_slide(prs, "结构模板使用说明", [
    "结构模板是脚手架，不是要求机械套用",
    "讲师需要根据每个学员的故事特点灵活调整",
    "有些故事可能天然就是'总-分-总'结构",
    "关键是：让听众跟着你的节奏走，最终自己得出结论"
])

# 第四步详解
add_section_slide(prs, 4, "第四步：语言转译", "把业务语言转化成讲台语言")

add_content_slide(prs, "业务语言 vs 讲台语言", [
    "业务干部平时习惯用：数据、指标、专业术语",
    "这套语言直接搬上党课讲台会显得生硬冰冷",
    "需要引导其加入更多'画面感'和'情绪细节'",
    "比如不说'团队压力很大'，而是说具体是哪个深夜、谁说了哪句话"
])

add_content_slide(prs, "语言转译的核心训练", [
    "反复追问：'当时具体是什么场景？'",
    "追问：'谁说了什么？'",
    "追问：'你当时的第一反应是什么？'",
    "逼学员讲出细节而非概括，让故事变得生动"
])

# 四步法总结
add_content_slide(prs, "案例转化四步法总结", [
    "故事盘点：找到有冲突、有抉择的素材",
    "主题锚定：建立故事与主题的真实连接",
    "结构搭建：五段式叙事让故事有节奏感",
    "语言转译：增加画面感和情绪细节",
    "四步法的核心：让故事自己'长出'主题，而不是'扣上'主题"
])

# ===== 第三章：登台表达 =====
add_section_slide(prs, 3, "登台表达训练", "分层递进的练习体系")

add_quote_slide(prs,
    "业务干部真正的痛点往往不是不知道讲什么，\n而是'写好了也不敢讲、讲的时候放不开'。\n本课程必须投入至少与内容设计同等甚至更多的时间在实际登台练习上。")

# 心理阻力
add_two_col_slide(prs, "登台心理阻力拆解与应对",
    "阻力一：怕讲错", [
        "怕被人说不专业",
        "深层恐惧：在同事面前丢脸",
        "应对：强调党课不是政治理论考试",
        "允许有个人风格和不完美",
        "讲师给出具体、真实的肯定反馈"
    ],
    "阻力二：放不开", [
        "觉得讲自己的故事很'肉麻'",
        "觉得在同事面前表达感情很奇怪",
        "应对：先从小范围（2-3人小组）练习开始",
        "逐步过渡到全班展示",
        "降低一开始面对大场面的心理压力"
    ])

add_content_slide(prs, "阻力三：照着稿子念", [
    "表现：脱稿就忘词，必须看着稿子才能讲",
    "原因：训练的是背诵，而不是讲故事",
    "应对：训练'记住故事的骨架和几个关键画面'",
    "用讲故事的自然节奏代替背诵的机械感",
    "核心：不是背稿子，是讲故事"
])

# 分层练习
add_section_slide(prs, 3, "分层递进练习体系", "从低压力到高压力，循序渐进")

add_process_slide(prs, "四层练习体系", [
    (1, "小组讲述", "2-3人一组，低压力环境先把内容捋顺"),
    (2, "录制回看", "让学员看到自己实际讲述时的状态"),
    (3, "全班展示", "每位学员完整讲述，讲师现场点评"),
    (4, "模拟真实", "有讲台、有完整时长的正式试讲")
])

add_content_slide(prs, "第一层：小组内讲述练习", [
    "目的：在低压力环境下先把内容捋顺",
    "习惯'说出来'这件事本身",
    "同伴反馈聚焦：",
    "'我作为听众，哪个瞬间让我有感觉？'",
    "'哪里我走神了？'",
    "不做内容对错评判"
])

add_content_slide(prs, "第二层：录制与回看", [
    "目的：让学员看到自己实际讲述时的状态",
    "包括：语速、眼神、肢体语言",
    "很多人从未看过自己讲话的样子",
    "这个环节的冲击力很大",
    "回看后引导：'你觉得自己讲得最自然的是哪一段？'"
])

add_content_slide(prs, "第三层：全班展示与点评", [
    "这是全天最重要的环节",
    "每位学员完整讲述一次（3-5分钟精简版）",
    "讲师现场给出具体、可操作的改进建议",
    "点评原则：先肯定具体的亮点细节",
    "再给1-2条最关键的改进建议",
    "不做面面俱到式的挑刺"
])

add_content_slide(prs, "第四层：模拟真实场景（选配）", [
    "如果时间允许，安排一次模拟真实党课场景",
    "有讲台、有完整时长",
    "让学员提前适应真实场景的紧张感",
    "这是从练习到实战的关键过渡"
])

# 表达技巧
add_section_slide(prs, 4, "表达技巧训练", "穿插在练习环节中逐一带过")

add_card_slide(prs, "四大表达技巧", [
    ("开场设计", ["前3句话的设计", "用问题或画面抓住注意力", "避免'今天我要给大家讲一下...'"]),
    ("停顿使用", ["关键转折处的停顿", "比语速快更有感染力", "让听众有时间消化"]),
    ("眼神交流", ["有意识地扫视不同区域", "不要死盯一个人或稿子", "与听众建立连接"]),
    ("手势自然", ["克制使用手势", "避免汇报PPT时的习惯", "讲故事时手势更自然放松"])
])

# ===== 第四章：实战演练 =====
add_section_slide(prs, 4, "实战演练", "全班展示与点评")

add_content_slide(prs, "全班展示环节设计", [
    "目的：让每位学员都有登台机会并获得讲师点评",
    "人数控制：建议每场不超过16-18人",
    "时间：每位学员3-5分钟精简版讲述",
    "点评时间：每位学员1-2分钟"
])

add_content_slide(prs, "讲师点评的四大原则", [
    "原则一：具体肯定优先于笼统鼓励",
    "原则二：关键建议少而精优于面面俱到",
    "原则三：先说亮点，再说建议",
    "原则四：保护学员的信心和积极性"
])

add_content_slide(prs, "点评话术范例", [
    "\"你刚才讲到'那个深夜'的时候，我作为听众一下子就被带进去了\"",
    "\"你在转折处的停顿用得特别好，让我想知道接下来发生了什么\"",
    "\"建议你在手势上可以更放松一些，你现在有点端着\"",
    "\"整体结构很清晰，继续练，你会越来越自然的\""
])

add_content_slide(prs, "避免的点评方式", [
    "\"你的内容不错\"（太笼统）",
    "\"这里不对，那里也有问题\"（打击信心）",
    "\"你应该这样讲才对\"（否定学员风格）",
    "在全班面前让某位学员显得难堪（产生连锁防御心理）"
])

# ===== 第五章：工具与转化 =====
add_section_slide(prs, 5, "工具与转化", "带走可迁移的方法论")

add_content_slide(prs, "核心工具包", [
    "《党课素材转化卡》：从原始故事到讲稿框架的转化记录表",
    "《党课基础叙事结构模板》：五段式叙事结构，附示范案例",
    "《登台自查清单》：开场设计、停顿使用、眼神交流等要点速查",
    "《同伴反馈卡》：规范反馈只聚焦'感受层面'而非内容对错"
])

add_content_slide(prs, "学员离场时带走", [
    "一份经过四步法转化的《党课讲稿框架》",
    "经过录像回看、现场点评三轮打磨的最终版本",
    "一份《登台自查清单》",
    "一套可迁移的党课内容设计方法论"
])

add_content_slide(prs, "后续转化路径", [
    "回去后对着镜子练习讲3遍",
    "用自己的手机录下来，回看找问题",
    "正式讲党课前，在小范围先试讲一遍",
    "每次讲完后对照《登台自查清单》自检"
])

# ===== 课程总结 =====
add_section_slide(prs, 6, "课程总结", "")

add_content_slide(prs, "课程核心收获", [
    "认知：业务干部能讲好党课，因为有真实故事",
    "方法：案例转化四步法 — 故事盘点→主题锚定→结构搭建→语言转译",
    "训练：分层递进练习 — 从小组到全班，循序渐进",
    "工具：带走可迁移的党课设计方法论",
    "目标：敢讲、能讲、讲得自然"
])

add_content_slide(prs, "课程特色总结", [
    "不代写讲稿，只做赋能",
    "现场练到能自然讲出来为止",
    "16-18人小班，确保充分练习",
    "国际版权课标准的内容质量"
])

# ===== 结尾 =====
add_cover_slide(prs)

# 保存
os.makedirs(SLIDES_DIR, exist_ok=True)
output_path = os.path.join(SLIDES_DIR, '经营者讲党课-授课PPT.pptx')
prs.save(output_path)
print(f'PPT已保存到: {output_path}')
print(f'共 {len(prs.slides)} 页')
