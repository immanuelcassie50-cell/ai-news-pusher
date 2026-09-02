# -*- coding: utf-8 -*-
"""
《经营者讲党课》PPT生成脚本
红灰配色，浅底背景
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
THEME = {
    'primary': RGBColor(0xC4, 0x12, 0x30),   # 党建红
    'secondary': RGBColor(0x4A, 0x4A, 0x4A),  # 深灰
    'accent': RGBColor(0x8B, 0x45, 0x13),     # 红褐
    'light': RGBColor(0xF5, 0xF5, 0xF5),      # 浅灰白
    'bg': RGBColor(0xFF, 0xFF, 0xFF),         # 白色背景
    'text': RGBColor(0x33, 0x33, 0x33),       # 深色文字
    'gold': RGBColor(0xD4, 0xAF, 0x37),       # 金色
}

SLIDES_DIR = 'D:/新课开发/党业融合/经营者讲党课/完整课程包/003-授课PPT'
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['primary']
    background.line.fill.background()

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME['secondary']
    top_bar.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = THEME['light']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, section_num, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 浅色背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 左侧红色装饰条
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = THEME['primary']
    left_bar.line.fill.background()

    # 章节号
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(2), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = THEME['primary']

    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(3), Inches(7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = THEME['secondary']

    return slide

def add_content_slide(prs, title, content_items, has_left_bar=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 标题栏背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = THEME['primary']
    title_bg.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 左侧装饰条（可选）
    if has_left_bar:
        left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(0.08), prs.slide_height - Inches(1.1))
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = THEME['secondary']
        left_bar.line.fill.background()

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = THEME['text']
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = THEME['primary']
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 左栏
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(4.5), Inches(5.8))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
    left_box.line.color.rgb = THEME['primary']

    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = THEME['primary']

    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.1), Inches(4.8))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = THEME['text']
        p.space_before = Pt(8)

    # 右栏
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.8))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    right_box.line.color.rgb = THEME['secondary']

    right_title_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(4), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = THEME['secondary']

    right_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.1), Inches(4.1), Inches(4.8))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = THEME['text']
        p.space_before = Pt(8)

    return slide

def add_process_slide(prs, title, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = THEME['bg']
    background.line.fill.background()

    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = THEME['primary']
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 流程步骤
    step_width = Inches(2.8)
    for i, (step_num, step_title, step_desc) in enumerate(steps):
        x = Inches(0.5 + i * 3.1)

        # 步骤圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(1.5), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = THEME['primary']
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(x + Inches(0.9), Inches(1.55), Inches(0.8), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(step_num)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # 步骤标题
        title_box = slide.shapes.add_textbox(x, Inches(2.4), step_width, Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = THEME['secondary']
        p.alignment = PP_ALIGN.CENTER

        # 步骤描述
        desc_box = slide.shapes.add_textbox(x, Inches(2.9), step_width, Inches(3.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(12)
        p.font.color.rgb = THEME['text']
        p.alignment = PP_ALIGN.CENTER

        # 连接箭头
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.6), Inches(1.7), Inches(0.4), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = THEME['secondary']
            arrow.line.fill.background()

    return slide

# ============ 开始生成幻灯片 ============

# 封面
add_title_slide(prs,
    "讲党课",
    "业务干部的登台表达赋能工作坊")

# 课程定位
add_content_slide(prs, "课程定位", [
    "内容设计 + 登台表达",
    "不代写讲稿，只做一件事：",
    "教业务口干部把自己真实经历的管理故事，",
    "转化成一堂能讲、敢讲、学员爱听的党课，",
    "并且现场练到能自然讲出来为止。"
])

# 目标学员
add_two_column_slide(prs, "目标学员画像",
    "核心人群", [
        "\"一岗双责\"制度下必须自己上台讲党课的经营管理者/业务口中高层",
        "通常本人并非党务出身",
        "平时讲业务汇报没问题",
        "一到讲党课就发怵"
    ],
    "次要人群", [
        "即将被要求承担党课任务的后备中高层",
        "提前赋能，避免第一次上台就出丑"
    ])

# 常见困境
add_content_slide(prs, "常见困境", [
    "完全不会讲：只会念PPT或复述文件",
    "讲得像业务汇报：学员听得昏昏欲睡",
    "深层恐惧：怕在同事面前显得不专业",
    "表达焦虑：比\"不知道讲什么\"更影响实际表现"
])

# 课程边界
add_content_slide(prs, "课程边界", [
    "不代写讲稿、不代做PPT",
    "不介入党课内容的具体思想政治内核",
    "不做党史知识讲授、不做党建理论培训",
    "不涉及具体政策文件的解读争议",
    "AI工具不作为课程卖点"
])

# 课程特色
add_content_slide(prs, "课程特色", [
    "独家\"案例转化四步法\"方法论",
    "\"内容设计+登台表达\"双核训练",
    "现场练到能自然讲出来为止",
    "16-18人小班，确保每位学员充分练习"
])

# 章节一：案例转化四步法
add_section_slide(prs, 1, "案例转化四步法")

# 四步法总览
add_process_slide(prs, "案例转化四步法", [
    (1, "故事盘点", "从干部的真实经历中找出有戏剧张力的素材"),
    (2, "主题锚定", "找到故事与党课主题之间真实的连接点"),
    (3, "结构搭建", "用讲故事的方式重新组织内容"),
    (4, "语言转译", "把业务语言转化成适合讲台表达的语言")
])

# 第一步：故事盘点
add_content_slide(prs, "第一步：故事盘点", [
    "引导提问：\"你带团队这些年，有没有一次真的很难抉择、纠结到睡不着觉的时刻？\"",
    "关键筛选标准：好的党课素材必须有\"冲突\"和\"抉择\"",
    "常见误区：学员第一反应往往会讲\"团队多么努力\"这类苦劳型素材",
    "讲师需要引导往\"当时面临的真实两难选择是什么\"方向深挖"
])

# 第二步：主题锚定
add_content_slide(prs, "第二步：主题锚定", [
    "不是生硬地把故事\"扣\"到主题上",
    "而是问：\"这个故事让你自己悟到了什么道理？\"",
    "这个道理和今天要讲的主题，有没有真实的呼应？",
    "如果连接生硬，宁可换一个故事，也不要强行嫁接"
])

# 第三步：结构搭建
add_content_slide(prs, "第三步：结构搭建", [
    "五段式叙事结构：",
    "开场设置悬念/抛出真实困境（30秒内让听众进入情境）",
    "还原当时的纠结和几种可能的选择",
    "揭示当时的真实选择和过程中的转折",
    "结果与感悟（引导听众自己感受到那个道理）",
    "回扣主题，用一句简短有力的话收尾"
])

# 第四步：语言转译
add_content_slide(prs, "第四步：语言转译", [
    "业务干部平时习惯用数据、指标、专业术语表达",
    "需要引导其加入更多\"画面感\"和\"情绪细节\"的描述",
    "反复追问\"当时具体是什么场景、谁说了什么、你当时的第一反应是什么\"",
    "逼学员讲出细节而非概括"
])

# 章节二：登台表达训练
add_section_slide(prs, 2, "登台表达训练")

# 为什么表达训练更重要
add_content_slide(prs, "为什么表达训练更重要", [
    "业务干部真正的痛点往往不是不知道讲什么，",
    "而是\"写好了也不敢讲、讲的时候放不开\"",
    "本课程必须投入至少与内容设计同等甚至更多的时间在实际登台练习上",
    "否则内容设计做得再好，学员回去实战时依然会紧张卡壳"
])

# 登台心理阻力
add_two_column_slide(prs, "登台心理阻力拆解与应对",
    "阻力一：怕讲错", [
        "\"怕被人说不专业\"",
        "应对：强调党课不是政治理论考试，允许有个人风格和不完美",
        "讲师需要反复给出具体、真实的肯定反馈"
    ],
    "阻力二：放不开", [
        "觉得讲自己的故事很\"肉麻\"",
        "应对：先从小范围（2-3人小组）练习开始",
        "逐步过渡到全班展示，降低心理压力"
    ])

# 阻力三
add_content_slide(prs, "阻力三：照着稿子念", [
    "\"照着稿子念，脱稿就忘词\"",
    "训练的不是背诵稿子",
    "而是训练\"记住故事的骨架和几个关键画面\"",
    "用讲故事的自然节奏代替背诵的机械感"
])

# 分层递进练习
add_process_slide(prs, "分层递进练习体系", [
    (1, "小组内讲述", "2-3人一组，低压力环境先把内容捋顺"),
    (2, "录制回看", "让学员看到自己实际讲述时的状态"),
    (3, "全班展示", "每位学员完整讲述，讲师现场点评"),
    (4, "模拟真实场景", "有讲台、有完整时长的正式试讲")
])

# 第一层练习
add_content_slide(prs, "第一层：小组内讲述练习", [
    "目的：在低压力环境下先把内容捋顺，习惯\"说出来\"这件事本身",
    "同伴反馈聚焦：只反馈\"我作为听众，哪个瞬间让我有感觉、哪里我走神了\"",
    "不做内容对错评判"
])

# 第二层练习
add_content_slide(prs, "第二层：录制与回看", [
    "目的：让学员看到自己实际讲述时的状态（语速、眼神、肢体）",
    "很多人从未看过自己讲话的样子，这个环节冲击力很大",
    "回看录像后的引导提问：",
    "\"你觉得自己讲得最自然的是哪一段？最紧张的是哪一段？\""
])

# 第三层练习
add_content_slide(prs, "第三层：全班展示与点评", [
    "每位学员完整讲述一次（3-5分钟精简版）",
    "讲师现场给出具体、可操作的改进建议",
    "点评原则：先肯定具体的亮点细节",
    "再给1-2条最关键的改进建议",
    "不做面面俱到式的挑刺，避免打击学员积极性"
])

# 表达技巧
add_content_slide(prs, "表达技巧训练点", [
    "开场前3句话的设计：如何用一个问题或一个画面瞬间抓住听众注意力",
    "停顿的使用：关键转折处的停顿比语速快更有感染力",
    "眼神交流的基础训练：有意识地扫视不同区域",
    "手势的克制使用：讲故事时手势更自然放松"
])

# 章节三：课程安排
add_section_slide(prs, 3, "课程安排")

# 时间分配
add_two_column_slide(prs, "课程时间分配（1天6-7课时）",
    "上午（3.5课时）", [
        "开场导入（0.5课时）",
        "四步法讲解示范（1课时）",
        "学员分组练习（1.5课时）",
        "结构搭建初稿（0.5课时）"
    ],
    "下午（3.5课时）", [
        "语言转译打磨（0.5课时）",
        "分层练习第一层（0.5课时）",
        "分层练习第二层（1课时）",
        "分层练习第三层（1.5课时）"
    ])

# 人数控制
add_content_slide(prs, "人数控制", [
    "建议每场不超过16-18人",
    "原因：全班展示环节需要保证每位学员都有登台机会并获得讲师点评",
    "人数过多会导致点评质量下降或时间严重超支"
])

# 章节四：核心工具
add_section_slide(prs, 4, "核心工具包")

# 工具清单
add_content_slide(prs, "核心工具包清单", [
    "《党课素材转化卡》：从原始故事到讲稿框架的转化记录表",
    "《党课基础叙事结构模板》：五段式叙事结构，附示范案例",
    "《登台自查清单》：开场设计、停顿使用、眼神交流等要点速查",
    "《同伴反馈卡》：规范反馈只聚焦\"感受层面\"而非内容对错评判"
])

# 章节五：讲师要求
add_section_slide(prs, 5, "讲师要求")

# 讲师定位
add_content_slide(prs, "讲师定位", [
    "需要同时具备两种能力：",
    "既懂课程设计和叙事结构方法论",
    "又有足够的现场亲和力和点评的分寸感",
    "因为这门课直接触及学员\"怕丢面子\"的敏感心理"
])

# 讲师示范
add_content_slide(prs, "讲师现场示范的必要性", [
    "建议讲师在开场讲解四步法时，必须用自己的真实故事完整示范一遍转化过程",
    "而不是只讲理论框架",
    "这样才能让学员相信\"这套方法真的能把我的普通经历变成一个好故事\"",
    "单纯讲方法论而没有示范，学员的信任度和参与度会明显打折扣"
])

# 点评分寸
add_content_slide(prs, "点评环节的分寸把控", [
    "\"具体肯定优先于笼统鼓励、关键建议少而精优于面面俱到\"",
    "尤其要避免在全班展示环节让某位学员因为点评方式不当而当众显得难堪",
    "一旦有一个人因为点评受挫而回避，会在班级中产生连锁的防御心理"
])

# 学员收获
add_content_slide(prs, "学员收获", [
    "一份经过四步法转化、录像回看、现场点评三轮打磨的《党课讲稿框架》",
    "一份《登台自查清单》",
    "一套可迁移的党课内容设计方法论",
    "敢讲、能讲、讲得自然的自信"
])

# 结语
add_title_slide(prs,
    "谢谢",
    "祝大家讲出一堂好党课！")

# 保存
os.makedirs(SLIDES_DIR, exist_ok=True)
output_path = os.path.join(SLIDES_DIR, '经营者讲党课-授课PPT.pptx')
prs.save(output_path)
print(f'PPT已保存到: {output_path}')
print(f'共 {len(prs.slides)} 页')
