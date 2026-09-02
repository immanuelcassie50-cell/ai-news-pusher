#!/usr/bin/env python3
"""Merge all preview PPTX files into one final presentation."""

import os
import re
from pptx import Presentation
from pptx.util import Inches

slides_dir = r"D:\新课开发\HR\薪酬\10.全面薪酬新叙事：弹性福利与非物质回报的AI个性化设计\授课PPT\slides"
output_pptx = r"D:\新课开发\HR\薪酬\10.全面薪酬新叙事：弹性福利与非物质回报的AI个性化设计\授课PPT\全面薪酬新叙事_授课PPT.pptx"

def get_slide_number(filename):
    match = re.search(r'slide-(\d+)', filename)
    if match:
        return int(match.group(1))
    return 999999

# Find all preview PPTX files
pptx_files = []
for f in os.listdir(slides_dir):
    if f.endswith('.pptx') and '-preview' in f:
        pptx_files.append(f)

pptx_files.sort(key=get_slide_number)
print(f"Found {len(pptx_files)} preview PPTX files")

# Create merged presentation
merged_prs = Presentation()
merged_prs.slide_width = Inches(10)
merged_prs.slide_height = Inches(5.625)

# Use blank layout
blank_layout = merged_prs.slide_layouts[6]

added_count = 0
for i, pptx_file in enumerate(pptx_files):
    pptx_path = os.path.join(slides_dir, pptx_file)
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            new_slide = merged_prs.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    try:
                        new_slide.shapes.add_textbox(
                            shape.left, shape.top, shape.width, shape.height
                        ).text = shape.text
                    except:
                        pass
        added_count += 1
        if added_count % 20 == 0:
            print(f"[{added_count}] Added: {pptx_file}")
    except Exception as e:
        print(f"[ERROR] {pptx_file}: {e}")

print(f"\nAdded {added_count} slide groups to merged presentation")
print(f"Total slides: {len(merged_prs.slides)}")

# Save
merged_prs.save(output_pptx)
print(f"\nSaved: {output_pptx}")