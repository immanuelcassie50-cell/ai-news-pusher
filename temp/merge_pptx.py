#!/usr/bin/env python3
"""Merge multiple PPTX files into one presentation in order."""

from pptx import Presentation
from pptx.util import Inches
import os
import re

def get_slide_number(filename):
    """Extract slide number from filename like slide-05.pptx"""
    match = re.search(r'slide-(\d+)', filename)
    if match:
        return int(match.group(1))
    return 999999

def merge_pptx_files(input_dir, output_file):
    """Merge all PPTX files in input_dir into a single presentation."""

    # Get all pptx files (exclude preview files and .js.pptx files)
    pptx_files = []
    for f in os.listdir(input_dir):
        if f.endswith('.pptx') and f.startswith('slide-'):
            if 'preview' in f or '.js.' in f:
                continue
            pptx_files.append(f)

    # Sort by slide number
    pptx_files.sort(key=get_slide_number)

    print(f"Found {len(pptx_files)} PPTX files to merge")

    # Create new presentation
    merged_prs = Presentation()
    merged_prs.slide_width = Inches(10)
    merged_prs.slide_height = Inches(5.625)

    for i, pptx_file in enumerate(pptx_files):
        pptx_path = os.path.join(input_dir, pptx_file)
        try:
            prs = Presentation(pptx_path)
            for slide in prs.slides:
                if len(merged_prs.slide_layouts) > 0:
                    slide_layout = merged_prs.slide_layouts[0]
                else:
                    slide_layout = merged_prs.slide_layouts[0]

                new_slide = merged_prs.slides.add_slide(slide_layout)

                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        try:
                            new_slide.shapes.add_textbox(
                                shape.left, shape.top, shape.width, shape.height
                            ).text = shape.text
                        except:
                            pass

            print(f"[{i+1}/{len(pptx_files)}] Added: {pptx_file}")
        except Exception as e:
            print(f"[{i+1}/{len(pptx_files)}] Error with {pptx_file}: {e}")

    merged_prs.save(output_file)
    print(f"\nMerged presentation saved to: {output_file}")

if __name__ == "__main__":
    input_dir = r"D:\新课开发\HR\薪酬\10.全面薪酬新叙事：弹性福利与非物质回报的AI个性化设计\授课PPT\slides"
    output_file = r"D:\新课开发\HR\薪酬\10.全面薪酬新叙事：弹性福利与非物质回报的AI个性化设计\授课PPT\全面薪酬新叙事_授课PPT.pptx"

    merge_pptx_files(input_dir, output_file)
