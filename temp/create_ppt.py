from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2F, 0x54, 0x96)
    shape.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    sb = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    p = sb.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = RGBColor(0x2F, 0x54, 0x96)
    hdr.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    cb = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)

add_title(prs, "Hypothesis Status Report", "Huadong Manufacturing Group ERP Upgrade Project")

items1 = ["1. Project Initiation Phase Hypothesis List", "2. Hypothesis Verification Status", "3. Triggered Hypotheses and Response", "4. Key Decision Points", "5. Project Lessons Learned"]
add_content(prs, "Report Contents", items1)

items2 = ["H001: Data migration success rate assumption", "H005: Subsidiaries accept unified process assumption", "H006: Key users available full-time assumption", "H013: 12-month timeline is sufficient assumption", "H019: Counterpart has sufficient decision authority"]
add_content(prs, "High-Risk Hypotheses at Initiation", items2)

output_path = 'D:/CC/新课开发/工作手册/假设管理：项目经理的风险前置手册/完整课程包/07-成果demo/假设状态追踪报告-演示.pptx'
prs.save(output_path)
print(f"OK: {output_path}")
