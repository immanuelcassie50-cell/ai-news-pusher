# -*- coding: utf-8 -*-
"""
Training ROI Course PPT Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

RED = RGBColor(0xC0, 0x00, 0x00)
DARK_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF8, 0xF8)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, left, top, width, height, size=44, color=RED, bold=True):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_text(slide, text, left, top, width, height, size=18, color=DARK_GRAY):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.line_spacing = 1.5
        p.space_after = Pt(6)
    return shape

def add_bullet(slide, items, left, top, width, height, size=16):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.line_spacing = 1.5
        p.space_after = Pt(8)
    return shape

def add_table(slide, data, left, top, width, height, col_widths=None, font_size=12):
    rows = len(data)
    cols = len(data[0]) if data else 0
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for ri, row in enumerate(data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.name = "Microsoft YaHei"

            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = RED
            else:
                p.font.color.rgb = DARK_GRAY
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG

            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table

def section_divider(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RED)

    shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
    return slide

def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    add_title(slide, "培训闭环与效果证明", Inches(0.8), Inches(2), Inches(11.5), Inches(1), size=48)
    add_title(slide, "从完课率到ROI的数据驱动实操路径", Inches(0.8), Inches(3), Inches(11.5), Inches(0.8), size=32, color=DARK_GRAY, bold=False)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4), Inches(4), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    add_text(slide, "授课讲师：罗宏伟", Inches(0.8), Inches(4.5), Inches(5), Inches(0.5), size=18)
    add_text(slide, "课程时长：2天（每天6小时，共12小时）", Inches(0.8), Inches(5.1), Inches(6), Inches(0.5), size=16)
    return slide

def page_header(slide, prs):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

def section_header(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    add_title(slide, title, Inches(0.5), Inches(2.8), Inches(12), Inches(1.2), size=44)

    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(4.2), Inches(12), Inches(0.8), size=20)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.2), Inches(3), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    return slide

def content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    page_header(slide, prs)
    add_title(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7), size=32)

    y = 1.1
    for item in items:
        if isinstance(item, tuple):
            add_title(slide, "• " + item[0], Inches(0.5), Inches(y), Inches(12), Inches(0.4), size=18, color=DARK_GRAY, bold=True)
            y += 0.4
            add_text(slide, "  " + item[1], Inches(0.9), Inches(y), Inches(11.5), Inches(0.5), size=14)
            y += 0.6
        else:
            add_bullet(slide, [item], Inches(0.5), Inches(y), Inches(12), Inches(0.5), size=16)
            y += 0.55
    return slide

def table_slide(prs, title, headers, data, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    page_header(slide, prs)
    add_title(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7), size=32)

    all_data = [headers] + data
    h = min(5.5, 0.4 * len(all_data))
    add_table(slide, all_data, Inches(0.3), Inches(1.1), Inches(12.5), Inches(h), col_widths=col_widths, font_size=13)
    return slide

def two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    page_header(slide, prs)
    add_title(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7), size=32)

    add_title(slide, left_title, Inches(0.5), Inches(1.1), Inches(5.5), Inches(0.5), size=20)
    y = 1.7
    for item in left_items:
        add_bullet(slide, [item], Inches(0.5), Inches(y), Inches(5.8), Inches(0.5), size=14)
        y += 0.55

    add_title(slide, right_title, Inches(6.8), Inches(1.1), Inches(5.5), Inches(0.5), size=20)
    y = 1.7
    for item in right_items:
        add_bullet(slide, [item], Inches(6.8), Inches(y), Inches(5.8), Inches(0.5), size=14)
        y += 0.55
    return slide

def quote_slide(prs, quote, source=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

    qm = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(1), Inches(1))
    qtf = qm.text_frame
    qp = qtf.paragraphs[0]
    qp.text = '"'
    qp.font.size = Pt(120)
    qp.font.color.rgb = RED
    qp.font.name = "Georgia"

    add_text(slide, quote, Inches(1), Inches(2.5), Inches(11), Inches(3), size=24)

    if source:
        add_text(slide, "- " + source, Inches(8), Inches(5.5), Inches(4), Inches(0.5), size=16)
    return slide

def key_point_slide(prs, title, point, exp=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    page_header(slide, prs)
    add_title(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7), size=32)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box.line.color.rgb = RED
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = point
    p.font.size = Pt(22)
    p.font.color.rgb = RED
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    if exp:
        p2 = tf.add_paragraph()
        p2.text = exp
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = "Microsoft YaHei"
    return slide

def summary_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    page_header(slide, prs)
    add_title(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.7), size=32)

    y = 1.2
    for i, item in enumerate(items):
        nb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.35), Inches(0.35))
        nb.fill.solid()
        nb.fill.fore_color.rgb = RED
        nb.line.fill.background()

        tf = nb.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text(slide, item, Inches(1), Inches(y), Inches(11.5), Inches(0.5), size=16)
        y += 0.6
    return slide

def kano_table_slide(prs, title, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    page_header(slide, prs)
    add_title(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.6), size=32)

    headers = data[0]
    add_table(slide, data, Inches(0.3), Inches(1.0), Inches(12.5), Inches(4), font_size=14)
    return slide

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Cover
    cover(prs)

    # Course Objectives
    content_slide(prs, "课程学习目标", [
        ("认知重构", "理解从\"完课率\"到\"ROI\"的全链路评估逻辑，建立数据驱动思维"),
        ("体系建设", "能够设计一套覆盖学习行为数据、业务关联数据的采集方案"),
        ("指标设计", "能够根据培训目标，设计从反应层到成果层的完整指标体系"),
        ("ROI计算", "能够计算培训的货币价值，向管理层证明培训的投资回报率"),
        ("闭环运营", "能够基于数据，持续优化培训计划并推动迭代改进"),
        ("成果汇报", "能够向管理层呈现专业的培训效果报告，驱动资源支持"),
    ])

    # Course Outline
    content_slide(prs, "课程学习地图", [
        "模块一：培训效果评估的认知升级 — 为什么做了培训却说不清价值",
        "模块二：数据采集体系建设 — 数据从哪来、怎么采、采什么",
        "模块三：培训指标体系建设 — 指标那么多，哪些才有用",
        "模块四：培训ROI计算实操 — ROI到底怎么算",
        "模块五：数据驱动闭环运营 — 数据有了，怎么用起来",
        "模块六：成果展示与汇报 — 怎么让管理层相信培训有价值",
        "综合演练：培训效果证明手册 — 把所有产出整合成一套完整方案",
    ])

    # Opening Quote
    quote_slide(prs, "企业培训行业有一个被反复验证的数据差距：建立了完整培训效果评估体系的企业，其培训ROI是仅有基础线上学习追踪企业的2.3倍。", "ATD State of the Industry Report, 2023")

    key_point_slide(prs, "核心问题", "培训做了，课也上了，学员满意度也不错——但老板问'培训到底带来了什么业务价值'的时候，你拿不出能让管理层信服的数据。",
        "这个困境的背后，是大多数培训管理者只掌握了\"完课率\"这一层指标，而没有建立起从\"学习行为\"到\"业务成果\"的完整数据链条。")

    # ===== Part 1 =====
    section_divider(prs, "第一部分", "培训效果评估的认知升级")

    section_header(prs, "为什么\"完课率85%\"不是一个答案", "完课率是学习行为指标，不是效果指标")

    content_slide(prs, "2.3倍差距的真正含义", [
        ("不是\"花钱多效果好\"", "两类企业的培训预算中位数没有显著差异"),
        ("而是\"知道如何追踪和优化培训价值\"的企业，和\"只知道完课率\"的企业之间的差距", ""),
        "知道哪些培训真的有效，从而优化资源分配",
        "知道哪些培训无效，从而停止浪费",
        "能够用数据证明培训价值，从而获得更多资源",
    ])

    table_slide(prs, "柯氏四级评估框架", ["层级", "名称", "核心问题", "衡量内容", "常见指标"], [
        ["L1", "反应层", "学员喜欢这个培训吗？", "学员对培训的即时感受", "满意度评分、净推荐值（NPS）"],
        ["L2", "学习层", "学员学到了什么？", "知识、技能、态度的习得程度", "考试/测试成绩、技能考核得分"],
        ["L3", "行为层", "学员回到岗位后，行为改变了吗？", "培训内容在工作中的应用程度", "行为观察合格率、标准化操作执行率"],
        ["L4", "成果层", "培训对业务产生了什么影响？", "业务结果的改善情况", "质量合格率、客诉率、离职率、ROI"],
    ])

    key_point_slide(prs, "关键认知", "四个层级不是\"四选一\"，而是层层递进的关系。",
        "没有前一层级打好的基础，就不可能有后一层级的结果。没有L2（学到知识），不可能有L3（行为改变）；没有L3（行为改变），不可能有L4（业务成果）。")

    content_slide(prs, "完课率在哪一级", [
        "完课率不属于任何一级评估",
        "完课率是一个学习行为指标，反映的是\"学员有没有完成学习过程\"",
        "完课率可以辅助L1和L2的评估，但不能替代任何一个层级的评估",
    ])

    table_slide(prs, "行业现状：大多数企业停在哪个层级", ["评估层级", "实施比例"], [
        ["L1 反应层（有系统性的满意度调查）", "约75%的企业"],
        ["L2 学习层（有考核或测评）", "约50%的企业"],
        ["L3 行为层（有结构化的行为追踪）", "约20%的企业"],
        ["L4 成果层（有业务指标关联）", "约10%的企业"],
        ["L4+ ROI计算（货币化收益）", "不到5%的企业"],
    ])

    section_header(prs, "评估升级的三大障碍")

    two_col_slide(prs, "评估升级的三大障碍",
        "障碍一：数据缺失", [
            "L3行为层和L4成果层的数据没有被系统性采集",
            "LMS只记录\"有没有完成\"，不记录\"用没用\"",
            "业务数据在业务系统里，培训部没有权限",
        ],
        "障碍二：指标模糊", [
            "没有把\"培训效果\"转化为可量化的业务指标",
            "说\"提升管理能力\"，但没有定义用什么数字衡量",
            "\"减少质量事故\"，但没有明确\"减少多少算有效\"",
        ])

    content_slide(prs, "障碍三：因果归因", [
        "无法证明业务指标的改善是培训造成的，而非其他因素",
        "\"这个季度销售额提升了，是培训的效果吗？还是市场好了？\"",
        "\"客诉率下降了，是培训的效果？还是换了客服主管？\"",
        "业务指标的改善可能来自多种因素，培训只是其中一种",
    ])

    section_header(prs, "企业培训评估层级自测")

    content_slide(prs, "第一步：对照各层级标准，判断你企业目前做到了哪个层级", [
        "L1 反应层：有没有系统性的培训满意度调查",
        "L2 学习层：有没有考核或测评机制，能衡量学员的知识/技能习得程度",
        "L3 行为层：有没有结构化的机制追踪学员培训后的行为改变",
        "L4 成果层：能不能将培训与业务指标关联，能看到业务指标的改善",
        "ROI计算：能不能将培训收益货币化，计算出ROI",
    ])

    content_slide(prs, "第二步：识别你企业向L3/L4升级的核心障碍", [
        "数据缺失 — 行为层和成果层的数据没有系统性采集",
        "指标模糊 — \"培训效果\"没有转化为可量化的业务指标",
        "因果归因 — 无法证明业务改善是培训带来的",
        "系统断连 — 学习数据（培训系统）和业务数据（业务系统）没有打通",
        "部门壁垒 — 业务部门不配合提供数据",
    ])

    section_header(prs, "本部分小结")

    content_slide(prs, "第一部分：知识框架", [
        "核心问题：为什么做了培训却说不清价值",
        "模块一：2.3倍差距的真正含义 — 不是\"花钱多效果好\"",
        "模块二：柯氏四级评估框架 — L1→L2→L3→L4层层递进",
        "模块三：评估升级的三大障碍 — 数据缺失/指标模糊/因果归因",
    ])

    # ===== Part 2 =====
    section_divider(prs, "第二部分", "数据采集体系建设")

    key_point_slide(prs, "开场思考", "没有数据，评估就是空中楼阁。",
        "ROI计算需要业务数据，行为追踪需要过程数据，而这些数据不会凭空出现——需要你主动去设计、去推动、去建立机制。")

    section_header(prs, "数据采集的三个原则")

    table_slide(prs, "数据采集三原则", ["原则", "含义", "没做到的后果"], [
        ["可追踪", "数据能够被系统性地记录，而不是靠人工回忆或临时填表", "数据不完整、不真实"],
        ["可量化", "数据必须转化为数字形式，才能进行后续的计算和比较", "无法做统计分析"],
        ["可关联", "学习数据必须能够和业务数据打通，否则永远只是\"学习数据\"", "无法做ROI分析"],
    ])

    key_point_slide(prs, "关键认知", "三个原则的关系是层层递进的——不可追踪的数据，必然不可量化；不可量化的数据，必然不可关联。",
        "可关联是最高标准，也是ROI计算的前提。")

    section_header(prs, "学习行为数据的五层采集框架")

    table_slide(prs, "学习行为数据五层框架", ["数据层", "描述", "具体数据项", "采集难度"], [
        ["L0 基础层", "学员的静态信息", "部门、岗位、司龄、职级、学历等", "⭐ 容易"],
        ["L1 登录层", "学员有没有来学习", "登录次数、登录时间分布、最后登录时间", "⭐ 容易"],
        ["L2 完课层", "学员有没有学完", "课程完成率、章节完成率、学习进度排名", "⭐ 容易"],
        ["L3 互动层", "学员有没有深度参与", "讨论区发言、答疑提问、作业提交、笔记记录", "⭐⭐ 中等"],
        ["L4 掌握层", "学员真的掌握了吗", "测评得分、实操练习得分、场景模拟成绩", "⭐⭐⭐ 较难"],
    ], col_widths=[1.2, 1.5, 3.5, 1.0])

    section_header(prs, "业务关联数据的四种采集方法")

    table_slide(prs, "业务数据采集方法", ["方法", "适用场景", "操作要点"], [
        ["业务系统对接", "有标准化CRM/OA/ERP/MES系统的企业", "与IT部门协作，通过API或数据报表获取"],
        ["业务部门定期填报", "业务数据难以自动采集的场景", "设计简化的业务数据填报模板，每月/每季度收集"],
        ["绩效数据关联", "培训与绩效评估体系挂钩的企业", "从绩效管理系统获取培训前后的绩效对比数据"],
        ["问卷/访谈调查", "缺乏系统数据支持的企业", "设计标准化的业务改善感知问卷，定期收集"],
    ])

    section_header(prs, "数据采集的常见坑与应对")

    table_slide(prs, "四大常见坑与应对", ["坑", "具体表现", "后果", "应对方法"], [
        ["数据孤岛", "学习数据和业务数据存在不同系统，无法关联", "无法做ROI分析", "建立统一的学员ID映射，或推动系统集成"],
        ["数据质量差", "完课数据依赖学员\"自点击\"，可以刷课", "数据失真", "引入深度学习行为指标（停留时长、互动频率）"],
        ["采集负担重", "业务部门抱怨填报数据太麻烦，配合度低", "数据中断或不完整", "最小化数据采集字段，只采集真正需要的"],
        ["数据隐私合规", "采集员工个人行为数据未经过合规审批", "法律风险", "提前做好数据使用协议和隐私合规审查"],
    ], col_widths=[1.2, 2.2, 1.8, 3.5])

    section_header(prs, "本部分小结")

    content_slide(prs, "第二部分：知识框架", [
        "模块一：数据采集三原则 — 可追踪→可量化→可关联",
        "模块二：学习行为数据五层框架 — L0基础层→L4掌握层",
        "模块三：业务数据四种采集方法 — 系统对接/定期填报/绩效关联/问卷访谈",
        "模块四：四大常见坑 — 数据孤岛/数据质量差/采集负担重/隐私合规",
    ])

    # ===== Part 3 =====
    section_divider(prs, "第三部分", "培训指标体系建设")

    key_point_slide(prs, "开场思考", "指标设计，不是从\"能采集什么\"开始。",
        "\"能采集什么\"不等于\"应该衡量什么\"。真正有效的指标设计，逻辑应该是\"目标导向\"——从业务目标开始，逆向推导需要衡量什么指标。")

    section_header(prs, "指标设计的三个起点")

    table_slide(prs, "指标设计三起点", ["起点", "问题", "输出"], [
        ["业务需求", "这个培训要解决什么业务问题？", "业务目标（定性）"],
        ["行为目标", "为了解决这个问题，学员需要做出什么行为改变？", "行为目标"],
        ["学习目标", "学员需要学习什么内容、投入多少时间？", "学习目标"],
    ])

    section_header(prs, "三层指标体系的构建方法")

    table_slide(prs, "三层指标体系", ["层级", "指标类型", "回答的问题", "时间维度"], [
        ["反应层", "满意度/体验指标", "学员喜欢这个培训吗？学得愉快吗？", "培训结束时"],
        ["行为层", "应用/行为指标", "学员回到岗位后，有没有用上培训所学？", "培训后1-3个月"],
        ["成果层", "业务/成效指标", "业务目标达成了吗？问题解决了吗？", "培训后3-12个月"],
    ])

    section_header(prs, "行为层指标设计——BEO框架")

    table_slide(prs, "BEO框架", ["要素", "全称", "含义", "示例"], [
        ["B", "Behavior（行为）", "具体、可观察的学员行为", "\"使用新的话术与客户沟通\""],
        ["E", "Evidence（证据）", "如何证明这个行为发生了", "\"销售会话记录中出现了新话术的关键词\""],
        ["O", "Observation（观察条件）", "在什么条件下观察", "\"在模拟电话场景中\""],
    ])

    section_header(prs, "成果层指标设计")

    table_slide(prs, "成果层指标类型", ["类型", "特点", "示例", "计算难度"], [
        ["直接业务指标", "可以直接归因到培训效果", "质检合格率从92%提升到96%", "⭐⭐⭐ 较难"],
        ["间接业务指标", "受多因素影响，但可作为参考", "客户满意度提升", "⭐⭐⭐ 较难"],
        ["财务指标", "可以货币化呈现", "废品率降低节省成本X元", "⭐⭐⭐⭐ 难"],
    ])

    section_header(prs, "指标权重设计")

    table_slide(prs, "权重设计原则", ["层级", "权重原则", "说明"], [
        ["反应层", "权重较低", "体验层是基础，但非核心目标"],
        ["行为层", "权重中等", "连接学习和成果的关键"],
        ["成果层", "权重最高", "业务成效是最终目标"],
    ])

    key_point_slide(prs, "常见误区", "把大量权重放在反应层（满意度占70%以上）。",
        "这会导致\"追求满意度\"而非\"追求效果\"——为了让学生说好话，降低考核难度、减少课程内容。")

    section_header(prs, "本部分小结")

    content_slide(prs, "第三部分：知识框架", [
        "模块一：指标设计的三个起点 — 业务需求→行为目标→学习目标→指标体系",
        "模块二：三层指标体系构建 — L1反应层/L3行为层/L4成果层",
        "模块三：指标权重设计 — 行为层权重≥反应层权重，成果层权重≥行为层权重",
    ])

    # ===== Part 4 =====
    section_divider(prs, "第四部分", "培训ROI计算实操")

    key_point_slide(prs, "开场思考", "为什么ROI是终极答案？",
        "满意度、完课率、测评成绩——这些数字都回答不了\"值不值\"这个问题。ROI把培训的收益和成本都换算成\"钱\"，让你可以直接回答\"这个培训赚了多少、亏了多少\"。")

    section_header(prs, "ROI公式与底层逻辑")

    content_slide(prs, "培训ROI的标准公式", [
        "ROI（%）= [（培训总收益 - 培训总成本） ÷ 培训总成本] × 100%",
        "",
        "变量说明：",
        "• 培训总成本：实施这个培训项目所花费的全部资金（直接+间接）",
        "• 培训总收益：培训带来的货币化价值",
        "• 净货币收益：总收益减去总成本的差值",
    ])

    table_slide(prs, "ROI值解读", ["ROI值", "含义", "解读"], [
        ["ROI = 200%", "每投入1元，赚了2元", "绩效显著，非常理想"],
        ["ROI = 100%", "每投入1元，赚了1元", "收回成本，还赚1元，达标"],
        ["ROI = 50%", "每投入1元，赚了0.5元", "有正收益，但未收回成本"],
        ["ROI = 0%", "每投入1元，不赔不赚", "刚好收回成本"],
        ["ROI = -30%", "每投入1元，亏了0.3元", "亏损，需要分析原因"],
    ])

    section_header(prs, "成本计算的完整框架")

    table_slide(prs, "成本分类总览", ["成本类型", "明细项", "特点"], [
        ["直接成本", "讲师费/课酬、教材/讲义印刷、场地/设备租赁、外部培训/认证费、差旅费", "可以直接计入该项目"],
        ["间接成本", "学员参加培训的时间成本、培训管理人员的工资分摊、培训系统/平台的运维成本、机会成本", "需要分摊计算，容易被忽视"],
    ])

    key_point_slide(prs, "最常见的错误", "只计算直接成本，忽略间接成本（尤其是学员时间成本），导致成本被严重低估，ROI被虚高。")

    section_header(prs, "货币收益的计算方法")

    two_col_slide(prs, "三种收益计算方法",
        "方法一：直接法（最准确）", [
            "适用：生产质量指标、废品率、事故率等",
            "步骤：确定业务指标→获取基线值→获取改善后值→计算改善幅度→乘以业务量→乘以单位价值",
        ],
        "方法二：间接法（专家判断）", [
            "适用：难以直接量化的业务改善",
            "通过专家判断或业务负责人估算来确定收益",
            "必须明确说明这是\"基于专家判断的估算\"",
        ])

    section_header(prs, "效果隔离——证明因果关系的技术")

    key_point_slide(prs, "为什么效果隔离这么难", "业务指标的改善，几乎永远是\"多因素\"的。",
        "业务指标改善 = 培训效果 + 市场因素 + 团队因素 + 管理层因素 + 运气 + ...")

    table_slide(prs, "四种效果隔离方法", ["方法", "可靠性", "适用条件"], [
        ["对照组设计", "⭐⭐⭐⭐⭐ 最高", "可以随机分配学员到实验组和对照组"],
        ["前后对比", "⭐⭐⭐⭐ 较高", "无法设置对照组，但有历史数据"],
        ["趋势分析", "⭐⭐⭐ 中等", "业务指标有明显的趋势或季节性"],
        ["业务负责人判断", "⭐⭐ 参考使用", "其他方法都不可行时"],
    ])

    table_slide(prs, "归因比例参考", ["归因比例", "含义", "适用场景"], [
        ["80%-100%", "大部分归因于培训", "培训与业务指标高度相关，且无明显其他因素"],
        ["50%-80%", "多因素共同作用，培训是主要因素", "最常见的归因区间"],
        ["20%-50%", "多因素共同作用，培训是辅助因素", "当业务指标受多种因素显著影响时"],
        ["<20%", "培训贡献较小", "应坦诚披露，并说明其他因素"],
    ])

    key_point_slide(prs, "关键原则", "归因比例必须有依据、有逻辑、保守。宁可低估，不要高估。")

    section_header(prs, "ROI计算五步骤")

    content_slide(prs, "ROI计算步骤", [
        "第一步：收集数据 — 采集培训项目各层级指标的数据",
        "第二步：换算货币收益 — 将业务指标的改善换算为货币价值",
        "第三步：计算培训成本 — 全面计算培训项目的所有成本",
        "第四步：隔离培训效果 — 通过方法分离培训带来的效果",
        "第五步：计算ROI — 套用公式，计算ROI",
    ])

    section_header(prs, "ROI结果的解读与沟通")

    table_slide(prs, "不同ROI结果的沟通策略", ["ROI值", "含义", "沟通策略"], [
        ["> 150%", "非常理想", "重点说明成功因素，可以作为最佳实践推广"],
        ["100%-150%", "达标", "客观呈现，说明ROI计算方法和假设"],
        ["50%-100%", "有正收益但未达标准", "坦诚分析制约因素，提出改进建议"],
        ["< 50%", "ROI较低", "深入分析原因，判断是否继续投入"],
        ["负值", "亏损", "诚实说明，分析原因，考虑停止或大幅调整"],
    ])

    section_header(prs, "本部分小结")

    content_slide(prs, "第四部分：知识框架", [
        "模块一：ROI公式 — ROI=（净货币收益÷总成本）×100%",
        "模块二：成本计算 — 直接成本vs间接成本，学员时间成本最易被低估",
        "模块三：收益计算 — 直接法（最准）→间接法→影子价值",
        "模块四：效果隔离 — 对照组→前后对比→趋势分析→专家判断",
    ])

    # ===== Part 5 =====
    section_divider(prs, "第五部分", "数据驱动闭环运营")

    key_point_slide(prs, "开场思考", "为什么大多数培训\"无效\"？",
        "培训的真正价值，不是在培训结束那一刻释放的，而是在培训之后的工作中持续体现的。但大多数培训管理者的精力，都花在了\"培训前\"和\"培训中\"——培训一结束，工作就结束了。")

    section_header(prs, "闭环思维的本质")

    content_slide(prs, "培训闭环的核心逻辑", [
        "计划（Plan）→ 实施（Do）→ 检查（Check）→ 改进（Act）→ 计划（Plan'）→ ...",
        "",
        "这不是一个线性流程，而是一个不断螺旋上升的循环。",
    ])

    table_slide(prs, "闭环为什么重要", ["对比维度", "没有闭环", "有闭环"], [
        ["培训价值", "一次性释放", "持续释放"],
        ["问题发现", "很久之后才发现", "实时监控，快速响应"],
        ["改进行动", "凭经验、拍脑袋", "凭数据、有依据"],
        ["资源分配", "平均分配或看领导脸色", "精准投放给高ROI项目"],
        ["培训团队价值", "\"花钱的部门\"", "\"创造价值的部门\""],
    ])

    section_header(prs, "PDCA循环在培训中的应用")

    table_slide(prs, "PDCA四阶段", ["阶段", "核心任务", "关键动作"], [
        ["Plan（计划）", "明确目标、设计指标、规划数据采集", "业务需求分析、目标设定、指标设计"],
        ["Do（实施）", "按计划执行，同时采集过程数据", "课程设计、教学实施、学习支持"],
        ["Check（检查）", "分析数据、评估效果、发现问题", "数据收集、指标分析、ROI计算"],
        ["Act（改进）", "基于检查结果，优化下一轮计划", "问题识别、改进措施、更新计划"],
    ])

    section_header(prs, "数据驱动决策的三个层级")

    table_slide(prs, "数据分析三层级", ["层级", "描述", "决策类型", "示例"], [
        ["描述性分析", "发生了什么", "回顾性判断", "\"这个培训项目的完课率是75%\""],
        ["诊断性分析", "为什么会发生", "原因分析", "\"完课率低是因为第三章的实操练习太难\""],
        ["预测性分析", "接下来会发生什么", "前瞻性决策", "\"如果把第三章拆成两节，完课率预计能到85%\""],
    ])

    key_point_slide(prs, "关键认知", "大多数培训管理者停留在\"描述性分析\"阶段——知道发生了什么，但不知道为什么、也不知道下一步怎么做。",
        "闭环运营要求至少进入\"诊断性分析\"阶段，而理想状态是达到\"预测性分析\"阶段。")

    section_header(prs, "培训运营仪表盘设计")

    table_slide(prs, "仪表盘核心模块", ["模块", "展示内容", "关键指标", "更新频率"], [
        ["学习概览", "整体学习情况", "完课率、平均学习时长、激活率", "实时"],
        ["效果概览", "培训效果汇总", "各层级指标达成率、ROI估算", "月度"],
        ["项目追踪", "各培训项目进展", "单项目完课率、通过率、满意度", "周度"],
        ["异常预警", "需要关注的问题", "完课率低于阈值、满意度下降", "实时"],
    ])

    section_header(prs, "本部分小结")

    content_slide(prs, "第五部分：知识框架", [
        "模块一：闭环思维 — Plan→Do→Check→Act循环",
        "模块二：PDCA四阶段 — 计划/实施/检查/改进",
        "模块三：数据驱动决策三层级 — 描述性→诊断性→预测性",
    ])

    # ===== Part 6 =====
    section_divider(prs, "第六部分", "成果展示与汇报")

    key_point_slide(prs, "开场思考", "为什么汇报比培训本身更重要？",
        "培训效果汇报，是培训价值的最终证明时刻。你前期做得再好，如果在汇报时不能让管理层信服，所有的努力都会打折扣。")

    section_header(prs, "管理层真正想听什么")

    table_slide(prs, "管理层三问", ["问题", "潜台词", "汇报重点"], [
        ["\"花了多少钱？\"", "投资是否合理", "成本结构、投入产出比"],
        ["\"带来了什么？\"", "收益是否值得", "业务指标改善、ROI"],
        ["\"和其他选择比呢？\"", "机会成本", "对比分析、替代方案"],
    ])

    key_point_slide(prs, "关键认知", "管理层不是不想听数据，而是想听和钱有关的数据。",
        "\"满意度、完课率\"这些\"过程数据\"对他们来说是噪音，\"ROI是多少\"才是信号。")

    section_header(prs, "不同受众的信息需求")

    table_slide(prs, "四类典型受众", ["受众", "最关心的点", "汇报重点", "呈现方式"], [
        ["CEO/高管", "战略价值、ROI", "ROI数字、业务影响力", "一页纸执行摘要、数字+故事"],
        ["业务部门负责人", "业务指标改善", "具体指标的前后对比", "数据对比图表、业务语言"],
        ["HR同行/培训同行", "方法论、数据质量", "评估方法论、数据可靠性", "详细数据、工具模板"],
        ["培训学员", "学习收获", "自己的成长和收获", "个人学习报告、证书"],
    ])

    section_header(prs, "七步汇报结构")

    table_slide(prs, "汇报七步", ["步骤", "内容", "时间占比", "注意事项"], [
        ["1. 开场", "培训项目背景和目标", "10%", "30秒内切入，不要从历史讲起"],
        ["2. 方法", "评估框架和数据来源", "10%", "建立可信度"],
        ["3. 投入", "成本结构说明", "10%", "透明呈现，建立信任"],
        ["4. 产出-反应", "满意度数据", "10%", "不要花太多时间"],
        ["5. 产出-学习", "学习效果数据", "15%", "重点展示技能/知识提升"],
        ["6. 产出-行为/成果", "业务指标改善+ROI", "35%", "核心重点，用数据说话"],
        ["7. 结论与建议", "下一步行动", "10%", "提出具体建议"],
    ])

    section_header(prs, "数据可视化五项原则")

    content_slide(prs, "选对图表类型", [
        "比较：用柱状图",
        "趋势：用折线图",
        "构成：用饼图/堆叠柱状图",
        "关联：用散点图",
        "分布：用直方图/箱线图",
    ])

    table_slide(prs, "五项原则", ["原则", "说明", "常见错误"], [
        ["选对图表类型", "比较用柱状图、趋势用折线图、构成用饼图", "用错图表类型导致信息传达错误"],
        ["去掉不必要的装饰", "3D效果、阴影、过多颜色都是干扰", "华而不实的图表反而降低可信度"],
        ["突出关键数字", "用箭头、色块标注关键数字", "听众需要找很久才能看到重点"],
        ["注明数据来源", "每个图表下方标注数据来源和时间", "没有来源的数据=没有说服力"],
        ["结论先行", "先说结论，再说图表", "不要让听众自己解读"],
    ])

    section_header(prs, "汇报现场应对技巧")

    table_slide(prs, "六类典型质疑与应对", ["问题类型", "典型问题", "应对思路"], [
        ["质疑数据可信度", "\"这个数据是怎么来的？\"", "提前准备数据来源说明"],
        ["质疑因果归因", "\"怎么证明是培训的效果？\"", "说明效果隔离方法和局限性"],
        ["质疑方法论", "\"为什么用这个评估框架？\"", "解释选择原因和行业通用性"],
        ["关注成本", "\"花这么多钱，值得吗？\"", "直接展示ROI数字和对标数据"],
        ["要求更多信息", "\"能再详细讲讲吗？\"", "准备附录材料"],
        ["提出反对意见", "\"我觉得没什么用\"", "先倾听，再以数据回应"],
    ])

    section_header(prs, "本部分小结")

    content_slide(prs, "第六部分：知识框架", [
        "模块一：培训效果汇报的本质 — 回答\"花了多少钱/带来了什么/值不值\"三个问题",
        "模块二：不同受众需求 — CEO/业务部门/HR同行/学员",
        "模块三：七步汇报结构 — 开场→方法→投入→产出→结论",
    ])

    # ===== Part 7 =====
    section_divider(prs, "第七部分", "综合演练与课程总结")

    key_point_slide(prs, "核心目标", "将前六个模块的产出整合成一套完整、连贯、可直接在实际工作中使用的\"培训效果证明手册\"。",
        "这不是一个孤立的练习，而是把前六天所有的产出串联成一条完整的证据链。")

    section_header(prs, "演练流程")

    table_slide(prs, "综合演练五阶段", ["阶段", "时间", "内容"], [
        ["回顾与整合", "20分钟", "回顾前六个模块的产出，确认每个模块的产出都已完成"],
        ["串联与补全", "40分钟", "将各模块产出串联成完整链路，识别缺失环节并补全"],
        ["模拟汇报", "30分钟", "每人用3分钟向组内汇报自己的培训效果证明手册"],
        ["反馈与优化", "20分钟", "组内互评，提出改进建议"],
        ["总结与承诺", "10分钟", "每人写下课程结束后的第一个行动承诺"],
    ])

    section_header(prs, "培训效果证明手册模板")

    content_slide(prs, "手册结构", [
        "1. 培训项目概述（项目名称、业务背景、培训目标、目标学员、培训形式、时间）",
        "2. 评估框架与指标体系（柯氏四级、指标体系）",
        "3. 数据采集方案（采集字段清单、来源和责任人、采集频率）",
        "4. 数据分析结果（反应层、学习层、行为层、成果层数据）",
        "5. ROI计算（成本明细、收益明细、ROI计算结果、效果隔离说明）",
        "6. 结论与建议（培训效果总结、成功因素分析、待改进领域、下一步行动计划）",
        "7. 附件（原始数据、数据采集工具、评估量表）",
    ])

    section_header(prs, "课程全景回顾")

    content_slide(prs, "完整学习路径", [
        "第一部分：认知升级 — 理解为什么\"完课率\"不等于\"培训效果\"",
        "第二部分：数据采集 — 知道数据从哪来、怎么采、采什么",
        "第三部分：指标设计 — 把业务目标转化为可衡量的指标体系",
        "第四部分：ROI计算 — 用货币语言证明培训价值",
        "第五部分：闭环运营 — 让培训从一次性活动变成持续优化的系统",
        "第六部分：成果汇报 — 让管理层看见培训的价值",
        "综合演练 — 把以上所有产出整合成一套完整的证据链",
    ])

    summary_slide(prs, "三个核心认知", [
        "完课率不是培训效果，行为改变才是。培训效果的终极证明，不是学员\"学完了\"，而是学员\"用上了\"。",
        "ROI不是目的，是副产品。计算ROI是为了让培训管理者真正理解\"培训带来了什么价值\"。",
        "闭环比评估更重要。评估的目的是改进。如果评估完了没有后续动作，那评估就只是\"交差\"，而不是\"赋能\"。",
    ])

    section_header(prs, "三十天行动计划")

    table_slide(prs, "行动计划", ["时间", "行动", "产出"], [
        ["第1周", "选择一个正在做的培训项目，用F1判断当前评估层级", "评估层级诊断"],
        ["第2周", "为这个项目设计数据采集方案（F2）", "数据采集方案"],
        ["第3周", "设计这个项目的指标体系（F3）", "指标体系"],
        ["第4周", "向老板/同事做一次3分钟的效果汇报", "汇报材料"],
    ])

    # ===== Appendix =====
    section_divider(prs, "附录", "术语速查与工具索引")

    table_slide(prs, "术语速查", ["术语", "定义"], [
        ["柯氏四级评估", "全球最广泛使用的培训效果评估框架，包含反应层、学习层、行为层、成果层"],
        ["ROI", "投资回报率，衡量投入产出效率的通用商业语言"],
        ["WSDF", "学习数据采集规划表（What/Scale/Data Source/First Time）"],
        ["PDCA循环", "计划-实施-检查-改进的持续优化循环"],
        ["BEO框架", "行为指标设计的三个要素：Behavior、Evidence、Observation"],
    ])

    table_slide(prs, "工具表单速查", ["工具编号", "工具名称", "主要用途"], [
        ["F1", "培训效果评估层级判断卡", "快速判断培训需求属于哪一级评估"],
        ["F2", "学习数据采集规划表（WSDF）", "规划培训项目的数据采集方案"],
        ["F3", "培训指标体系设计表", "设计从反应层到成果层的完整指标体系"],
        ["F4", "柯氏四级评估设计卡", "为每一级评估设计具体评估工具"],
        ["F5", "培训ROI计算工作表", "计算培训项目的投资回报率"],
        ["F6", "行为改变追踪表", "追踪训后行为改变"],
        ["F7", "业务指标关联分析表", "关联培训与业务数据"],
        ["F8", "培训效果报告模板", "撰写效果报告"],
        ["F9", "培训闭环运营检查清单", "闭环运营检查"],
        ["F10", "培训数据仪表盘模板", "设计数据仪表盘"],
    ])

    # ===== Closing =====
    section_divider(prs, "课程结束", "致出发者")

    quote_slide(prs, "当你下次被老板问\"这个培训有什么效果\"的时候，希望你能自信地回答：\n\n\"让我用数据给你讲讲这个培训的价值。\"")

    content_slide(prs, "这不是一门教你\"怎么做培训\"的课，而是一门教你\"怎么证明培训有价值\"的课。", [
        "而你，已经掌握了这项能力。",
    ])

    # Save
    output_dir = "D:/CC/temp"
    os.makedirs(output_dir, exist_ok=True)
    output_path = "D:/CC/temp/培训闭环与效果证明_授课PPT.pptx"
    prs.save(output_path)
    print(f"PPT generated: {output_path}")
    print(f"Total pages: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    build_presentation()
