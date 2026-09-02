# -*- coding: utf-8 -*-
"""
《经营者讲党课》完整PPT生成脚本 - 140页版
红灰配色，浅底背景
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ============ 配色方案 ============
RED = RGBColor(0xC4, 0x12, 0x30)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_RED = RGBColor(0xFF, 0xF0, 0xF0)
ACCENT = RGBColor(0x8B, 0x45, 0x13)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)

OUT_DIR = 'D:/新课开发/党业融合/经营者讲党课/完整课程包/003-授课PPT'

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============ 辅助函数 ============
def bg_white(s):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = BG; b.line.fill.background()

def title_bar(s, title, h=Inches(1.1)):
    tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, h)
    tb.fill.solid(); tb.fill.fore_color.rgb = RED; tb.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    p = t.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE

def left_bar(s, color=GRAY):
    lb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(0.08), prs.slide_height - Inches(1.1))
    lb.fill.solid(); lb.fill.fore_color.rgb = color; lb.line.fill.background()

def content_text(s, items, x=Inches(0.5), y=Inches(1.4), w=Inches(9), fs=18):
    cb = s.shapes.add_textbox(x, y, w, Inches(5.5))
    ct = cb.text_frame; ct.word_wrap = True
    for i, item in enumerate(items):
        p = ct.paragraphs[0] if i == 0 else ct.add_paragraph()
        p.text = item; p.font.size = Pt(fs); p.font.color.rgb = TEXT
        p.space_before = Pt(12); p.space_after = Pt(6)

def page_num(s, n):
    pn = s.shapes.add_textbox(Inches(9.3), Inches(6.8), Inches(0.5), Inches(0.3))
    pp = pn.text_frame.paragraphs[0]; pp.text = str(n).zfill(2)
    pp.font.size = Pt(10); pp.font.color.rgb = ACCENT; pp.alignment = PP_ALIGN.RIGHT

# ============ 幻灯片类型函数 ============
def cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = RED; b.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = GRAY; bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
    p = t.text_frame.paragraphs[0]; p.text = '讲党课'
    p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    s2 = s.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    p2 = s2.text_frame.paragraphs[0]; p2.text = '业务干部的登台表达赋能工作坊'
    p2.font.size = Pt(32); p2.font.color.rgb = LIGHT; p2.alignment = PP_ALIGN.CENTER
    b2 = s.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    pb = b2.text_frame.paragraphs[0]; pb.text = '一天学会 当众讲好一个故事'
    pb.font.size = Pt(18); pb.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); pb.alignment = PP_ALIGN.CENTER

def section(n, title, sub=''):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    nb = s.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(3), Inches(2))
    np = nb.text_frame.paragraphs[0]; np.text = str(n)
    np.font.size = Pt(120); np.font.bold = True; np.font.color.rgb = RED
    tb = s.shapes.add_textbox(Inches(3.5), Inches(2.5), Inches(6), Inches(1.5))
    tp = tb.text_frame.paragraphs[0]; tp.text = title
    tp.font.size = Pt(40); tp.font.bold = True; tp.font.color.rgb = GRAY
    if sub:
        sb = s.shapes.add_textbox(Inches(3.5), Inches(4), Inches(6), Inches(0.8))
        sp = sb.text_frame.paragraphs[0]; sp.text = sub
        sp.font.size = Pt(20); sp.font.color.rgb = ACCENT

def content(title, items, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    left_bar(s)
    content_text(s, items)
    if num: page_num(s, num)

def two_col(title, lt, li, rt, ri, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    lb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.4), Inches(4.5), Inches(5.6))
    lb.fill.solid(); lb.fill.fore_color.rgb = LIGHT_RED; lb.line.color.rgb = RED
    ltb = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4), Inches(0.5))
    ltp = ltb.text_frame.paragraphs[0]; ltp.text = lt
    ltp.font.size = Pt(20); ltp.font.bold = True; ltp.font.color.rgb = RED
    lcb = s.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.1), Inches(4.6))
    lct = lcb.text_frame; lct.word_wrap = True
    for i, item in enumerate(li):
        p = lct.paragraphs[0] if i == 0 else lct.add_paragraph()
        p.text = '* ' + item; p.font.size = Pt(14); p.font.color.rgb = TEXT; p.space_before = Pt(10)
    rb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.4), Inches(4.5), Inches(5.6))
    rb.fill.solid(); rb.fill.fore_color.rgb = LIGHT; rb.line.color.rgb = GRAY
    rtb = s.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(4), Inches(0.5))
    rtp = rtb.text_frame.paragraphs[0]; rtp.text = rt
    rtp.font.size = Pt(20); rtp.font.bold = True; rtp.font.color.rgb = GRAY
    rcb = s.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(4.1), Inches(4.6))
    rct = rcb.text_frame; rct.word_wrap = True
    for i, item in enumerate(ri):
        p = rct.paragraphs[0] if i == 0 else rct.add_paragraph()
        p.text = '* ' + item; p.font.size = Pt(14); p.font.color.rgb = TEXT; p.space_before = Pt(10)
    if num: page_num(s, num)

def process(title, steps, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    n = len(steps)
    sw = Inches(2.2); gap = Inches(0.3); sx = Inches(0.5)
    for i, (n_, st, sd) in enumerate(steps):
        x = sx + i * (sw + gap)
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.5), Inches(0.8), Inches(0.8))
        c.fill.solid(); c.fill.fore_color.rgb = RED; c.line.fill.background()
        nb = s.shapes.add_textbox(x + Inches(0.7), Inches(1.55), Inches(0.8), Inches(0.7))
        np = nb.text_frame.paragraphs[0]; np.text = str(n_)
        np.font.size = Pt(24); np.font.bold = True; np.font.color.rgb = WHITE; np.alignment = PP_ALIGN.CENTER
        crd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), sw, Inches(4.5))
        crd.fill.solid(); crd.fill.fore_color.rgb = LIGHT_RED; crd.line.color.rgb = RED
        tb = s.shapes.add_textbox(x + Inches(0.1), Inches(2.7), sw - Inches(0.2), Inches(0.6))
        tpp = tb.text_frame.paragraphs[0]; tpp.text = st
        tpp.font.size = Pt(14); tpp.font.bold = True; tpp.font.color.rgb = RED; tpp.alignment = PP_ALIGN.CENTER
        db = s.shapes.add_textbox(x + Inches(0.1), Inches(3.3), sw - Inches(0.2), Inches(3.5))
        dp = db.text_frame.paragraphs[0]; dp.text = sd
        dp.font.size = Pt(11); dp.font.color.rgb = TEXT; dp.alignment = PP_ALIGN.CENTER
        if i < n - 1:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + sw + Inches(0.05), Inches(1.7), Inches(0.2), Inches(0.3))
            a.fill.solid(); a.fill.fore_color.rgb = GRAY; a.line.fill.background()
    if num: page_num(s, num)

def quote(q, a='', num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = LIGHT_RED; b.line.fill.background()
    qm = s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(2), Inches(2))
    qmp = qm.text_frame.paragraphs[0]; qmp.text = '"'
    qmp.font.size = Pt(120); qmp.font.color.rgb = RED
    qb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    qp = qb.text_frame.paragraphs[0]; qp.text = q
    qp.font.size = Pt(28); qp.font.color.rgb = GRAY; qp.alignment = PP_ALIGN.CENTER
    if a:
        ab = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
        ap = ab.text_frame.paragraphs[0]; ap.text = '-- ' + a
        ap.font.size = Pt(16); ap.font.color.rgb = ACCENT; ap.alignment = PP_ALIGN.RIGHT
    if num: page_num(s, num)

def cards(title, card_list, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    n = len(card_list)
    cw = Inches(2.8) if n <= 3 else Inches(2.2)
    gap = Inches(0.2); sx = Inches(0.5)
    for i, (ct, ci) in enumerate(card_list):
        x = sx + i * (cw + gap)
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), cw, Inches(5.5))
        c.fill.solid(); c.fill.fore_color.rgb = LIGHT_RED; c.line.color.rgb = RED
        ctb = s.shapes.add_textbox(x + Inches(0.1), Inches(1.7), cw - Inches(0.2), Inches(0.5))
        ctp = ctb.text_frame.paragraphs[0]; ctp.text = ct
        ctp.font.size = Pt(14); ctp.font.bold = True; ctp.font.color.rgb = RED; ctp.alignment = PP_ALIGN.CENTER
        ccb = s.shapes.add_textbox(x + Inches(0.1), Inches(2.3), cw - Inches(0.2), Inches(4.5))
        cct = ccb.text_frame; cct.word_wrap = True
        for j, item in enumerate(ci):
            p = cct.paragraphs[0] if j == 0 else cct.add_paragraph()
            p.text = '* ' + item; p.font.size = Pt(11); p.font.color.rgb = TEXT; p.space_before = Pt(8)
    if num: page_num(s, num)

def timeline(title, items, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = RED; line.line.fill.background()
    n = len(items)
    for i, (t, d) in enumerate(items):
        x = Inches(0.8 + i * 8.4 / max(n-1, 1))
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.15), Inches(3.35), Inches(0.3), Inches(0.3))
        dot.fill.solid(); dot.fill.fore_color.rgb = RED; dot.line.fill.background()
        tb = s.shapes.add_textbox(x - Inches(1), Inches(1.8), Inches(2), Inches(0.5))
        tp = tb.text_frame.paragraphs[0]; tp.text = t
        tp.font.size = Pt(14); tp.font.bold = True; tp.font.color.rgb = RED; tp.alignment = PP_ALIGN.CENTER
        db = s.shapes.add_textbox(x - Inches(1), Inches(4), Inches(2), Inches(2))
        dp = db.text_frame.paragraphs[0]; dp.text = d
        dp.font.size = Pt(11); dp.font.color.rgb = TEXT; dp.alignment = PP_ALIGN.CENTER
    if num: page_num(s, num)

def big_num(title, num_text, desc, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    nb = s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    np = nb.text_frame.paragraphs[0]; np.text = num_text
    np.font.size = Pt(100); np.font.bold = True; np.font.color.rgb = RED; np.alignment = PP_ALIGN.CENTER
    db = s.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    dp = db.text_frame.paragraphs[0]; dp.text = desc
    dp.font.size = Pt(20); dp.font.color.rgb = GRAY; dp.alignment = PP_ALIGN.CENTER
    if num: page_num(s, num)

def three_col(title, items, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    w = Inches(2.9); gap = Inches(0.2); sx = Inches(0.5)
    for i, (t, itms) in enumerate(items):
        x = sx + i * (w + gap)
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), w, Inches(5.6))
        c.fill.solid(); c.fill.fore_color.rgb = LIGHT_RED; c.line.color.rgb = RED
        tb = s.shapes.add_textbox(x + Inches(0.1), Inches(1.6), w - Inches(0.2), Inches(0.5))
        tp = tb.text_frame.paragraphs[0]; tp.text = t
        tp.font.size = Pt(16); tp.font.bold = True; tp.font.color.rgb = RED; tp.alignment = PP_ALIGN.CENTER
        cb = s.shapes.add_textbox(x + Inches(0.1), Inches(2.2), w - Inches(0.2), Inches(4.6))
        ct = cb.text_frame; ct.word_wrap = True
        for j, item in enumerate(itms):
            p = ct.paragraphs[0] if j == 0 else ct.add_paragraph()
            p.text = '* ' + item; p.font.size = Pt(12); p.font.color.rgb = TEXT; p.space_before = Pt(8)
    if num: page_num(s, num)

def four_col(title, items, num=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_white(s)
    title_bar(s, title)
    w = Inches(2.15); gap = Inches(0.13); sx = Inches(0.4)
    for i, (t, itms) in enumerate(items):
        x = sx + i * (w + gap)
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), w, Inches(5.6))
        c.fill.solid(); c.fill.fore_color.rgb = LIGHT_RED; c.line.color.rgb = RED
        tb = s.shapes.add_textbox(x + Inches(0.1), Inches(1.6), w - Inches(0.2), Inches(0.5))
        tp = tb.text_frame.paragraphs[0]; tp.text = t
        tp.font.size = Pt(13); tp.font.bold = True; tp.font.color.rgb = RED; tp.alignment = PP_ALIGN.CENTER
        cb = s.shapes.add_textbox(x + Inches(0.1), Inches(2.2), w - Inches(0.2), Inches(4.6))
        ct = cb.text_frame; ct.word_wrap = True
        for j, item in enumerate(itms):
            p = ct.paragraphs[0] if j == 0 else ct.add_paragraph()
            p.text = '* ' + item; p.font.size = Pt(10); p.font.color.rgb = TEXT; p.space_before = Pt(6)
    if num: page_num(s, num)

def ending():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = RED; b.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.15), prs.slide_width, Inches(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = GRAY; bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    p = t.text_frame.paragraphs[0]; p.text = '谢谢'
    p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    s2 = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    p2 = s2.text_frame.paragraphs[0]; p2.text = '祝大家讲出一堂好党课！'
    p2.font.size = Pt(28); p2.font.color.rgb = LIGHT; p2.alignment = PP_ALIGN.CENTER

# ============ 生成全部幻灯片 ============
slide_num = 1

# ==================== 封面 ====================
cover(); slide_num += 1

# ==================== 目录 ====================
content('课程目录', [
    '> 第一章：认知重塑 -- 业务干部为什么能讲好党课',
    '> 第二章：内容设计 -- 案例转化四步法',
    '> 第三章：登台表达 -- 分层递进练习体系',
    '> 第四章：实战演练 -- 全班展示与点评',
    '> 第五章：工具与转化 -- 带走可迁移的方法论',
    '> 第六章：课程总结 -- 核心收获回顾'
], slide_num); slide_num += 1

# ==================== 第一章：认知重塑 ====================
section(1, '认知重塑', '业务干部为什么能讲好党课'); slide_num += 1

content('课程核心定位', [
    '> 本课程不代写讲稿，不代做PPT',
    '> 只做一件事：教业务口干部把自己真实经历的管理故事，',
    '> 转化成一堂能讲、敢讲、学员爱听的党课，',
    '> 并且现场练到能自然讲出来为止。'
], slide_num); slide_num += 1

content('一句话定位', [
    '不代写讲稿，只做一件事：',
    '教业务口干部把自己真实经历的管理故事，',
    '转化成一堂能讲、敢讲、学员爱听的党课，',
    '并且现场练到能自然讲出来为止。'
], slide_num); slide_num += 1

two_col('业务干部讲党课的两大困境', '困境一：不会讲', [
    '完全不会讲党课',
    '只会念PPT或复述文件',
    '不知道如何开头',
    '不知道如何让学员愿意听',
    '讲得像业务汇报，学员昏昏欲睡'
], '困境二：不敢讲', [
    '怕在同事面前显得不专业',
    '怕被议论讲得差',
    '觉得自己不是党务出身',
    '表达焦虑比内容焦虑更严重',
    '深层恐惧：在同事面前丢脸'
], slide_num); slide_num += 1

big_num('深层恐惧的本质', '怕丢脸', '表达焦虑往往比"不知道讲什么"更影响实际表现', slide_num); slide_num += 1

content('目标学员画像', [
    '核心人群："一岗双责"制度下必须自己上台讲党课的经营管理者/业务口中高层',
    '通常本人并非党务出身',
    '平时讲业务汇报没问题，一到讲党课就发怵',
    '次要人群：即将被要求承担党课任务的后备中高层（提前赋能）',
    '行业覆盖：国企各业务条线，尤其是制造、金融、能源等业务性强的中高层群体'
], slide_num); slide_num += 1

content('学员典型心态', [
    '心态A：完全不会讲，只会念PPT或复述文件',
    '心态B：讲得像业务汇报，学员听得昏昏欲睡',
    '共同深层恐惧：怕在同事面前显得不专业、怕被议论讲得差',
    '这种表达焦虑往往比"不知道讲什么内容"更影响实际表现'
], slide_num); slide_num += 1

two_col('业务干部讲党课的独特优势', '优势一：有真实故事', [
    '带领团队克服困难的故事',
    '做出艰难抉择的故事',
    '团队成长的故事',
    '失败与挫折的故事',
    '这些经历本身就是最好的党课素材'
], '优势二：有说服力', [
    '因为真实，所以可信',
    '因为亲身经历，所以有细节',
    '因为有感情，所以能打动人',
    '这恰恰是党课最需要的',
    '代笔稿永远做不到这一点'
], slide_num); slide_num += 1

quote('你能讲好一门业务课，\\n就能讲好一堂党课。\\n区别只在于内容，不在于能力。', '', slide_num); slide_num += 1

content('本课程与其他课程的本质区别', [
    '不是"经验萃取转化"类课程（解决书记材料转化问题）',
    '不是"判断力/决策训练"类课程（政治判断力、廉政风险）',
    '而是"内容设计+登台表达"类课程',
    '解决的是"有故事，但不知道怎么讲、也不敢讲"',
    '效果要体现在"真的能站上台从容讲完一堂课"'
], slide_num); slide_num += 1

content('课程边界：五个不做', [
    '不代写讲稿、不代做PPT -- 学员必须自己参与转化',
    '不介入党课内容的具体思想政治内核 -- 学员自行对照组织部门要求',
    '不做党史知识讲授、不做党建理论培训 -- 假设学员已确定主题',
    '不涉及具体政策文件的解读争议 -- 引导对照本单位口径',
    'AI工具不作为课程卖点 -- 仅用于辅助整理讲稿逻辑框架草稿'
], slide_num); slide_num += 1

content('课程唯一聚焦点', [
    '用课程设计方法论，把业务干部真实的管理故事转化成有结构、有共鸣点的党课内容',
    '再用反复的登台练习，把这份内容内化成敢讲、能讲、讲得自然的能力',
    '离开时带着自己的党课框架',
    '核心承诺：一天学会 当众讲好一个故事'
], slide_num); slide_num += 1

# ==================== 第二章：内容设计 ====================
section(2, '内容设计', '案例转化四步法'); slide_num += 1

process('案例转化四步法', [
    (1, '故事盘点', '从真实经历中找出有戏剧张力的素材'),
    (2, '主题锚定', '找到故事与主题的真实连接点'),
    (3, '结构搭建', '用讲故事的方式重新组织内容'),
    (4, '语言转译', '把业务语言转化成讲台语言')
], slide_num); slide_num += 1

quote('党课内容本质上是"有立场的故事"\\n用真实故事让听众自己得出结论，\\n而不是先讲道理再配一个故事。', '', slide_num); slide_num += 1

# 第一步：故事盘点
section(2, '第一步：故事盘点', '从真实经历中找出有戏剧张力的素材'); slide_num += 1

content('什么是好故事素材', [
    '好的党课素材必须有"冲突"和"抉择"',
    '平铺直叙的"我们完成了任务"没有戏剧张力',
    '真正的张力来自：当时面临的真实两难选择是什么？',
    '引导学员回忆那个"纠结到睡不着觉"的时刻'
], slide_num); slide_num += 1

content('故事盘点的引导提问', [
    '"你带团队这些年，有没有一次真的很难抉择、纠结到睡不着觉的时刻？"',
    '"有没有一次你原本想放弃，但最后咬牙坚持下来的事？"',
    '"当时最让你睡不着觉的是什么？"',
    '"如果重新来过一次，你会做不同的选择吗？为什么？"'
], slide_num); slide_num += 1

content('学员常见反应与应对', [
    '第一反应：讲"团队多么努力、多么辛苦"（苦劳型素材）',
    '这类素材只有苦劳，没有冲突，没有抉择',
    '讲师需要引导其往"两难选择"方向深挖：',
    '"那当时有没有想过放弃？是什么让你坚持下来的？"'
], slide_num); slide_num += 1

content('故事盘点的关键标准', [
    '必须有"冲突"：理想与现实的差距、目标与阻力的对抗',
    '必须有"抉择"：在多个选项之间的权衡与取舍',
    '必须有"行动"：做出了什么选择，采取了什么行动',
    '平铺直叙三段论（背景-行动-结果）缺少张力，不适合党课'
], slide_num); slide_num += 1

quote('没有冲突的故事就像没有波折的电影，\\n观众从一开始就知道结局。', '', slide_num); slide_num += 1

content('故事盘点练习', [
    '每人写下3个让你印象深刻的带团队经历',
    '从中选出最有冲突感的一个',
    '描述当时面临的两难选择是什么',
    '同伴反馈：这个故事让你有感觉的地方在哪里？'
], slide_num); slide_num += 1

content('冲突的常见类型', [
    '类型一：资源冲突 -- 资源有限，必须取舍',
    '类型二：目标冲突 -- 多重目标无法同时兼顾',
    '类型三：人际冲突 -- 团队内部或与外部的利益矛盾',
    '类型四：自我冲突 -- 内心的挣扎与煎熬',
    '好的故事通常包含多种类型的冲突叠加'
], slide_num); slide_num += 1

content('故事盘点自查清单', [
    '这个故事中，最大的冲突是什么？',
    '当时我面临哪几种选择？',
    '最终我做了什么选择？为什么？',
    '这个选择带来的代价和收获分别是什么？',
    '有没有一个让我至今记忆犹新的细节或场景？'
], slide_num); slide_num += 1

# 第二步：主题锚定
section(2, '第二步：主题锚定', '找到故事与主题的真实连接点'); slide_num += 1

content('主题锚定的原则', [
    '不是生硬地把故事"扣"到主题上',
    '而是问：这个故事让你自己悟到了什么道理？',
    '再问：这个道理和今天要讲的主题，有没有真实的呼应？',
    '如果有真实的呼应，故事自然会和主题产生共鸣'
], slide_num); slide_num += 1

quote('如果学员选的故事和主题连接生硬，\\n宁可换一个故事，也不要强行嫁接。\\n牵强的连接是党课让人听着尴尬的主要原因。', '', slide_num); slide_num += 1

content('主题锚定的三个问题', [
    '问题1：这个故事让你自己悟到了什么道理？（不是主题让你想到什么）',
    '问题2：这个道理和今天要讲的主题之间，有什么真实的连接点？',
    '问题3：这个连接是故事"长出来"的，还是硬"扣"上去的？'
], slide_num); slide_num += 1

content('主题锚定练习', [
    '给出党课主题，让学员思考',
    '这个主题让你想起自己经历中的哪个故事？',
    '为什么这个故事让你想起这个主题？',
    '它们之间最真实的连接点是什么？'
], slide_num); slide_num += 1

content('常见错误：强行嫁接', [
    '为了扣主题，生硬地改编故事结局',
    '用主题反过来筛选和裁剪故事细节',
    '让故事失去原有的真实感和感染力',
    '结果：学员觉得假，听众觉得尴尬'
], slide_num); slide_num += 1

content('好主题与好故事的连接方式', [
    '方式一：主题是故事的"意义" -- 故事本身有力量，主题是对故事意义的提炼',
    '方式二：故事是主题的"证据" -- 主题是先有的，故事来印证它',
    '方式三：主题和故事是"对话"关系 -- 两者相互丰富，相互深化',
    '关键：连接必须是自然的、有机的'
], slide_num); slide_num += 1

# 第三步：结构搭建
section(2, '第三步：结构搭建', '用讲故事的方式重新组织内容'); slide_num += 1

content('五段式叙事结构', [
    '第一段：开场 -- 设置悬念/抛出真实困境（30秒内让听众进入情境）',
    '第二段：抉择 -- 还原当时的纠结和几种可能的选择',
    '第三段：转折 -- 揭示当时的真实选择和过程中的转折',
    '第四段：感悟 -- 引导听众自己感受到那个道理（不是替听众总结）',
    '第五段：收尾 -- 回扣主题，用一句简短有力的话收尾'
], slide_num); slide_num += 1

content('开场设计：30秒抓住注意力', [
    '不要：从"今天我要给大家讲一下XX主题"开始',
    '要：用一个问题、一个画面、或一个惊人事实抓住听众',
    '好的开场让听众立刻产生"然后呢？"的好奇',
    '开场设计是整个党课成功的一半'
], slide_num); slide_num += 1

content('开场白类型一：抛出问题', [
    '"大家有没有遇到过这种情况——"',
    '通过一个普遍性的问题让听众产生共鸣',
    '然后说："今天我想分享的，就是我们团队曾经经历过的..."',
    '问题要具体，不要太泛泛'
], slide_num); slide_num += 1

content('开场白类型二：描述画面', [
    '"那是去年冬天的一个晚上，凌晨两点，我收到了一条消息..."',
    '用一个具体的、带有感官细节的画面开场',
    '让听众仿佛能看到当时的场景',
    '画面要真实，不要编造'
], slide_num); slide_num += 1

content('开场白类型三：惊人事实', [
    '"我们团队差点就解散了。"',
    '用一个出人意料的事实抓住听众注意力',
    '然后慢慢展开故事',
    '事实必须是真的，不能为了效果编造'
], slide_num); slide_num += 1

content('抉择段：让听众代入', [
    '还原当时的几种可能选择',
    '不是简单罗列，而是让听众感受到每个选择的代价',
    '好的抉择段让听众不自觉地想"如果是我，我会怎么选？"',
    '这是让听众从旁观者变成参与者的关键'
], slide_num); slide_num += 1

content('转折段：揭示真实选择', [
    '你最终做了什么选择？为什么？',
    '过程中有什么意想不到的转折？',
    '有没有人反对？有没有什么转机？',
    '转折段让故事有起伏，有节奏感'
], slide_num); slide_num += 1

content('感悟段：让道理自己"长出来"', [
    '不是讲师替听众总结"所以我们要坚持党的领导"',
    '而是引导听众自己得出结论',
    '技巧：讲完故事后问"大家觉得这个故事告诉我们什么？"',
    '听众自己说出来的道理，比你直接讲出来更有说服力'
], slide_num); slide_num += 1

content('收尾：一句有力的话', [
    '不要长篇大论重复道理',
    '用一句简短有力的话回扣主题',
    '可以是一句发人深省的反问',
    '可以是一个意味深长的画面',
    '好的收尾让听众回味无穷'
], slide_num); slide_num += 1

content('收尾句范例', [
    '反问式："所以，当我们再次面对这样的选择时，我们会怎么做？"',
    '画面式："那个冬天的夜晚，永远刻在了我的心里。"',
    '行动式："从那天起，我们团队再也没有抱怨过困难。"',
    '启示式："这就是我想和大家分享的——"'
], slide_num); slide_num += 1

content('结构模板使用说明', [
    '结构模板是脚手架，不是要求机械套用',
    '有些故事天然是"总-分-总"，有些是"起承转合"',
    '关键是：让听众跟着你的节奏走',
    '最终引导听众自己得出结论'
], slide_num); slide_num += 1

content('结构灵活调整原则', [
    '如果故事本身有清晰的时间线，可以用时间顺序',
    '如果故事有强烈的对比（如前后变化），可以用对比结构',
    '如果故事有多个层次，可以层层递进',
    '核心原则：让故事清晰、有力、易懂'
], slide_num); slide_num += 1

# 第四步：语言转译
section(2, '第四步：语言转译', '把业务语言转化成讲台语言'); slide_num += 1

content('业务语言 vs 讲台语言', [
    '业务干部平时习惯用：数据、指标、专业术语',
    '这套语言直接搬上党课讲台会显得生硬冰冷',
    '需要引导其加入更多"画面感"和"情绪细节"',
    '比如不说"团队压力很大"，而是说具体是哪个深夜、谁说了哪句话'
], slide_num); slide_num += 1

content('语言转译的核心训练：追问法', [
    '反复追问：当时具体是什么场景？',
    '追问：谁说了什么？当时的气氛是怎样的？',
    '追问：你当时的第一反应是什么？',
    '追问：那个时刻你心里在想什么？',
    '逼学员讲出细节而非概括，让故事变得生动'
], slide_num); slide_num += 1

quote('一个细节胜过一千个形容词。\\n"那天晚上加班到凌晨两点"比"团队非常努力"更有力量。', '', slide_num); slide_num += 1

content('画面感训练', [
    '好的讲述应该有"画面感"',
    '听众仿佛能看到当时的场景',
    '而不是听到一堆抽象的形容词',
    '练习：把"我们克服了很大困难"转化成具体的场景描述'
], slide_num); slide_num += 1

content('画面感转化范例', [
    '原文："团队压力很大，大家都很辛苦"',
    '转化："连续一个月，每天晚上十一二点，办公室的灯都亮着。老张的爱人打电话来问怎么还不回家，他说快了快了，结果又到凌晨..."',
    '关键是：具体的时间、具体的场景、具体的人'
], slide_num); slide_num += 1

content('情绪细节训练', [
    '不仅要讲事实，还要讲感受',
    '"当时我真的很想放弃"比"我们遇到了困难"更打动人',
    '适当的脆弱感让讲述者更真实、更可亲',
    '但不要过度煽情，保持分寸'
], slide_num); slide_num += 1

content('情绪细节转化范例', [
    '原文："我们克服了困难，完成了任务"',
    '转化："那一刻，我真的想放弃了。但看到同事们还在咬牙坚持，我又把那句话咽了回去..."',
    '关键是：呈现真实的内心挣扎，而不是包装过的结果'
], slide_num); slide_num += 1

content('案例转化四步法总结', [
    '故事盘点：找到有冲突、有抉择的素材',
    '主题锚定：建立故事与主题的真实连接',
    '结构搭建：五段式叙事让故事有节奏感',
    '语言转译：增加画面感和情绪细节',
    '核心：让故事自己"长出"主题，而不是"扣上"主题'
], slide_num); slide_num += 1

# ==================== 第三章：登台表达 ====================
section(3, '登台表达训练', '分层递进的练习体系'); slide_num += 1

quote('业务干部真正的痛点往往不是不知道讲什么，\\n而是"写好了也不敢讲、讲的时候放不开"。\\n本课程必须投入与内容设计同等甚至更多的时间在实际登台练习上。', '', slide_num); slide_num += 1

content('为什么表达训练更重要', [
    '很多同类课程只做到"帮你把内容写好"就结束了',
    '但真正的痛点是"写好了也不敢讲、讲的时候放不开"',
    '内容设计做得再好，学员回去实战时依然会紧张卡壳',
    '达不到赋能效果'
], slide_num); slide_num += 1

two_col('登台心理阻力拆解与应对', '阻力一：怕讲错', [
    '深层恐惧：在同事面前丢脸',
    '应对：强调党课不是政治理论考试',
    '允许有个人风格和不完美',
    '讲师需要反复给出具体、真实的肯定反馈',
    '而不是泛泛表扬'
], '阻力二：放不开', [
    '觉得讲自己的故事很"肉麻"',
    '觉得在同事面前表达感情很奇怪',
    '应对：先从小范围（2-3人小组）练习开始',
    '逐步过渡到全班展示',
    '降低面对大场面的心理压力'
], slide_num); slide_num += 1

content('阻力三：照着稿子念', [
    '表现：脱稿就忘词，必须看着稿子才能讲',
    '原因：训练的是背诵，而不是讲故事',
    '应对：训练"记住故事的骨架和几个关键画面"',
    '用讲故事的自然节奏代替背诵的机械感',
    '核心：不是背稿子，是讲故事'
], slide_num); slide_num += 1

content('阻力四：语速和停顿失控', [
    '紧张时语速会不自觉加快',
    '或者相反——卡壳、停顿过多',
    '应对：通过录像回看让学员意识到自己的语速问题',
    '停顿训练：刻意在关键转折处停顿',
    '停顿是自信的表现，比语速快更有感染力'
], slide_num); slide_num += 1

content('阻力五：眼神躲闪', [
    '不敢看听众，低头看稿或看地面',
    '眼神只盯着一两个人，不敢扫视全场',
    '应对：刻意训练眼神交流，有意识地扫视不同区域',
    '眼神交流让人感觉自信、有诚意',
    '眼神躲闪会显得不自信、不诚实'
], slide_num); slide_num += 1

content('阻力六：手势僵硬', [
    '业务干部常见问题：带着汇报PPT的手势习惯',
    '比如：手心向下压、手指比划数字、背手等',
    '这些手势在党课场景下显得生硬、不亲切',
    '应对：放松手臂，自然手势，不要刻意控制'
], slide_num); slide_num += 1

# 分层递进练习体系
section(3, '分层递进练习体系', '从低压力到高压力，循序渐进'); slide_num += 1

process('四层练习体系', [
    (1, '小组讲述', '2-3人一组，低压力环境先把内容捋顺'),
    (2, '录制回看', '让学员看到自己实际讲述时的状态'),
    (3, '全班展示', '每位学员完整讲述，讲师现场点评'),
    (4, '模拟真实', '有讲台、有完整时长的正式试讲')
], slide_num); slide_num += 1

content('第一层：小组内讲述练习', [
    '目的：在低压力环境下先把内容捋顺',
    '习惯"说出来"这件事本身',
    '同伴反馈聚焦：我作为听众，哪个瞬间让我有感觉？',
    '哪里我走神了？',
    '不做内容对错评判'
], slide_num); slide_num += 1

content('第一层练习要点', [
    '2-3人一组，互相讲述自己的故事',
    '讲述者：完整讲述自己选定的故事',
    '听众：只反馈感受，不做评价',
    '反馈句式：这个瞬间让我有感觉，那个地方我走神了'
], slide_num); slide_num += 1

content('第一层练习流程', [
    '步骤1：学员分组（2-3人一组）',
    '步骤2：第一位学员讲述自己的故事（5分钟）',
    '步骤3：其他成员反馈感受',
    '步骤4：轮换，第二位学员讲述',
    '步骤5：讲师巡场，解答疑问'
], slide_num); slide_num += 1

content('第二层：录制与回看', [
    '目的：让学员看到自己实际讲述时的状态',
    '包括：语速、眼神、肢体语言',
    '很多人从未看过自己讲话的样子',
    '这个环节的冲击力很大',
    '回看后引导：你觉得自己讲得最自然的是哪一段？'
], slide_num); slide_num += 1

content('第二层练习要点', [
    '用手机录制自己的讲述过程（3-5分钟）',
    '回看录像，标记：语速最快的段落、最紧张的时刻',
    '讲师引导：看录像时不要关注内容对错，只关注呈现状态',
    '同伴互相看，提出具体反馈'
], slide_num); slide_num += 1

content('第三层：全班展示与点评', [
    '这是全天最重要的环节',
    '每位学员完整讲述一次（3-5分钟精简版）',
    '讲师现场给出具体、可操作的改进建议',
    '点评原则：先肯定具体的亮点细节',
    '再给1-2条最关键的改进建议',
    '不做面面俱到式的挑刺'
], slide_num); slide_num += 1

content('第三层练习流程', [
    '步骤1：学员依次上台（每人3-5分钟）',
    '步骤2：讲师在每人结束后即时点评',
    '步骤3：其他学员可以举手补充反馈',
    '步骤4：全部结束后，讲师做整体总结',
    '步骤5：每人获得一份个性化改进建议'
], slide_num); slide_num += 1

content('第四层：模拟真实场景（选配）', [
    '如果时间允许，安排一次模拟真实党课场景',
    '有讲台、有完整时长',
    '让学员提前适应真实场景的紧张感',
    '这是从练习到实战的关键过渡'
], slide_num); slide_num += 1

# 讲师点评
section(3, '讲师点评技巧', '保护信心，促进成长'); slide_num += 1

content('讲师点评的四大原则', [
    '原则一：具体肯定优先于笼统鼓励',
    '原则二：关键建议少而精优于面面俱到',
    '原则三：先说亮点，再说建议',
    '原则四：保护学员的信心和积极性'
], slide_num); slide_num += 1

content('点评话术范例：肯定部分', [
    '"你刚才讲到那个深夜的时候，我一下子就被带进去了"',
    '"你在转折处的停顿用得特别好，让我想知道接下来发生了什么"',
    '"你描述的那个画面非常生动，我仿佛看到了当时的场景"',
    '"整体结构很清晰，开场的那个问题问得很有吸引力"'
], slide_num); slide_num += 1

content('点评话术范例：建议部分', [
    '"建议你在手势上可以更放松一些，现在有点刻意"',
    '"语速在中间部分稍微快了一点，可以刻意放慢一些"',
    '"眼神可以多扫视一下全场，不要只盯着中间的位置"',
    '"收尾那句话可以更有力量一些，现在稍微有点平"'
], slide_num); slide_num += 1

content('避免的点评方式', [
    '"你的内容不错"（太笼统，没有具体点）',
    '"这里不对，那里也有问题"（打击信心）',
    '"你应该这样讲才对"（否定学员风格）',
    '在全班面前让某位学员显得难堪',
    '"你的故事不够感人"（否定学员的真实感受）'
], slide_num); slide_num += 1

content('讲师现场示范的必要性', [
    '建议讲师在开场讲解四步法时，必须用自己的真实故事完整示范一遍转化过程',
    '而不是只讲理论框架',
    '这样才能让学员相信"这套方法真的能把我的普通经历变成一个好故事"',
    '单纯讲方法论而没有示范，学员的信任度和参与度会明显打折扣'
], slide_num); slide_num += 1

content('讲师定位', [
    '需要同时具备两种能力：',
    '既懂课程设计和叙事结构方法论',
    '又有足够的现场亲和力和点评的分寸感',
    '因为这门课直接触及学员"怕丢面子"的敏感心理',
    '讲师引导方式如果偏生硬，会导致学员不敢真实练习'
], slide_num); slide_num += 1

content('点评分寸把控', [
    '"具体肯定优先于笼统鼓励、关键建议少而精优于面面俱到"',
    '尤其要避免在全班展示环节让某位学员因为点评方式不当而当众显得难堪',
    '一旦有一个人因为点评受挫而回避，会在班级中产生连锁的防御心理',
    '点评的目标是让学员更有信心去讲，而不是证明讲师更专业'
], slide_num); slide_num += 1

# 表达技巧
section(3, '表达技巧训练', '穿插在练习环节中逐一带过'); slide_num += 1

cards('四大表达技巧', [
    ('开场设计', ['前3句话的设计', '用问题或画面抓住注意力', '避免"今天我要给大家讲一下..."', '好的开场让听众产生好奇心']),
    ('停顿使用', ['关键转折处的停顿', '比语速快更有感染力', '让听众有时间消化', '停顿是自信的表现']),
    ('眼神交流', ['有意识地扫视不同区域', '不要死盯一个人或稿子', '与听众建立连接', '眼神躲闪会显得不自信']),
    ('手势自然', ['克制使用手势', '避免汇报PPT时的习惯', '讲故事时手势更自然放松', '手放在身体两侧比抱胸更好'])
], slide_num); slide_num += 1

content('开场技巧详解', [
    '禁忌开场："各位领导、各位同事，大家好，今天我要给大家讲一下XX主题"',
    '好的开场类型1：抛出问题——"大家有没有遇到过这种情况..."',
    '好的开场类型2：描述画面——"那是去年冬天的一个晚上..."',
    '好的开场类型3：惊人事实——"我们团队差点就..."'
], slide_num); slide_num += 1

content('开场技巧练习', [
    '每人准备3种不同类型的开场白',
    '在小组内试讲，比较哪种开场效果最好',
    '选择最能引发好奇心的开场方式',
    '记录下来，用于正式党课'
], slide_num); slide_num += 1

content('停顿技巧详解', [
    '为什么停顿重要：停顿让听众有时间消化，给讲述者时间调整',
    '关键停顿点1：开场后第一句话说完，稍作停顿',
    '关键停顿点2：转折处（"然后，我做了一个决定..."）',
    '关键停顿点3：高潮前（"那一刻，我想到了..."）',
    '关键停顿点4：收尾前（"所以，我想说..."）'
], slide_num); slide_num += 1

content('停顿技巧练习', [
    '在关键转折处刻意停顿2-3秒',
    '停顿时看着听众，不要躲闪',
    '停顿后继续讲，不要急着填补空白',
    '回看录像，感受停顿的效果'
], slide_num); slide_num += 1

content('眼神交流技巧详解', [
    '不要低头看稿——这会显得不自信',
    '不要死盯一个人——会让对方不自在',
    '正确做法：有意识地扫视不同区域（左侧、中间、右侧）',
    '每个区域停留2-3秒，让每位听众都觉得你在和他们说话'
], slide_num); slide_num += 1

content('眼神交流技巧练习', [
    '对着镜子练习，观察自己的眼神',
    '或者用手机录像，回看时注意眼神位置',
    '练习有意识地转移视线，而不是一直盯着一个方向',
    '如果怕看听众的眼睛，可以看他们的额头或鼻梁'
], slide_num); slide_num += 1

content('手势技巧详解', [
    '业务干部常见问题：带着汇报PPT的手势习惯',
    '比如：手心向下压、手指比划数字、背手等',
    '党课讲故事时，手势应该更自然放松',
    '建议：手势在腰部到肩部之间，不要过高或过低',
    '适度的手势可以增加感染力，过度会分散注意力'
], slide_num); slide_num += 1

content('手势技巧练习', [
    '放松手臂，自然下垂',
    '讲故事时让手势跟随情绪自然流动',
    '避免刻意比划或控制手势',
    '回看录像，检查手势是否自然'
], slide_num); slide_num += 1

# ==================== 第四章：实战演练 ====================
section(4, '实战演练', '全班展示与点评'); slide_num += 1

content('全班展示环节设计', [
    '目的：让每位学员都有登台机会并获得讲师点评',
    '人数控制：建议每场不超过16-18人',
    '时间：每位学员3-5分钟精简版讲述',
    '点评时间：每位学员1-2分钟'
], slide_num); slide_num += 1

content('展示环节流程', [
    '1. 学员依次上台（每人3-5分钟）',
    '2. 讲师在每人结束后即时点评',
    '3. 其他学员可以举手补充反馈',
    '4. 全部结束后，讲师做整体总结',
    '5. 每人获得一份个性化改进建议'
], slide_num); slide_num += 1

content('展示环节注意事项', [
    '如果学员人数过多，可以采用抽签展示+其余提交录像由讲师课后点评的方案',
    '如果学员在某个点上有共性问题，可以暂停做一次集体示范',
    '全程保持轻松、被支持的氛围',
    '让学员知道：在这个教室里犯错是安全的'
], slide_num); slide_num += 1

content('讲师整体总结要点', [
    '总结本次练习的整体亮点',
    '指出最常见的1-2个改进点',
    '强调最重要的1-2个技巧',
    '鼓励学员回去后继续练习'
], slide_num); slide_num += 1

# ==================== 第五章：工具与转化 ====================
section(5, '工具与转化', '带走可迁移的方法论'); slide_num += 1

content('核心工具包清单', [
    '《党课素材转化卡》：从原始故事到讲稿框架的转化记录表',
    '《党课基础叙事结构模板》：五段式叙事结构，附示范案例',
    '《登台自查清单》：开场设计、停顿使用、眼神交流等要点速查',
    '《同伴反馈卡》：规范反馈只聚焦"感受层面"而非内容对错评判'
], slide_num); slide_num += 1

three_col('工具一：党课素材转化卡', [
    ('用途', ['记录从原始故事到讲稿框架的转化全过程']),
    ('内容', ['原始故事描述', '故事中的冲突与抉择', '与主题的连接点', '初步框架']),
    ('价值', ['方便课后复盘', '可重复使用这套方法论', '形成自己的素材库'])
], slide_num); slide_num += 1

three_col('工具二：叙事结构模板', [
    ('五段式结构', ['开场悬念', '抉择还原', '转折揭示', '感悟引导', '回扣收尾']),
    ('使用说明', ['模板是脚手架不是枷锁', '根据故事特点灵活调整', '不要求机械套用']),
    ('示范案例', ['含一个完整示范案例', '展示如何将普通素材', '转化为有感染力的党课'])
], slide_num); slide_num += 1

content('工具三：登台自查清单', [
    '开场设计：我的前3句话能抓住注意力吗？',
    '眼神交流：我有意识地扫视不同区域吗？',
    '停顿使用：我在关键转折处有停顿吗？',
    '手势自然：我的手势是否自然放松？',
    '语速控制：我的语速是否有快有慢？'
], slide_num); slide_num += 1

content('工具四：同伴反馈卡', [
    '反馈原则：只反馈感受层面，不做内容对错评判',
    '反馈句式：这个瞬间让我有感觉',
    '反馈句式：那个地方我走神了',
    '反馈句式：如果我是听众，我想听到更多细节的是...',
    '避免：不要说你应该这样讲才对'
], slide_num); slide_num += 1

content('学员离场时带走', [
    '一份经过四步法转化的《党课讲稿框架》',
    '经过录像回看、现场点评三轮打磨的最终版本',
    '一份《登台自查清单》',
    '一套可迁移的党课内容设计方法论'
], slide_num); slide_num += 1

content('后续转化路径', [
    '回去后对着镜子练习讲3遍',
    '用自己的手机录下来，回看找问题',
    '正式讲党课前，在小范围先试讲一遍',
    '每次讲完后对照《登台自查清单》自检',
    '持续迭代，不断完善自己的党课'
], slide_num); slide_num += 1

# ==================== 第六章：课程总结 ====================
section(6, '课程总结', ''); slide_num += 1

content('课程核心收获', [
    '认知：业务干部能讲好党课，因为有真实故事',
    '方法：案例转化四步法（故事盘点-主题锚定-结构搭建-语言转译）',
    '训练：分层递进练习（小组讲述-录制回看-全班展示-模拟真实）',
    '工具：带走可迁移的党课设计方法论和工具包',
    '目标：敢讲、能讲、讲得自然'
], slide_num); slide_num += 1

content('课程特色总结', [
    '不代写讲稿，只做赋能',
    '现场练到能自然讲出来为止',
    '16-18人小班，确保充分练习',
    '国际版权课标准的内容质量',
    '一天学会 当众讲好一个故事'
], slide_num); slide_num += 1

timeline('课程时间分配（1天6-7课时）', [
    ('上午', '开场导入+四步法讲解+学员分组练习+结构搭建初稿'),
    ('下午', '语言转译打磨+分层递进练习+全班展示与点评')
], slide_num); slide_num += 1

content('上午时间安排', [
    '开场导入（0.5课时）：为什么业务干部能讲好党课',
    '四步法讲解+示范（1课时）：讲师现场示范完整转化过程',
    '学员分组练习（1.5课时）：故事盘点与主题锚定',
    '结构搭建初稿（0.5课时）：第三步结构搭建'
], slide_num); slide_num += 1

content('下午时间安排', [
    '语言转译打磨（0.5课时）：第四步语言转译',
    '分层练习第一层（0.5课时）：小组内讲述',
    '分层练习第二层（1课时）：录制与回看',
    '分层练习第三层（1.5课时）：全班展示与点评'
], slide_num); slide_num += 1

content('人数控制要求', [
    '建议每场不超过16-18人',
    '原因：全班展示环节需要保证每位学员都有登台机会',
    '人数过多会导致点评质量下降或时间严重超支',
    '如果超过18人，建议分批或采用抽签展示方案'
], slide_num); slide_num += 1

content('讲师带课要求', [
    '既懂课程设计和叙事结构方法论',
    '又有足够的现场亲和力和点评的分寸感',
    '必须现场示范用自己的故事走一遍四步法',
    '点评遵循"具体肯定优先于笼统鼓励"原则',
    '保护学员信心，避免当众难堪'
], slide_num); slide_num += 1

quote('离开时带着你的党课框架，\\n回去后对着镜子练3遍，\\n正式讲党课前先在小范围试讲一遍。\\n——这是课程结束后的行动承诺', '', slide_num); slide_num += 1

big_num('课程核心理念', '真实', '有真实故事，用真实感受，讲真实的话', slide_num); slide_num += 1

four_col('案例转化四步法详解', [
    ('故事盘点', ['找到有冲突的素材', '追问两难选择', '挖掘内心挣扎', '记录关键时刻']),
    ('主题锚定', ['找到故事感悟', '寻找主题连接', '判断连接真实', '自然长出主题']),
    ('结构搭建', ['开场设置悬念', '抉择还原代入', '转折揭示选择', '收尾有力回扣']),
    ('语言转译', ['追问场景细节', '增加画面感', '表达情绪感受', '让故事生动'])
], slide_num); slide_num += 1

four_col('分层递进练习详解', [
    ('第一层小组', ['2-3人一组', '低压力环境', '习惯说出来', '反馈聚焦感受']),
    ('第二层录像', ['手机录像', '回看状态', '看语速眼神', '冲击力大']),
    ('第三层全班', ['完整展示', '讲师点评', '同伴反馈', '即时改进']),
    ('第四层模拟', ['真实场景', '有讲台时长', '提前适应', '关键过渡'])
], slide_num); slide_num += 1

content('与其他课程的关系', [
    '与双带头人课：转化方向相反（业务经历 vs 党建动作）',
    '与政治判断力课、廉政风险课：完全不同类型（内容设计+表达 vs 情景决策）',
    '分属不同能力赋能类型，不建议打包销售',
    '应根据客户具体痛点精准匹配对应课程'
], slide_num); slide_num += 1

content('待后续补充事项', [
    '讲师现场示范用的真实故事素材',
    '《叙事结构模板》中的完整示范案例正文',
    '全班展示环节的抽签展示+录像点评备选方案',
    '不同层级（部门负责人vs分管领导）的难度版本设计'
], slide_num); slide_num += 1

content('课后行动承诺', [
    '承诺一：回去后对着镜子练习讲3遍',
    '承诺二：用自己的手机录下来，回看找问题',
    '承诺三：正式讲党课前在小范围先试讲一遍',
    '承诺四：每次讲完后对照清单自检',
    '持续迭代，讲出属于你自己的好党课'
], slide_num); slide_num += 1

# ==================== 结束页 ====================
ending(); slide_num += 1

# ============ 保存文件 ============
import os
os.makedirs(OUT_DIR, exist_ok=True)
out = os.path.join(OUT_DIR, '经营者讲党课-完整PPT.pptx')
prs.save(out)
print(f'Saved: {out}')
print(f'Total slides: {len(prs.slides)}')
