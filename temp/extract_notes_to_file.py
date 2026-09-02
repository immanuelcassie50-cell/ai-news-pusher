# -*- coding: utf-8 -*-
from pptx import Presentation
import codecs

pptx_path = r'D:\CC\temp\volvo_work.pptx'

prs = Presentation(pptx_path)

output = []
output.append("总幻灯片数: {}".format(len(prs.slides)))
output.append("="*80)

notes_count = 0
has_content_count = 0

for i, slide in enumerate(prs.slides, 1):
    has_notes = False
    notes_text = ""

    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        if text_frame and text_frame.text.strip():
            has_notes = True
            notes_count += 1
            notes_text = text_frame.text.strip()

            if len(notes_text) > 0:
                has_content_count += 1

    output.append("\n--- 第 {} 页 ---".format(i))
    output.append("有备注: {}".format('是' if has_notes else '否'))
    if has_notes:
        output.append("备注长度: {} 字符".format(len(notes_text)))
        output.append("备注内容:\n{}".format(notes_text))
    output.append("-"*80)

output.append("\n\n===== 统计结果 =====")
output.append("总页数: {}".format(len(prs.slides)))
output.append("有备注的页数: {}".format(notes_count))
output.append("备注有实质内容的页数: {}".format(has_content_count))
output.append("备注为空白或无备注的页数: {}".format(len(prs.slides) - notes_count))

with codecs.open(r'D:\CC\temp\notes_output.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(output))

print("Output written to D:\\CC\\temp\\notes_output.txt")
