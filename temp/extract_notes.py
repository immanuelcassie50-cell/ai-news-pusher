# -*- coding: utf-8 -*-
from pptx import Presentation

# 使用重命名后的文件
pptx_path = r'D:\CC\temp\volvo_work.pptx'

prs = Presentation(pptx_path)

print("总幻灯片数: {}".format(len(prs.slides)))
print("="*80)

notes_count = 0
has_content_count = 0

for i, slide in enumerate(prs.slides, 1):
    has_notes = False
    notes_text = ""

    # 提取备注
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        if text_frame and text_frame.text.strip():
            has_notes = True
            notes_count += 1
            notes_text = text_frame.text.strip()

            if len(notes_text) > 0:
                has_content_count += 1

    print("\n--- 第 {} 页 ---".format(i))
    print("有备注: {}".format('是' if has_notes else '否'))
    if has_notes:
        print("备注长度: {} 字符".format(len(notes_text)))
        print("备注内容:\n{}".format(notes_text))
    print("-"*80)

print("\n\n===== 统计结果 =====")
print("总页数: {}".format(len(prs.slides)))
print("有备注的页数: {}".format(notes_count))
print("备注有实质内容的页数: {}".format(has_content_count))
print("备注为空白或无备注的页数: {}".format(len(prs.slides) - notes_count))
