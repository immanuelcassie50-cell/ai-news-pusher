#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
expand_ppt.py - 创新领导力PPT扩展到120页+
Day 1 PM + Day 2 全部 + Closing
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ========== Theme ==========
THEME = {
    'primary':  '1A1A1A',
    'secondary':'4A4748',
    'accent':   'B81025',
    'light':    'EAE6E4',
    'bg':       'FFFFFF',
    'gray':     'B8B4B5',
}

FONT = 'Microsoft YaHei'

def RGB(hex_color):
    return RGBColor.from_string(hex_color)


class PPTBuilder:
    def __init__(self, prs):
        self.prs = prs
        self.page = 33

    def add_text(self, slide, text, x, y, w, h, size=12, bold=False, color=None,
                 align='left', valign='top', italic=False, fill_color=None, line_color=None):
        if isinstance(x, float): x = Inches(x)
        if isinstance(y, float): y = Inches(y)
        if isinstance(w, float): w = Inches(w)
        if isinstance(h, float): h = Inches(h)

        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)

        if valign == 'middle':
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif valign == 'bottom':
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM

        if fill_color:
            tb.fill.solid()
            tb.fill.fore_color.rgb = RGB(fill_color)
        else:
            tb.fill.background()

        if line_color:
            tb.line.color.rgb = RGB(line_color)
            tb.line.width = Pt(0.5)
        else:
            tb.line.fill.background()

        p = tf.paragraphs[0]
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)

        lines = text.split('\n') if isinstance(text, str) else [text]
        for i, line in enumerate(lines):
            if i == 0:
                paragraph = p
            else:
                paragraph = tf.add_paragraph()
                paragraph.alignment = align_map.get(align, PP_ALIGN.LEFT)

            run = paragraph.add_run()
            run.text = line
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            if color:
                run.font.color.rgb = RGB(color)

        return tb

    def add_rect(self, slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5):
        if isinstance(x, float): x = Inches(x)
        if isinstance(y, float): y = Inches(y)
        if isinstance(w, float): w = Inches(w)
        if isinstance(h, float): h = Inches(h)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGB(fill_color)
        else:
            shape.fill.background()

        if line_color:
            shape.line.color.rgb = RGB(line_color)
            shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()

        if shape.has_text_frame:
            shape.text_frame.text = ''
        return shape

    def add_oval(self, slide, x, y, w, h, fill_color=None, line_color=None):
        if isinstance(x, float): x = Inches(x)
        if isinstance(y, float): y = Inches(y)
        if isinstance(w, float): w = Inches(w)
        if isinstance(h, float): h = Inches(h)

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGB(fill_color)
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = RGB(line_color)
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()
        return shape

    def content_slide(self, title, section):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        self.add_rect(slide, 0, 0, 10, 5.625, fill_color=THEME['bg'])
        self.add_rect(slide, 0, 0, 10, 0.55, fill_color=THEME['primary'])
        self.add_rect(slide, 0, 0.55, 10, 0.04, fill_color=THEME['accent'])

        self.add_text(slide, title, 0.4, 0.05, 7.5, 0.45,
                     size=18, bold=True, color='FFFFFF', valign='middle')
        self.add_text(slide, section, 7.5, 0.05, 2.4, 0.45,
                     size=11, color='FFFFFF', align='right', valign='middle')

        self.page += 1
        page_str = str(self.page).zfill(2)
        self.add_oval(slide, 9.3, 5.1, 0.4, 0.4, fill_color=THEME['accent'])
        self.add_text(slide, page_str, 9.3, 5.1, 0.4, 0.4,
                     size=11, bold=True, color='FFFFFF', align='center', valign='middle')

        return slide

    def divider_slide(self, label, main, sub):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        self.add_rect(slide, 0, 0, 10, 5.625, fill_color=THEME['primary'])
        self.add_rect(slide, 0, 2.3, 10, 0.04, fill_color=THEME['accent'])

        self.add_text(slide, label, 0.5, 0.8, 9, 0.5,
                     size=14, bold=True, color=THEME['accent'])
        self.add_text(slide, main, 0.5, 1.5, 9, 1.0,
                     size=38, bold=True, color='FFFFFF')
        self.add_text(slide, sub, 0.5, 2.8, 9, 0.4,
                     size=14, color=THEME['gray'])

        self.page += 1
        page_str = str(self.page).zfill(2)
        self.add_oval(slide, 9.3, 5.1, 0.4, 0.4, fill_color=THEME['accent'])
        self.add_text(slide, page_str, 9.3, 5.1, 0.4, 0.4,
                     size=11, bold=True, color='FFFFFF', align='center', valign='middle')

        return slide


def build_day1_pm(b):
    """Day 1 PM: Parts 4-5 + summary + assignment"""

    b.divider_slide(
        "Day 1 · Afternoon",
        "从认知到行动：三大要素与挑战卡",
        "第四至五部分：客户洞察 · 交互涌现 · Day 1 总结"
    )

    # ---- 4.1: 进入第四部分 ----
    slide = b.content_slide("第四部分：要素一·客户洞察", "Day 1 · 第四部分")
    b.add_text(slide, "为什么客户洞察是创新第一要素？", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide,
        "· 90% 的创新失败不是因为技术不行\n"
        "· 而是因为没真正理解客户在做什么、为什么做、卡在哪里\n"
        "· 客户洞察不是问客户想要什么——他们也不知道\n"
        "· 是用一套穿透方法，进入客户没说出口的层面",
        0.4, 1.4, 9.2, 2.0, size=14, color=THEME['primary'])

    b.add_text(slide, "这一部分你要带走：", 0.4, 3.6, 9.2, 0.3,
               size=13, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "✓ 一套四层洞察框架（行为/目标/情感/身份）\n"
        "✓ 一个真实案例的四层穿透分析\n"
        "✓ 你自己的一次穿越四层练习\n"
        "✓ 一项明天可以启动的验证行动",
        0.4, 4.0, 9.2, 1.4, size=12, color=THEME['primary'])

    # ---- 4.2: 四层洞察框架 ----
    slide = b.content_slide("四层洞察框架", "Day 1 · 第四部分")
    b.add_text(slide, "从客户行为表层，到身份认同深层：四层穿透", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    layers = [
        ("第1层", "行为", "客户在做什么？频率/场景/工具/路径", THEME['accent']),
        ("第2层", "目标", "客户想达成什么？短期/中期/长期", "4A4748"),
        ("第3层", "情感", "客户的纠结/不满/骄傲/焦虑是什么？", "B81025"),
        ("第4层", "身份", "客户想成为什么样的自己？", "1A1A1A"),
    ]
    for i, (num, name, desc, color) in enumerate(layers):
        y = 1.5 + i * 0.85
        b.add_rect(slide, 0.4, y, 1.0, 0.75, fill_color=color)
        b.add_text(slide, num, 0.4, y, 1.0, 0.75, size=14, bold=True,
                   color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, name, 1.5, y + 0.05, 1.8, 0.65,
                   size=18, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 3.4, y + 0.05, 6.2, 0.65,
                   size=11, color=THEME['secondary'], valign='middle')

    # ---- 4.3: 第1层-行为 ----
    slide = b.content_slide("第1层 · 行为：看见客户真正做什么", "Day 1 · 第四部分")
    b.add_text(slide, '"客户说想要" ≠ "客户真的做"', 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "经典案例：\n"
        "· 客户说想要更快的马——亨利·福特\n"
        "· 客户说想要功能更多的手机——苹果\n"
        "· 客户说想要更便宜的网约车——滴滴",
        0.4, 1.4, 9.2, 1.2, size=12, color=THEME['primary'])

    b.add_text(slide, "如何观察行为层：", 0.4, 2.8, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "→ 不靠访谈，靠陪伴观察——陪客户去干他想干的事\n"
        "→ 不用问卷，用一日追踪——记录他这一天到底怎么过的\n"
        "→ 不问他怎么看，问他上一次做这事是什么时候\n"
        "→ 重点找抱怨与妥协——那里藏着真实需求",
        0.4, 3.2, 9.2, 1.8, size=12, color=THEME['primary'])

    # ---- 4.4: 第2层-目标 ----
    slide = b.content_slide("第2层 · 目标：客户想达成什么", "Day 1 · 第四部分")
    b.add_text(slide, "行为是表层动作，目标是动作背后的想要", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    b.add_text(slide, "三层目标结构：", 0.4, 1.4, 9.2, 0.3,
               size=13, bold=True, color=THEME['accent'])
    goals = [
        ("短期目标", "今天/本周想完成什么？", "例：今天要把方案做完"),
        ("中期目标", "未来3-6个月想达成什么？", "例：季度业绩达标"),
        ("长期目标", "3-5年想成为什么？", "例：成为业内认可的专家"),
    ]
    for i, (name, q, ex) in enumerate(goals):
        y = 1.85 + i * 0.95
        b.add_rect(slide, 0.4, y, 9.2, 0.85,
                   fill_color=THEME['light'] if i % 2 == 0 else THEME['bg'],
                   line_color=THEME['light'])
        b.add_text(slide, name, 0.6, y + 0.1, 1.8, 0.65,
                   size=13, bold=True, color=THEME['accent'], valign='middle')
        b.add_text(slide, q, 2.5, y + 0.1, 2.5, 0.65,
                   size=11, color=THEME['primary'], valign='middle')
        b.add_text(slide, ex, 5.0, y + 0.1, 4.5, 0.65,
                   size=11, color=THEME['secondary'], valign='middle', italic=True)

    # ---- 4.5: 第3层-情感 ----
    slide = b.content_slide("第3层 · 情感：客户的纠结、焦虑、骄傲", "Day 1 · 第四部分")
    b.add_text(slide, "客户买的不只是功能，是情感解决方案", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    b.add_text(slide, "情感层次结构：", 0.4, 1.4, 9.2, 0.3,
               size=13, bold=True, color=THEME['accent'])
    emotions = [
        ("痛点", "客户怕什么？担心什么？", "例：怕错过机会、怕被团队看不起"),
        ("痒点", "客户想要但不好意思说？", "例：渴望被认可、渴望轻松"),
        ("爽点", "客户用完会偷偷开心？", "例：用了你的产品后被同事夸"),
    ]
    for i, (name, q, ex) in enumerate(emotions):
        y = 1.85 + i * 0.95
        b.add_rect(slide, 0.4, y, 9.2, 0.85,
                   fill_color=THEME['light'] if i % 2 == 0 else THEME['bg'],
                   line_color=THEME['light'])
        b.add_text(slide, name, 0.6, y + 0.1, 1.8, 0.65,
                   size=13, bold=True, color=THEME['accent'], valign='middle')
        b.add_text(slide, q, 2.5, y + 0.1, 2.5, 0.65,
                   size=11, color=THEME['primary'], valign='middle')
        b.add_text(slide, ex, 5.0, y + 0.1, 4.5, 0.65,
                   size=11, color=THEME['secondary'], valign='middle', italic=True)

    # ---- 4.6: 第4层-身份 ----
    slide = b.content_slide("第4层 · 身份：客户想成为谁", "Day 1 · 第四部分")
    b.add_text(slide, "身份是动机的最终驱动", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "· 同一行为可以服务于不同身份\n"
        "· 客户买特斯拉 = 买环保先锋的身份\n"
        "· 客户用顶配笔记本 = 买专业设计师的身份\n"
        "· 客户买奢侈品 = 买经济自由的身份",
        0.4, 1.4, 9.2, 1.4, size=12, color=THEME['primary'])

    b.add_text(slide, "如何洞察身份层：", 0.4, 3.0, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "→ 问客户：你最不想被看成什么样的人？\n"
        "→ 问客户：你想让你的同事/朋友知道你在用什么吗？\n"
        "→ 观察客户如何使用产品——是为了炫耀，还是为了私下使用？\n"
        "→ 看客户朋友圈/微博/小红书——他在扮演什么角色？",
        0.4, 3.4, 9.2, 1.6, size=12, color=THEME['primary'])

    # ---- 4.7: 案例-星巴克第三空间 ----
    slide = b.content_slide("案例：星巴克第三空间", "Day 1 · 第四部分")
    b.add_text(slide, "一家咖啡店 vs 一种身份归属", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    layers = [
        ("行为", "买一杯咖啡，坐在店里待2小时"),
        ("目标", "找到工作之外的舒适空间"),
        ("情感", "逃离办公室的窒息感，假装我是自由的"),
        ("身份", "我是懂生活的城市中产"),
    ]
    for i, (name, desc) in enumerate(layers):
        y = 1.4 + i * 0.8
        b.add_rect(slide, 0.4, y, 4.5, 0.7,
                   fill_color=THEME['accent'] if i == 3 else (THEME['light'] if i % 2 == 0 else THEME['bg']),
                   line_color=THEME['light'])
        b.add_text(slide, name, 0.5, y, 1.0, 0.7,
                   size=14, bold=True,
                   color='FFFFFF' if i == 3 else THEME['accent'],
                   align='center', valign='middle')
        b.add_text(slide, desc, 1.6, y, 3.2, 0.7,
                   size=11, color=THEME['primary'], valign='middle')

    b.add_rect(slide, 5.2, 1.4, 4.4, 3.3, fill_color=THEME['primary'])
    b.add_text(slide, "→ 关键洞察", 5.4, 1.5, 4.0, 0.4,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "· 如果只看见第1层（行为）\n"
        "你会做：更便宜的咖啡\n\n"
        "· 如果看见第4层（身份）\n"
        "你会做：场景化的归属仪式\n"
        "——装潢、爵士乐、沙发、WiFi\n\n"
        "→ 后者比前者贵10倍，\n"
        "但毛利率高2倍。",
        5.4, 1.95, 4.0, 2.8, size=10, color='FFFFFF')

    # ---- 4.8: 练习引导 ----
    slide = b.content_slide("练习：穿越四层 · 15 分钟", "Day 1 · 第四部分")
    b.add_text(slide, "请小组选择一个真实客户/用户场景，做一次四层穿透", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    steps = [
        ("1", "5 分钟", "选场景", "选定一个你团队服务/想服务的客户场景"),
        ("2", "5 分钟", "穿透四层", "逐层讨论，写在便利贴上"),
        ("3", "5 分钟", "找穿透点", "找出从第1层到第4层最关键的转变点"),
    ]
    for i, (num, time, name, desc) in enumerate(steps):
        y = 1.4 + i * 1.2
        b.add_oval(slide, 0.4, y, 0.6, 0.6, fill_color=THEME['accent'])
        b.add_text(slide, num, 0.4, y, 0.6, 0.6,
                   size=20, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, time, 1.2, y + 0.05, 1.2, 0.5,
                   size=12, bold=True, color=THEME['accent'], valign='middle')
        b.add_text(slide, name, 2.4, y, 2.0, 0.6,
                   size=14, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 2.4, y + 0.3, 7.2, 0.5,
                   size=11, color=THEME['secondary'], valign='middle')

    b.add_text(slide, "⚠ 关键提醒：别停在第1层。穿透到第3层才刚开始。", 0.4, 4.95, 9.2, 0.4,
               size=12, bold=True, color=THEME['accent'], align='center')

    # ---- 4.9: 验证行动 ----
    slide = b.content_slide("产出：你的验证行动", "Day 1 · 第四部分")
    b.add_text(slide, "回到团队后，第一周可以启动的验证行动", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    actions = [
        "① 选 3 个最铁的老客户，预约一次陪伴观察，时长 2 小时",
        "② 找 5 个客户做 30 分钟的五层问（不是问卷）",
        "③ 看自己团队最近发布的客户洞察报告——是否停留在行为层",
        "④ 把你服务的客户画一幅身份地图——他到底想成为谁",
    ]
    for i, action in enumerate(actions):
        y = 1.5 + i * 0.8
        b.add_rect(slide, 0.4, y, 9.2, 0.7, fill_color=THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, action, 0.6, y, 9.0, 0.7,
                   size=13, color=THEME['primary'], valign='middle')

    # ---- Part 5 divider ----
    b.divider_slide(
        "Day 1 · 第五部分",
        "要素二：交互涌现",
        "团队不只是人在一起，是知识在流动"
    )

    # ---- 5.1: 进入第五部分 ----
    slide = b.content_slide("第五部分：要素二·交互涌现", "Day 1 · 第五部分")
    b.add_text(slide, "为什么团队在一起，却没涌现出新东西？", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide,
        "· 团队 ≠ 群体\n"
        "· 一群人坐在一起不等于在涌现\n"
        "· 涌现 = 不同想法碰撞 → 产生新的想法\n"
        "· 大部分团队的问题：知识只在一个小圈子里流动",
        0.4, 1.4, 9.2, 1.6, size=14, color=THEME['primary'])

    b.add_text(slide, "这一部分你要带走：", 0.4, 3.2, 9.2, 0.3,
               size=13, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "✓ 一套三堵墙诊断（部门墙/层级墙/信任墙）\n"
        "✓ 一套三种机制设计（轮岗/碰撞/外溢）\n"
        "✓ 你团队的知识流通审计练习\n"
        "✓ 一个你设计的碰撞机制——明天就能启动",
        0.4, 3.6, 9.2, 1.6, size=12, color=THEME['primary'])

    # ---- 5.2: 三堵墙-部门墙 ----
    slide = b.content_slide("三堵墙 1/3 · 部门墙", "Day 1 · 第五部分")
    b.add_text(slide, "墙 1：部门墙——信息在部门之间不流通", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "典型症状：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    symptoms = [
        "· 销售说：客户要的功能研发做不出来——研发不知道",
        "· 研发说：我们做了很多新功能没人用——销售没传达",
        "· 运营说：活动效果差——产品不知道哪里有问题",
        "· 部门周报制度存在，但只是形式",
    ]
    for i, s in enumerate(symptoms):
        b.add_text(slide, s, 0.4, 1.8 + i * 0.45, 9.2, 0.4,
                   size=11, color=THEME['primary'])

    b.add_text(slide, "为什么会有这堵墙？", 0.4, 3.8, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "→ 部门 KPI 设计相互独立（销售看业绩，研发看交付）\n"
        "→ 信息被视为部门资源——不愿轻易分享\n"
        "→ 跨部门沟通没有翻译者——专业语境不同\n"
        "→ 管理者担心暴露问题——影响部门评价",
        0.4, 4.2, 9.2, 1.4, size=11, color=THEME['primary'])

    # ---- 5.3: 三堵墙-层级墙 ----
    slide = b.content_slide("三堵墙 2/3 · 层级墙", "Day 1 · 第五部分")
    b.add_text(slide, "墙 2：层级墙——基层想法到不了决策层", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "真实案例：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "· 一线员工发现客户退货率异常，但报告流程要 5 层签字\n"
        "· 等问题到达高层时，已经过了 3 个月\n"
        "· 等高层知道时，已经不是问题，是危机",
        0.4, 1.8, 9.2, 1.2, size=12, color=THEME['primary'])

    b.add_text(slide, "为什么会形成层级墙？", 0.4, 3.2, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "→ 管理者习惯自上而下决策，不愿听到下面的反对\n"
        "→ 中层过滤信息——老板不爱听这个\n"
        "→ 基层表达方式粗糙——被翻译成 PPT 后失真\n"
        "→ 决策权集中在少数人——信息流动路径单向",
        0.4, 3.6, 9.2, 1.6, size=12, color=THEME['primary'])

    # ---- 5.4: 三堵墙-信任墙 ----
    slide = b.content_slide("三堵墙 3/3 · 信任墙", "Day 1 · 第五部分")
    b.add_text(slide, "墙 3：信任墙——说了真话会付出代价", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "故事：赵建设的沉默", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "那个三个月前曾说提案平台没人认真看的骨干员工\n"
        "在赵建设再次问我们哪里出了问题时，沉默了。\n\n"
        "他为什么沉默？\n"
        "· 他之前说过，但没人真的改\n"
        "· 他担心被认为不配合\n"
        "· 他觉得反正也没用\n"
        "· 他已经学会了保护自己——不说真话",
        0.4, 1.8, 9.2, 3.0, size=12, color=THEME['primary'])

    b.add_text(slide, "💡 信任墙是最危险的墙——它让人不再相信说出来有用", 0.4, 5.0, 9.2, 0.4,
               size=12, bold=True, color=THEME['accent'], align='center')

    # ---- 5.5: 三种机制-轮岗 ----
    slide = b.content_slide("三种机制 1/3 · 轮岗：让知识跨界流动", "Day 1 · 第五部分")
    b.add_text(slide, "机制 1：让员工在不同岗位/项目/客户之间流动", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "具体做法：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    practices = [
        ("短期轮岗", "2-4 周", "参与其他部门的关键项目，结束后做分享"),
        ("中期轮岗", "3-6 个月", "借调到其他岗位/项目组，回流后带回视角"),
        ("反向轮岗", "1-2 周", "高管到一线体验——不是视察，是真干"),
        ("客户轮岗", "持续", "客服、产品、研发、销售轮流跟客户接触"),
    ]
    for i, (name, time, desc) in enumerate(practices):
        y = 1.85 + i * 0.75
        b.add_rect(slide, 0.4, y, 9.2, 0.65,
                   fill_color=THEME['light'] if i % 2 == 0 else THEME['bg'],
                   line_color=THEME['light'])
        b.add_text(slide, name, 0.6, y, 1.6, 0.65,
                   size=12, bold=True, color=THEME['accent'], valign='middle')
        b.add_text(slide, time, 2.3, y, 1.0, 0.65,
                   size=10, color=THEME['secondary'], valign='middle')
        b.add_text(slide, desc, 3.4, y, 6.0, 0.65,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 5.6: 三种机制-碰撞 ----
    slide = b.content_slide("三种机制 2/3 · 碰撞：让想法高频碰撞", "Day 1 · 第五部分")
    b.add_text(slide, "机制 2：让不同背景的人高频次产生连接", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    collisions = [
        ("跨界午餐会", "每月 1 次", "8 人随机抽签，包括不同部门/职级"),
        ("案例分享夜", "每月 2 次", "员工分享自己最近学到的/撞见的"),
        ("挑战者机制", "持续", "每个项目指定一个挑战者角色——专门提问"),
        ("客户面对面", "每周", "团队成员轮流陪同拜访客户，回来汇报"),
        ("失败复盘会", "每月 1 次", "专门讲失败——不找责任人，找原因"),
    ]
    for i, (name, time, desc) in enumerate(collisions):
        y = 1.4 + i * 0.7
        b.add_rect(slide, 0.4, y, 9.2, 0.6,
                   fill_color=THEME['accent'] if i == 0 else (THEME['light'] if i % 2 == 0 else THEME['bg']),
                   line_color=THEME['light'])
        b.add_text(slide, name, 0.6, y, 2.2, 0.6,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, time, 2.9, y, 1.4, 0.6,
                   size=10, color=THEME['secondary'], valign='middle')
        b.add_text(slide, desc, 4.4, y, 5.0, 0.6,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 5.7: 三种机制-外溢 ----
    slide = b.content_slide("三种机制 3/3 · 外溢：让知识主动外溢", "Day 1 · 第五部分")
    b.add_text(slide, "机制 3：让团队与外部世界高频碰撞", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "外溢路径：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    paths = [
        "→ 跨行业参观：每季度组织 1 次，去不同行业/不同规模的公司看",
        "→ 顾问机制：请外部专家作为月度顾问，提供不同视角",
        "→ 行业会议：员工轮流参加，回来必须做团队分享",
        "→ 客户开放日：让客户到公司坐一天，看团队怎么工作",
        "→ 公开写作：鼓励员工写公众号/知乎/小红书——倒逼思考",
    ]
    for i, p in enumerate(paths):
        y = 1.85 + i * 0.6
        b.add_text(slide, p, 0.6, y, 9.0, 0.5,
                   size=12, color=THEME['primary'], valign='middle')

    # ---- 5.8: 知识流通审计 ----
    slide = b.content_slide("练习：知识流通审计", "Day 1 · 第五部分")
    b.add_text(slide, "用一张图，看清你的团队知识在哪流动、在哪卡住", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    audit = [
        ("① 盘点关键节点", "你的团队里，关键知识掌握在谁手里？（不超过 8 个人）"),
        ("② 画流通图", "在团队成员之间画箭头：谁向谁传递什么知识？"),
        ("③ 找断点", "箭头稀疏/单向/缺失的地方——就是涌现的卡点"),
        ("④ 诊断墙类型", "断点背后是部门墙、层级墙还是信任墙？"),
        ("⑤ 设计机制", "针对断点选一种或多种机制，让知识开始流动"),
    ]
    for i, (step, desc) in enumerate(audit):
        y = 1.4 + i * 0.7
        b.add_rect(slide, 0.4, y, 9.2, 0.6,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, step, 0.6, y, 3.0, 0.6,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 3.8, y, 5.8, 0.6,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 5.9: 创新挑战卡作业 ----
    slide = b.content_slide("今晚作业：填一份创新挑战卡", "Day 1 · 第五部分")
    b.add_text(slide, "把模糊的痛变成具体的卡", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "卡片的 6 个字段：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    fields = [
        ("挑战是什么", "用一句话描述你团队里最想解决的那个创新相关挑战"),
        ("为什么是挑战", "它卡在哪里？是行为问题/认知问题/结构问题？"),
        ("已经尝试过什么", "过去 6 个月团队尝试过的方法、结果"),
        ("卡在哪里了", "做了很多但没结果——卡点是？"),
        ("最想突破的", "如果只能突破一点，那是什么？"),
        ("30 天内能做什么", "明天开始的具体动作"),
    ]
    for i, (name, desc) in enumerate(fields):
        y = 1.85 + i * 0.55
        b.add_rect(slide, 0.4, y, 0.6, 0.45, fill_color=THEME['accent'])
        b.add_text(slide, str(i+1), 0.4, y, 0.6, 0.45,
                   size=14, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, name, 1.1, y, 2.4, 0.45,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 3.5, y, 6.0, 0.45,
                   size=10, color=THEME['secondary'], valign='middle')

    # ---- Day 1 Summary ----
    slide = b.content_slide("Day 1 · 我们走过的路", "Day 1 · 总结")
    b.add_text(slide, "从看见真相到开始行动", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    journey = [
        ("开场", "赵建设的沉默", "找到真问题"),
        ("Part 1", "八个真相测试", "打破认知盲区"),
        ("Part 2", "五个关键因素", "诊断团队状态"),
        ("Part 3", "创新型 vs 运营管理", "识别自身行为"),
        ("Part 4", "要素一·客户洞察", "穿透到第4层"),
        ("Part 5", "要素二·交互涌现", "拆三堵墙，建三机制"),
        ("作业", "创新挑战卡", "今晚 15-20 分钟"),
    ]
    for i, (part, name, desc) in enumerate(journey):
        y = 1.4 + i * 0.5
        b.add_rect(slide, 0.4, y, 1.5, 0.4, fill_color=THEME['accent'] if i == 6 else THEME['light'])
        b.add_text(slide, part, 0.4, y, 1.5, 0.4,
                   size=10, bold=True,
                   color='FFFFFF' if i == 6 else THEME['accent'],
                   align='center', valign='middle')
        b.add_text(slide, name, 2.0, y, 3.5, 0.4,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 5.6, y, 4.0, 0.4,
                   size=11, color=THEME['secondary'], valign='middle')

    # ---- Day 1 关键认知 ----
    slide = b.content_slide("Day 1 关键认知 · 5 句话", "Day 1 · 总结")
    b.add_text(slide, "如果今天只能带走 5 句话", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    quotes = [
        "1. 团队不创新的原因，大部分不在团队——在管理者身上",
        "2. 你以为给了空间，可能恰恰是你用了 9 种方式压制了空间",
        "3. 客户洞察不是问客户想要什么——是看见他没说出口的",
        "4. 团队的真正问题不是没有人，是知识在墙里不流动",
        "5. 改变的起点不是大动作，是第一个具体的卡——把它写下来",
    ]
    for i, q in enumerate(quotes):
        y = 1.4 + i * 0.7
        b.add_rect(slide, 0.4, y, 9.2, 0.6, fill_color=THEME['light'])
        b.add_text(slide, q, 0.6, y, 9.0, 0.6,
                   size=12, color=THEME['primary'], valign='middle')


def build_day2(b):
    """Day 2 AM + PM"""
    b.divider_slide(
        "Day 2 · 从认知到行动",
        "看见 × 行动 × 承诺",
        "第六至十二部分：敏捷迭代 · 角色转换 · 完整分析 · 情景模拟 · 承诺"
    )

    # ---- Day 2 Opening ----
    slide = b.content_slide("Day 2 开场：昨晚发生了什么", "Day 2 · 开场")
    b.add_text(slide, "两人对话 · 5 分钟", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "请找到一位昨天没怎么交流的学员，两两一组：", 0.4, 1.4, 9.2, 0.4,
               size=14, color=THEME['primary'])

    questions = [
        ("第 1 问", "3 分钟", "你昨晚填创新挑战卡时，最让你停顿的那个瞬间是什么？"),
        ("第 2 问", "2 分钟", "你写完挑战卡后，有什么跟之前想法不一样的发现？"),
    ]
    for i, (q, time, desc) in enumerate(questions):
        y = 2.0 + i * 1.0
        b.add_rect(slide, 0.4, y, 9.2, 0.9,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, q, 0.6, y, 1.5, 0.9,
                   size=14, bold=True,
                   color='FFFFFF' if i == 0 else THEME['accent'],
                   valign='middle')
        b.add_text(slide, time, 2.2, y, 1.0, 0.9,
                   size=12, color=THEME['secondary'], valign='middle')
        b.add_text(slide, desc, 3.3, y, 6.0, 0.9,
                   size=12, color=THEME['primary'], valign='middle')

    b.add_text(slide, "⚠ 规则：只听，不评判，不建议", 0.4, 4.2, 9.2, 0.4,
               size=13, bold=True, color=THEME['accent'], align='center')
    b.add_text(slide, "听完只需说：谢谢你告诉我这些", 0.4, 4.65, 9.2, 0.4,
               size=12, color=THEME['primary'], align='center')

    # ---- Part 6 divider ----
    b.divider_slide(
        "Day 2 · 第六部分",
        "要素三：敏捷迭代",
        "客户洞察×交互涌现×敏捷迭代 = 创新型团队三角"
    )

    # ---- 6.1: 进入第六部分 ----
    slide = b.content_slide("第六部分：要素三·敏捷迭代", "Day 2 · 第六部分")
    b.add_text(slide, "为什么大多数创新项目做着做着就死了？", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "真相：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "· 大多数团队没有迭代——只有计划→执行\n"
        "· 计划完美，执行变形，结果偏差，回不到原点\n"
        "· 创新不是想清楚再做，是做着做着才想清楚",
        0.4, 1.85, 9.2, 1.4, size=12, color=THEME['primary'])

    b.add_text(slide, "这一部分你要带走：", 0.4, 3.4, 9.2, 0.3,
               size=13, bold=True, color=THEME['accent'])
    b.add_text(slide,
        "✓ 四个关键动作（假设/最小验证/预定义改变/显性学习）\n"
        "✓ 管理者在迭代中的五种角色转换\n"
        "✓ 把你昨晚的挑战卡改写成一份最小可学习实验",
        0.4, 3.85, 9.2, 1.5, size=12, color=THEME['primary'])

    # ---- 6.2: 敏捷迭代 vs 传统计划 ----
    slide = b.content_slide("敏捷迭代 vs 传统计划", "Day 2 · 第六部分")
    b.add_text(slide, "两种创新范式的根本差异", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    compare = [
        ("目标", "完美方案 → 一次执行到位", "最小可行方案 → 学到东西"),
        ("起点", "已知一切，可推演", "承认未知，需要探索"),
        ("节奏", "周/月 → 一次大检查", "天/周 → 多次小检查"),
        ("失败", "尽量避免", "必须发生，是学习来源"),
        ("决策", "基于预测", "基于观察"),
        ("资源", "前期大投入", "前期小投入，按学习追加"),
    ]
    for i, (aspect, a, b_text) in enumerate(compare):
        y = 1.4 + i * 0.55
        b.add_rect(slide, 0.4, y, 1.5, 0.5, fill_color=THEME['accent'])
        b.add_text(slide, aspect, 0.4, y, 1.5, 0.5,
                   size=11, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_rect(slide, 2.0, y, 3.7, 0.5, fill_color=THEME['light'])
        b.add_text(slide, a, 2.1, y, 3.6, 0.5,
                   size=10, color=THEME['primary'], valign='middle')
        b.add_rect(slide, 5.8, y, 3.7, 0.5, fill_color=THEME['primary'])
        b.add_text(slide, b_text, 5.9, y, 3.6, 0.5,
                   size=10, color='FFFFFF', valign='middle')

    b.add_text(slide, "→ 不是左和右的对错，是什么场景用哪种", 0.4, 4.85, 9.2, 0.4,
               size=12, bold=True, color=THEME['accent'], align='center')

    # ---- 6.3: 四步循环 1-假设 ----
    slide = b.content_slide("四步循环 1/4 · 提出假设", "Day 2 · 第六部分")
    b.add_text(slide, "第 1 步：把模糊的想做什么变成可被证伪的假设", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "模糊 vs 假设", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    compare = [
        ("× 模糊", "我们要做一个让客户喜欢的产品", "无法被验证"),
        ("× 模糊", "我们想做一个新功能", "无法被验证"),
        ("✓ 假设", "假设：客户愿意为节省每周 3 小时付费 30 元/月", "可被验证"),
        ("✓ 假设", "假设：30% 的老客户会试用我们设计的会员体系", "可被验证"),
    ]
    for i, (tag, text, comment) in enumerate(compare):
        y = 1.85 + i * 0.6
        is_good = '✓' in tag
        color = THEME['accent'] if is_good else THEME['secondary']
        b.add_rect(slide, 0.4, y, 1.0, 0.5, fill_color=color)
        b.add_text(slide, tag, 0.4, y, 1.0, 0.5,
                   size=14, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, text, 1.5, y, 6.5, 0.5,
                   size=11, color=THEME['primary'], valign='middle')
        b.add_text(slide, comment, 8.1, y, 1.5, 0.5,
                   size=10, italic=True, color=color, valign='middle')

    # ---- 6.4: 四步循环 2-最小验证 ----
    slide = b.content_slide("四步循环 2/4 · 最小验证", "Day 2 · 第六部分")
    b.add_text(slide, "第 2 步：用最小成本、最快速度验证假设", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, '"最小"≠"简单"，而是"刚好能验证假设"', 0.4, 1.4, 9.2, 0.4,
               size=13, color=THEME['secondary'], italic=True)

    b.add_text(slide, "验证工具箱：", 0.4, 1.9, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    tools = [
        ("访谈", "5-10 个目标客户的深度访谈，验证需求是否真存在"),
        ("原型", "做一个低保真原型，看客户是否能理解/接受"),
        ("落地页", "做一个落地页，看转化率是否支持假设"),
        ("小流量测试", "用 5% 的用户做 A/B 测试"),
        ("预售", "在没做产品前，看是否有人愿意预付"),
    ]
    for i, (tool, desc) in enumerate(tools):
        y = 2.3 + i * 0.55
        b.add_rect(slide, 0.4, y, 9.2, 0.5,
                   fill_color=THEME['light'] if i % 2 == 0 else THEME['bg'],
                   line_color=THEME['light'])
        b.add_text(slide, tool, 0.6, y, 1.5, 0.5,
                   size=12, bold=True, color=THEME['accent'], valign='middle')
        b.add_text(slide, desc, 2.2, y, 7.4, 0.5,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 6.5: 四步循环 3-预定义改变条件 ----
    slide = b.content_slide("四步循环 3/4 · 预定义改变条件", "Day 2 · 第六部分")
    b.add_text(slide, "第 3 步：在开始之前，先定义什么情况下我们要改变方向", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    b.add_text(slide, '"预定义改变条件"避免"沉没成本谬误"', 0.4, 1.4, 9.2, 0.4,
               size=13, color=THEME['secondary'], italic=True)

    b.add_text(slide, "示范框架：", 0.4, 1.9, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    framework = [
        ("绿灯", "3 周内访谈 8 个客户，6 个表达愿意付费，进入下一阶段"),
        ("黄灯", "3 周内访谈 8 个客户，4-5 个犹豫，需要调整方案后再试"),
        ("红灯", "3 周内访谈 8 个客户，少于 3 个愿意，立刻转向或停止"),
    ]
    for i, (signal, desc) in enumerate(framework):
        y = 2.3 + i * 0.85
        color = ['0EA34A', 'F5A623', 'B81025'][i]
        b.add_rect(slide, 0.4, y, 1.5, 0.75, fill_color=color)
        b.add_text(slide, signal, 0.4, y, 1.5, 0.75,
                   size=14, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, desc, 2.0, y, 7.6, 0.75,
                   size=12, color=THEME['primary'], valign='middle')

    # ---- 6.6: 四步循环 4-显性学习 ----
    slide = b.content_slide("四步循环 4/4 · 显性学习", "Day 2 · 第六部分")
    b.add_text(slide, "第 4 步：把我们学到了什么显性化，进入下一轮", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    b.add_text(slide, "显性学习的 3 个关键：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    keys = [
        ("假设被验证/证伪", "我们原以为 X，实际是 Y"),
        ("下轮的具体动作", "基于学到，下一轮我们做 Z"),
        ("新假设/新问题", "新的问题是 W，下一轮验证"),
    ]
    for i, (k, v) in enumerate(keys):
        y = 1.85 + i * 0.95
        b.add_rect(slide, 0.4, y, 9.2, 0.85,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, k, 0.6, y + 0.1, 2.8, 0.65,
                   size=13, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, v, 3.5, y + 0.1, 6.0, 0.65,
                   size=12, color=THEME['primary'], valign='middle')

    b.add_text(slide, "⚠ 不做显性学习 = 之前的投入浪费", 0.4, 4.7, 9.2, 0.4,
               size=12, bold=True, color=THEME['accent'], align='center')

    # ---- 6.7: 五种角色 ----
    slide = b.content_slide("管理者在迭代中的五种角色", "Day 2 · 第六部分")
    b.add_text(slide, "迭代中管理者不是旁观者——是五种角色的动态切换", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    roles = [
        ("战略设定者", "定义问题边界，确保方向对齐"),
        ("资源守护者", "保护团队不被外部打断，集中精力迭代"),
        ("挑战者", "持续追问假设，问我们凭什么相信"),
        ("学习教练", "每次迭代后强制显性化学习"),
        ("边界打破者", "必要时打破部门墙/层级墙，让团队有空间"),
    ]
    for i, (role, desc) in enumerate(roles):
        y = 1.4 + i * 0.7
        b.add_rect(slide, 0.4, y, 2.8, 0.6, fill_color=THEME['accent'])
        b.add_text(slide, role, 0.4, y, 2.8, 0.6,
                   size=12, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_rect(slide, 3.3, y, 6.3, 0.6, fill_color=THEME['light'])
        b.add_text(slide, desc, 3.5, y, 6.0, 0.6,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 6.8: 最小可学习实验 ----
    slide = b.content_slide("练习：把挑战卡改写成最小可学习实验", "Day 2 · 第六部分")
    b.add_text(slide, "从想做，到下次能学到什么", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    template = [
        ("原挑战", "我团队的老客户流失率上升"),
        ("核心假设", "假设：流失是因为我们的服务节奏没跟上客户需求变化"),
        ("最小验证", "访谈 8 个流失客户，问：是什么让你决定不续约"),
        ("改变条件", "绿灯=6/8 客户说是；黄灯=4-5 个；红灯=少于 4 个"),
        ("预计学习", "我们将学到：流失的真实原因是 X，不是我们以为的 Y"),
        ("下一轮", "基于学习，调整服务节奏或改变定位"),
    ]
    for i, (k, v) in enumerate(template):
        y = 1.4 + i * 0.55
        b.add_rect(slide, 0.4, y, 1.8, 0.5, fill_color=THEME['accent'] if i == 0 else THEME['light'])
        b.add_text(slide, k, 0.4, y, 1.8, 0.5,
                   size=11, bold=True,
                   color='FFFFFF' if i == 0 else THEME['primary'],
                   align='center', valign='middle')
        b.add_text(slide, v, 2.3, y, 7.3, 0.5,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- Part 7 divider ----
    b.divider_slide(
        "Day 2 · 第七部分",
        "管理者在迭代中的角色",
        "不是推动者，是五种角色的动态切换"
    )

    # ---- 7.1: Part 7 引导 ----
    slide = b.content_slide("第七部分：管理者在迭代中的角色", "Day 2 · 第七部分")
    b.add_text(slide, "大多数管理者的本能反应——都在阻碍迭代", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "本能反应 → 真实后果", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    instincts = [
        ("本能", "汇报详细计划、等批准再执行"),
        ("后果", "决策周期长，错过最佳学习窗口"),
        ("", ""),
        ("本能", "看到偏差立即纠正"),
        ("后果", "团队失去试错机会，只敢做安全的事"),
        ("", ""),
        ("本能", "把进度当成功指标"),
        ("后果", "团队赶进度，忽视学习——下一轮继续跑偏"),
    ]
    for i, (a, c) in enumerate(instincts):
        y = 1.85 + i * 0.45
        if a:
            b.add_text(slide, a, 0.4, y, 4.4, 0.4,
                       size=11, color=THEME['accent'], valign='middle')
            b.add_text(slide, c, 5.0, y, 4.6, 0.4,
                       size=11, color=THEME['primary'], valign='middle')

    # ---- 7.2 to 7.6: 五种角色深度 ----
    role_details = [
        ("角色深度 1/5 · 战略设定者", "你的核心动作：定义什么问题值得探索",
         "· 不参与细节——但定方向\n"
         "· 用我们 30 天想学到什么代替 30 天我们要达成什么\n"
         "· 拒绝做一个能用的产品——改为验证 3 个关键假设\n"
         "· 当团队失去焦点时，把问题收回\n"
         "· 当团队陷入细节时，问：这个对客户价值是什么",
         "产品经理说我们做这个，你说等一下，我们假设是什么？验证一下再说"),

        ("角色深度 2/5 · 资源守护者", "你的核心动作：保护团队不被外部打断",
         "· 会议不是越多越好——取消不必要的会议\n"
         "· 突发任务有缓冲——等等，给团队 3 天\n"
         "· 设定无干扰时段——团队成员每天有 4 小时不被任何人打扰\n"
         "· 你自己也不要在他们工作时顺路问一下\n"
         "· 当公司 KPI 压力来时，你是墙而不是放大器",
         "CEO 来电话让你汇报下季度规划——你挡回去，给我 1 周，我先跟团队对齐"),

        ("角色深度 3/5 · 挑战者", "你的核心动作：持续追问我们凭什么相信",
         "5 个关键追问：",
         "CEO 来电话让你汇报下季度规划——你挡回去，给我 1 周，我先跟团队对齐"),

        ("角色深度 4/5 · 学习教练", "你的核心动作：把我们学到了什么显性化",
         "学习教练的固定动作：",
         "每周一次 30 分钟复盘会——只问我们学到了什么，不讲进度"),

        ("角色深度 5/5 · 边界打破者", "你的核心动作：必要时打破部门墙/层级墙",
         "· 当团队需要其他部门资源时——你亲自协调\n"
         "· 当层级阻碍信息流动时——你亲自下沉\n"
         "· 当部门墙阻碍跨团队合作时——你主动打破\n"
         "· 当团队陷入自己人思维时——你引入外部视角\n"
         "· 当制度阻碍创新时——你敢于绕开——然后向上说明",
         "你的团队需要借调一个工程师，但 HR 流程要 2 周——你亲自去要，1 天解决"),
    ]

    for title, subtitle, content, scenario in role_details:
        slide = b.content_slide(title, "Day 2 · 第七部分")
        b.add_text(slide, subtitle, 0.4, 0.85, 9.2, 0.4,
                   size=16, bold=True, color=THEME['accent'])

        if "5 个关键追问" in content:
            questions = [
                ("①", "你做这件事的核心假设是什么？"),
                ("②", "为什么我们相信这个假设？"),
                ("③", "要验证这个假设，最小成本的方法是什么？"),
                ("④", "如果假设是错的，我们多快能知道？"),
                ("⑤", "这次迭代结束，我们具体能学到什么？"),
            ]
            for i, (num, q) in enumerate(questions):
                y = 1.4 + i * 0.5
                b.add_oval(slide, 0.4, y, 0.4, 0.4, fill_color=THEME['accent'])
                b.add_text(slide, num, 0.4, y, 0.4, 0.4,
                           size=12, bold=True, color='FFFFFF', align='center', valign='middle')
                b.add_text(slide, q, 1.0, y, 8.6, 0.4,
                           size=12, color=THEME['primary'], valign='middle')
        else:
            b.add_text(slide, content, 0.4, 1.4, 9.2, 2.5,
                       size=12, color=THEME['primary'])

        b.add_text(slide, "真实场景：", 0.4, 4.1, 9.2, 0.3,
                   size=14, bold=True, color=THEME['accent'])
        b.add_text(slide, scenario, 0.4, 4.5, 9.2, 0.5,
                   size=12, italic=True, color=THEME['primary'])

    # ---- 7.7: 角色切换自检 ----
    slide = b.content_slide("练习：角色切换自检", "Day 2 · 第七部分")
    b.add_text(slide, "在你的迭代方案中，你当前承担了哪种角色？", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    self_check = [
        ("战略设定者", "我清楚我们 30 天想学到什么"),
        ("资源守护者", "我为团队挡掉过 3 件以上干扰"),
        ("挑战者", "我每周至少追问团队 5 个假设"),
        ("学习教练", "每次迭代我们都有显性学习文档"),
        ("边界打破者", "我亲自为团队打通过跨部门障碍"),
    ]
    for i, (role, q) in enumerate(self_check):
        y = 1.4 + i * 0.65
        b.add_rect(slide, 0.4, y, 2.5, 0.55, fill_color=THEME['accent'])
        b.add_text(slide, role, 0.4, y, 2.5, 0.55,
                   size=12, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_rect(slide, 3.0, y, 6.6, 0.55, fill_color=THEME['light'])
        b.add_text(slide, q, 3.2, y, 6.4, 0.55,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- Part 8 divider ----
    b.divider_slide(
        "Day 2 · 第八部分",
        "三大要素完整分析",
        "用三要素框架，给自己的挑战卡做完整诊断"
    )

    # ---- 8.1: Part 8 引导 ----
    slide = b.content_slide("第八部分：三大要素完整分析", "Day 2 · 第八部分")
    b.add_text(slide, "客户洞察 × 交互涌现 × 敏捷迭代 = 创新三角", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_oval(slide, 4.0, 1.5, 2.0, 2.0, fill_color=THEME['accent'])
    b.add_text(slide, "创新型团队", 4.0, 1.5, 2.0, 2.0,
               size=16, bold=True, color='FFFFFF', align='center', valign='middle')

    b.add_rect(slide, 1.0, 1.8, 2.8, 0.6, fill_color=THEME['light'])
    b.add_text(slide, "① 客户洞察", 1.0, 1.8, 2.8, 0.6,
               size=14, bold=True, color=THEME['primary'], align='center', valign='middle')
    b.add_text(slide, "看见真实需求", 1.0, 2.4, 2.8, 0.4,
               size=10, color=THEME['secondary'], align='center')

    b.add_rect(slide, 6.2, 1.8, 2.8, 0.6, fill_color=THEME['light'])
    b.add_text(slide, "② 交互涌现", 6.2, 1.8, 2.8, 0.6,
               size=14, bold=True, color=THEME['primary'], align='center', valign='middle')
    b.add_text(slide, "让知识流动", 6.2, 2.4, 2.8, 0.4,
               size=10, color=THEME['secondary'], align='center')

    b.add_rect(slide, 3.6, 4.0, 2.8, 0.6, fill_color=THEME['light'])
    b.add_text(slide, "③ 敏捷迭代", 3.6, 4.0, 2.8, 0.6,
               size=14, bold=True, color=THEME['primary'], align='center', valign='middle')
    b.add_text(slide, "持续学到东西", 3.6, 4.6, 2.8, 0.4,
               size=10, color=THEME['secondary'], align='center')

    b.add_text(slide, "⚠ 三要素缺一不可，缺任何一个 = 创新不可持续", 0.4, 5.05, 9.2, 0.4,
               size=12, bold=True, color=THEME['accent'], align='center')

    # ---- 8.2: 三要素自诊表 ----
    slide = b.content_slide("个人诊断报告：你的三要素得分", "Day 2 · 第八部分")
    b.add_text(slide, "请根据昨天的练习，给你的团队打分（0-10）", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    diag = [
        ("① 客户洞察", "5.0", "我们能否穿透到第4层？", "4.5"),
        ("② 交互涌现", "5.0", "我们是否有让知识流动的机制？", "4.0"),
        ("③ 敏捷迭代", "5.0", "我们能否最小验证、快速学习？", "4.5"),
    ]
    for i, (name, score, q, suggestion) in enumerate(diag):
        y = 1.4 + i * 1.2
        b.add_rect(slide, 0.4, y, 9.2, 1.1,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, name, 0.6, y + 0.1, 2.5, 0.9,
                   size=14, bold=True,
                   color='FFFFFF' if i == 0 else THEME['primary'],
                   valign='middle')
        b.add_text(slide, score, 3.2, y + 0.1, 1.0, 0.9,
                   size=22, bold=True,
                   color='FFFFFF' if i == 0 else THEME['accent'],
                   align='center', valign='middle')
        b.add_text(slide, q, 4.4, y + 0.1, 3.0, 0.9,
                   size=11, color=THEME['primary'], valign='middle')
        b.add_text(slide, "建议：" + suggestion,
                   7.4, y + 0.1, 2.2, 0.9,
                   size=10, italic=True, color=THEME['secondary'], valign='middle')

    # ---- 8.3: 互诊练习 ----
    slide = b.content_slide("两人互诊 · 15 分钟", "Day 2 · 第八部分")
    b.add_text(slide, "找一位挑战卡场景相似的学员，互相诊断", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    pair = [
        ("A 分享", "5 分钟", "读你的挑战卡给 B 听"),
        ("B 提问", "3 分钟", "问 3 个问题：哪个要素最弱？为什么？"),
        ("B 建议", "3 分钟", "基于自己的经验，给一个具体动作建议"),
        ("B 分享", "5 分钟", "读你的挑战卡给 A 听"),
        ("A 提问", "3 分钟", "同上"),
        ("A 建议", "3 分钟", "同上"),
    ]
    for i, (step, time, desc) in enumerate(pair):
        y = 1.4 + i * 0.55
        b.add_rect(slide, 0.4, y, 9.2, 0.5,
                   fill_color=THEME['accent'] if i % 2 == 0 else THEME['light'])
        b.add_text(slide, step, 0.6, y, 1.8, 0.5,
                   size=11, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, time, 2.5, y, 1.0, 0.5,
                   size=10, color=THEME['secondary'], valign='middle')
        b.add_text(slide, desc, 3.6, y, 6.0, 0.5,
                   size=11, color=THEME['primary'], valign='middle')

    b.add_text(slide, "⚠ 关键：建议要明天就能做的具体动作，不是泛泛方向", 0.4, 4.7, 9.2, 0.4,
               size=12, bold=True, color=THEME['accent'], align='center')

    # ---- 8.4: 挑战卡完整分析 ----
    slide = b.content_slide("挑战卡完整分析：你的最终诊断", "Day 2 · 第八部分")
    b.add_text(slide, "基于互诊结果，完善你的挑战卡", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    analysis = [
        ("挑战是什么", "（保留原版）"),
        ("最弱要素诊断", "（从三要素中选一项：洞察/涌现/迭代）"),
        ("最弱的原因", "（是 9 种抑制行为中的哪种？还是机制缺失？）"),
        ("30 天具体动作", "（基于互诊建议，写 3 条具体动作）"),
        ("验证学习方式", "（怎么知道学到了？最小验证动作是什么？）"),
        ("支持资源", "（需要谁支持？跨部门协调？预算？）"),
    ]
    for i, (k, v) in enumerate(analysis):
        y = 1.4 + i * 0.55
        b.add_rect(slide, 0.4, y, 2.5, 0.5, fill_color=THEME['accent'] if i == 0 else THEME['light'])
        b.add_text(slide, k, 0.4, y, 2.5, 0.5,
                   size=11, bold=True,
                   color='FFFFFF' if i == 0 else THEME['primary'],
                   align='center', valign='middle')
        b.add_text(slide, v, 3.0, y, 6.6, 0.5,
                   size=11, italic=True, color=THEME['secondary'], valign='middle')

    # ---- 亮界科技背景 1 ----
    slide = b.content_slide("情景案例：亮界科技", "Day 2 · 亮界背景")
    b.add_text(slide, "一家 B2B SaaS 公司 · 真实困境", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "公司画像：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    profile = [
        ("行业", "企业 SaaS · 客户关系管理（CRM）赛道"),
        ("规模", "120 人 · 其中产品研发 60 人 · 销售 25 人"),
        ("成立", "6 年 · B 轮融资 8000 万"),
        ("现状", "老产品 500+ 客户 · 续约率 92% · 但 ARPU 停滞"),
        ("挑战", "3 家新对手用 AI 切入，市场份额 6 个月掉 8%"),
    ]
    for i, (k, v) in enumerate(profile):
        y = 1.85 + i * 0.55
        b.add_rect(slide, 0.4, y, 1.2, 0.5, fill_color=THEME['accent'])
        b.add_text(slide, k, 0.4, y, 1.2, 0.5,
                   size=11, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, v, 1.7, y, 8.0, 0.5,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 亮界科技背景 2 ----
    slide = b.content_slide("关键人物：张力", "Day 2 · 亮界背景")
    b.add_text(slide, "产品副总裁 · 45 岁 · 工号 008", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide,
        "· 公司联合创始人 · 主导老产品的成功\n"
        "· 行业 18 年经验 · 性格强势 · 决策快\n"
        "· 业务直觉极强 · 但不擅长反思\n"
        "· 团队内部评价：跟着张力干能成事，但很累\n"
        "· 最近半年开始焦虑——市场变了，过去的方法失灵了",
        0.4, 1.4, 9.2, 2.5, size=13, color=THEME['primary'])

    b.add_text(slide, "关键时刻：", 0.4, 4.0, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    b.add_text(slide, "CEO 上周说：张力，市场份额连续 3 个月下滑，你必须在 30 天内拿出方案", 0.4, 4.4, 9.2, 0.5,
               size=12, italic=True, color=THEME['accent'])

    # ---- 亮界科技背景 3 ----
    slide = b.content_slide("亮界的具体困境", "Day 2 · 亮界背景")
    b.add_text(slide, "张力过去 6 个月做了什么？为什么没起作用？", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    actions = [
        ("行动 1", "招聘 3 个 AI 工程师 · 启动AI升级项目"),
        ("结果", "3 个月过去了，产品还在 demo 阶段"),
        ("行动 2", "组织 4 次全员创新提案会"),
        ("结果", "收到 60 条提案，没有一条落地"),
        ("行动 3", "引进外部咨询公司做创新战略"),
        ("结果", "PPT 80 页，团队不认——不懂业务"),
        ("行动 4", "个人每周跟 5 个老客户深聊"),
        ("结果", "发现客户确实在用 AI 替代——但他自己没产品"),
    ]
    for i, (k, v) in enumerate(actions):
        y = 1.4 + i * 0.55
        is_action = '行动' in k
        b.add_rect(slide, 0.4, y, 1.5, 0.5,
                   fill_color=THEME['accent'] if is_action else THEME['light'])
        b.add_text(slide, k, 0.4, y, 1.5, 0.5,
                   size=11, bold=True,
                   color='FFFFFF' if is_action else THEME['primary'],
                   align='center', valign='middle')
        b.add_text(slide, v, 2.0, y, 7.6, 0.5,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- Part 9 divider ----
    b.divider_slide(
        "Day 2 · 下午",
        "第九部分：情景模拟 · 诊断亮界",
        "用今天的框架，诊断一个真实案例"
    )

    # ---- 9.1: Part 9 引导 ----
    slide = b.content_slide("第九部分：诊断亮界科技", "Day 2 · 第九部分")
    b.add_text(slide, "小组任务：3 小时内给出一份诊断报告", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    tasks = [
        ("第 1 步 · 25 分钟", "五因素诊断", "亮界在 5 个关键因素上各得几分？"),
        ("第 2 步 · 25 分钟", "三要素失效分析", "亮界的客户洞察/交互涌现/敏捷迭代哪里失效？"),
        ("第 3 步 · 25 分钟", "行为链分析", "张力做了哪些对的事却没结果？为什么？"),
        ("第 4 步 · 15 分钟", "关键发现", "找出张力团队最核心的 1-2 个问题"),
    ]
    for i, (step, name, desc) in enumerate(tasks):
        y = 1.4 + i * 0.85
        b.add_rect(slide, 0.4, y, 9.2, 0.75,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, step, 0.6, y, 2.5, 0.75,
                   size=12, bold=True,
                   color='FFFFFF' if i == 0 else THEME['primary'],
                   valign='middle')
        b.add_text(slide, name, 3.2, y, 2.5, 0.75,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 5.8, y, 3.8, 0.75,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 9.2: 五因素诊断 ----
    slide = b.content_slide("五因素诊断 · 小组讨论", "Day 2 · 第九部分")
    b.add_text(slide, "请评估亮界在 5 个关键因素上的得分（0-10）", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    factors = [
        ("心理安全感", "张力强势，员工不敢说真话"),
        ("认知多样性", "团队背景相似（同质化）"),
        ("探索空间", "KPI 压力大，没空间探索"),
        ("学习速度", "试错文化弱，怕失败"),
        ("领导者信号", "张力言行不一致——说一套做一套"),
    ]
    for i, (factor, note) in enumerate(factors):
        y = 1.4 + i * 0.65
        b.add_rect(slide, 0.4, y, 2.5, 0.55, fill_color=THEME['accent'])
        b.add_text(slide, factor, 0.4, y, 2.5, 0.55,
                   size=12, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_rect(slide, 3.0, y, 6.6, 0.55, fill_color=THEME['light'])
        b.add_text(slide, note, 3.2, y, 6.4, 0.55,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 9.3: 三要素失效分析 ----
    slide = b.content_slide("三要素失效分析", "Day 2 · 第九部分")
    b.add_text(slide, "哪个要素失效最严重？为什么？", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    failures = [
        ("客户洞察", "5", "张力每周跟客户聊——但只聊我们想做的产品有没有需求，没问客户真正想要什么"),
        ("交互涌现", "3", "信息只在张力一个人手里。60 条提案没人看——因为没人敢问张力为什么不做"),
        ("敏捷迭代", "4", "AI 升级项目 3 个月没 demo——因为没人敢告诉张力这个方向可能不对"),
    ]
    for i, (e, score, reason) in enumerate(failures):
        y = 1.4 + i * 1.2
        b.add_rect(slide, 0.4, y, 9.2, 1.1,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'])
        b.add_text(slide, e, 0.6, y + 0.1, 1.8, 0.9,
                   size=14, bold=True,
                   color='FFFFFF' if i == 0 else THEME['primary'],
                   valign='middle')
        b.add_text(slide, score, 2.5, y + 0.1, 0.8, 0.9,
                   size=28, bold=True,
                   color='FFFFFF' if i == 0 else THEME['accent'],
                   align='center', valign='middle')
        b.add_text(slide, reason, 3.5, y + 0.1, 6.1, 0.9,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 9.4: 行为链分析 ----
    slide = b.content_slide("行为链分析 · 张力的因果链", "Day 2 · 第九部分")
    b.add_text(slide, '"对的行为"如何一步步导致"失败"', 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    chain = [
        "张力强势决策 → 员工不敢说真话",
        "员工不说真话 → 张力看不到真实信息",
        "看不到真实信息 → 张力继续按直觉决策",
        "按直觉决策 → 决策越来越准（短期）",
        "决策短期有效 → 张力更相信自己",
        "张力更自信 → 更强势",
        "（循环强化）",
    ]
    for i, c in enumerate(chain):
        y = 1.4 + i * 0.5
        is_loop = '循环' in c
        b.add_rect(slide, 0.4, y, 9.2, 0.45,
                   fill_color=THEME['accent'] if is_loop else THEME['light'])
        b.add_text(slide, c, 0.6, y, 9.0, 0.45,
                   size=12, bold=is_loop,
                   color='FFFFFF' if is_loop else THEME['primary'],
                   valign='middle')

    b.add_text(slide, "→ 这是创新型团队→运营型团队演化的典型路径", 0.4, 5.0, 9.2, 0.4,
               size=12, bold=True, color=THEME['accent'], align='center')

    # ---- 9.5: 关键发现 ----
    slide = b.content_slide("亮界诊断的关键发现", "Day 2 · 第九部分")
    b.add_text(slide, "5 个最核心问题", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    findings = [
        ("1", "张力的强势决策风格是亮界创新的最大障碍"),
        ("2", "员工的心理安全感几乎为零——60 条提案无人问津"),
        ("3", "知识只在张力脑子里——缺乏跨部门/跨层级的流动机制"),
        ("4", "迭代速度慢——3 个月没看到 demo，本质是不敢失败"),
        ("5", "言行不一致——张力说大家多提意见，但提了没采纳，下次没人提"),
    ]
    for i, (n, f) in enumerate(findings):
        y = 1.4 + i * 0.65
        b.add_oval(slide, 0.4, y, 0.5, 0.5, fill_color=THEME['accent'])
        b.add_text(slide, n, 0.4, y, 0.5, 0.5,
                   size=18, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, f, 1.1, y, 8.5, 0.5,
                   size=12, color=THEME['primary'], valign='middle')

    # ---- 9.6: 小组汇报 ----
    slide = b.content_slide("小组汇报 · 5 分钟/组", "Day 2 · 第九部分")
    b.add_text(slide, "每个小组 5 分钟，1 分钟提问", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    b.add_text(slide, "汇报结构（建议）：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    structure = [
        ("1 分钟", "诊断结论：亮界最核心的 1-2 个问题"),
        ("2 分钟", "分析支撑：基于 5 因素 + 3 要素 + 行为链"),
        ("1 分钟", "如果你来：作为张力，第一周会做什么"),
        ("1 分钟", "给张力 1 个具体的明天就做的动作"),
    ]
    for i, (time, desc) in enumerate(structure):
        y = 1.85 + i * 0.65
        b.add_rect(slide, 0.4, y, 9.2, 0.55,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'])
        b.add_text(slide, time, 0.6, y, 1.5, 0.55,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 2.2, y, 7.4, 0.55,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- Part 10 divider ----
    b.divider_slide(
        "Day 2 · 第十部分",
        "如果你是张力：方案设计",
        "三件事原则 · 全场展示 · 互相挑战"
    )

    # ---- 10.1: Part 10 引导 ----
    slide = b.content_slide("第十部分：如果你是张力", "Day 2 · 第十部分")
    b.add_text(slide, "诊断完还不够——你要拿出能落地的方案", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    b.add_text(slide, "关键问题：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    questions = [
        "· 如果你是张力，明天回公司第一件事做什么？",
        "· 30 天内，最想做的 3 件事是什么？",
        "· 90 天后，怎么判断自己在改变方向？",
    ]
    for i, q in enumerate(questions):
        y = 1.85 + i * 0.55
        b.add_rect(slide, 0.4, y, 9.2, 0.5,
                   fill_color=THEME['light'] if i % 2 == 0 else THEME['bg'],
                   line_color=THEME['light'])
        b.add_text(slide, q, 0.6, y, 9.0, 0.5,
                   size=13, color=THEME['primary'], valign='middle')

    # ---- 10.2: 三件事原则 ----
    slide = b.content_slide("三件事原则", "Day 2 · 第十部分")
    b.add_text(slide, "不是 10 件事——只做 3 件", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    principles = [
        ("事 1", "自我改变", "张力本人做出一个反差极大的动作——让团队看到"),
        ("事 2", "机制建立", "建立一个新机制——让知识流动的规则变了"),
        ("事 3", "最小验证", "选 1 个最小项目——快速试错、快速学习"),
    ]
    for i, (n, t, desc) in enumerate(principles):
        y = 1.4 + i * 1.15
        b.add_rect(slide, 0.4, y, 9.2, 1.05,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, n, 0.6, y + 0.1, 1.2, 0.85,
                   size=18, bold=True,
                   color='FFFFFF' if i == 0 else THEME['accent'],
                   align='center', valign='middle')
        b.add_text(slide, t, 1.9, y + 0.1, 2.0, 0.85,
                   size=14, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 4.0, y + 0.1, 5.5, 0.85,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 10.3: 张力示例方案 ----
    slide = b.content_slide("示范方案 · 张力的三件事", "Day 2 · 第十部分")
    b.add_text(slide, "一个可参考的张力式方案", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    plan = [
        ("事 1", "自我改变", "下一次全员会，我主动承认过去的失误——并宣布新规则"),
        ("事 2", "机制建立", "每周五下午 4 点全员开放 1 小时——任何人都可以问我任何问题"),
        ("事 3", "最小验证", "选 1 个 AI 小功能，3 周内访谈 8 个客户，做最小可用版本——不上大项目"),
    ]
    for i, (n, t, desc) in enumerate(plan):
        y = 1.4 + i * 1.15
        b.add_rect(slide, 0.4, y, 9.2, 1.05,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'])
        b.add_text(slide, n, 0.6, y + 0.1, 1.2, 0.85,
                   size=16, bold=True,
                   color='FFFFFF' if i == 0 else THEME['accent'],
                   align='center', valign='middle')
        b.add_text(slide, t, 1.9, y + 0.1, 1.8, 0.85,
                   size=13, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, desc, 3.8, y + 0.1, 5.7, 0.85,
                   size=10, italic=True, color=THEME['primary'], valign='middle')

    # ---- 10.4: 全场展示 ----
    slide = b.content_slide("全场展示 · 互相挑战", "Day 2 · 第十部分")
    b.add_text(slide, "每个小组 3 分钟方案展示 + 2 分钟互相挑战", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    challenges = [
        "你这三件事，张力能坚持 30 天吗？为什么？",
        "如果张力明天就做，他会遇到什么最大阻力？",
        "你假设团队会怎么反应？如果他们不当回事，怎么办？",
        "这件事不成功，你怎么知道？最坏情况是什么？",
    ]
    for i, c in enumerate(challenges):
        y = 1.4 + i * 0.85
        b.add_rect(slide, 0.4, y, 9.2, 0.75,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, c, 0.6, y, 9.0, 0.75,
                   size=12, color=THEME['primary'], valign='middle')

    # ---- Part 11 divider ----
    b.divider_slide(
        "Day 2 · 第十一部分",
        "连接自己的团队",
        "从亮界回到你的真实战场"
    )

    # ---- 11.1: 对照练习 ----
    slide = b.content_slide("对照练习：亮界 vs 我的团队", "Day 2 · 第十一部分")
    b.add_text(slide, "15 分钟独立思考 + 5 分钟小组分享", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    comparisons = [
        ("相似点", "我的团队有没有跟亮界一样的问题？"),
        ("不同点", "我的团队最不一样的特征是什么？"),
        ("可以借鉴", "亮界的方案里，哪些可以带回我的团队？"),
        ("必须调整", "亮界的方案在我团队不能直接用的部分是什么？"),
    ]
    for i, (k, v) in enumerate(comparisons):
        y = 1.4 + i * 0.85
        b.add_rect(slide, 0.4, y, 9.2, 0.75,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, k, 0.6, y, 2.0, 0.75,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, v, 2.7, y, 6.9, 0.75,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 11.2: 个人反思 ----
    slide = b.content_slide("个人反思 · 10 分钟", "Day 2 · 第十一部分")
    b.add_text(slide, "请找一个安静的地方，独自思考", 0.4, 0.85, 9.2, 0.4,
               size=16, bold=True, color=THEME['accent'])

    b.add_text(slide, "反思 5 个问题：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    reflections = [
        "① 两天下来，最触动你的 1 句话是什么？",
        "② 你的挑战卡里最让你焦虑的部分是什么？",
        "③ 你作为管理者，过去 6 个月做过最像张力的事是什么？",
        "④ 你愿意改变吗？什么阻碍你改变？",
        "⑤ 回到岗位第一周，你最想做的 1 个动作是什么？",
    ]
    for i, q in enumerate(reflections):
        y = 1.85 + i * 0.6
        b.add_text(slide, q, 0.6, y, 9.0, 0.55,
                   size=12, color=THEME['primary'], valign='middle')

    # ---- Part 12 divider ----
    b.divider_slide(
        "Day 2 · 第十二部分",
        "行动承诺 + 课程收尾",
        "不是结束——是开始"
    )

    # ---- 12.1: 行动承诺卡 ----
    slide = b.content_slide("30 天行动承诺卡", "Day 2 · 第十二部分")
    b.add_text(slide, "写给未来的你 · 30 天后回看", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_text(slide, "承诺卡的 6 个字段：", 0.4, 1.4, 9.2, 0.3,
               size=14, bold=True, color=THEME['accent'])
    commitments = [
        ("我的 30 天挑战", "一句话：我要解决的具体问题是什么"),
        ("第一周做什么", "Day 1-7 必须完成的具体动作"),
        ("谁会知道", "谁会成为你的问责伙伴——他知道你在做这件事"),
        ("30 天后学到什么", "30 天后我希望验证/学到的具体内容"),
        ("90 天后变成什么", "90 天后这个改变是否已经融入我的日常"),
        ("我会怎么知道", "具体的可观察的信号——不是我感觉"),
    ]
    for i, (k, v) in enumerate(commitments):
        y = 1.85 + i * 0.55
        b.add_rect(slide, 0.4, y, 2.6, 0.5, fill_color=THEME['accent'] if i == 0 else THEME['light'])
        b.add_text(slide, k, 0.4, y, 2.6, 0.5,
                   size=11, bold=True,
                   color='FFFFFF' if i == 0 else THEME['primary'],
                   align='center', valign='middle')
        b.add_text(slide, v, 3.1, y, 6.5, 0.5,
                   size=10, color=THEME['primary'], valign='middle')

    # ---- 12.2: 两人分享 ----
    slide = b.content_slide("两人分享承诺 · 10 分钟", "Day 2 · 第十二部分")
    b.add_text(slide, "找一个伙伴，把你的承诺卡读给他听", 0.4, 0.85, 9.2, 0.4,
               size=15, bold=True, color=THEME['accent'])

    share = [
        ("5 分钟", "A 读给 B 听", "A 完整读自己的承诺卡"),
        ("3 分钟", "B 反馈", "B 问 1 个问题 + 给 1 个具体建议"),
        ("5 分钟", "B 读给 A 听", "B 完整读自己的承诺卡"),
        ("3 分钟", "A 反馈", "同上"),
        ("2 分钟", "互相交换承诺卡", "写在承诺卡上，30 天后回看"),
    ]
    for i, (t, a, desc) in enumerate(share):
        y = 1.4 + i * 0.65
        b.add_rect(slide, 0.4, y, 9.2, 0.55,
                   fill_color=THEME['accent'] if i % 2 == 0 else THEME['light'])
        b.add_text(slide, t, 0.6, y, 1.0, 0.55,
                   size=11, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, a, 1.7, y, 2.0, 0.55,
                   size=11, bold=True, color=THEME['accent'], valign='middle')
        b.add_text(slide, desc, 3.8, y, 5.8, 0.55,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 12.3: 全场承诺 ----
    slide = b.content_slide("全场承诺仪式", "Day 2 · 第十二部分")
    b.add_text(slide, "一起说一句承诺", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    b.add_rect(slide, 1.0, 1.6, 8.0, 1.4, fill_color=THEME['accent'])
    b.add_text(slide, "我承诺：从明天起，", 1.0, 1.7, 8.0, 0.5,
               size=20, bold=True, color='FFFFFF', align='center', valign='middle')
    b.add_text(slide, "做一件让我的团队更敢创新的事。", 1.0, 2.3, 8.0, 0.5,
               size=20, bold=True, color='FFFFFF', align='center', valign='middle')

    b.add_text(slide, "请全体起立，一起读一遍", 0.4, 3.3, 9.2, 0.4,
               size=14, color=THEME['secondary'], align='center')

    b.add_text(slide, "⚠ 不是给别人承诺——是给自己承诺", 0.4, 3.9, 9.2, 0.4,
               size=14, bold=True, color=THEME['accent'], align='center')
    b.add_text(slide, "⚠ 不是要做到完美——是开始做", 0.4, 4.3, 9.2, 0.4,
               size=14, bold=True, color=THEME['accent'], align='center')

    # ---- 12.4: 课程核心收获 ----
    slide = b.content_slide("课程核心收获 · 12 句话", "Day 2 · 收尾")
    b.add_text(slide, "两天下来最值得带走的 12 句话", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    quotes = [
        ("1", "团队不创新的原因，大部分在管理者身上"),
        ("2", "给了空间可能恰恰是你用了 9 种方式压制空间"),
        ("3", "客户洞察不是问客户想要什么——是看见他没说出口的"),
        ("4", "团队的真正问题不是没有人——是知识在墙里不流动"),
        ("5", "迭代不是计划→执行，是假设→验证→学习"),
        ("6", "管理者的角色不是推动者——是 5 种角色的动态切换"),
        ("7", "改变不是大动作——是第一个具体的卡"),
        ("8", "失败不是要避免——是要让它发生，并显性学习"),
        ("9", "你给的空间是否真的存在——团队心里最清楚"),
        ("10", "你不是缺方法——你缺的是第一个具体的动作"),
        ("11", "30 天后不是结束——是 30 天行动承诺的开始"),
        ("12", "真正的改变，是从你愿意承认我可能错了开始"),
    ]
    for i, (n, q) in enumerate(quotes):
        y = 1.4 + i * 0.32
        b.add_rect(slide, 0.4, y, 0.5, 0.28, fill_color=THEME['accent'])
        b.add_text(slide, n, 0.4, y, 0.5, 0.28,
                   size=10, bold=True, color='FFFFFF', align='center', valign='middle')
        b.add_text(slide, q, 1.0, y, 8.6, 0.28,
                   size=10, color=THEME['primary'], valign='middle')

    # ---- 12.5: 30 天支持 ----
    slide = b.content_slide("30 天跟进支持", "Day 2 · 收尾")
    b.add_text(slide, "你不是一个人——30 天后还有对话", 0.4, 0.85, 9.2, 0.4,
               size=18, bold=True, color=THEME['accent'])

    support = [
        ("第 7 天", "小组群打卡", "每人 1 句话：本周最触动你的一件事"),
        ("第 14 天", "学习伙伴对话", "2 人组 30 分钟视频——聊做对了什么、卡在哪里"),
        ("第 30 天", "公开复盘", "线上小组会 60 分钟——每人分享 30 天学习"),
        ("第 90 天", "深度复盘", "回到工作坊所在城市，1 天深度复盘"),
    ]
    for i, (time, name, desc) in enumerate(support):
        y = 1.4 + i * 0.85
        b.add_rect(slide, 0.4, y, 9.2, 0.75,
                   fill_color=THEME['accent'] if i == 0 else THEME['light'],
                   line_color=THEME['light'])
        b.add_text(slide, time, 0.6, y, 1.5, 0.75,
                   size=12, bold=True, color=THEME['primary'], valign='middle')
        b.add_text(slide, name, 2.2, y, 2.5, 0.75,
                   size=12, bold=True, color=THEME['accent'], valign='middle')
        b.add_text(slide, desc, 4.8, y, 4.8, 0.75,
                   size=11, color=THEME['primary'], valign='middle')

    # ---- 12.6: 致谢页 ----
    slide = b.content_slide("致谢 · 期待你的故事", "Day 2 · 收尾")
    b.add_rect(slide, 0, 0, 10, 5.625, fill_color=THEME['primary'])
    b.add_rect(slide, 0, 2.0, 10, 0.04, fill_color=THEME['accent'])

    b.add_text(slide, "THANK YOU", 0.5, 0.8, 9, 1.0,
               size=44, bold=True, color='FFFFFF', align='center')

    b.add_text(slide, "创新领导力", 0.5, 1.8, 9, 0.6,
               size=24, bold=True, color=THEME['accent'], align='center')

    b.add_text(slide, "——打造创新型团队", 0.5, 2.4, 9, 0.5,
               size=18, color=THEME['gray'], align='center')

    b.add_text(slide, "做出那个改变——不是为了你的团队，", 0.5, 3.0, 9, 0.4,
               size=14, color='FFFFFF', align='center')
    b.add_text(slide, "是为了让你 5 年后还能在管理岗位上继续做事", 0.5, 3.4, 9, 0.4,
               size=14, color='FFFFFF', align='center')

    b.add_text(slide, "罗老师团队 · 竞越课程研发中心", 0.5, 4.5, 9, 0.4,
               size=12, color=THEME['gray'], align='center')
    b.add_text(slide, "30 天后，我们再见", 0.5, 4.9, 9, 0.4,
               size=12, color=THEME['accent'], align='center')

    b.add_oval(slide, 9.3, 5.1, 0.4, 0.4, fill_color=THEME['accent'])
    b.add_text(slide, str(b.page + 1), 9.3, 5.1, 0.4, 0.4,
               size=11, bold=True, color='FFFFFF', align='center', valign='middle')


def main():
    src = 'D:/2026年课程/竞越/创新领导力：打造创新型团队/完整课程表/05-授课PPT/创新领导力_授课PPT_v1.0_33页.pptx'
    dst = 'D:/2026年课程/竞越/创新领导力：打造创新型团队/完整课程表/05-授课PPT/创新领导力_授课PPT_v2.0_120页.pptx'

    prs = Presentation(src)
    print("原始 PPT 页数:", len(prs.slides))

    builder = PPTBuilder(prs)
    print("构建 Day 1 PM...")
    build_day1_pm(builder)
    print("构建 Day 2...")
    build_day2(builder)
    print("最终页数:", builder.page)

    prs.save(dst)
    print("已保存:", dst)


if __name__ == '__main__':
    main()