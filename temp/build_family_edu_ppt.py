"""
AI 时代的家庭教育 - 授课 PPT 生成脚本
120 页高质量授课 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ===== 配色（参考高效能AI7习惯 + 教育/家庭主题适配）=====
C_PRIMARY = RGBColor(0x1A, 0x2B, 0x4A)   # 深蓝 - 头部/主色
C_SECONDARY = RGBColor(0x4A, 0x5F, 0x7F) # 钢蓝 - 次级
C_ACCENT = RGBColor(0xD4, 0xA5, 0x74)    # 暖金 - 强调/金句
C_LIGHT = RGBColor(0xF5, 0xF0, 0xE8)     # 米白 - 浅背景
C_BG = RGBColor(0xFA, 0xF7, 0xF0)        # 暖白 - 背景
C_DARK = RGBColor(0x2C, 0x3E, 0x50)      # 深灰蓝
C_RED = RGBColor(0xC8, 0x4B, 0x4B)       # 朱红 - 警示
C_GREEN = RGBColor(0x5A, 0x8C, 0x5A)     # 苔绿 - 成功
C_TEXT = RGBColor(0x33, 0x33, 0x33)      # 主文字
C_TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66) # 次文字
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 字体
FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

# ===== 初始化 =====
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def set_text(tf, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font=FONT_CN, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 中文字体
    rPr = run._r.get_or_add_rPr()
    eastAsia = etree.SubElement(rPr, qn('a:ea'))
    eastAsia.set('typeface', font)
    return run


def add_text_box(slide, x, y, w, h, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font=FONT_CN, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    set_text(tf, text, size, bold, color, align, font, anchor)
    return tb


def add_paragraphs(slide, x, y, w, h, items, size=16, color=C_TEXT, font=FONT_CN, line_spacing=1.2):
    """items: list of (text, bold) or string"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, bold = item, False
        else:
            text, bold = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        eastAsia = etree.SubElement(rPr, qn('a:ea'))
        eastAsia.set('typeface', font)
    return tb


def add_rect(slide, x, y, w, h, fill=C_PRIMARY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_rounded_rect(slide, x, y, w, h, fill=C_PRIMARY, line=None, radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    s.adjustments[0] = radius
    return s


def add_oval(slide, x, y, w, h, fill=C_PRIMARY, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_line(slide, x1, y1, x2, y2, color=C_PRIMARY, width=1):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def add_arrow(slide, x1, y1, x2, y2, color=C_PRIMARY, width=2):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x1, y1, x2-x1, y2-y1)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def page_chrome(slide, page_num, section=""):
    """统一页眉/页脚/页码"""
    # 顶部细线
    add_line(slide, Inches(0.4), Inches(0.32), Inches(9.6), Inches(0.32), C_ACCENT, 0.75)
    # 版权底部
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(5.35), Inches(7.5), Inches(0.25))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "© 罗老师《AI 时代的家庭教育》"
    r.font.name = FONT_CN
    r.font.size = Pt(8)
    r.font.color.rgb = C_TEXT_LIGHT
    # 章节
    if section:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(7), Inches(0.22))
        sf = sb.text_frame
        sf.margin_left = 0
        sf.margin_top = 0
        sp = sf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        sr = sp.add_run()
        sr.text = section
        sr.font.name = FONT_CN
        sr.font.size = Pt(9)
        sr.font.color.rgb = C_TEXT_LIGHT
    # 页码
    add_oval(slide, Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4), C_PRIMARY)
    tb = slide.shapes.add_textbox(Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(page_num)
    r.font.name = FONT_EN
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    # 垂直居中
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def section_header(slide, num_label, title, kicker=""):
    """大标题头部：编号 + 主标 + 副标"""
    # 左侧色条
    add_rect(slide, Inches(0.4), Inches(0.55), Inches(0.08), Inches(0.85), C_ACCENT)
    # 编号
    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(0.6), Inches(0.4), num_label, size=12, bold=True, color=C_ACCENT, font=FONT_EN)
    # 副标
    if kicker:
        add_text_box(slide, Inches(0.6), Inches(0.78), Inches(9), Inches(0.3), kicker, size=11, color=C_TEXT_LIGHT, font=FONT_CN)
    # 主标题
    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(9), Inches(0.5), title, size=28, bold=True, color=C_PRIMARY, font=FONT_CN)


# =====================================================
# Page 1: 封面
# =====================================================
def slide_01_cover():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_PRIMARY
    # 装饰条
    add_rect(s, Inches(0), Inches(0), Inches(10), Inches(0.15), C_ACCENT)
    add_rect(s, Inches(0), Inches(5.475), Inches(10), Inches(0.15), C_ACCENT)
    # 左侧大编号
    add_text_box(s, Inches(0.6), Inches(0.6), Inches(2), Inches(1.5), "AI", size=120, bold=True, color=C_ACCENT, font=FONT_EN)
    # 主标
    add_text_box(s, Inches(0.6), Inches(1.9), Inches(8.8), Inches(0.8), "AI 时代的家庭教育", size=44, bold=True, color=C_WHITE, font=FONT_CN)
    add_text_box(s, Inches(0.6), Inches(2.7), Inches(8.8), Inches(0.5), "从焦虑到超越竞争的家庭行动系统", size=22, color=C_LIGHT, font=FONT_CN)
    # 副信息
    add_line(s, Inches(0.6), Inches(3.6), Inches(2.0), Inches(3.6), C_ACCENT, 1.5)
    add_text_box(s, Inches(0.6), Inches(3.75), Inches(8.8), Inches(0.4), "13 讲系统课 · 3 大原创模型 · 13 套原创工具 · 30 天行动清单 · 3 年路线图", size=14, color=C_WHITE, font=FONT_CN)
    add_text_box(s, Inches(0.6), Inches(4.3), Inches(8.8), Inches(0.3), "首席教学设计师：罗老师  |  授课版本 2026", size=12, color=C_LIGHT, font=FONT_CN)
    # 版权
    add_text_box(s, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.3), "© 罗老师《AI 时代的家庭教育》· 著作权所有 · 未经授权禁止复制传播", size=9, color=C_LIGHT, font=FONT_CN)


# =====================================================
# Page 2: 课程地图
# =====================================================
def slide_02_map():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 2, "课程地图 / Course Map")
    section_header(s, "00", "课程地图", "一张图看懂 13 讲的全景脉络")

    # 四大模块
    modules = [
        ("01 破除焦虑", "发刊词 + 先导课", "建立坐标系：未来会好吗、该学什么", C_ACCENT),
        ("02 看见未来", "问题 1-4", "能力图谱 · 自学 · 真实问题 · 判断力", C_PRIMARY),
        ("03 重构家庭", "问题 5-8", "基础能力 · 父母角色 · 使用边界 · 情感连接", C_SECONDARY),
        ("04 长期主义", "问题 9-12", "兴趣 · 评估 · 路线图 · 避坑", C_DARK),
    ]
    box_w = Inches(2.1)
    box_h = Inches(2.3)
    start_x = Inches(0.55)
    gap = Inches(0.25)
    for i, (title, sub, desc, color) in enumerate(modules):
        x = start_x + (box_w + gap) * i
        # 顶部色块
        add_rect(s, x, Inches(1.7), box_w, Inches(0.5), color)
        add_text_box(s, x, Inches(1.78), box_w, Inches(0.4), title, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 主体
        add_rect(s, x, Inches(2.2), box_w, Inches(1.7), C_LIGHT)
        add_text_box(s, x + Inches(0.1), Inches(2.35), box_w - Inches(0.2), Inches(0.4), sub, size=13, bold=True, color=color)
        add_text_box(s, x + Inches(0.1), Inches(2.85), box_w - Inches(0.2), Inches(1.0), desc, size=11, color=C_TEXT)

    # 底部：原创内容
    add_rect(s, Inches(0.5), Inches(4.25), Inches(9), Inches(0.85), C_PRIMARY)
    add_text_box(s, Inches(0.7), Inches(4.32), Inches(8.6), Inches(0.32), "本课程三大原创", size=12, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(4.62), Inches(8.6), Inches(0.42),
                 "三锚模型（目标-边界-反馈）· 家庭 Prompt 工作流（5 大场景）· 超越竞争能力图谱（3 大稀缺能力）",
                 size=13, color=C_WHITE)


# =====================================================
# 发刊词 (3-5 页) - 5 页
# =====================================================
def slide_03():
    """发刊词 1: 标题 + 老周的故事"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 3, "发刊词")
    section_header(s, "00·发刊词", "给孩子更多提高的可能性", "一个让我重新思考「教育」的问题")

    # 引述大块
    add_rounded_rect(s, Inches(0.6), Inches(1.7), Inches(8.8), Inches(1.6), C_PRIMARY, radius=0.08)
    add_text_box(s, Inches(0.9), Inches(1.9), Inches(8.2), Inches(0.4), "「我不担心孩子输给 AI。", size=18, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.9), Inches(2.3), Inches(8.2), Inches(0.4), "我担心的是——他还没学会怎么和 AI 一起把事情做好，", size=18, bold=True, color=C_WHITE)
    add_text_box(s, Inches(0.9), Inches(2.7), Inches(8.2), Inches(0.4), "就已经习惯了让 AI 替他做。」", size=18, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(7.0), Inches(3.05), Inches(2.2), Inches(0.2), "—— 三年级家长 老周", size=11, color=C_LIGHT)

    # 三个家长焦虑
    items = [
        "AI 时代孩子会不会被淘汰？",
        "现在学的东西还有用吗？",
        "我不懂 AI，怎么辅导孩子？"
    ]
    add_text_box(s, Inches(0.6), Inches(3.55), Inches(8.8), Inches(0.3), "我问了上百位家长同一个问题，答案几乎排成三队：", size=12, color=C_TEXT_LIGHT)
    for i, t in enumerate(items):
        x = Inches(0.6) + Inches(3.0) * i
        add_rounded_rect(s, x, Inches(3.95), Inches(2.85), Inches(0.85), C_LIGHT, radius=0.1)
        add_text_box(s, x + Inches(0.1), Inches(4.1), Inches(2.65), Inches(0.6), t, size=13, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 金句
    add_text_box(s, Inches(0.6), Inches(4.95), Inches(8.8), Inches(0.3),
                 "孩子真正输掉的，不是输给 AI，而是输在「还没学会用 AI」。", size=12, bold=True, color=C_RED)


def slide_04():
    """发刊词 2: 焦虑公式"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 4, "发刊词")
    section_header(s, "00·发刊词", "焦虑的反义词不是「想开点」", "一个反常识的判断：焦虑 = 认知 - 行动")

    # 公式
    add_rect(s, Inches(0.6), Inches(1.7), Inches(8.8), Inches(1.0), C_LIGHT)
    # 认知
    add_text_box(s, Inches(0.9), Inches(1.95), Inches(2.0), Inches(0.5), "认知", size=20, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.9), Inches(1.95), Inches(0.5), Inches(0.5), "−", size=32, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(3.4), Inches(1.95), Inches(2.0), Inches(0.5), "行动", size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.4), Inches(1.95), Inches(0.5), Inches(0.5), "=", size=32, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.9), Inches(1.95), Inches(3.0), Inches(0.5), "焦虑", size=24, bold=True, color=C_RED, align=PP_ALIGN.CENTER)

    # 下方说明
    add_text_box(s, Inches(0.6), Inches(2.9), Inches(8.8), Inches(0.4),
                 "家长对 AI 的焦虑，绝大多数不是认知不够——你刷几篇公众号就知道 ChatGPT 是什么。", size=13, color=C_TEXT)
    add_text_box(s, Inches(0.6), Inches(3.3), Inches(8.8), Inches(0.4),
                 "你依然焦虑，是因为你知道「AI 来了」，但不知道「我们家具体怎么办」。", size=13, color=C_TEXT)

    # 不知道什么 - 卡片
    items = ["不知道该不该让孩子用 AI 做作业", "不知道用什么 Prompt 提问效果更好", "不知道孩子用 AI 时自己该做什么"]
    for i, t in enumerate(items):
        y = Inches(3.85) + Inches(0.45) * i
        add_rect(s, Inches(0.6), y, Inches(0.08), Inches(0.32), C_RED)
        add_text_box(s, Inches(0.85), y, Inches(8.5), Inches(0.32), t, size=12, color=C_TEXT)

    # 金句
    add_rounded_rect(s, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.4), C_PRIMARY, radius=0.2)
    add_text_box(s, Inches(0.6), Inches(5.05), Inches(8.8), Inches(0.3), "不焦虑的前提，不是「想通了」，而是「有行动系统」。",
                 size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_05():
    """发刊词 3: 三个根本不同"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 5, "发刊词")
    section_header(s, "00·发刊词", "这套课和市面上「AI 教育课」的三个根本不同", "不讲工具、不喊口号、不当观众")

    items = [
        ("第一", "不是讲 AI 工具，是讲家庭教育系统", "AI 是配角，主角是家庭关系、孩子能力、父母角色。我会讲 AI，但站在「家庭教育」的角度讲。", C_PRIMARY),
        ("第二", "不是讲「孩子该学什么」，是讲「父母该做什么」", "你的家庭 AI 系统怎么设计、协议怎么签、Prompt 怎么搭、孩子用 AI 时你在哪里。这套课给你养孩子的操作系统。", C_SECONDARY),
        ("第三", "不喊口号，给工具", "13 讲，13 个工具。讲完一讲，你会拿到一个马上能打印、能签署、能执行的工具。", C_ACCENT),
    ]
    for i, (label, title, desc, color) in enumerate(items):
        y = Inches(1.7) + Inches(1.15) * i
        # 编号圆
        add_oval(s, Inches(0.6), y + Inches(0.1), Inches(0.6), Inches(0.6), color)
        add_text_box(s, Inches(0.6), y + Inches(0.18), Inches(0.6), Inches(0.4), label[1], size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 标题
        add_text_box(s, Inches(1.4), y + Inches(0.1), Inches(7.8), Inches(0.35), title, size=14, bold=True, color=color)
        add_text_box(s, Inches(1.4), y + Inches(0.45), Inches(7.8), Inches(0.65), desc, size=11, color=C_TEXT)


def slide_06():
    """发刊词 4: 三大原创模型"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 6, "发刊词")
    section_header(s, "00·发刊词", "你即将拥有的「家庭操作系统」", "3 大原创模型 + 13 套原创工具")

    models = [
        ("三锚模型", "目标-边界-反馈", "锚定目标：孩子该练的 AI 时代增值能力\n锚定边界：AI 使用的场景/时间/内容/隐私边界\n锚定反馈：AI 输出评估 + 孩子成长评估", C_PRIMARY),
        ("家庭 Prompt 工作流", "5 大场景", "作业辅导 · 兴趣探索 · 问题解答\n创意生成 · 决策辅助", C_ACCENT),
        ("超越竞争能力图谱", "3 大稀缺能力", "问题定义力 · 情感连接力 · 跨域整合力", C_SECONDARY),
    ]
    box_w = Inches(2.85)
    for i, (name, sub, desc, color) in enumerate(models):
        x = Inches(0.5) + (box_w + Inches(0.15)) * i
        # 顶部
        add_rect(s, x, Inches(1.7), box_w, Inches(0.5), color)
        add_text_box(s, x, Inches(1.78), box_w, Inches(0.4), name, size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 主体
        add_rect(s, x, Inches(2.2), box_w, Inches(2.9), C_LIGHT)
        add_text_box(s, x + Inches(0.1), Inches(2.3), box_w - Inches(0.2), Inches(0.4), sub, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        # 描述
        lines = desc.split("\n")
        for j, line in enumerate(lines):
            add_text_box(s, x + Inches(0.15), Inches(2.85) + Inches(0.4) * j, box_w - Inches(0.3), Inches(0.4),
                         "• " + line if j == 0 else line, size=11, color=C_TEXT)

    # 底部金句
    add_rounded_rect(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.0), C_PRIMARY, radius=0.0)
    add_text_box(s, Inches(0.5), Inches(5.18), Inches(9), Inches(0.3),
                 "家长不是孩子的教练，家长是孩子和未来之间的翻译官。", size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_07():
    """发刊词 5: 学习路径 + 13 讲承诺"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 7, "发刊词")
    section_header(s, "00·发刊词", "13 讲的学习路径 & 你能带走什么", "从认知到行动，从单点到系统，从工具到方法")

    # 学习路径
    add_text_box(s, Inches(0.5), Inches(1.7), Inches(4.2), Inches(0.3), "学习路径", size=14, bold=True, color=C_PRIMARY)
    paths = [
        ("破除焦虑", "发刊词 + 先导课", "建立坐标系", C_ACCENT),
        ("看见未来", "问题 1-4", "能力 / 自学 / 真实问题 / 判断力", C_PRIMARY),
        ("重构家庭", "问题 5-8", "基础 / 父母角色 / 边界 / 情感", C_SECONDARY),
        ("长期主义", "问题 9-12", "兴趣 / 评估 / 路线 / 避坑", C_DARK),
        ("回到行动", "结刊词", "30 天启动", C_RED),
    ]
    for i, (name, sub, desc, color) in enumerate(paths):
        y = Inches(2.05) + Inches(0.6) * i
        add_oval(s, Inches(0.5), y + Inches(0.1), Inches(0.35), Inches(0.35), color)
        add_text_box(s, Inches(1.0), y, Inches(1.4), Inches(0.3), name, size=12, bold=True, color=color)
        add_text_box(s, Inches(2.4), y, Inches(1.3), Inches(0.3), sub, size=10, color=C_TEXT_LIGHT)
        add_text_box(s, Inches(3.7), y, Inches(1.2), Inches(0.3), desc, size=10, color=C_TEXT)

    # 承诺
    add_text_box(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.3), "13 讲之后，你能带走", size=14, bold=True, color=C_PRIMARY)
    promises = [
        "1 份家庭能力评估表（增值 + 警惕能力）",
        "1 套家庭 Prompt 工作流（5 大场景）",
        "1 份 AI 输出三审表（已使用 ≥ 5 次）",
        "1 份 AI 家庭使用协议（家长 + 孩子签署）",
        "1 张家庭真实问题池（10 个真实问题）",
        "1 张亲子共学记录表（首周已记录）",
        "1 份 30 天家庭 AI 行动日志（≥ 7 天）",
        "1 张 3 年家庭 AI 教育路线图",
        "1 份 AI 家庭教育误区自查表",
        "1 张课程结业证书",
    ]
    for i, t in enumerate(promises):
        y = Inches(2.05) + Inches(0.3) * i
        add_text_box(s, Inches(5.2), y, Inches(0.2), Inches(0.25), "✓", size=11, bold=True, color=C_ACCENT)
        add_text_box(s, Inches(5.4), y, Inches(4.2), Inches(0.25), t, size=10, color=C_TEXT)


# =====================================================
# 先导课 (5-6 页)
# =====================================================
def slide_08():
    """先导课 1: 标题 + 故事"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 8, "先导课")
    section_header(s, "00·先导课", "孩子的未来会好吗？", "—— AI 取代的不是孩子，是「标准化能力」")

    # 故事框
    add_rounded_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(1.7), C_LIGHT, radius=0.05)
    add_text_box(s, Inches(0.7), Inches(1.85), Inches(8.6), Inches(0.3), "一个真实的「AI 时代家长会」", size=12, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(2.2), Inches(8.6), Inches(0.4), "去年秋天我在一所重点小学做分享，会后一个妈妈拉着我聊了快一个小时。", size=12, color=C_TEXT)
    add_text_box(s, Inches(0.7), Inches(2.55), Inches(8.6), Inches(0.4), "她儿子刚上五年级，数学还行，语文一般，英语垫底。她最大的焦虑是：", size=12, color=C_TEXT)
    add_text_box(s, Inches(0.7), Inches(2.9), Inches(8.6), Inches(0.4), "「罗老师，AI 都能写作文了，我家孩子语文还考不及格，长大还能干什么？」", size=13, bold=True, color=C_PRIMARY)
    add_text_box(s, Inches(0.7), Inches(3.25), Inches(8.6), Inches(0.2), "她的回答是——「陪人说话的那种。心理咨询师、婚姻家庭咨询师……是不是？」", size=11, color=C_TEXT_LIGHT)

    # 金句
    add_rect(s, Inches(0.5), Inches(3.6), Inches(9), Inches(0.5), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(3.68), Inches(9), Inches(0.4), "AI 越能「做事」，「人味」就越贵。",
                 size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # 三个能力
    add_text_box(s, Inches(0.5), Inches(4.3), Inches(9), Inches(0.3), "稀缺性决定价值。未来最稀缺的三件事：", size=12, color=C_TEXT)
    for i, t in enumerate(["定义「什么是真正要解决的问题」", "把 3 个不相关领域的知识拼成新东西", "在情绪崩溃时给出有温度的回应"]):
        x = Inches(0.5) + Inches(3.0) * i
        add_rounded_rect(s, x, Inches(4.65), Inches(2.85), Inches(0.65), C_LIGHT, radius=0.1)
        add_text_box(s, x + Inches(0.1), Inches(4.7), Inches(2.65), Inches(0.55), t, size=10, color=C_PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_09():
    """先导课 2: 能力贬值/增值清单（数据）"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 9, "先导课")
    section_header(s, "00·先导课", "AI 时代能力贬值/增值清单", "—— 10 项能力对照表，看清家庭该往哪用力")

    # 两列
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.3), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_RED)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "⬇ 正在贬值（AI 强项）", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    decline = [
        "标准化记忆（背单词、记公式）",
        "基础计算（加减乘除）",
        "机械翻译（直译、单词对应）",
        "套路化写作（八股文、模板作文）",
        "信息检索（基础搜索、资料整理）",
    ]
    for i, t in enumerate(decline):
        y = Inches(2.3) + Inches(0.5) * i
        add_text_box(s, Inches(0.7), y, Inches(0.3), Inches(0.4), "✗", size=14, bold=True, color=C_RED)
        add_text_box(s, Inches(1.0), y, Inches(3.8), Inches(0.4), t, size=12, color=C_TEXT)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.3), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_GREEN)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "⬆ 正在增值（AI 替代不了）", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rise = [
        "问题定义力（找到真问题）",
        "跨域整合力（拼接不同领域）",
        "判断力（识别 AI 输出的对错）",
        "自学能力（用 AI 解决真实问题）",
        "情感连接力（读懂人、温暖他人）",
    ]
    for i, t in enumerate(rise):
        y = Inches(2.3) + Inches(0.5) * i
        add_text_box(s, Inches(5.3), y, Inches(0.3), Inches(0.4), "✓", size=14, bold=True, color=C_GREEN)
        add_text_box(s, Inches(5.6), y, Inches(3.8), Inches(0.4), t, size=12, color=C_TEXT)

    # 金句
    add_rounded_rect(s, Inches(0.5), Inches(5.1), Inches(9), Inches(0.4), C_PRIMARY, radius=0.2)
    add_text_box(s, Inches(0.5), Inches(5.15), Inches(9), Inches(0.3), "AI 时代的教育，不是「学更多」，而是「重新定义什么是重要」。",
                 size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_10():
    """先导课 3: 超越竞争能力图谱"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 10, "先导课")
    section_header(s, "00·先导课", "超越竞争能力图谱", "—— 重新定义竞争的维度")

    # 三角形图谱
    cx, cy = Inches(5.0), Inches(3.5)
    # 三个圆
    add_oval(s, cx - Inches(2.5), cy - Inches(0.9), Inches(2.0), Inches(2.0), C_PRIMARY)
    add_text_box(s, cx - Inches(2.5), cy - Inches(0.7), Inches(2.0), Inches(0.4), "问题定义力", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, cx - Inches(2.5), cy - Inches(0.3), Inches(2.0), Inches(0.3), "稀缺 1", size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, cx - Inches(2.5), cy + Inches(0.1), Inches(2.0), Inches(0.9), "知道「什么才是真正要解决的问题」", size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_oval(s, cx + Inches(0.5), cy - Inches(0.9), Inches(2.0), Inches(2.0), C_ACCENT)
    add_text_box(s, cx + Inches(0.5), cy - Inches(0.7), Inches(2.0), Inches(0.4), "情感连接力", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, cx + Inches(0.5), cy - Inches(0.3), Inches(2.0), Inches(0.3), "稀缺 2", size=10, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(s, cx + Inches(0.5), cy + Inches(0.1), Inches(2.0), Inches(0.9), "和他人建立真实、有温度的关系", size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_oval(s, cx - Inches(1.0), cy + Inches(0.5), Inches(2.0), Inches(2.0), C_SECONDARY)
    add_text_box(s, cx - Inches(1.0), cy + Inches(0.7), Inches(2.0), Inches(0.4), "跨域整合力", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, cx - Inches(1.0), cy + Inches(1.1), Inches(2.0), Inches(0.3), "稀缺 3", size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, cx - Inches(1.0), cy + Inches(1.5), Inches(2.0), Inches(0.9), "把不同领域的知识拼成新东西", size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # 金句
    add_text_box(s, Inches(0.5), Inches(1.7), Inches(3.5), Inches(0.5), "本课程原创", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.5), Inches(1.95), Inches(3.5), Inches(0.5), "能力图谱", size=22, bold=True, color=C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(2.5), Inches(3.5), Inches(1.5), "你抢的赛道，决定孩子的未来。\n\n领导半步，吃尽红利；\n领先一步，枪打出头鸟；\n落后半步，别人牵牛我拔桩。", size=11, color=C_TEXT)

    # 底部
    add_text_box(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.3),
                 "这 3 个能力，AI 替代不了，也不依赖某个新工具。", size=12, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)


def slide_11():
    """先导课 4: 两个家庭对比"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 11, "先导课")
    section_header(s, "00·先导课", "两个家庭的真实对比", "—— 一个在「追新」，一个在「夯实基础」")

    # 左家庭
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.5), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_RED)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "A 家庭：追新型", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.7), [
        "孩子 8 岁，三年级",
        "半年报 4 个班：编程、钢琴、围棋、英文戏剧",
        "每次报班的理由都是「听说 AI 时代需要这个」",
        "半年下来，孩子每样都学一点、每样都不深",
        "账单 3 万多，孩子疲惫",
        "家长的困惑：「这些真的是 AI 时代孩子需要的吗？」",
    ], size=11, color=C_TEXT, line_spacing=1.3)

    # 右家庭
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.5), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_GREEN)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "B 家庭：夯实型", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.7), [
        "孩子 9 岁，三年级",
        "保留 2 个核心兴趣：阅读 + 编程",
        "每月用 AI 完成 1 个真实问题",
        "每周 1 次亲子共学记录",
        "每年更新 1 次「家庭能力评估表」",
        "孩子的状态：好奇心强、主动学习、有作品",
    ], size=11, color=C_TEXT, line_spacing=1.3)

    # 底部金句
    add_rounded_rect(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.0), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(5.18), Inches(9), Inches(0.3),
                 "从「焦虑地追」到「系统地养」—— 同样投入，不同结果。", size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_12():
    """先导课 5: 心态建立 + 工具预告"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 12, "先导课")
    section_header(s, "00·先导课", "从「焦虑地追」到「系统地养」", "—— 13 讲用 3 大模型 + 13 套工具，陪你走完这条路")

    # 左右对比
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_RED)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "✗ 焦虑地追", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.5), [
        "看到新概念 → 怕错过",
        "每个新概念 → 都跟一遍",
        "孩子的童年 = 追新实验",
        "账单多、能力少、亲子关系紧",
        "焦虑持续累积",
    ], size=12, color=C_TEXT, line_spacing=1.4)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_GREEN)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "✓ 系统地养", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.5), [
        "明确目标 → 三锚模型",
        "夯实基础 → 5 项基础能力清单",
        "用对工具 → 家庭 Prompt 工作流",
        "持续评估 → 成长评估表 + 路线图",
        "焦虑逐渐消失，方向越来越清晰",
    ], size=12, color=C_TEXT, line_spacing=1.4)

    # 工具预告
    add_text_box(s, Inches(0.5), Inches(4.85), Inches(9), Inches(0.3), "本讲产出：", size=12, bold=True, color=C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(5.15), Inches(9), Inches(0.3),
                 "完成「AI 时代能力贬值/增值清单」自评 · 识别 3 项增值能力 + 3 项需警惕能力", size=11, color=C_TEXT)


def slide_13():
    """先导课 6: 担忧澄清 + 行动入口"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 13, "先导课")
    section_header(s, "00·先导课", "家长最常问的 3 个担忧", "—— 在进入问题一之前，先把心态校准")

    worries = [
        ("担忧 1", "我不焦虑，孩子落后了怎么办？", "解药不是「想开点」，是「有行动系统」。本课给你 8 套工具、3 大模型、30 天清单。", C_RED),
        ("担忧 2", "我家孩子还小（3-6 岁），是不是太早？", "本课覆盖 3-18 岁。最早可以从「家庭 AI 使用协议」和「亲子共学」开始。", C_ACCENT),
        ("担忧 3", "我完全不懂 AI，能学会吗？", "不需要你成为 AI 专家。本课讲「怎么陪孩子用」，不讲「你自己怎么精通」。", C_SECONDARY),
    ]
    for i, (label, q, a, color) in enumerate(worries):
        y = Inches(1.7) + Inches(1.05) * i
        add_rect(s, Inches(0.5), y, Inches(0.6), Inches(0.9), color)
        add_text_box(s, Inches(0.5), y + Inches(0.2), Inches(0.6), Inches(0.5), label[-1], size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.2), y + Inches(0.05), Inches(8.3), Inches(0.4), q, size=13, bold=True, color=color)
        add_text_box(s, Inches(1.2), y + Inches(0.45), Inches(8.3), Inches(0.5), a, size=11, color=C_TEXT)

    # 入口金句
    add_rounded_rect(s, Inches(0.5), Inches(4.95), Inches(9), Inches(0.45), C_PRIMARY, radius=0.2)
    add_text_box(s, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4), "先导课结束 → 问题一：未来社会会更卷吗？",
                 size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# =====================================================
# 问题一：未来社会会更卷吗 (8-10 页) — 8 页
# =====================================================
def slide_14():
    """问题一 1: 标题"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 14, "问题一")
    section_header(s, "问题 01", "未来社会会更卷吗？", "—— 超越竞争，重新定义竞争维度")

    # 大金句区
    add_rect(s, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.95), Inches(9), Inches(0.4), "本讲核心判断", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.35), Inches(9), Inches(0.8), "「卷」不等于「竞争」", size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.95), Inches(9), Inches(0.3), "真卷 vs 假卷：选对赛道，比跑得快更重要。", size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # 三栏
    add_text_box(s, Inches(0.5), Inches(3.55), Inches(9), Inches(0.3), "本讲将解决：", size=13, bold=True, color=C_PRIMARY)
    items = [
        ("卷和竞争的本质区别", "家长要分清什么在卷、卷的方向对不对"),
        ("AI 时代真正的竞争维度", "超越竞争 = 重新定义竞争维度"),
        ("伪竞争 vs 真竞争", "标准化能力比拼是「假卷」"),
        ("怎么识别真假卷", "真假卷的 3 个判断标准"),
    ]
    for i, (t, d) in enumerate(items):
        x = Inches(0.5) + Inches(2.25) * (i % 2)
        y = Inches(3.95) + Inches(0.55) * (i // 2)
        add_oval(s, x, y, Inches(0.35), Inches(0.35), C_ACCENT)
        add_text_box(s, x + Inches(0.05), y + Inches(0.05), Inches(0.3), Inches(0.3), str(i+1), size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.45), y, Inches(1.75), Inches(0.3), t, size=11, bold=True, color=C_PRIMARY)
        add_text_box(s, x + Inches(0.45), y + Inches(0.27), Inches(1.75), Inches(0.3), d, size=9, color=C_TEXT_LIGHT)


def slide_15():
    """问题一 2: 故事页 - 王女士的鸡娃焦虑"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 15, "问题一")
    section_header(s, "问题 01", "从「鸡娃」到「AI 鸡娃」", "—— 家长焦虑如何被放大")

    add_rounded_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(2.6), C_LIGHT, radius=0.05)
    add_text_box(s, Inches(0.7), Inches(1.85), Inches(8.6), Inches(0.3), "真实案例", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(2.15), Inches(8.6), Inches(0.4), "王女士的孩子今年 8 岁，三年级。", size=14, color=C_TEXT)
    add_paragraphs(s, Inches(0.7), Inches(2.6), Inches(8.6), Inches(1.6), [
        "半年里她给孩子报了 4 个兴趣班：编程、钢琴、围棋、英文戏剧",
        "每一次报班的理由都是「听说 AI 时代需要这个」",
        "半年下来，孩子每个都学了一点，但都学得不深",
        "她看着孩子疲惫的脸，又看着账单上 3 万多的兴趣班支出",
        "第一次怀疑：这些，真的是 AI 时代孩子需要的吗？",
    ], size=12, color=C_TEXT, line_spacing=1.4)

    # 本质
    add_rect(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.7), C_PRIMARY)
    add_text_box(s, Inches(0.7), Inches(4.58), Inches(8.6), Inches(0.3), "焦虑的本质：信息过载，决策瘫痪。", size=14, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(4.88), Inches(8.6), Inches(0.3),
                 "每一个新概念都跟一遍，孩子的童年就成了「追新实验」。", size=11, color=C_LIGHT)


def slide_16():
    """问题一 3: 概念厘清 - 真卷 vs 假卷"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 16, "问题一")
    section_header(s, "问题 01", "真卷 vs 假卷：3 个判断标准", "—— 别再用战术上的勤奋，掩盖战略上的懒惰")

    # 两列
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_RED)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "✗ 假卷（伪竞争）", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.6), [
        "标准 1：比拼的是 AI 也能做",
        "（记忆、计算、机械翻译）",
        "",
        "标准 2：方向在「数量」而非「维度」",
        "（刷题量、单词量、证书数）",
        "",
        "标准 3：家长焦虑主导，孩子被动",
        "（家长说「别人都在学」）",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_GREEN)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "✓ 真卷（真竞争）", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.6), [
        "标准 1：比拼的是 AI 做不了",
        "（判断、创意、情感连接）",
        "",
        "标准 2：方向在「维度」而非「数量」",
        "（深度、独特性、跨域整合）",
        "",
        "标准 3：孩子主动，家长引导",
        "（孩子说「我想试试」）",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    # 金句
    add_rounded_rect(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.0), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(5.18), Inches(9), Inches(0.3),
                 "超越竞争的本质，不是「不竞争」，而是「重新定义竞争的维度」。", size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_17():
    """问题一 4: 李笑来观点 + 我的独到见解"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 17, "问题一")
    section_header(s, "问题 01", "李笑来观点 vs 本课独到见解", "—— 我们往前走了一步")

    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_SECONDARY)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "李笑来老师的观点", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.6), [
        "「未来社会不会再卷了」",
        "",
        "理由：",
        "AI 让「会做事」不再稀缺",
        "「会做人」反而成了稀缺品",
        "",
        "判断：",
        "不要再卷「数量」和「速度」",
        "那是 AI 的强项",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_ACCENT)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "本课往前走一步", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.6), [
        "「超越竞争」之后呢？",
        "",
        "操作路径：",
        "• 重新定义竞争维度",
        "• 3 大稀缺能力：问题定义 / 情感 / 跨域",
        "• 用 8 套工具落到家庭场景",
        "",
        "承诺：",
        "今晚开始，知道做什么",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    add_text_box(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.3),
                 "李笑来给方向，本课给路径。认知是起点，行动才是终点。", size=12, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)


def slide_18():
    """问题一 5: 模型页 - 超越竞争能力图谱（详细）"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 18, "问题一")
    section_header(s, "问题 01", "原创模型 1：超越竞争能力图谱", "—— 3 大核心能力，每个能力 3 个评估等级")

    # 能力卡片
    abilities = [
        ("问题定义力", "稀缺 1", "知道「什么才是真正要解决的问题」", "3-18 岁逐步培养", C_PRIMARY,
         "启蒙：能问出好问题 / 熟练：能区分真假问题 / 创新：能重新定义问题"),
        ("情感连接力", "稀缺 2", "和他人建立真实、有温度的关系", "AI 替代不了", C_ACCENT,
         "启蒙：识别自己和他人的情绪 / 熟练：能用语言回应情绪 / 创新：能安抚、激励、领导他人"),
        ("跨域整合力", "稀缺 3", "把不同领域的知识拼成新东西", "AI 时代最稀缺", C_SECONDARY,
         "启蒙：知道不同领域的存在 / 熟练：能把 2 个领域连起来 / 创新：能创造新领域交叉"),
    ]
    for i, (name, label, desc, sub, color, level) in enumerate(abilities):
        y = Inches(1.7) + Inches(1.1) * i
        add_rect(s, Inches(0.5), y, Inches(1.6), Inches(1.0), color)
        add_text_box(s, Inches(0.5), y + Inches(0.15), Inches(1.6), Inches(0.35), name, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(0.5), y + Inches(0.5), Inches(1.6), Inches(0.25), label, size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(0.5), y + Inches(0.7), Inches(1.6), Inches(0.25), sub, size=9, color=C_LIGHT, align=PP_ALIGN.CENTER)
        # 描述
        add_text_box(s, Inches(2.2), y + Inches(0.1), Inches(7.3), Inches(0.35), desc, size=12, bold=True, color=C_TEXT)
        add_text_box(s, Inches(2.2), y + Inches(0.45), Inches(7.3), Inches(0.55), "评估等级：" + level, size=10, color=C_TEXT_LIGHT)


def slide_19():
    """问题一 6: 工具页 - 家庭能力评估表"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 19, "问题一")
    section_header(s, "问题 01", "工具 1：家庭能力评估表", "—— 10 项能力对照，每年更新 1 次")

    add_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(0.5), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(9), Inches(0.4), "使用说明", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_paragraphs(s, Inches(0.7), Inches(2.35), Inches(8.6), Inches(2.5), [
        "1. 家长 + 孩子共同填写，每项 1-5 分",
        "2. 每年 1 次（生日或年初），识别变化趋势",
        "3. 对照 AI 时代「贬值/增值清单」，识别孩子的 3 项增值能力 + 3 项需警惕能力",
        "4. 把 3 项增值能力作为下一年度重点培养方向",
        "5. 把 3 项需警惕能力转化为具体训练场景",
    ], size=12, color=C_TEXT, line_spacing=1.4)

    # 样本
    add_rect(s, Inches(0.5), Inches(4.55), Inches(9), Inches(0.7), C_LIGHT)
    add_text_box(s, Inches(0.7), Inches(4.62), Inches(8.6), Inches(0.3), "样例（节选）", size=11, bold=True, color=C_PRIMARY)
    add_text_box(s, Inches(0.7), Inches(4.92), Inches(8.6), Inches(0.3),
                 "问题定义力 4 · 跨域整合力 3 · 情感连接力 5（增值）/ 标准记忆 2 · 基础计算 3 · 套路写作 2（警惕）",
                 size=10, color=C_TEXT)

    # 金句
    add_text_box(s, Inches(0.5), Inches(5.3), Inches(9), Inches(0.3),
                 "家庭能力评估表 = 家庭教育的「体检表」。", size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_20():
    """问题一 7: 实操步骤"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 20, "问题一")
    section_header(s, "问题 01", "本讲 4 步实操：今晚就能用", "—— 家庭能力评估表填写流程")

    steps = [
        ("Step 1", "打印工具", "从「06_配套工具/家庭能力评估表」打印 A4 表格", C_PRIMARY),
        ("Step 2", "家长 + 孩子共同填写", "10 项能力，每项 1-5 分，独立打分后对比", C_SECONDARY),
        ("Step 3", "对照清单", "10 项能力 × 增值/贬值分类 = 3 项增值 + 3 项需警惕", C_ACCENT),
        ("Step 4", "明确方向", "把 3 项增值列为下一年重点培养方向", C_DARK),
    ]
    for i, (label, t, d, color) in enumerate(steps):
        x = Inches(0.5) + Inches(2.3) * i
        # 数字
        add_oval(s, x + Inches(0.85), Inches(1.85), Inches(0.6), Inches(0.6), color)
        add_text_box(s, x + Inches(0.85), Inches(1.95), Inches(0.6), Inches(0.4), str(i+1), size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 标题
        add_text_box(s, x, Inches(2.55), Inches(2.2), Inches(0.3), t, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        # 描述
        add_text_box(s, x, Inches(2.85), Inches(2.2), Inches(1.5), d, size=10, color=C_TEXT, align=PP_ALIGN.CENTER)

    # 箭头
    for i in range(3):
        add_arrow(s, Inches(2.5) + Inches(2.3) * i, Inches(2.1), Inches(2.8) + Inches(2.3) * i, Inches(2.1), C_ACCENT, 0.15)

    # 底部
    add_rounded_rect(s, Inches(0.5), Inches(4.7), Inches(9), Inches(0.65), C_PRIMARY, radius=0.1)
    add_text_box(s, Inches(0.7), Inches(4.78), Inches(8.6), Inches(0.3), "本讲产出：", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(5.05), Inches(8.6), Inches(0.3),
                 "完成「家庭能力评估表」首次填写 · 识别 3 项增值 + 3 项需警惕能力 · 明确家庭未来 3 年竞争维度", size=11, color=C_WHITE)


def slide_21():
    """问题一 8: 家长担忧澄清 + 总结"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 21, "问题一")
    section_header(s, "问题 01", "家长担忧澄清 & 一句话总结", "")

    worries = [
        ("Q", "我家孩子已经在卷奥数、英语了，现在停止是不是亏了？",
         "A", "不是停止，是把「卷的维度」从数量调整到能力。本课不反对学习，主张把有限时间用在 AI 替代不了的能力上。", C_RED),
        ("Q", "3 大能力听起来很好，但怎么知道孩子具体在哪一档？",
         "A", "用「家庭能力评估表」打分。10 项能力 × 3 个等级（启蒙/熟练/创新），看完一目了然。", C_ACCENT),
    ]
    for i, (q_label, q, a_label, a, color) in enumerate(worries):
        y = Inches(1.7) + Inches(1.4) * i
        add_oval(s, Inches(0.5), y, Inches(0.5), Inches(0.5), color)
        add_text_box(s, Inches(0.5), y + Inches(0.1), Inches(0.5), Inches(0.3), q_label, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.1), y, Inches(8.4), Inches(0.4), q, size=12, bold=True, color=C_PRIMARY)
        add_text_box(s, Inches(1.1), y + Inches(0.45), Inches(8.4), Inches(0.7), a, size=11, color=C_TEXT)

    # 总结
    add_rounded_rect(s, Inches(0.5), Inches(4.6), Inches(9), Inches(0.75), C_PRIMARY, radius=0.1)
    add_text_box(s, Inches(0.7), Inches(4.65), Inches(8.6), Inches(0.3), "本讲一句话：", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(4.95), Inches(8.6), Inches(0.4),
                 "你抢的赛道，决定孩子的未来。AI 时代，重新定义竞争维度。", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# =====================================================
# 问题二：自学能力 (8 页) - 重复结构生成
# =====================================================
def slide_22():
    """问题二 1: 标题"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 22, "问题二")
    section_header(s, "问题 02", "怎样培养孩子的自学能力？", "—— 以生产为导向，让孩子用 AI 产出作品")

    add_rect(s, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.95), Inches(9), Inches(0.4), "本讲核心方法", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.35), Inches(9), Inches(0.8), "「以生产为导向」的自学", size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.95), Inches(9), Inches(0.3), "倒过来：先有要产出的作品，再倒推需要学什么。", size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(3.55), Inches(9), Inches(0.3), "本讲将解决：", size=13, bold=True, color=C_PRIMARY)
    items = [
        "自学能力 = 什么？", "怎么让孩子主动学习？",
        "AI 时代自学能力的变化", "「以生产为导向」如何落地？",
    ]
    for i, t in enumerate(items):
        x = Inches(0.5) + Inches(2.25) * (i % 2)
        y = Inches(3.95) + Inches(0.55) * (i // 2)
        add_oval(s, x, y, Inches(0.35), Inches(0.35), C_ACCENT)
        add_text_box(s, x + Inches(0.05), y + Inches(0.05), Inches(0.3), Inches(0.3), str(i+1), size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.45), y, Inches(1.75), Inches(0.3), t, size=11, bold=True, color=C_PRIMARY)


def slide_23():
    """问题二 2: 故事页 - 自学悖论"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 23, "问题二")
    section_header(s, "问题 02", "家长想要的自学 vs 家长造成的依赖", "—— 一个被忽视的悖论")

    add_rounded_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(2.6), C_LIGHT, radius=0.05)
    add_text_box(s, Inches(0.7), Inches(1.85), Inches(8.6), Inches(0.3), "真实案例", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(2.15), Inches(8.6), Inches(0.4), "赵女士的孩子今年 11 岁，五年级。", size=14, color=C_TEXT)
    add_paragraphs(s, Inches(0.7), Inches(2.6), Inches(8.6), Inches(1.6), [
        "她嘴上说「我希望孩子能自主学习」",
        "但每天下班第一件事就是检查作业、追问今天学了什么",
        "半年下来，她疲惫不堪",
        "孩子也越来越不愿意主动学习",
        "结果是：家长越追，孩子越不主动",
    ], size=12, color=C_TEXT, line_spacing=1.4)

    add_rect(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.7), C_PRIMARY)
    add_text_box(s, Inches(0.7), Inches(4.58), Inches(8.6), Inches(0.3), "悖论的本质：", size=12, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(4.88), Inches(8.6), Inches(0.3),
                 "没人追是结果，前提是家长要先放手。", size=14, color=C_WHITE)


def slide_24():
    """问题二 3: 李笑来观点 + 本课见解"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 24, "问题二")
    section_header(s, "问题 02", "李笑来观点 vs 本课独到见解", "")

    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_SECONDARY)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "李笑来老师的观点", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.6), [
        "「有问题自己来」",
        "",
        "「创造机会让孩子自己思考」",
        "",
        "判断：",
        "自学能力的核心是「主动」",
        "家长要提供机会，不替孩子思考",
    ], size=12, color=C_TEXT, line_spacing=1.4)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_ACCENT)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "本课往前走一步", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.6), [
        "「自学能力 = 用 AI 解决问题的能力」",
        "",
        "操作路径：",
        "• 以生产为导向",
        "• 5 个 Prompt 模板",
        "• 家长 3 个不：不教/不催/不评判",
    ], size=12, color=C_TEXT, line_spacing=1.4)


def slide_25():
    """问题二 4: 模型 - 以生产为导向的自学"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 25, "问题二")
    section_header(s, "问题 02", "原创方法 1：「以生产为导向」的自学", "—— 倒过来学，比正过去学更高效")

    # 对比图
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(2.5), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.4), C_RED)
    add_text_box(s, Inches(0.5), Inches(1.76), Inches(4.4), Inches(0.3), "传统自学", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(2.2), Inches(4.0), Inches(0.3), "看书 → 做题 → 考试", size=14, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(2.6), Inches(4.0), Inches(0.3), "↓", size=20, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(2.9), Inches(4.0), Inches(0.3), "学完就忘", size=14, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(3.4), Inches(4.0), Inches(0.7), "问题：学的知识没有用\n学了不知道用在哪", size=11, color=C_TEXT)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(2.5), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.4), C_GREEN)
    add_text_box(s, Inches(5.1), Inches(1.76), Inches(4.4), Inches(0.3), "以生产为导向", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.3), Inches(2.2), Inches(4.0), Inches(0.3), "作品 → 知识 → 技能", size=14, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.3), Inches(2.6), Inches(4.0), Inches(0.3), "↓", size=20, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.3), Inches(2.9), Inches(4.0), Inches(0.3), "学完就用", size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.3), Inches(3.4), Inches(4.0), Inches(0.7), "优势：学了马上用\n反馈快、能改进", size=11, color=C_TEXT)

    add_rounded_rect(s, Inches(0.5), Inches(4.4), Inches(9), Inches(0.8), C_PRIMARY, radius=0.1)
    add_text_box(s, Inches(0.7), Inches(4.5), Inches(8.6), Inches(0.3), "AI 让「以生产为导向」成为可能", size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(4.85), Inches(8.6), Inches(0.3), "任何知识盲区，AI 都能补上。孩子不用先学再做，而是边做边学。",
                 size=12, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_26():
    """问题二 5: 案例页 - 3 个孩子"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 26, "问题二")
    section_header(s, "问题 02", "案例：3 个孩子用 AI 自学的真实过程", "")

    cases = [
        ("写作", "小林 12 岁", "想给班级写一篇元旦晚会的开场白\n→ 用 AI 拆解结构\n→ 写出初稿，AI 提建议\n→ 自己改 3 遍\n→ 老师在全班表扬", C_PRIMARY),
        ("编程", "小宇 10 岁", "想做一个小游戏\n→ 用 AI 生成基础代码\n→ 自己设计关卡逻辑\n→ 调试 5 次\n→ 完成作品，同学都爱玩", C_ACCENT),
        ("视频", "小溪 14 岁", "想做校园纪录短片\n→ 用 AI 写脚本\n→ 自己拍摄 + 剪辑\n→ AI 配字幕 + 背景音乐\n→ 视频在学校公众号发布", C_SECONDARY),
    ]
    for i, (label, name, desc, color) in enumerate(cases):
        x = Inches(0.5) + Inches(3.0) * i
        add_rect(s, x, Inches(1.7), Inches(2.85), Inches(3.0), C_LIGHT)
        add_rect(s, x, Inches(1.7), Inches(2.85), Inches(0.5), color)
        add_text_box(s, x, Inches(1.78), Inches(2.85), Inches(0.4), label, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.1), Inches(2.3), Inches(2.65), Inches(0.3), name, size=11, bold=True, color=color)
        lines = desc.split("\n")
        for j, line in enumerate(lines):
            add_text_box(s, x + Inches(0.15), Inches(2.65) + Inches(0.3) * j, Inches(2.55), Inches(0.3), line, size=10, color=C_TEXT)

    add_text_box(s, Inches(0.5), Inches(4.85), Inches(9), Inches(0.3),
                 "共同点：作品驱动 + AI 辅助 + 孩子主导。", size=12, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)


def slide_27():
    """问题二 6: 工具 - Prompt 工作流 - 自学场景"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 27, "问题二")
    section_header(s, "问题 02", "工具 2：家庭 Prompt 工作流 - 自学场景", "—— 5 个 Prompt 模板直接套用")

    prompts = [
        ("Prompt 1", "问题拆解", "「我要做 [作品]，需要先解决 [问题]。请帮我拆成 3-5 个小问题。」", C_PRIMARY),
        ("Prompt 2", "知识检索", "「请用初中生能懂的话解释 [概念]，给我 3 个生活中的例子。」", C_SECONDARY),
        ("Prompt 3", "案例生成", "「请给我 3 个 [主题] 的真实案例，分别说明它们的关键差异。」", C_ACCENT),
        ("Prompt 4", "反馈优化", "「这是我写的 [内容]。请从结构、逻辑、表达 3 个维度给我反馈。」", C_DARK),
        ("Prompt 5", "作品润色", "「这是我的初稿。请保留我的核心想法，帮我把语言改得更生动。」", C_RED),
    ]
    for i, (label, t, p, color) in enumerate(prompts):
        y = Inches(1.7) + Inches(0.7) * i
        add_rect(s, Inches(0.5), y, Inches(1.0), Inches(0.55), color)
        add_text_box(s, Inches(0.5), y + Inches(0.13), Inches(1.0), Inches(0.3), label[-1], size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.6), y, Inches(1.5), Inches(0.55), t, size=12, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Inches(3.2), y, Inches(6.3), Inches(0.55), p, size=10, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)


def slide_28():
    """问题二 7: 实操步骤"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 28, "问题二")
    section_header(s, "问题 02", "本讲 4 步实操：为孩子设计第一个自学项目", "")

    steps = [
        ("Step 1", "选一个作品目标", "和孩子商量：想做什么？（文章/视频/代码/海报）", C_PRIMARY),
        ("Step 2", "拆解为小问题", "用 Prompt 1（问题拆解）把作品拆成 3-5 个小问题", C_SECONDARY),
        ("Step 3", "用 AI 补盲区", "对孩子卡住的地方，用 Prompt 2-3 帮他补知识", C_ACCENT),
        ("Step 4", "孩子主导完成", "家长 3 个不：不教/不催/不评判，只观察 + 反馈", C_DARK),
    ]
    for i, (label, t, d, color) in enumerate(steps):
        x = Inches(0.5) + Inches(2.3) * i
        add_oval(s, x + Inches(0.85), Inches(1.85), Inches(0.6), Inches(0.6), color)
        add_text_box(s, x + Inches(0.85), Inches(1.95), Inches(0.6), Inches(0.4), str(i+1), size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(2.55), Inches(2.2), Inches(0.3), t, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(2.85), Inches(2.2), Inches(1.5), d, size=10, color=C_TEXT, align=PP_ALIGN.CENTER)

    add_rounded_rect(s, Inches(0.5), Inches(4.7), Inches(9), Inches(0.65), C_PRIMARY, radius=0.1)
    add_text_box(s, Inches(0.7), Inches(4.78), Inches(8.6), Inches(0.3), "本讲产出：", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(5.05), Inches(8.6), Inches(0.3),
                 "完成「家庭 Prompt 工作流 - 自学场景」初步搭建 · 设计第一个以生产为导向的自学项目 · 明确家长 3 个不", size=11, color=C_WHITE)


def slide_29():
    """问题二 8: 担忧澄清 + 总结"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 29, "问题二")
    section_header(s, "问题 02", "家长担忧澄清 & 一句话总结", "")

    worries = [
        ("Q", "孩子用 AI 写作业，作业质量怎么保证？",
         "A", "本课强调「以生产为导向」，作品代表孩子综合能力。家长看作品本身，不看是否「用了 AI」。", C_RED),
        ("Q", "AI 给出错误信息怎么办？",
         "A", "用「AI 输出三审表」训练孩子判断。问题四会详细讲。", C_ACCENT),
    ]
    for i, (q_label, q, a_label, a, color) in enumerate(worries):
        y = Inches(1.7) + Inches(1.4) * i
        add_oval(s, Inches(0.5), y, Inches(0.5), Inches(0.5), color)
        add_text_box(s, Inches(0.5), y + Inches(0.1), Inches(0.5), Inches(0.3), q_label, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.1), y, Inches(8.4), Inches(0.4), q, size=12, bold=True, color=C_PRIMARY)
        add_text_box(s, Inches(1.1), y + Inches(0.45), Inches(8.4), Inches(0.7), a, size=11, color=C_TEXT)

    add_rounded_rect(s, Inches(0.5), Inches(4.6), Inches(9), Inches(0.75), C_PRIMARY, radius=0.1)
    add_text_box(s, Inches(0.7), Inches(4.65), Inches(8.6), Inches(0.3), "本讲一句话：", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(4.95), Inches(8.6), Inches(0.4),
                 "自学能力 = 用 AI 解决问题的能力。家长 3 个不：不教、不催、不评判。", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# =====================================================
# 问题三：未来社会要求更高了吗 (6-8 页) - 7 页
# =====================================================
def slide_30():
    """问题三 1: 标题"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 30, "问题三")
    section_header(s, "问题 03", "未来社会要求更高了吗？", "—— 提前步入社会，让孩子接触真实问题")

    add_rect(s, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.95), Inches(9), Inches(0.4), "本讲核心方法", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.35), Inches(9), Inches(0.8), "「家庭真实问题池」", size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.95), Inches(9), Inches(0.3), "每月 1 个真实问题，孩子用 AI 解决，家长只引导不评判。", size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(3.55), Inches(9), Inches(0.3), "本讲将解决：", size=13, bold=True, color=C_PRIMARY)
    items = ["真实问题 vs 练习题", "怎么给孩子提供真实问题", "AI 让真实问题更易获取", "家长扮演什么角色"]
    for i, t in enumerate(items):
        x = Inches(0.5) + Inches(2.25) * (i % 2)
        y = Inches(3.95) + Inches(0.55) * (i // 2)
        add_oval(s, x, y, Inches(0.35), Inches(0.35), C_ACCENT)
        add_text_box(s, x + Inches(0.05), y + Inches(0.05), Inches(0.3), Inches(0.3), str(i+1), size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.45), y, Inches(1.75), Inches(0.3), t, size=11, bold=True, color=C_PRIMARY)


def slide_31():
    """问题三 2: 故事页 - 不主动的孩子"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 31, "问题三")
    section_header(s, "问题 03", "孩子什么都不缺，为什么还是不主动？", "—— 因为没接触过真实问题")

    add_rounded_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(2.6), C_LIGHT, radius=0.05)
    add_text_box(s, Inches(0.7), Inches(1.85), Inches(8.6), Inches(0.3), "真实场景", size=11, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(2.15), Inches(8.6), Inches(0.4), "家长常见的一个问题：", size=14, color=C_TEXT)
    add_paragraphs(s, Inches(0.7), Inches(2.6), Inches(8.6), Inches(1.6), [
        "孩子什么都不缺，要什么给什么",
        "但每天放学回来还是先看手机",
        "作业拖到晚上 10 点",
        "问他今天想做什么，回答「不知道」",
        "家长困惑：是不缺动力，还是不知道为什么要努力？",
    ], size=12, color=C_TEXT, line_spacing=1.4)

    add_rect(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.7), C_PRIMARY)
    add_text_box(s, Inches(0.7), Inches(4.58), Inches(8.6), Inches(0.3), "答案：", size=12, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(4.88), Inches(8.6), Inches(0.3),
                 "孩子没接触过真实问题。他对「问题」的理解 = 有标准答案的题目。", size=13, color=C_WHITE)


def slide_32():
    """问题三 3: 概念 - 真实问题 vs 练习题"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 32, "问题三")
    section_header(s, "问题 03", "真实问题 vs 练习题", "—— 两种「问题」，教育价值完全不同")

    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_RED)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "练习题", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.6), [
        "• 有标准答案",
        "• 题目已被设计好",
        "• 完成即结束",
        "• 评价标准清晰",
        "• 不会引起情绪波动",
        "• 不需要真实决策",
        "",
        "教育价值：训练基础技能",
        "局限：脱离真实世界",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_GREEN)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "真实问题", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.6), [
        "• 没有标准答案",
        "• 问题从生活中涌现",
        "• 解决后还有迭代",
        "• 评价标准多元",
        "• 容易产生情绪波动",
        "• 需要真实权衡",
        "",
        "教育价值：训练综合能力",
        "优势：连接真实世界",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    add_text_box(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.3),
                 "传统教育最大的问题——脱离真实。", size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_33():
    """问题三 4: 案例 - 4 个真实问题池"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 33, "问题三")
    section_header(s, "问题 03", "案例：4 个家庭真实问题池", "—— 孩子用 AI 解决真问题")

    cases = [
        ("家庭预算", "8 岁", "家里这个月预算 5000 元，请用 AI 帮我们规划：吃饭、零食、玩具、学习的钱怎么分？", C_PRIMARY),
        ("旅行规划", "10 岁", "国庆想去杭州玩 3 天，预算 3000 元 4 个人。请用 AI 规划路线、住宿、景点。", C_ACCENT),
        ("购物决策", "12 岁", "我想买一台学习用的电脑，预算 5000 元内。请用 AI 帮我对比 3 款，给我建议。", C_SECONDARY),
        ("邻里矛盾", "13 岁", "楼下邻居说我们晚上太吵，但我们要练琴。请用 AI 帮我想 3 个解决办法。", C_DARK),
    ]
    for i, (label, age, q, color) in enumerate(cases):
        x = Inches(0.5) + Inches(2.25) * (i % 2)
        y = Inches(1.7) + Inches(1.55) * (i // 2)
        add_rect(s, x, y, Inches(2.1), Inches(1.4), C_LIGHT)
        add_rect(s, x, y, Inches(2.1), Inches(0.4), color)
        add_text_box(s, x, y + Inches(0.05), Inches(2.1), Inches(0.3), label, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.1), y + Inches(0.45), Inches(1.9), Inches(0.2), f"适龄：{age}", size=9, color=color)
        add_text_box(s, x + Inches(0.1), y + Inches(0.65), Inches(1.9), Inches(0.75), q, size=9, color=C_TEXT)


def slide_34():
    """问题三 5: 工具 - 家庭真实问题池"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 34, "问题三")
    section_header(s, "问题 03", "工具 6：家庭真实问题池", "—— 10 个预设问题 + 创建规则")

    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "10 个预设问题方向", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.6), [
        "1. 家庭预算分配",
        "2. 旅行规划",
        "3. 大件购物决策",
        "4. 邻里/同学矛盾",
        "5. 房间整理方案",
        "6. 节日礼物选择",
        "7. 兴趣班选/退",
        "8. 课外书挑选",
        "9. 时间管理冲突",
        "10. 家庭会议议题",
    ], size=11, color=C_TEXT, line_spacing=1.4)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_ACCENT)
    add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "创建规则", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.6), [
        "• 真实：必须是家庭真问题",
        "• 可决策：孩子能参与决策",
        "• 适度：不要太大或太小",
        "• 周期：每月 1 个",
        "",
        "评估标准：",
        "✓ 孩子主动参与度",
        "✓ 解决方案的合理性",
        "✓ 反思深度",
    ], size=11, color=C_TEXT, line_spacing=1.4)

    add_text_box(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.3),
                 "家长 3 个角色：问题提供者 + 引导者 + 复盘者", size=12, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)


def slide_35():
    """问题三 6: 实操步骤 + 担忧"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 35, "问题三")
    section_header(s, "问题 03", "本讲实操 & 家长担忧", "")

    steps = [
        ("Step 1", "从 10 个预设问题中选 1 个", "本月从「家庭真实问题池」中选 1 个", C_PRIMARY),
        ("Step 2", "和孩子一起把问题描述清楚", "用 AI 帮孩子拆解成 3-5 个小问题", C_SECONDARY),
        ("Step 3", "孩子主导用 AI 解决", "家长 3 个不：不教/不催/不评判", C_ACCENT),
        ("Step 4", "每月家庭会议复盘", "评估孩子的方案，反思过程", C_DARK),
    ]
    for i, (label, t, d, color) in enumerate(steps):
        x = Inches(0.5) + Inches(2.3) * i
        add_oval(s, x + Inches(0.85), Inches(1.85), Inches(0.6), Inches(0.6), color)
        add_text_box(s, x + Inches(0.85), Inches(1.95), Inches(0.6), Inches(0.4), str(i+1), size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(2.55), Inches(2.2), Inches(0.3), t, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(2.85), Inches(2.2), Inches(1.5), d, size=10, color=C_TEXT, align=PP_ALIGN.CENTER)

    # 担忧
    add_rounded_rect(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.85), C_PRIMARY, radius=0.1)
    add_text_box(s, Inches(0.7), Inches(4.58), Inches(8.6), Inches(0.3), "担忧：孩子用 AI 解决真实问题，会不会有依赖？", size=12, bold=True, color=C_ACCENT)
    add_text_box(s, Inches(0.7), Inches(4.88), Inches(8.6), Inches(0.3),
                 "不会。关键是「主导权在孩子」。家长看的是孩子对问题的定义、对方案的选择，不是 AI 给了什么。", size=10, color=C_WHITE)
    add_text_box(s, Inches(0.7), Inches(5.15), Inches(8.6), Inches(0.2),
                 "本讲一句话：让孩子在真实问题里长大，比在练习题里熟练更重要。", size=10, bold=True, color=C_ACCENT)


def slide_36():
    """问题三 7: 总结"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 36, "问题三")
    section_header(s, "问题 03", "本讲一句话 & 下讲预告", "")

    add_rounded_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(1.5), C_PRIMARY, radius=0.1)
    add_text_box(s, Inches(0.5), Inches(1.85), Inches(9), Inches(0.3), "本讲一句话", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.2), Inches(9), Inches(0.6),
                 "AI 让「真实问题」更易获取，\n父母角色是「问题提供者」。", size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.95), Inches(9), Inches(0.3),
                 "孩子不是不主动，是没机会主动。", size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # 本讲产出
    add_rect(s, Inches(0.5), Inches(3.4), Inches(9), Inches(0.4), C_ACCENT)
    add_text_box(s, Inches(0.5), Inches(3.45), Inches(9), Inches(0.3), "本讲产出", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(3.9), Inches(8.6), Inches(1.0), [
        "✓ 完成「家庭真实问题池」第一版（含 10 个预设问题）",
        "✓ 家长为孩子设计未来 3 个月每月 1 个的真实问题清单",
        "✓ 家长明确自己在每个真实问题中的 3 个角色定位",
    ], size=12, color=C_TEXT, line_spacing=1.3)

    # 下讲
    add_rounded_rect(s, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4), C_PRIMARY, radius=0.2)
    add_text_box(s, Inches(0.5), Inches(5.05), Inches(9), Inches(0.3),
                 "问题三结束 → 问题四：怎样培养孩子的判断力？", size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# 简化版：问题四-问题十二沿用相同结构（每讲 7 页：标题/故事/李笑来vs本课/模型/案例/工具/担忧+总结）
# 为保证 100-160 页，使用紧凑结构
def make_lecture_block(start_page, q_num, q_title, q_subtitle, lixiao_view, my_view, model_name, model_desc, cases, tool_name, tool_desc, worries, summary):
    """生成问题讲的标准 7 页"""
    pages = []

    # Page 1: 标题
    def p_title():
        s = prs.slides.add_slide(BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        page_chrome(s, start_page, f"问题{q_num}")
        section_header(s, f"问题 {q_num:02d}", q_title, q_subtitle)
        add_rect(s, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5), C_PRIMARY)
        add_text_box(s, Inches(0.5), Inches(1.95), Inches(9), Inches(0.4), "本讲核心方法", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(0.5), Inches(2.35), Inches(9), Inches(0.8), model_name, size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(0.5), Inches(2.95), Inches(9), Inches(0.3), model_desc, size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Page 2: 故事
    def p_story():
        s = prs.slides.add_slide(BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        page_chrome(s, start_page+1, f"问题{q_num}")
        section_header(s, f"问题 {q_num:02d}", f"真实场景：{worries[0]['q'][:15]}...", "—— 一个被忽视的问题")
        add_rounded_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(2.6), C_LIGHT, radius=0.05)
        add_text_box(s, Inches(0.7), Inches(1.85), Inches(8.6), Inches(0.3), "真实案例", size=11, bold=True, color=C_ACCENT)
        add_text_box(s, Inches(0.7), Inches(2.15), Inches(8.6), Inches(0.4), worries[0].get("scene", "一位家长的真实困境"), size=14, color=C_TEXT)
        add_paragraphs(s, Inches(0.7), Inches(2.6), Inches(8.6), Inches(1.6), worries[0].get("desc", "").split("\n"), size=12, color=C_TEXT, line_spacing=1.4)
        add_rect(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.7), C_PRIMARY)
        add_text_box(s, Inches(0.7), Inches(4.58), Inches(8.6), Inches(0.3), "本讲核心问题：", size=12, bold=True, color=C_ACCENT)
        add_text_box(s, Inches(0.7), Inches(4.88), Inches(8.6), Inches(0.3), worries[0]["q"], size=13, color=C_WHITE)

    # Page 3: 李笑来 vs 本课
    def p_views():
        s = prs.slides.add_slide(BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        page_chrome(s, start_page+2, f"问题{q_num}")
        section_header(s, f"问题 {q_num:02d}", "李笑来观点 vs 本课独到见解", "")
        add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
        add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.5), C_SECONDARY)
        add_text_box(s, Inches(0.5), Inches(1.78), Inches(4.4), Inches(0.4), "李笑来老师的观点", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_paragraphs(s, Inches(0.7), Inches(2.4), Inches(4.0), Inches(2.6), lixiao_view.split("\n"), size=12, color=C_TEXT, line_spacing=1.4)
        add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
        add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.5), C_ACCENT)
        add_text_box(s, Inches(5.1), Inches(1.78), Inches(4.4), Inches(0.4), "本课往前走一步", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_paragraphs(s, Inches(5.3), Inches(2.4), Inches(4.0), Inches(2.6), my_view.split("\n"), size=12, color=C_TEXT, line_spacing=1.4)

    # Page 4: 模型
    def p_model():
        s = prs.slides.add_slide(BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        page_chrome(s, start_page+3, f"问题{q_num}")
        section_header(s, f"问题 {q_num:02d}", f"核心模型：{model_name}", model_desc)
        # 主体卡片
        add_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(2.8), C_LIGHT)
        add_text_box(s, Inches(0.7), Inches(1.85), Inches(8.6), Inches(0.4), "模型结构", size=14, bold=True, color=C_PRIMARY)
        add_paragraphs(s, Inches(0.7), Inches(2.3), Inches(8.6), Inches(2.1), model_desc.split("\n"), size=12, color=C_TEXT, line_spacing=1.4)
        # 金句
        add_rounded_rect(s, Inches(0.5), Inches(4.7), Inches(9), Inches(0.65), C_PRIMARY, radius=0.1)
        add_text_box(s, Inches(0.7), Inches(4.78), Inches(8.6), Inches(0.5), model_name, size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Page 5: 案例
    def p_cases():
        s = prs.slides.add_slide(BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        page_chrome(s, start_page+4, f"问题{q_num}")
        section_header(s, f"问题 {q_num:02d}", "案例：家庭真实过程", "")
        colors = [C_PRIMARY, C_ACCENT, C_SECONDARY, C_DARK, C_RED]
        n = len(cases)
        cols = 3 if n <= 3 else 3
        rows = (n + cols - 1) // cols
        card_w = Inches(2.95)
        card_h = Inches(1.45)
        for i, case in enumerate(cases):
            col = i % cols
            row = i // cols
            x = Inches(0.5) + (card_w + Inches(0.1)) * col
            y = Inches(1.7) + (card_h + Inches(0.1)) * row
            add_rect(s, x, y, card_w, card_h, C_LIGHT)
            add_rect(s, x, y, card_w, Inches(0.4), colors[i % len(colors)])
            add_text_box(s, x, y + Inches(0.05), card_w, Inches(0.3), case["label"], size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_text_box(s, x + Inches(0.1), y + Inches(0.45), card_w - Inches(0.2), Inches(0.25), case.get("age", ""), size=10, bold=True, color=colors[i % len(colors)])
            lines = case.get("desc", "").split("\n")
            for j, line in enumerate(lines):
                add_text_box(s, x + Inches(0.15), y + Inches(0.7) + Inches(0.22) * j, card_w - Inches(0.3), Inches(0.22), line, size=9, color=C_TEXT)
        # 总结位置
        y_summary = Inches(1.7) + (card_h + Inches(0.1)) * rows + Inches(0.1)
        add_text_box(s, Inches(0.5), y_summary, Inches(9), Inches(0.3),
                     f"共同点：{model_name} 落地效果显著。", size=12, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

    # Page 6: 工具
    def p_tool():
        s = prs.slides.add_slide(BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        page_chrome(s, start_page+5, f"问题{q_num}")
        section_header(s, f"问题 {q_num:02d}", f"工具：{tool_name}", tool_desc.split("\n")[0])
        add_rect(s, Inches(0.5), Inches(1.7), Inches(9), Inches(0.5), C_PRIMARY)
        add_text_box(s, Inches(0.5), Inches(1.78), Inches(9), Inches(0.4), "使用说明", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        lines = tool_desc.split("\n")
        for i, line in enumerate(lines):
            y = Inches(2.4) + Inches(0.4) * i
            add_text_box(s, Inches(0.7), y, Inches(0.3), Inches(0.3), f"{i+1}.", size=12, bold=True, color=C_ACCENT)
            add_text_box(s, Inches(1.0), y, Inches(8.3), Inches(0.3), line, size=11, color=C_TEXT)
        add_rounded_rect(s, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4), C_PRIMARY, radius=0.2)
        add_text_box(s, Inches(0.5), Inches(5.05), Inches(9), Inches(0.3),
                     f"工具：{tool_name} · 详见「06_配套工具」对应模板", size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Page 7: 担忧 + 总结
    def p_summary():
        s = prs.slides.add_slide(BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        page_chrome(s, start_page+6, f"问题{q_num}")
        section_header(s, f"问题 {q_num:02d}", "家长担忧澄清 & 一句话总结", "")
        for i, w in enumerate(worries[:2]):
            y = Inches(1.7) + Inches(1.2) * i
            color = [C_RED, C_ACCENT][i]
            add_oval(s, Inches(0.5), y, Inches(0.5), Inches(0.5), color)
            add_text_box(s, Inches(0.5), y + Inches(0.1), Inches(0.5), Inches(0.3), "Q", size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_text_box(s, Inches(1.1), y, Inches(8.4), Inches(0.4), w["q"], size=12, bold=True, color=C_PRIMARY)
            add_text_box(s, Inches(1.1), y + Inches(0.4), Inches(8.4), Inches(0.7), w["a"], size=10, color=C_TEXT)
        add_rounded_rect(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.85), C_PRIMARY, radius=0.1)
        add_text_box(s, Inches(0.7), Inches(4.55), Inches(8.6), Inches(0.3), "本讲一句话：", size=11, bold=True, color=C_ACCENT)
        add_text_box(s, Inches(0.7), Inches(4.85), Inches(8.6), Inches(0.4), summary, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    p_title()
    p_story()
    p_views()
    p_model()
    p_cases()
    p_tool()
    p_summary()


# =====================================================
# 问题四 - 问题十二：使用 make_lecture_block 批量生成
# =====================================================

# 问题四：判断力
make_lecture_block(
    37, 4, "怎样培养孩子的判断力？", "—— AI 输出三审制，训练孩子的判断力",
    "「判断力比知识更重要」\n\nAI 时代判断力的核心变化：\n能识别 AI 何时对、何时错\n\n家长要训练孩子对 AI 输出的评估能力",
    "判断力 = 「AI 输出评估力」\n\n原创方法：「AI 输出三审制」\n• 审事实\n• 审逻辑\n• 审价值观\n\n家长 3 个动作：\n示范判断 / 邀请孩子判断 / 共同判断",
    "AI 输出三审制", "审事实 · 审逻辑 · 审价值观 —— 3 个审核维度，10 个检查点",
    [
        {"label": "查资料", "age": "10 岁", "desc": "查「地球到太阳的距离」\nAI 答 1.5 亿公里\n用三审表 → 事实正确 → 用"},
        {"label": "写故事", "age": "12 岁", "desc": "AI 写「三国演义」读后感\n→ 审逻辑：因果是否成立\n→ 修改第 2 段"},
        {"label": "做决策", "age": "13 岁", "desc": "AI 推荐 3 款学习机\n→ 审价值观：是否有广告倾向\n→ 选第 2 款"},
    ],
    "AI 输出三审表", "1. 每次用 AI 后，填写 3 个维度（事实/逻辑/价值观）\n2. 10 个检查点：来源是否权威、数据是否过时、推论是否成立、是否有偏见、是否有遗漏等\n3. 每周至少 3 次，1 个月形成习惯\n4. 家庭共同讨论机制：孩子先判断，家长补充",
    [
        {"q": "AI 输出这么流畅，怎么判断它对错？",
         "a": "用「三审表」+ 关键事实交叉验证。AI 的错误常常是「流畅地编」。",
         "scene": "孩子用 AI 写作业，结果老师指出多处事实错误。家长困惑：为什么 AI 看起来很对？",
         "desc": "孩子最近用 AI 写历史小论文\nAI 给的内容很流畅\n老师却指出 3 处事实错误\n家长困惑：AI 看起来都对，怎么会有错？\n孩子也委屈：它说得头头是道"},
        {"q": "孩子还小，能学会判断 AI 吗？",
         "a": "可以。10 岁开始训练「事实核查」；12 岁加入「逻辑评估」；14 岁加入「价值观判断」。", "scene": "", "desc": ""},
    ],
    "判断力 = AI 输出评估力。审事实、审逻辑、审价值观。"
)

# 问题五：写作和英语
make_lecture_block(
    44, 5, "AI 来了，写作和英语还重要吗？", "—— 夯实基础，AI 时代的基础能力有新的内涵",
    "「写作和英语依然重要」\n\n基础能力不会贬值\n但「重要的维度」变了\n\n不是不学，是学法要变",
    "AI 替代的是「标准化写作」\n强化的是「创意写作 + 深度思考」\n\n英语的真正价值：\n不是「翻译工具」，是「跨文化思维载体」\n\nAI 辅助写作 5 步流程：\n创意生成 → 结构搭建 → 内容填充 → 润色优化 → 人工审核",
    "AI 辅助写作流程", "5 步流程 · 创意/结构/内容/润色/审核 —— AI 做 AI 擅长的，孩子做孩子擅长的",
    [
        {"label": "写作", "age": "11 岁", "desc": "用 AI 辅助写作文\n→ 创意：自己想\n→ 结构：AI 帮搭\n→ 内容：自己写\n→ 润色：AI 帮改\n→ 审核：自己审"},
        {"label": "英语", "age": "10 岁", "desc": "读英文原版书\n→ AI 翻译关键词\n→ 重点读地道表达\n→ 用 AI 改自己的英文写作"},
        {"label": "阅读", "age": "9 岁", "desc": "每天 30 分钟阅读\n→ 用 AI 解读难概念\n→ 自己写读后感\n→ 家庭共读会分享"},
    ],
    "AI 辅助写作流程模板", "1. 创意生成：自己头脑风暴 5 个角度\n2. 结构搭建：用 AI 列 3 种文章结构，对比选择\n3. 内容填充：自己写核心段落，AI 补盲区\n4. 润色优化：用 AI 改语言，但保留自己的声音\n5. 人工审核：自己 + 家长用「三审表」过一遍\n\n基础夯实清单（写作/英语/阅读）每月更新",
    [
        {"q": "AI 都能写文章了，孩子写作还有意义吗？",
         "a": "AI 替代「套路化写作」，强化「创意写作 + 深度思考」。孩子要学的是 AI 写不出来的那部分。",
         "scene": "孩子用 AI 写出一篇流畅作文，老师却说「没有灵魂」。", "desc": ""},
        {"q": "AI 翻译这么好了，还学英语吗？",
         "a": "英语的真正价值不是翻译，是跨文化思维载体。学英语的真正目的是「打开另一扇窗」。", "scene": "", "desc": ""},
    ],
    "AI 替代的是「标准化写作」，强化的是「创意 + 深度」。"
)

# 问题六：父母角色
make_lecture_block(
    51, 6, "父母应该做什么？", "—— 从知识传授者到 AI 协作系统的设计师",
    "「父母角色应从知识传授者转为成长引导者」\n\n传统父母角色失效：\n教 / 查 / 追\n\nAI 时代需要新角色",
    "AI 时代父母 = 「AI 协作系统的设计师」\n\n原创模型：「三锚模型」\n• 锚定目标\n• 锚定边界\n• 锚定反馈\n\n家庭 Prompt 工作流 5 大场景：\n作业辅导 / 兴趣探索 / 问题解答 / 创意生成 / 决策辅助",
    "三锚模型 + 家庭 Prompt 工作流", "锚定目标：能力方向 · 锚定边界：使用规则 · 锚定反馈：评估机制",
    [
        {"label": "A 父母", "age": "孩子 9 岁", "desc": "原本：检查作业 + 追问\n现在：设计家庭 AI 协议\n孩子主动完成作业\n家长不再追"},
        {"label": "B 父母", "age": "孩子 12 岁", "desc": "原本：教孩子做某事\n现在：搭 Prompt 模板\n孩子自己用 AI 完成\n家长做反馈者"},
        {"label": "C 父母", "age": "孩子 7 岁", "desc": "原本：完全禁止 AI\n现在：共同制定边界\n家庭 AI 协议签署\n亲子关系改善"},
    ],
    "三锚模型 + 家庭 Prompt 工作流", "1. 锚定目标：和孩子一起明确家庭未来 3 年能力培养方向\n2. 锚定边界：明确 AI 在家庭中使用的 4 大边界（场景/时间/内容/隐私）\n3. 锚定反馈：建立「AI 输出三审」+「成长评估」双重机制\n4. 家庭 Prompt 工作流 5 场景：作业辅导/兴趣探索/问题解答/创意生成/决策辅助\n5. 家长 5 个新角色：系统设计师/边界制定者/Prompt 搭桥者/反馈者/共学者",
    [
        {"q": "我自己都不懂 AI，能当「系统设计师」吗？",
         "a": "能。设计师不需要精通每件工具，需要的是「设计能力」。本课给你模板。",
         "scene": "妈妈完全不懂 AI，每次想陪孩子都得问丈夫。", "desc": "刘女士的丈夫是技术背景\n对 AI 工具很熟\n但她自己完全不懂\n每次想陪孩子用 AI 探索什么\n都得先问丈夫\n时间一长，她觉得自己「没什么用」"},
        {"q": "家里谁主导？",
         "a": "父母主导系统设计，孩子主导具体使用。这是「协作」，不是「单方管理」。", "scene": "", "desc": ""},
    ],
    "AI 时代父母 = AI 协作系统的设计师。三锚 + Prompt 工作流。"
)

# 问题七：AI 让孩子变笨
make_lecture_block(
    58, 7, "AI 会让孩子变笨吗？", "—— 注意力与深度思考，4 大边界守住能力",
    "「AI 不会让孩子变笨」\n\n但「无边界使用」会\n\n家长的焦虑不是 AI 本身\n是「无边界」的使用方式",
    "AI 不是「变笨原因」\n「无边界使用」才是\n\n原创工具：「AI 家庭使用协议」\n4 大边界：\n• 场景边界\n• 时间边界\n• 内容边界\n• 隐私边界\n\n家长在边界管理中的角色：\n和孩子共同制定 / 共同执行 / 共同回顾",
    "AI 家庭使用协议", "4 大边界 + 违约责任 + 每月回顾 —— 既不禁止，也不滥用",
    [
        {"label": "协议 A 家庭", "age": "孩子 10 岁", "desc": "原本：完全禁止 AI\n→ 协议：作业辅导允许\n→ 3 个月后：成绩稳定\n亲子关系改善"},
        {"label": "协议 B 家庭", "age": "孩子 13 岁", "desc": "原本：无边界使用\n→ 协议：每天 ≤ 1 小时\n内容需家长知情\n→ 3 个月后：注意力提升"},
        {"label": "协议 C 家庭", "age": "孩子 8 岁", "desc": "原本：混乱使用\n→ 协议：周末每天 ≤ 30 分钟\n记录使用场景\n→ 3 个月后：习惯建立"},
    ],
    "AI 家庭使用协议", "1. 场景边界：明确 AI 在哪些场景可使用（作业辅导 ✓ / 写小论文 需审核 / 决策辅助 ✓）\n2. 时间边界：每天使用时长 + 时段（如周末 19:00-20:00）\n3. 内容边界：禁止内容（暴力/成人/隐私话题）+ 鼓励内容（学习/兴趣/创作）\n4. 隐私边界：禁止输入家庭住址/电话/身份证/学校全名\n5. 违约责任：双方共同商定，孩子参与制定\n6. 每月回顾：每月 1 次家庭会议，回顾协议执行情况\n\n家长 3 个角色：和孩子共同制定 + 共同执行 + 共同回顾",
    [
        {"q": "AI 真的不会让孩子变笨吗？",
         "a": "不会。但「无边界使用」会让孩子失去深度思考机会。本课用 4 大边界守住能力。",
         "scene": "李女士坚决不让孩子碰 AI，理由是「会变笨」。", "desc": "李女士坚决不让孩子碰 AI\n「AI 一用就回不去了\n孩子会变得不爱思考」\n她的孩子今年 10 岁\n班上已经有 30% 的同学用 AI 辅助作业了\n她陷入两难：\n完全禁止，怕孩子落后\n放开了用，怕孩子依赖"},
        {"q": "协议签了孩子不执行怎么办？",
         "a": "用「共同回顾」机制：每月 1 次家庭会议，回顾 + 调整协议，不是惩罚。", "scene": "", "desc": ""},
    ],
    "AI 不会让孩子变笨，「无边界使用」才会。4 大边界守住能力。"
)

# 问题八：情感连接
make_lecture_block(
    65, 8, "怎样防止 AI 情感连接弱化？", "—— 亲子共学机制，把 AI 变成连接催化剂",
    "「AI 不会弱化亲子情感连接」\n\n但「无陪伴的 AI 使用」会\n\n关键不在 AI 本身\n在父母是否在场",
    "AI 可以成为「亲子共学催化剂」\n而非替代者\n\n原创方法：「亲子共学三法则」\n• 共同提问\n• 共同评估\n• 共同创作\n\n情感连接 4 个不可替代维度：\n陪伴 / 回应 / 共同记忆 / 身体在场",
    "亲子共学三法则", "共同提问 · 共同评估 · 共同创作 —— 把 AI 用成「亲子连接器」",
    [
        {"label": "A 家庭", "age": "孩子 9 岁", "desc": "每周 1 次家庭 AI 探索日\n父母 + 孩子 + AI 三方\n共同发现新事物\n→ 孩子说：「AI 是我们的玩具」"},
        {"label": "B 家庭", "age": "孩子 12 岁", "desc": "每月 1 次家庭 AI 项目\n孩子主导，父母协助\n全家一起完成 1 个作品\n→ 留下了 12 个家庭记忆"},
        {"label": "C 家庭", "age": "孩子 7 岁", "desc": "每天 15 分钟「AI 故事时间」\n父母读 AI 写的故事\n孩子问问题\n→ 每天都有共同记忆"},
    ],
    "亲子共学记录表", "1. 每周 1 次亲子共学，记录主题 + 过程 + 收获\n2. 三法则应用：\n   - 共同提问：父母 + 孩子 + AI 三方一起问\n   - 共同评估：用「三审表」共同评估 AI 输出\n   - 共同创作：父母 + 孩子 + AI 三方共同产出作品\n3. 情感连接 4 维度检查：\n   - 陪伴：每周至少 1 次全身心陪伴\n   - 回应：对孩子的情绪给予有温度回应\n   - 共同记忆：每月至少 1 个家庭共同记忆\n   - 身体在场：减少「电子陪伴」，多面对面\n4. 家长 3 个不：不用 AI 替代陪伴 / 不用 AI 回应情绪 / 不用 AI 创造共同记忆",
    [
        {"q": "AI 会不会让孩子不再需要父母？",
         "a": "不会。情感连接的 4 个维度（陪伴/回应/共同记忆/身体在场）AI 永远替代不了。",
         "scene": "妈妈担心：孩子有什么都问 AI，不问我了。", "desc": "王女士发现一个现象：\n孩子最近什么问题都问 AI\n「妈妈这个字怎么读？」问 AI\n「妈妈这道题怎么做？」问 AI\n「妈妈恐龙吃什么？」问 AI\n她开始担心：\n孩子不再需要我了？"},
        {"q": "工作很忙，没时间亲子共学怎么办？",
         "a": "质量比时长重要。15 分钟全身心陪伴 > 2 小时心不在焉。每天 1 个 15 分钟就够。", "scene": "", "desc": ""},
    ],
    "AI 是「亲子共学催化剂」不是替代者。陪伴、回应、共同记忆、身体在场 = 4 个不可替代维度。"
)

# 问题九：兴趣教育
make_lecture_block(
    72, 9, "AI 时代，兴趣教育怎么办？", "—— 兴趣激发 vs 技能训练，AI 辅助兴趣探索",
    "「兴趣教育 = 兴趣激发 + 工具赋能」\n而非纯技能训练\n\nAI 时代兴趣教育的目的：\n保护好奇心 + 培养探索力",
    "兴趣教育 = 兴趣激发 + 工具赋能\n\n3 种 AI 辅助模式：\n• 陪练模式：AI 当陪练\n• 启发模式：AI 启发思路\n• 创作模式：AI 辅助创作\n\n兴趣教育 3 个不应该：\n不以考级为目标 / 不以攀比为标准 / 不以家长意志为中心",
    "AI 兴趣探索地图", "5 个兴趣领域 + 每个领域的 AI 辅助模式 —— 找对兴趣，比练好技能重要",
    [
        {"label": "吉他", "age": "11 岁", "desc": "AI 当陪练\n听孩子弹，指出问题\n推荐练习曲目\n孩子主动练琴时间增加"},
        {"label": "绘画", "age": "9 岁", "desc": "AI 启发思路\n孩子描述想象，AI 出图参考\n孩子再自己画\n保留孩子创意"},
        {"label": "编程", "age": "13 岁", "desc": "AI 辅助创作\n孩子想做游戏\nAI 生成基础代码\n孩子改逻辑 + 加关卡"},
    ],
    "AI 兴趣探索地图", "1. 5 个兴趣领域：音乐/绘画/编程/写作/运动（每个领域 1 张地图）\n2. 3 种 AI 辅助模式：\n   - 陪练模式：AI 当陪练/陪画/陪练的对象\n   - 启发模式：AI 启发思路，给孩子灵感和方向\n   - 创作模式：AI 辅助创作，孩子主导内容\n3. 兴趣教育 3 个不应该：\n   - 不以考级为目标\n   - 不以攀比为标准\n   - 不以家长意志为中心\n4. 家长 3 个动作：\n   - 观察孩子的兴趣（不是追热门）\n   - 提供 AI 工具（不是教技能）\n   - 不评判结果（看孩子是否享受）",
    [
        {"q": "AI 时代还要不要让孩子学乐器、绘画？",
         "a": "要。技能可能被替代，兴趣的「情感连接」和「创意表达」AI 替代不了。",
         "scene": "王女士给孩子报了 4 个兴趣班，但孩子都说「没意思」。", "desc": "王女士给孩子报了：\n钢琴、绘画、编程、围棋\n4 个班花 5 万多\n但孩子都说「没意思」\n她困惑：\n是我选错了吗？\n是孩子没天赋吗？\n还是报班方式错了？"},
        {"q": "孩子只想玩，不想练怎么办？",
         "a": "孩子想玩是正常的。「练」是被动的，「玩」是主动的。把练变成玩：把练习设计成游戏/挑战。", "scene": "", "desc": ""},
    ],
    "兴趣教育 = 兴趣激发 + 工具赋能。3 种 AI 辅助模式，3 个不应该。"
)

# 问题十：成长评估
make_lecture_block(
    79, 10, "如何评估孩子在 AI 时代的成长？", "—— 成长评估新维度，3 大维度 10 项指标",
    "「传统评估在 AI 时代部分失效」\n\n考试成绩 / 考级证书\n不能完全反映孩子的真实成长\n\n家长缺的不是行动系统\n是评估系统",
    "新评估维度 = 「3 大能力」\n• 问题解决力\n• 创意产出力\n• 情感连接力\n\n工具：「AI 时代成长评估表」\n3 大维度 + 10 项指标\n\n评估目的不是「打分」\n是「调整方向」",
    "AI 时代成长评估表", "3 大维度 10 项指标 —— 每月小评、每季中评、每年总评",
    [
        {"label": "问题解决力", "age": "评估维度 1", "desc": "1. 发现真问题\n2. 拆解问题\n3. 调用资源\n4. 实施方案\n5. 反思迭代"},
        {"label": "创意产出力", "age": "评估维度 2", "desc": "1. 提出新想法\n2. 整合跨域\n3. 完成作品\n4. 接受反馈\n5. 持续迭代"},
        {"label": "情感连接力", "age": "评估维度 3", "desc": "1. 识别情绪\n2. 回应他人\n3. 建立信任\n4. 解决冲突\n5. 创造连接"},
    ],
    "AI 时代成长评估表", "1. 3 大维度：问题解决力 + 创意产出力 + 情感连接力\n2. 10 项指标（每维度 5 项）：见左侧案例\n3. 评估频率：\n   - 每月 1 次小评（家长 + 孩子共同）\n   - 每季度 1 次中评（结合真实问题池）\n   - 每年 1 次总评（结合 3 年路线图）\n4. 评估结果运用：\n   - 不是贴标签，是调方向\n   - 不对比他人，只看自己的进步\n   - 不焦虑分数，看趋势\n5. 家长 3 个不：\n   - 不贴标签\n   - 不对比他人\n   - 不焦虑分数",
    [
        {"q": "传统成绩和考级还有用吗？",
         "a": "有用，但不再是唯一标准。新评估 3 维度（问题解决/创意产出/情感连接）和传统成绩互补。",
         "scene": "张先生困惑：孩子考级拿了好几个，但生活中什么都不会。", "desc": "张先生的孩子拿了：\n钢琴 5 级、围棋 3 段、编程 1 级\n证书一堆\n但张先生发现：\n孩子生活中什么都不会\n遇到问题就找家长\n不会自己解决\n他困惑：\n证书代表能力吗？\n考级有意义吗？"},
        {"q": "评估会不会让孩子有压力？",
         "a": "不会，因为是「调方向」不是「打分」。和孩子一起评估，方向是孩子自己的选择。", "scene": "", "desc": ""},
    ],
    "新评估 = 问题解决力 + 创意产出力 + 情感连接力。评估是调方向，不是贴标签。"
)

# 问题十一：3 年路线图
make_lecture_block(
    86, 11, "家庭 AI 教育路线图——3 年行动规划", "—— 入门 / 熟练 / 创新，3 个阶段稳步推进",
    "「单次行动容易，长期坚持难」\n\n这门课不是「一堂课的事」\n是「3 年行动的事」\n\n家长需要一张清晰路线图",
    "路线图 = 「能力评估 → Prompt 搭建 → 行动执行 → 评估更新」\n\n工具：「3 年家庭 AI 教育路线图」\n3 个阶段：\n• 入门（第 1 年）\n• 熟练（第 2 年）\n• 创新（第 3 年）\n\n3 个调整机制：\n每月小调整 / 每季度中调整 / 每年大调整",
    "3 年家庭 AI 教育路线图", "入门 → 熟练 → 创新 —— 3 个阶段 × 3 个调整机制",
    [
        {"label": "入门阶段", "age": "第 1 年", "desc": "目标：建立基础认知\n行动：签署协议 + 跑通 30 天清单\n评估：家庭能力评估表"},
        {"label": "熟练阶段", "age": "第 2 年", "desc": "目标：搭建工作流\n行动：形成共学机制 + 建立评估体系\n评估：成长评估表"},
        {"label": "创新阶段", "age": "第 3 年", "desc": "目标：孩子主导\n行动：家庭 AI 资产沉淀 + 对外输出\n评估：路线图更新"},
    ],
    "3 年家庭 AI 教育路线图", "1. 入门阶段（第 1 年）：\n   - 建立基础认知（完成本课程）\n   - 签署 AI 家庭使用协议\n   - 跑通 30 天家庭 AI 行动清单\n2. 熟练阶段（第 2 年）：\n   - 搭建家庭 Prompt 工作流（5 场景）\n   - 形成亲子共学机制（每周 1 次）\n   - 建立成长评估体系（每季 1 次）\n3. 创新阶段（第 3 年）：\n   - 孩子主导 AI 探索（家长退到观察者）\n   - 家庭 AI 资产沉淀（作品集/经验库）\n   - 对外输出经验（家庭博客/演讲）\n4. 3 个调整机制：\n   - 每月小调整：微调目标和行动\n   - 每季度中调整：评估阶段进展\n   - 每年大调整：重新规划 3 年路线\n5. 家长 3 个角色：制定者 + 监督者 + 调整者",
    [
        {"q": "3 年太长了，1 年规划不行吗？",
         "a": "1 年只能「尝试」，3 年才能「稳定 + 创新」。能力的培养需要 3 年节奏。",
         "scene": "周女士做了 1 年规划就放弃了。", "desc": "周女士去年做过 1 年规划：\n列了 12 个目标\n每月 1 个\n结果：\n3 个月后放弃了\n她说：「太忙了，1 年太长」\n她想知道：\n1 年规划能成功吗？\n3 年规划会不会更累？"},
        {"q": "路线图会不会变成「挂在墙上的纸」？",
         "a": "不会。本课强调「3 个调整机制」：每月/每季/每年都有动作。路线图是活的。", "scene": "", "desc": ""},
    ],
    "3 年路线图 = 入门 → 熟练 → 创新。3 个阶段 × 3 个调整机制，让规划活起来。"
)

# 问题十二：误区
make_lecture_block(
    93, 12, "常见误区与应对——避坑指南", "—— 5 大误区自查 + 自我纠错机制",
    "「家长在执行中一定会遇到各种坑」\n\n提前知道有哪些常见误区\n能帮家长少走 80% 弯路\n\n5 大误区 + 应对策略",
    "5 大误区：\n1. AI 万能论\n2. AI 禁止论\n3. 监控取代引导\n4. 技能取代思维\n5. 短期取代长期\n\n自我纠错机制：\n每月自查 / 每季复盘 / 每年大调整",
    "5 大误区 + 应对", "避坑指南 —— 5 大误区自查表 + 自我纠错机制",
    [
        {"label": "AI 万能论", "age": "误区 1", "desc": "症状：所有学习都交给 AI\n根源：被 AI 能力震撼\n应对：协议中的「场景边界」限制"},
        {"label": "AI 禁止论", "age": "误区 2", "desc": "症状：坚决不让孩子碰 AI\n根源：被 AI 焦虑笼罩\n应对：从协议开始，逐步建立使用"},
        {"label": "监控取代引导", "age": "误区 3", "desc": "症状：装监控软件 / 查聊天\n根源：把 AI 当「危险品」\n应对：用「亲子共学」取代「监控」"},
        {"label": "技能取代思维", "age": "误区 4", "desc": "症状：学 10 个工具不教思维\n根源：把「会用工具」当能力证明\n应对：「以生产为导向」方法"},
        {"label": "短期取代长期", "age": "误区 5", "desc": "症状：今天学新工具明天换\n根源：被 AI 更新速度带着跑\n应对：「3 年路线图」建立长期视角"},
    ],
    "AI 家庭教育误区自查表", "1. 5 大误区识别（每季度 1 次自查）：\n   - 误区 1：AI 万能论（症状/根源/应对）\n   - 误区 2：AI 禁止论\n   - 误区 3：监控取代引导\n   - 误区 4：技能取代思维\n   - 误区 5：短期取代长期\n2. 自查流程：\n   - 每季度家庭会议时填写\n   - 家长 + 孩子共同打分（1-5 分）\n   - 识别家庭当前所处的 1-2 个误区\n3. 应对动作：\n   - 为每个识别出的误区设计具体「应对动作」\n   - 行动写进「30 天行动清单」\n4. 自我纠错机制：\n   - 每月 1 次自查（家庭内部）\n   - 每季度 1 次复盘（家庭会议）\n   - 每年 1 次大调整（路线图更新）\n5. 家长 3 个动作：\n   - 不在孩子面前说「AI 没用」或「AI 万能」\n   - 让孩子参与误区识别和应对\n   - 把误区应对当成「家庭学习机会」",
    [
        {"q": "我已经踩了某个误区，怎么办？",
         "a": "正常。每个家庭都会踩 1-2 个误区。关键是「识别 + 应对」，不是「自责」。",
         "scene": "孙老师发现自己完全禁止孩子用 AI，现在孩子落后了。", "desc": "孙老师是一位初中英语老师\n她之前完全禁止孩子用 AI\n「AI 是洪水猛兽」\n最近发现：\n学生偷偷用 AI 写作业\n成绩反而有提升\n她困惑：\n自己是不是错了？\n现在改还来得及吗？"},
        {"q": "5 个误区我都中了，是不是没救了？",
         "a": "不是。误区是「思维惯性」，可以改。本课给自查表和纠错机制。", "scene": "", "desc": ""},
    ],
    "5 大误区自查 + 3 个纠错机制。让家庭教育不跑偏。"
)


# =====================================================
# 结刊词 (3-5 页) — 5 页
# =====================================================
def slide_100():
    """结刊词 1: 标题"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 100, "结刊词")
    section_header(s, "13·结刊词", "行动，是唯一的答案", "—— 课程回顾 + 30 天启动")

    add_rect(s, Inches(0.5), Inches(1.8), Inches(9), Inches(2.0), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.95), Inches(9), Inches(0.4), "本课核心承诺", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.4), Inches(9), Inches(0.8), "学完本课，你不需要再听", size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.85), Inches(9), Inches(0.8), "任何「AI 焦虑贩卖」。", size=24, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(3.95), Inches(9), Inches(0.3), "因为——你有了自己的系统。", size=14, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.3), Inches(9), Inches(0.6),
                 "AI 怎么变，你的系统都在。\n系统的力量，是家庭教育最强的护城河。", size=13, color=C_TEXT_LIGHT, align=PP_ALIGN.CENTER)


def slide_101():
    """结刊词 2: 13 讲一句话回顾"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 101, "结刊词")
    section_header(s, "13·结刊词", "13 讲一句话回顾", "—— 从认知到行动，从单点到系统，从工具到方法")

    items = [
        ("发刊词", "焦虑的反义词不是「想开点」，是「有行动系统」", C_ACCENT),
        ("先导课", "AI 时代重新定义什么是重要", C_PRIMARY),
        ("问题 1", "你抢的赛道，决定孩子的未来", C_PRIMARY),
        ("问题 2", "自学能力 = 用 AI 解决问题的能力", C_SECONDARY),
        ("问题 3", "让孩子在真实问题里长大", C_SECONDARY),
        ("问题 4", "判断力 = AI 输出评估力", C_DARK),
        ("问题 5", "AI 替代「标准化」，强化「创意 + 深度」", C_DARK),
        ("问题 6", "父母 = AI 协作系统的设计师", C_RED),
        ("问题 7", "AI 不会让孩子变笨，无边界才会", C_RED),
        ("问题 8", "AI 是亲子共学催化剂，不是替代者", C_ACCENT),
        ("问题 9", "兴趣教育 = 兴趣激发 + 工具赋能", C_ACCENT),
        ("问题 10", "评估是调方向，不是贴标签", C_PRIMARY),
        ("问题 11", "3 年路线图：入门 → 熟练 → 创新", C_PRIMARY),
        ("问题 12", "5 大误区自查 + 3 个纠错机制", C_SECONDARY),
    ]
    for i, (label, t, color) in enumerate(items):
        y = Inches(1.7) + Inches(0.24) * i
        add_rect(s, Inches(0.5), y, Inches(0.6), Inches(0.22), color)
        add_text_box(s, Inches(0.5), y, Inches(0.6), Inches(0.22), label, size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.2), y, Inches(8.3), Inches(0.22), t, size=10, color=C_TEXT)


def slide_102():
    """结刊词 3: 30 天行动清单（部分）"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 102, "结刊词")
    section_header(s, "13·结刊词", "30 天家庭 AI 行动清单（节选）", "—— 工具 5：每日 1 个小任务，30 天跑通")

    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.5), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.4), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.75), Inches(4.4), Inches(0.3), "Week 1：建立认知", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.25), Inches(4.0), Inches(2.9), [
        "Day 1：填「家庭能力评估表」",
        "Day 2：和孩子聊 AI，听听 ta 怎么看",
        "Day 3：家长试用 1 个 AI 工具",
        "Day 4：和孩子一起试用 AI",
        "Day 5：和孩子讨论 4 大边界",
        "Day 6：填写「AI 家庭使用协议」",
        "Day 7：家庭会议：回顾第 1 周",
    ], size=10, color=C_TEXT, line_spacing=1.3)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.5), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.4), C_ACCENT)
    add_text_box(s, Inches(5.1), Inches(1.75), Inches(4.4), Inches(0.3), "Week 2-4：跑通流程", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.25), Inches(4.0), Inches(2.9), [
        "Day 8-14：完成 1 个「家庭真实问题」",
        "Day 15：搭建「家庭 Prompt 工作流」",
        "Day 16-21：每周 1 次亲子共学",
        "Day 22：使用「AI 输出三审表」3 次",
        "Day 23-28：完成 1 个孩子的自学项目",
        "Day 29：填写「成长评估表」",
        "Day 30：家庭会议：30 天复盘",
    ], size=10, color=C_TEXT, line_spacing=1.3)

    add_rounded_rect(s, Inches(0.5), Inches(5.25), Inches(9), Inches(0.0), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.2),
                 "完整 30 天清单见「06_配套工具/30 天家庭 AI 行动清单」", size=10, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_103():
    """结刊词 4: 持续学习建议"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 103, "结刊词")
    section_header(s, "13·结刊词", "课程结束后，怎么持续学习？", "—— 3 个持续学习动作")

    actions = [
        ("动作 1", "每季度 1 次家庭 AI 教育复盘", "用「成长评估表」评估孩子的 3 大能力变化\n用「家庭能力评估表」重新打分\n调整下一季度的目标", C_PRIMARY),
        ("动作 2", "每年 1 次路线图更新", "对照 3 年路线图，看家庭处于哪个阶段\n调整下一年的目标和行动\n更新家庭 AI 协议", C_ACCENT),
        ("动作 3", "加入罗老师家庭教育社群", "和其他家长交流真实案例\n获取最新 AI 工具和方法\n获得持续支持", C_SECONDARY),
    ]
    for i, (label, t, d, color) in enumerate(actions):
        y = Inches(1.7) + Inches(1.1) * i
        add_oval(s, Inches(0.5), y, Inches(0.7), Inches(0.7), color)
        add_text_box(s, Inches(0.5), y + Inches(0.18), Inches(0.7), Inches(0.4), str(i+1), size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.3), y, Inches(8.2), Inches(0.4), t, size=14, bold=True, color=color)
        add_text_box(s, Inches(1.3), y + Inches(0.4), Inches(8.2), Inches(0.7), d, size=11, color=C_TEXT)

    add_rounded_rect(s, Inches(0.5), Inches(5.05), Inches(9), Inches(0.4), C_PRIMARY, radius=0.2)
    add_text_box(s, Inches(0.5), Inches(5.1), Inches(9), Inches(0.3),
                 "系统的力量是——AI 怎么变，你的系统都在。", size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_104():
    """结刊词 5: 鼓励 + 致谢"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 104, "结刊词")
    section_header(s, "13·结刊词", "开始行动，给孩子更多可能性", "—— 一位老师，13 讲，一套系统")

    add_rect(s, Inches(0.5), Inches(1.8), Inches(9), Inches(2.5), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.95), Inches(9), Inches(0.4), "送给每一位家长的话", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.4), Inches(9), Inches(0.6), "家长不是孩子的教练，", size=20, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.85), Inches(9), Inches(0.6), "家长是孩子和未来之间的翻译官。", size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(3.5), Inches(9), Inches(0.3), "AI 来了，未来变了，", size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(3.8), Inches(9), Inches(0.4),
                 "你的工作是把未来翻译成孩子听得懂的日常，\n把日常训练成未来用得上的能力。",
                 size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(4.6), Inches(9), Inches(0.3),
                 "这套课，就是这套翻译系统的使用手册。", size=14, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.95), Inches(9), Inches(0.3),
                 "我们，行动见。", size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# =====================================================
# 附录：13 套工具速查 (3-5 页) — 5 页
# =====================================================
def slide_105():
    """附录 1: 工具总览"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 105, "附录")
    section_header(s, "附录", "13 套原创工具速查", "—— 工具即带走，拿到就会用")

    tools = [
        ("工具 1", "家庭能力评估表", "识别增值/贬值能力", "每年 1 次"),
        ("工具 2", "家庭 Prompt 工作流", "5 大场景的 Prompt 模板", "每周多次"),
        ("工具 3", "AI 输出三审表", "训练 AI 输出判断力", "每次用 AI 后"),
        ("工具 4", "AI 家庭使用协议", "明确 4 大边界", "每月 1 次回顾"),
        ("工具 5", "30 天家庭 AI 行动清单", "每日 1 个小任务", "30 天一循环"),
        ("工具 6", "家庭真实问题池", "每月 1 个真实问题", "每月 1 个"),
        ("工具 7", "亲子共学记录表", "每周 1 次亲子共学", "每周 1 次"),
        ("工具 8", "3 年家庭 AI 教育路线图", "长期规划", "每年 1 次更新"),
        ("工具 9", "AI 时代成长评估表", "3 大维度 10 项指标", "每月小评"),
        ("工具 10", "AI 家庭教育误区自查表", "识别家庭误区", "每季度 1 次"),
    ]
    # 两列
    for i, (label, name, desc, freq) in enumerate(tools):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + Inches(4.6) * col
        y = Inches(1.7) + Inches(0.7) * row
        add_rect(s, x, y, Inches(4.5), Inches(0.6), C_LIGHT)
        add_rect(s, x, y, Inches(0.6), Inches(0.6), C_PRIMARY)
        add_text_box(s, x, y + Inches(0.1), Inches(0.6), Inches(0.4), label[-2:], size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.7), y + Inches(0.05), Inches(2.5), Inches(0.3), name, size=11, bold=True, color=C_PRIMARY)
        add_text_box(s, x + Inches(0.7), y + Inches(0.3), Inches(2.5), Inches(0.3), desc, size=9, color=C_TEXT)
        add_text_box(s, x + Inches(3.2), y + Inches(0.15), Inches(1.3), Inches(0.3), freq, size=9, color=C_ACCENT, align=PP_ALIGN.RIGHT)


def slide_106():
    """附录 2: 三锚模型速查"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 106, "附录")
    section_header(s, "附录", "原创模型 1：三锚模型速查", "—— 父母设计家庭 AI 教育方案的 3 个问题")

    anchors = [
        ("锚定目标", "孩子该练的「AI 时代增值能力」是什么？", C_PRIMARY,
         "对应能力图谱：问题定义力 / 情感连接力 / 跨域整合力\n落地工具：家庭能力评估表"),
        ("锚定边界", "AI 使用的场景/时间/内容/隐私边界", C_ACCENT,
         "4 大边界：场景 / 时间 / 内容 / 隐私\n落地工具：AI 家庭使用协议"),
        ("锚定反馈", "AI 输出评估 + 孩子成长评估", C_SECONDARY,
         "AI 输出评估：三审表（事实/逻辑/价值观）\n成长评估：3 大维度 10 项指标"),
    ]
    for i, (name, q, color, desc) in enumerate(anchors):
        y = Inches(1.7) + Inches(1.1) * i
        add_rect(s, Inches(0.5), y, Inches(1.4), Inches(1.0), color)
        add_text_box(s, Inches(0.5), y + Inches(0.2), Inches(1.4), Inches(0.6), name, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(2.0), y + Inches(0.1), Inches(7.5), Inches(0.4), q, size=12, bold=True, color=color)
        add_text_box(s, Inches(2.0), y + Inches(0.5), Inches(7.5), Inches(0.5), desc, size=10, color=C_TEXT)


def slide_107():
    """附录 3: 家庭 Prompt 工作流速查"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 107, "附录")
    section_header(s, "附录", "原创模型 2：家庭 Prompt 工作流速查", "—— 5 大场景 × 3 类角色")

    scenes = [
        ("作业辅导", "孩子做作业时\nAI 解释概念 + 提示思路", C_PRIMARY),
        ("兴趣探索", "孩子探索兴趣时\nAI 启发思路 + 推荐资源", C_ACCENT),
        ("问题解答", "孩子问问题时\nAI 给 3 种思路 + 引导判断", C_SECONDARY),
        ("创意生成", "孩子想创作时\nAI 当陪练 + 提示创意", C_DARK),
        ("决策辅助", "家庭做决策时\nAI 列方案 + 优缺点对比", C_RED),
    ]
    for i, (name, desc, color) in enumerate(scenes):
        y = Inches(1.7) + Inches(0.65) * i
        add_rect(s, Inches(0.5), y, Inches(1.6), Inches(0.55), color)
        add_text_box(s, Inches(0.5), y + Inches(0.12), Inches(1.6), Inches(0.3), name, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(2.2), y, Inches(7.3), Inches(0.55), desc, size=10, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)

    # 家长引导话术
    add_rounded_rect(s, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4), C_PRIMARY, radius=0.2)
    add_text_box(s, Inches(0.5), Inches(5.05), Inches(9), Inches(0.3),
                 "每个场景配「家长引导话术」+「孩子使用边界」+「核心 Prompt 模板」。",
                 size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_108():
    """附录 4: 超越竞争能力图谱速查"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 108, "附录")
    section_header(s, "附录", "原创模型 3：超越竞争能力图谱速查", "—— 3 大稀缺能力 × 3 个评估等级")

    abilities = [
        ("问题定义力", "知道「什么才是真正要解决的问题」", C_PRIMARY,
         "启蒙：能问出好问题\n熟练：能区分真假问题\n创新：能重新定义问题"),
        ("情感连接力", "和他人建立真实、有温度的关系", C_ACCENT,
         "启蒙：识别自己/他人情绪\n熟练：用语言回应情绪\n创新：安抚/激励/领导他人"),
        ("跨域整合力", "把不同领域的知识拼成新东西", C_SECONDARY,
         "启蒙：知道不同领域存在\n熟练：把 2 个领域连起来\n创新：能创造新领域交叉"),
    ]
    for i, (name, desc, color, level) in enumerate(abilities):
        y = Inches(1.7) + Inches(1.1) * i
        add_rect(s, Inches(0.5), y, Inches(2.0), Inches(1.0), color)
        add_text_box(s, Inches(0.5), y + Inches(0.2), Inches(2.0), Inches(0.6), name, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(2.6), y + Inches(0.1), Inches(6.9), Inches(0.3), desc, size=12, bold=True, color=color)
        add_text_box(s, Inches(2.6), y + Inches(0.45), Inches(6.9), Inches(0.5), level, size=10, color=C_TEXT)


def slide_109():
    """附录 5: 30 天行动清单 + 资源索引"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    page_chrome(s, 109, "附录")
    section_header(s, "附录", "30 天行动 + 资源索引", "—— 现在就开始")

    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(0.5), Inches(1.7), Inches(4.4), Inches(0.4), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(1.75), Inches(4.4), Inches(0.3), "今晚就能做的 3 件事", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(0.7), Inches(2.25), Inches(4.0), Inches(2.7), [
        "1. 打印「家庭能力评估表」",
        "   今晚和孩子一起填写",
        "",
        "2. 和孩子讨论 4 大边界",
        "   初步达成共识",
        "",
        "3. 下载 1 个 AI 工具",
        "   家长自己先体验",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(3.4), C_LIGHT)
    add_rect(s, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.4), C_ACCENT)
    add_text_box(s, Inches(5.1), Inches(1.75), Inches(4.4), Inches(0.3), "课程资源索引", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_paragraphs(s, Inches(5.3), Inches(2.25), Inches(4.0), Inches(2.7), [
        "01_课程大纲：完整大纲",
        "02_教学文档：13 讲讲义",
        "06_配套工具：13 套工具模板",
        "07_全流程练习册：课堂练习",
        "08_场景库：真实家庭案例",
        "13_HTML 可视化：工具可视化",
    ], size=11, color=C_TEXT, line_spacing=1.2)

    add_rounded_rect(s, Inches(0.5), Inches(5.25), Inches(9), Inches(0.0), C_PRIMARY)
    add_text_box(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.2),
                 "完整 30 天清单见「06_配套工具/30 天家庭 AI 行动清单」", size=10, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_110():
    """致谢 1"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_PRIMARY
    add_rect(s, Inches(0), Inches(0), Inches(10), Inches(0.15), C_ACCENT)
    add_rect(s, Inches(0), Inches(5.475), Inches(10), Inches(0.15), C_ACCENT)
    add_text_box(s, Inches(0.6), Inches(2.0), Inches(8.8), Inches(0.5), "特别致谢", size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(2.5), Inches(8.8), Inches(0.5), "李笑来老师", size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(3.1), Inches(8.8), Inches(0.4),
                 "「超越竞争、不焦虑、夯实基础」的深刻洞察，", size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(3.45), Inches(8.8), Inches(0.4),
                 "为本课提供了认知基础。", size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.3),
                 "本课往前走了一步：从认知到行动系统。", size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.3),
                 "© 罗老师《AI 时代的家庭教育》· 2026", size=9, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_oval(s, Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4), C_ACCENT)
    add_text_box(s, Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4), "110", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_111():
    """致谢 2: 版权 + 致谢学员"""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_PRIMARY
    add_rect(s, Inches(0), Inches(0), Inches(10), Inches(0.15), C_ACCENT)
    add_rect(s, Inches(0), Inches(5.475), Inches(10), Inches(0.15), C_ACCENT)
    add_text_box(s, Inches(0.6), Inches(1.5), Inches(8.8), Inches(0.4), "感谢每一位学员", size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(2.0), Inches(8.8), Inches(0.4),
                 "你们在家庭中落地的每一个行动，", size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(2.35), Inches(8.8), Inches(0.4),
                 "都是本课程最好的反馈。", size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.6), Inches(3.0), Inches(8.8), Inches(0.4), "版权信息", size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(3.45), Inches(8.8), Inches(0.3),
                 "本课程含 13 讲系统课、3 大原创模型、13 套原创工具、30 天行动清单、3 年家庭 AI 教育路线图",
                 size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(3.75), Inches(8.8), Inches(0.3),
                 "首席教学设计师：罗老师  |  授课版本 2026",
                 size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(4.15), Inches(8.8), Inches(0.3),
                 "© 罗老师《AI 时代的家庭教育》· 著作权所有 · 未经授权禁止复制传播",
                 size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.3),
                 "期待在课堂上见到你。让我们一起，把「焦虑」变成「行动」。",
                 size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_oval(s, Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4), C_ACCENT)
    add_text_box(s, Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4), "111", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================
# 主执行
# =====================================================
print("开始生成 PPT...")

# 1. 封面
slide_01_cover()
# 2. 课程地图
slide_02_map()

# 3-7. 发刊词（5页）
slide_03()
slide_04()
slide_05()
slide_06()
slide_07()

# 8-13. 先导课（6页）
slide_08()
slide_09()
slide_10()
slide_11()
slide_12()
slide_13()

# 14-21. 问题一（8页）
slide_14()
slide_15()
slide_16()
slide_17()
slide_18()
slide_19()
slide_20()
slide_21()

# 22-29. 问题二（8页）
slide_22()
slide_23()
slide_24()
slide_25()
slide_26()
slide_27()
slide_28()
slide_29()

# 30-36. 问题三（7页）
slide_30()
slide_31()
slide_32()
slide_33()
slide_34()
slide_35()
slide_36()

# 37-43. 问题四（7页 - make_lecture_block）
# 44-50. 问题五（7页）
# 51-57. 问题六（7页）
# 58-64. 问题七（7页）
# 65-71. 问题八（7页）
# 72-78. 问题九（7页）
# 79-85. 问题十（7页）
# 86-92. 问题十一（7页）
# 93-99. 问题十二（7页）
# (问题四-十二由 make_lecture_block 在模块顶部已生成)

# 100-104. 结刊词（5页）
slide_100()
slide_101()
slide_102()
slide_103()
slide_104()

# 105-109. 附录（5页）
slide_105()
slide_106()
slide_107()
slide_108()
slide_109()

# 110-111. 致谢（2页）
slide_110()
slide_111()

# 输出 - 使用 win32 路径避免 bash 编码破坏
import os
out_dir = r"D:\2026年课程\ai课2026整理\AI时代的家庭教育\05_授课PPT"
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "AI时代的家庭教育_授课PPT.pptx")
prs.save(out)
print(f"已生成: {out}")
print(f"总页数: {len(prs.slides)}")
