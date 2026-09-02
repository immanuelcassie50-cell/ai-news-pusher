"""验证生成的PPT"""
import os, sys, io
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

out_dir = r"D:\2026年课程\ai课2026整理\AI时代的家庭教育\05_授课PPT"
out = os.path.join(out_dir, "AI时代的家庭教育_授课PPT.pptx")
p = Presentation(out)
print('总页数:', len(p.slides))
print('尺寸: %.2f x %.2f 英寸' % (p.slide_width/914400, p.slide_height/914400))
total_shapes = sum(len(s.shapes) for s in p.slides)
print('总shape数:', total_shapes)
print('平均shape/页:', total_shapes//len(p.slides))
print('文件大小: %.1f KB' % (os.path.getsize(out)/1024))
for i, s in enumerate(p.slides, 1):
    title = ""
    for shape in s.shapes:
        if shape.has_text_frame and shape.text_frame.text:
            title = shape.text_frame.text.split('\n')[0][:30]
            break
    if i <= 10 or i % 15 == 0 or i >= 105:
        print(f"  {i:3d}: {title}")
