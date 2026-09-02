#!/usr/bin/env python3
"""Compile all remaining slides and merge into one PPTX."""

import os
import subprocess
import re
import time

slides_dir = r"D:\新课开发\HR\薪酬\10.全面薪酬新叙事：弹性福利与非物质回报的AI个性化设计\授课PPT\slides"
output_pptx = r"D:\新课开发\HR\薪酬\10.全面薪酬新叙事：弹性福利与非物质回报的AI个性化设计\授课PPT\全面薪酬新叙事_授课PPT.pptx"

def get_slide_number(filename):
    match = re.search(r'slide-(\d+)', filename)
    if match:
        return int(match.group(1))
    return 999999

# Get all slide JS files
js_files = [f for f in os.listdir(slides_dir) if f.startswith('slide-') and f.endswith('.js')]
js_files.sort(key=get_slide_number)

print(f"Found {len(js_files)} slide JS files")

# Compile each slide
compiled = []
failed = []

for i, js_file in enumerate(js_files):
    js_path = os.path.join(slides_dir, js_file)
    try:
        result = subprocess.run(['node', js_path],
                              capture_output=True,
                              text=True,
                              timeout=30,
                              cwd=slides_dir)
        if result.returncode == 0:
            compiled.append(js_file)
            print(f"[{i+1}/{len(js_files)}] OK: {js_file}")
        else:
            failed.append(js_file)
            print(f"[{i+1}/{len(js_files)}] FAIL: {js_file} - {result.stderr[:100]}")
    except Exception as e:
        failed.append(js_file)
        print(f"[{i+1}/{len(js_files)}] ERROR: {js_file} - {str(e)[:100]}")

    # Small delay between compilations
    time.sleep(0.05)

print(f"\nCompiled: {len(compiled)}/{len(js_files)}")
print(f"Failed: {len(failed)}")

# Now merge all compiled PPTX files
from pptx import Presentation
from pptx.util import Inches

pptx_files = []
for f in os.listdir(slides_dir):
    if f.endswith('.pptx') and f.startswith('slide-') and 'preview' not in f:
        pptx_files.append(f)

pptx_files.sort(key=get_slide_number)
print(f"\nFound {len(pptx_files)} PPTX files to merge")

merged_prs = Presentation()
merged_prs.slide_width = Inches(10)
merged_prs.slide_height = Inches(5.625)

for i, pptx_file in enumerate(pptx_files):
    pptx_path = os.path.join(slides_dir, pptx_file)
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            slide_layout = merged_prs.slide_layouts[0]
            new_slide = merged_prs.slides.add_slide(slide_layout)
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    try:
                        new_slide.shapes.add_textbox(
                            shape.left, shape.top, shape.width, shape.height
                        ).text = shape.text
                    except:
                        pass
        print(f"[{i+1}/{len(pptx_files)}] Added: {pptx_file}")
    except Exception as e:
        print(f"[{i+1}/{len(pptx_files)}] Error with {pptx_file}: {e}")

merged_prs.save(output_pptx)
print(f"\nMerged presentation saved: {output_pptx}")
print(f"Total slides in merged presentation: {len(merged_prs.slides)}")
