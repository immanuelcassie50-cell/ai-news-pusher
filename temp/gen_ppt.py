# -*- coding: utf-8 -*-
"""PPT 生成脚本: 一线管理者的现代五项(150页)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ===== 配色 =====
NAVY = RGBColor(0x2E, 0x5C, 0x8A)        # 主深蓝
NAVY_DARK = RGBColor(0x1F, 0x3F, 0x5F)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)       # 暖橙
GREEN = RGBColor(0x27, 0xAE, 0x60)        # 亮绿
RED = RGBColor(0xC0, 0x39, 0x2B)          # 酒红
LIGHT = RGBColor(0xF5, 0xF5, 0xF0)        # 浅米
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
GRAY_LIGHT = RGBColor(0xD0, 0xD0, 0xD0)
GOLD = RGBColor(0xC8, 0xA1, 0x4A)

CN_TITLE = "微软雅黑"
CN_BODY = "微软雅轻"  # fallback to 微软雅黑 if not installed
EN_FONT = "Arial"

# ===== 创建 16:9 =====
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ===== 工具函数 =====
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    return shp


def add_round(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.15
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_oval(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=BLACK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=CN_BODY, line_h=None,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, list):
        for i, line in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if line_h is not None:
                p.line_spacing = line_h
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
            # east asian font
            rPr = r._r.get_or_add_rPr()
            ea = etree.SubElement(rPr, qn('a:ea'), typeface=CN_TITLE)
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        if line_h is not None:
            p.line_spacing = line_h
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        etree.SubElement(rPr, qn('a:ea'), typeface=CN_TITLE)
    return tb


def add_multi_text(slide, x, y, w, h, lines, font=CN_BODY, anchor=MSO_ANCHOR.TOP):
    """lines = [(text, size, color, bold, align), ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        etree.SubElement(rPr, qn('a:ea'), typeface=CN_TITLE)
        if i > 0:
            p.space_before = Pt(6)
    return tb


def page_chrome(slide, page_no, total=150, section=None, show_no=True):
    """顶部色条 + 页码 + 底部 logo 文字"""
    # 顶部细色条
    add_rect(slide, 0, 0, SW, Inches(0.18), fill=NAVY)
    add_rect(slide, 0, Inches(0.18), SW, Inches(0.05), fill=ORANGE)
    # 底部细线
    add_rect(slide, Inches(0.5), Inches(7.18), Inches(12.33), Emu(8000), fill=GRAY_LIGHT)
    # 底部 logo 文字
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.25),
             "一线管理者的现代五项  ·  国际版权课水准", size=9, color=GRAY)
    # 页码
    if show_no:
        add_text(slide, Inches(12.0), Inches(7.22), Inches(1.0), Inches(0.25),
                 f"{page_no:02d} / {total:02d}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)
    # 模块标签(左上角)
    if section:
        add_text(slide, Inches(0.5), Inches(0.05), Inches(8), Inches(0.18),
                 section, size=8, color=WHITE, bold=True, align=PP_ALIGN.LEFT)


def title_block(slide, title, subtitle=None, accent=ORANGE):
    """页内主标题块: 大字标题 + 副标题 + 左侧色条"""
    # 左侧色条
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.12), Inches(0.85), fill=accent)
    add_text(slide, Inches(0.75), Inches(0.40), Inches(11.5), Inches(0.7),
             title, size=28, color=NAVY_DARK, bold=True)
    if subtitle:
        add_text(slide, Inches(0.78), Inches(1.05), Inches(11.5), Inches(0.35),
                 subtitle, size=13, color=GRAY, font=CN_BODY)
    # 分隔线
    add_rect(slide, Inches(0.5), Inches(1.45), Inches(12.33), Emu(8000), fill=GRAY_LIGHT)


def bullet_block(slide, x, y, w, h, bullets, size=18, color=BLACK, line_h=1.4):
    """统一 bullet 列表"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, sc, co, bd = item
        else:
            text, sc, co, bd = item, size, color, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_h
        # 圆点
        r0 = p.add_run()
        r0.text = "■ "
        r0.font.size = Pt(sc)
        r0.font.color.rgb = ORANGE
        r0.font.bold = True
        r0.font.name = CN_TITLE
        rPr0 = r0._r.get_or_add_rPr()
        etree.SubElement(rPr0, qn('a:ea'), typeface=CN_TITLE)
        # 主文本
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sc)
        r.font.color.rgb = co
        r.font.bold = bd
        r.font.name = CN_TITLE
        rPr = r._r.get_or_add_rPr()
        etree.SubElement(rPr, qn('a:ea'), typeface=CN_TITLE)
        if i > 0:
            p.space_before = Pt(8)
    return tb


def quote_block(slide, x, y, w, h, text, author=None):
    """引用框: 橙色竖线 + 引文"""
    add_rect(slide, x, y, Emu(40000), h, fill=ORANGE)
    add_text(slide, x + Inches(0.18), y, w - Inches(0.2), h - Inches(0.1),
             f'"{text}"', size=14, color=NAVY_DARK, bold=False, font=CN_BODY)
    if author:
        add_text(slide, x + Inches(0.18), y + h - Inches(0.3), w - Inches(0.2), Inches(0.25),
                 f"— {author}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, big, small, color=NAVY):
    """大数字 + 小说明"""
    add_round(slide, x, y, w, h, fill=LIGHT, line=GRAY_LIGHT)
    add_text(slide, x, y + Inches(0.15), w, Inches(1.1),
             big, size=44, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.3), w, Inches(0.45),
             small, size=11, color=GRAY, align=PP_ALIGN.CENTER)


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


# ===================================================================
# 第 1 段: 封面 (1 页)
# ===================================================================
def slide_cover(idx):
    s = new_slide()
    # 左侧深色块
    add_rect(s, 0, 0, Inches(5.2), SH, fill=NAVY_DARK)
    # 橙色斜条
    add_rect(s, Inches(5.2), 0, Inches(0.1), SH, fill=ORANGE)
    # 大数字"五项"
    add_text(s, Inches(0.6), Inches(0.6), Inches(4.5), Inches(0.4),
             "FIVE  ·  ESSENTIALS", size=12, color=ORANGE, bold=True)
    add_text(s, Inches(0.6), Inches(1.5), Inches(4.5), Inches(2),
             "五项", size=120, color=WHITE, bold=True)
    add_text(s, Inches(0.6), Inches(4.0), Inches(4.5), Inches(0.5),
             "FOR FIRST-LINE MANAGERS", size=10, color=WHITE)
    # 右侧标题
    add_text(s, Inches(5.7), Inches(2.0), Inches(7.3), Inches(1.2),
             "一线管理者的", size=36, color=NAVY_DARK, bold=True)
    add_text(s, Inches(5.7), Inches(2.7), Inches(7.3), Inches(1.5),
             "现代五项", size=64, color=NAVY, bold=True)
    add_rect(s, Inches(5.7), Inches(4.3), Inches(1.0), Emu(30000), fill=ORANGE)
    add_text(s, Inches(5.7), Inches(4.45), Inches(7.3), Inches(0.5),
             "Five Essentials for First-Line Managers", size=14, color=GRAY)
    add_text(s, Inches(5.7), Inches(5.0), Inches(7.3), Inches(0.4),
             "复制成功 · 共谋抓手 · 应对难题 · 引领共创 · 前瞻思考", size=13, color=NAVY_DARK, bold=True)
    add_text(s, Inches(5.7), Inches(5.5), Inches(7.3), Inches(0.4),
             "12.5–13 小时  /  2 天线下工作坊", size=11, color=GRAY)
    # 底部信息
    add_text(s, Inches(5.7), Inches(6.6), Inches(7.3), Inches(0.4),
             "授课讲师: ___________________      日期: _____________", size=10, color=GRAY)
    add_text(s, Inches(5.7), Inches(6.9), Inches(7.3), Inches(0.4),
             "国际版权课水准  ·  DLP / Korn Ferry 级别", size=10, color=GRAY)
    # 左下小标
    add_text(s, Inches(0.6), Inches(6.6), Inches(4.5), Inches(0.4),
             "JINGYUE  /  LEADERSHIP  ACADEMY", size=10, color=WHITE)


# ===================================================================
# 第 2 段: 课程导览 (4 页)
# ===================================================================
def slide_01_course_map(idx):
    s = new_slide()
    page_chrome(s, idx, section="导览 · 课程全景")
    title_block(s, "课程全景: 你将获得的五件工具", "为什么是这五件? 因为它们正好回应了一线管理者最常遇到的三类困境")
    # 5 个工具卡片
    items = [
        ("M1", "螺旋深挖4问", "复制成功", "把'只有他会'变成'人人能学'", NAVY),
        ("M2", "花刺投票", "共谋抓手", "让集体共识真实发生", ORANGE),
        ("M3", "问题树 + 魔力提问", "应对难题", "把模糊难题拆成可下手方向", GREEN),
        ("M4", "高效脑暴双矩阵", "引领共创", "高质量发散 + 有依据聚焦", GOLD),
        ("M5", "推演双表格", "前瞻思考", "在行动前看清风险与机会", RED),
    ]
    card_w = Inches(2.35); card_h = Inches(2.7)
    gap = Inches(0.18)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (no, name, en, desc, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        add_rect(s, x, y, card_w, Inches(0.6), fill=color)
        add_text(s, x, y + Inches(0.1), card_w, Inches(0.4),
                 no, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.75), card_w, Inches(0.5),
                 name, size=15, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.2), card_w, Inches(0.4),
                 en, size=10, color=GRAY, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.6), y + Inches(1.65), card_w - Inches(1.2), Emu(15000), fill=color)
        add_text(s, x + Inches(0.15), y + Inches(1.75), card_w - Inches(0.3), Inches(0.9),
                 desc, size=11, color=BLACK, align=PP_ALIGN.CENTER)
    # 底部: 连接逻辑
    add_text(s, Inches(0.5), Inches(5.1), Inches(12.33), Inches(0.4),
             "连接逻辑", size=14, color=ORANGE, bold=True)
    add_text(s, Inches(0.5), Inches(5.45), Inches(12.33), Inches(0.5),
             "提炼成功原则  →  确定优先抓手  →  拆解关键难题  →  共创行动方案  →  做好前瞻准备",
             size=15, color=NAVY_DARK, bold=True)
    # 五箭头
    for i in range(4):
        x = Inches(2.7) + Inches(2.7) * i
        add_text(s, x, Inches(5.95), Inches(0.4), Inches(0.4), "→",
                 size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 最终产出
    add_round(s, Inches(0.5), Inches(6.45), Inches(12.33), Inches(0.55), fill=NAVY_DARK)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.45),
             "最终产出: 你的「管理行动地图」—— 带着 5 件工具的整合产出回到真实工作",
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def slide_02_three_pains(idx):
    s = new_slide()
    page_chrome(s, idx, section="导览 · 三大困境")
    title_block(s, "三类典型困境", "一线管理者普遍遇到 — 课程围绕它们各配了一件工具")
    # 三栏卡片
    items = [
        ("困境 1", "不擅长分析和解决新难题", "经验失灵\n用了以前有效的方法,结果和预期不符", NAVY),
        ("困境 2", "不擅长引导团队共创", "会议低效\n会议开了但什么也没决定", ORANGE),
        ("困境 3", "不擅长前瞻性思考", "救火式管理\n总在问题来时应对,而不是提前看到", GREEN),
    ]
    card_w = Inches(3.95); card_h = Inches(4.6)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(1.9)
    for i, (tag, title, desc, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        add_rect(s, x, y, card_w, Inches(0.7), fill=color)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.45),
                 tag, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.95), card_w, Inches(0.6),
                 title, size=18, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(1.2), y + Inches(1.65), card_w - Inches(2.4), Emu(25000), fill=color)
        for j, line in enumerate(desc.split("\n")):
            add_text(s, x + Inches(0.3), y + Inches(1.85) + Inches(0.4) * j,
                     card_w - Inches(0.6), Inches(0.4),
                     line, size=13, color=BLACK, align=PP_ALIGN.CENTER)
    # 对应工具
    tools = ["M1+M3 应对", "M2+M4 应对", "M5 应对"]
    for i, t in enumerate(tools):
        x = start_x + (card_w + gap) * i
        add_rect(s, x + Inches(0.4), y + card_h - Inches(0.7), card_w - Inches(0.8), Inches(0.5), fill=color)
        add_text(s, x + Inches(0.4), y + card_h - Inches(0.65), card_w - Inches(0.8), Inches(0.4),
                 t, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 底部金句
    quote_block(s, Inches(1.0), Inches(6.65), Inches(11.3), Inches(0.5),
                "可以学习的不是行为,而是行为背后的判断原则。",
                author="—— 第一部分核心信念")


def slide_03_yield(idx):
    s = new_slide()
    page_chrome(s, idx, section="导览 · 你将带走")
    title_block(s, "你将带走什么", "两天结束,你手里会有的 4 件具体产出")
    items = [
        ("01", "一个成功原则", "M1 螺旋深挖4问\n→ 可教给团队的成功经验", NAVY),
        ("02", "一份集体共识", "M2 花刺投票\n→ 团队共同认可的优先行动", ORANGE),
        ("03", "一套行动方案", "M3+M4 难题分析与共创\n→ 来自团队的高质量方案", GREEN),
        ("04", "一份前瞻推演", "M5 推演双表格\n→ 风险预案 + 机会准备", GOLD),
    ]
    card_w = Inches(2.95); card_h = Inches(4.2)
    gap = Inches(0.2)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (no, title, desc, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        # 圆形序号
        add_oval(s, x + Inches(1.0), y + Inches(0.3), Inches(1.0), Inches(1.0), fill=color)
        add_text(s, x + Inches(1.0), y + Inches(0.45), Inches(1.0), Inches(0.7),
                 no, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.5), card_w, Inches(0.5),
                 title, size=16, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.9), y + Inches(2.1), card_w - Inches(1.8), Emu(20000), fill=color)
        for j, line in enumerate(desc.split("\n")):
            t = line
            if t.startswith("→"):
                add_text(s, x + Inches(0.2), y + Inches(2.25) + Inches(0.4) * j,
                         card_w - Inches(0.4), Inches(0.4),
                         t, size=11, color=GRAY, align=PP_ALIGN.CENTER)
            else:
                add_text(s, x + Inches(0.2), y + Inches(2.25) + Inches(0.4) * j,
                         card_w - Inches(0.4), Inches(0.4),
                         t, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.5),
             "「 管理行动地图 」", size=22, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.4),
             "—— 五件工具的整合产出,可拍照放在手边,每隔两周看一眼",
             size=12, color=GRAY, align=PP_ALIGN.CENTER)


def slide_04_learning_roadmap(idx):
    s = new_slide()
    page_chrome(s, idx, section="导览 · 学习地图")
    title_block(s, "两天学习地图", "从认知 → 工具 → 案例 → 练习 → 整合")
    # 时间轴
    y_axis = Inches(4.0)
    add_rect(s, Inches(0.8), y_axis, Inches(11.7), Inches(0.06), fill=NAVY)
    # 5 个节点
    nodes = [
        ("Day 1 上午", "M1 复制成功", "螺旋深挖4问", NAVY),
        ("Day 1 下午", "M2 共谋抓手", "花刺投票", ORANGE),
        ("Day 1 傍晚", "M3 应对难题", "问题树+魔力提问", GREEN),
        ("Day 2 上午", "M4 引领共创", "高效脑暴双矩阵", GOLD),
        ("Day 2 下午", "M5 前瞻思考", "推演双表格", RED),
    ]
    n = len(nodes)
    for i, (t, name, sub, color) in enumerate(nodes):
        cx = Inches(0.8) + Inches(11.7 / (n - 1) * i)
        # 圆点
        add_oval(s, cx - Inches(0.22), y_axis - Inches(0.2), Inches(0.44), Inches(0.44), fill=color)
        add_oval(s, cx - Inches(0.12), y_axis - Inches(0.1), Inches(0.24), Inches(0.24), fill=WHITE)
        # 时间(下)
        add_text(s, cx - Inches(1.0), y_axis + Inches(0.4), Inches(2.0), Inches(0.4),
                 t, size=11, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
        # 名称(上)
        add_text(s, cx - Inches(1.0), y_axis - Inches(1.7), Inches(2.0), Inches(0.5),
                 name, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, cx - Inches(1.0), y_axis - Inches(1.15), Inches(2.0), Inches(0.4),
                 sub, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    # 顶部右侧金句框
    add_round(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.7), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, Inches(1.0), Inches(1.85), Inches(11.3), Inches(0.6),
             "从一个真实的管理场景出发 → 走完 5 件工具 → 带走可立即用的产出",
             size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    # 底部: 全程主线索
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.4),
             "全程主线索人物: 李明 (连锁电器零售店店长 · 12 人团队) + 助手小王 (35% 成交率)",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "—— 5 个模块全部用他的真实业务场景串起来",
             size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ===================================================================
# 第 3 段: 开场导入 (5 页)
# ===================================================================
def slide_05_self_check(idx):
    s = new_slide()
    page_chrome(s, idx, section="开场 · 学员心理诊断")
    title_block(s, "开场: 先给自己的管理水平做个诊断", "10 个问题 · 1 分钟 · 不评分,只为自我觉察")
    add_text(s, Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.4),
             "在以下情境里,你的第一反应更接近哪一个?",
             size=14, color=GRAY)
    items = [
        ("A", "团队里有人做得特别好,但其他人学不会", "我通常直接让ta分享感受", "我会用4问访谈提炼出原则"),
        ("B", "团队讨论了很久,散会后没人知道重点", "我通常最后总结一下", "我会用花刺投票聚焦共识"),
        ("C", "遇到了靠经验解决不了的新难题", "我会用过去最有效的办法再试", "我会先做问题树拆解再找新角度"),
        ("D", "团队会议经常冷场或乱场", "我尽量引导讨论", "我会用双矩阵结构化引导"),
        ("E", "总是在问题来了才应对", "救火本来就是管理者的工作", "我会用推演双表格提前准备"),
    ]
    y = Inches(2.2)
    row_h = Inches(0.85)
    for i, (no, sit, a, b) in enumerate(items):
        ry = y + row_h * i
        # 序号
        add_oval(s, Inches(0.5), ry + Inches(0.05), Inches(0.55), Inches(0.55), fill=NAVY)
        add_text(s, Inches(0.5), ry + Inches(0.1), Inches(0.55), Inches(0.45),
                 no, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 情境
        add_text(s, Inches(1.2), ry + Inches(0.05), Inches(4.5), Inches(0.7),
                 sit, size=12, color=BLACK)
        # A
        add_text(s, Inches(5.9), ry + Inches(0.05), Inches(3.6), Inches(0.7),
                 f"A  {a}", size=11, color=GRAY)
        # B
        add_text(s, Inches(9.6), ry + Inches(0.05), Inches(3.5), Inches(0.7),
                 f"B  {b}", size=11, color=NAVY_DARK, bold=True)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "选 A 多 = 习惯用经验 / 选 B 多 = 已经准备好升级工具 — 两种都是好起点",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def slide_06_five_tools_big(idx):
    s = new_slide()
    page_chrome(s, idx, section="开场 · 5 工具图示")
    title_block(s, "五件工具: 一个完整的链路", "M1→M2→M3→M4→M5 = 从成功提炼到行动前瞻的闭环")
    # 中央圆 + 五个外环
    cx, cy = Inches(6.67), Inches(4.5)
    add_oval(s, cx - Inches(1.0), cy - Inches(1.0), Inches(2.0), Inches(2.0), fill=NAVY_DARK)
    add_text(s, cx - Inches(1.0), cy - Inches(0.5), Inches(2.0), Inches(0.5),
             "管理", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx - Inches(1.0), cy, Inches(2.0), Inches(0.5),
             "行动地图", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 五个外环工具(圆形排列)
    radius = Inches(2.6)
    import math
    tools_pos = [
        ("M1", "复制成功", ORANGE),
        ("M2", "共谋抓手", GREEN),
        ("M3", "应对难题", GOLD),
        ("M4", "引领共创", RED),
        ("M5", "前瞻思考", NAVY),
    ]
    n = 5
    for i, (no, name, color) in enumerate(tools_pos):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        x = cx + radius * math.cos(angle) - Inches(0.7)
        y = cy + radius * math.sin(angle) - Inches(0.7)
        add_oval(s, x, y, Inches(1.4), Inches(1.4), fill=color)
        add_text(s, x, y + Inches(0.25), Inches(1.4), Inches(0.5),
                 no, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.75), Inches(1.4), Inches(0.4),
                 name, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 连线到中心
        add_line(s, cx, cy, x + Inches(0.7), y + Inches(0.7), color=GRAY_LIGHT)
    # 箭头
    for i in range(n):
        a1 = -math.pi / 2 + 2 * math.pi * i / n
        a2 = -math.pi / 2 + 2 * math.pi * ((i + 1) % n) / n
        ax1 = cx + radius * math.cos(a1) - Inches(0.7)
        ay1 = cy + radius * math.sin(a1) - Inches(0.7)
        ax2 = cx + radius * math.cos(a2) - Inches(0.7)
        ay2 = cy + radius * math.sin(a2) - Inches(0.7)
        add_line(s, ax1 + Inches(0.7), ay1 + Inches(0.7),
                 ax2 + Inches(0.7), ay2 + Inches(0.7), color=ORANGE, w=Inches(0.05))
    # 底部金句
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3),
             "记住一个名字: 这是一条从「 提炼已有经验 」到「 带着团队实现它 」的完整链路。",
             size=11, color=GRAY, align=PP_ALIGN.CENTER)


def add_line(slide, x1, y1, x2, y2, color=GRAY_LIGHT, w=None):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    if w is not None:
        line.line.width = w
    else:
        line.line.width = Pt(1.5)
    return line


def slide_07_li_ming_intro(idx):
    s = new_slide()
    page_chrome(s, idx, section="开场 · 主线索人物")
    title_block(s, "全程主线索: 李明的故事", "你将看到的所有工具,都会用这个真实场景演示一遍")
    # 左侧: 李明画像(几何)
    lx = Inches(0.8); ly = Inches(2.0)
    add_round(s, lx, ly, Inches(4.5), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_oval(s, lx + Inches(1.55), ly + Inches(0.4), Inches(1.4), Inches(1.4), fill=NAVY)
    add_text(s, lx + Inches(1.55), ly + Inches(0.7), Inches(1.4), Inches(0.8),
             "李明", size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, lx + Inches(0.4), ly + Inches(1.9), Inches(3.7), Inches(0.4),
             "连锁电器零售店 · 店长", size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, lx + Inches(0.4), ly + Inches(2.3), Inches(3.7), Inches(0.4),
             "团队规模: 12 人", size=12, color=GRAY, align=PP_ALIGN.CENTER)
    # 数据卡
    stats = [("21%", "门店成交率"), ("23%", "团队平均成交率"), ("35%", "小王成交率")]
    for i, (big, sm) in enumerate(stats):
        sx = lx + Inches(0.2) + Inches(1.45) * i
        sy = ly + Inches(2.85)
        add_round(s, sx, sy, Inches(1.25), Inches(1.4), fill=WHITE, line=ORANGE)
        add_text(s, sx, sy + Inches(0.25), Inches(1.25), Inches(0.6),
                 big, size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, sx, sy + Inches(0.85), Inches(1.25), Inches(0.4),
                 sm, size=9, color=GRAY, align=PP_ALIGN.CENTER)
    # 右侧: 故事线
    rx = Inches(5.7); ry = Inches(2.0)
    add_text(s, rx, ry, Inches(7.3), Inches(0.5),
             "故事起点", size=16, color=ORANGE, bold=True)
    add_text(s, rx, ry + Inches(0.4), Inches(7.3), Inches(1.0),
             "新竞争对手开业, 客流下滑。李明要回答一个问题:\n「 接下来这个季度,我们把力气放在哪里? 」",
             size=13, color=BLACK)
    # 工具路线
    steps = [
        ("M1", "向小王做4问访谈 — 提炼「 体验式接待 」原则", ORANGE),
        ("M2", "用花刺投票 — 团队确认优先机会和关键障碍", GREEN),
        ("M3", "问题树拆解 — 找到「 没有体验引导 」切入点", GOLD),
        ("M4", "双矩阵脑暴 — 共创出 3 条优先行动", RED),
        ("M5", "推演双表格 — 风险预案 + 机会准备", NAVY),
    ]
    for i, (no, desc, color) in enumerate(steps):
        sy = ry + Inches(1.6) + Inches(0.5) * i
        add_rect(s, rx, sy + Inches(0.08), Inches(0.55), Inches(0.35), fill=color)
        add_text(s, rx, sy + Inches(0.08), Inches(0.55), Inches(0.35),
                 no, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, rx + Inches(0.75), sy, Inches(6.5), Inches(0.5),
                 desc, size=12, color=BLACK)


def slide_08_promise(idx):
    s = new_slide()
    page_chrome(s, idx, section="开场 · 学习承诺")
    title_block(s, "我们对这门课的三条承诺", "不是教理论 — 是让你下周一回到办公室就能用")
    items = [
        ("01", "可立即用",
         "每个工具都配操作步骤\n模板清单可直接打印",
         "不是 PPT 上好看的图\n是你桌上能用的表", NAVY),
        ("02", "可被检验",
         "案例 + 练习 + 角色演练\n每个工具至少演练一次",
         "不是听一遍就完\n是用一次才能带走", ORANGE),
        ("03", "可带着走",
         "管理行动地图\n30 天 / 60 天 / 90 天节奏",
         "不是结束就结束\n是回到工作里真发生改变", GREEN),
    ]
    card_w = Inches(4.0); card_h = Inches(4.6)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(1.95)
    for i, (no, title, desc, foot, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        # 序号大圆
        add_oval(s, x + Inches(1.35), y + Inches(0.4), Inches(1.3), Inches(1.3), fill=color)
        add_text(s, x + Inches(1.35), y + Inches(0.55), Inches(1.3), Inches(0.9),
                 no, size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.85), card_w, Inches(0.5),
                 title, size=18, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        # 描述
        for j, line in enumerate(desc.split("\n")):
            add_text(s, x + Inches(0.3), y + Inches(2.45) + Inches(0.35) * j,
                     card_w - Inches(0.6), Inches(0.35),
                     line, size=11, color=BLACK, align=PP_ALIGN.CENTER)
        # 底部分隔
        add_rect(s, x + Inches(0.4), y + Inches(3.4), card_w - Inches(0.8), Emu(15000), fill=color)
        add_text(s, x + Inches(0.3), y + Inches(3.55), card_w - Inches(0.6), Inches(1.0),
                 foot, size=11, color=GRAY, align=PP_ALIGN.CENTER, font=CN_BODY)
    # 底部
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "今天下午, 你就会用第一个工具对你的真实场景做一次完整的演练。",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def slide_09_ground_rules(idx):
    s = new_slide()
    page_chrome(s, idx, section="开场 · 课堂公约")
    title_block(s, "课堂公约: 让两天价值最大化", "不是规矩 — 是对我们自己时间的尊重")
    items = [
        ("📵 手机", "静音或震动\n紧急电话请外出接听", NAVY),
        ("💬 真实场景", "练习用你自己的真实管理难题\n不是假设", ORANGE),
        ("⏰ 时间边界", "上午 9:00–12:30\n下午 14:00–17:30", GREEN),
        ("🤝 互助", "同桌 = 你的演练伙伴\n小组 = 你的对照镜子", GOLD),
    ]
    card_w = Inches(2.95); card_h = Inches(4.0)
    gap = Inches(0.2)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (icon_title, desc, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        # 顶部色块
        add_rect(s, x, y, card_w, Inches(0.8), fill=color)
        # 图标+标题
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.55),
                 icon_title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 描述
        for j, line in enumerate(desc.split("\n")):
            add_text(s, x + Inches(0.2), y + Inches(1.1) + Inches(0.45) * j,
                     card_w - Inches(0.4), Inches(0.45),
                     line, size=11, color=NAVY_DARK, align=PP_ALIGN.CENTER, line_h=1.3)
    # 底部
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.5),
             "一个简单承诺", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.4),
             "「 离开教室前, 至少有一件工具你在真实场景里用过一次。 」",
             size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


print("基础结构已就绪 — 共 10 页框架")


# ===================================================================
# 第 4 段补: 模块一 复制成功 (28 页)
# ===================================================================
def m1_section_divider(idx):
    s = new_slide()
    page_chrome(s, idx, section="模块 1 · 复制成功")
    add_rect(s, 0, Inches(1.5), SW, Inches(4.5), fill=ORANGE)
    add_text(s, Inches(0.5), Inches(1.7), Inches(5), Inches(3),
             "M1", size=180, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(2.3), Inches(8), Inches(1.0),
             "复制成功", size=44, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(3.5), Inches(8), Inches(0.6),
             "Replicate What Works", size=18, color=WHITE)
    quote_block(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.6),
                "把团队里「 只有一个人会 」的成功, 变成「 人人能学 」。",
                author="—— 本模块核心信念")


def slide_m1(idx, n, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section="模块 1 · 复制成功")
    title_block(s, title, subtitle, accent=ORANGE)
    return s


def m1_01(idx):
    s = slide_m1(idx, 11, "开场: 团队里那个「 只有他会 」的人",
                 "他到底做对了什么, 你能说清楚吗?")
    # 故事
    add_round(s, Inches(0.5), Inches(1.85), Inches(7.5), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, Inches(0.7), Inches(2.0), Inches(7.1), Inches(0.5),
             "📖 李明门店里的小王", size=14, color=ORANGE, bold=True)
    add_text(s, Inches(0.7), Inches(2.55), Inches(7.1), Inches(3.8),
             "小王是店里最年轻的新人, 但成交率是 35%, 是团队平均水平的两倍。\n\n李明不是没试过让其他人「 学小王 」, 晨会上让小王分享, 其他导购听的时候点头, 但执行起来就是不一样。\n\n李明的问题: 「 我让小王分享了, 大家也听了, 为什么还是只有小王能做到? 」",
             size=12, color=BLACK, line_h=1.5)
    # 右: 病根
    rx = Inches(8.2); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(4.65), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(4.0), Inches(0.5),
             "常见的误判", size=14, color=ORANGE, bold=True)
    items = [
        "他「 天生会卖东西 」",
        "他「 性格好, 客户喜欢他 」",
        "他「 认真用心 」",
        "他「 经验比我们多 」",
    ]
    for i, t in enumerate(items):
        yy = ry + Inches(0.9) + Inches(0.8) * i
        add_oval(s, rx + Inches(0.3), yy + Inches(0.1), Inches(0.3), Inches(0.3), fill=ORANGE)
        add_text(s, rx + Inches(0.7), yy, Inches(3.7), Inches(0.5),
             "✗ " + t, size=12, color=WHITE, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "→ 误判的背后, 是没有一套能挖掘出「 行为 - 动机 - 原则 - 路径 」的工具",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_02(idx):
    s = slide_m1(idx, 12, "为什么经验会失灵 — 三个根本原因",
                 "没有挖掘工具, 你看到的就只是「 表面 」")
    items = [
        ("原因 1", "成功者说不出自己在做什么", "「 成功直觉化 」\n他做了但说不出来\n你听了但抓不住", ORANGE),
        ("原因 2", "管理者听到的是抽象表扬", "「 他很认真 」\n「 他比较用心 」\n「 他爱动脑 」\n全是形容词, 不是动作", GOLD),
        ("原因 3", "复制变成了「 学精神 」", "不学动作, 学精神\n学完回去还是不知道\n第一步做什么", RED),
    ]
    card_h = Inches(1.45)
    for i, (no, name, sub, col) in enumerate(items):
        y = Inches(2.0) + (card_h + Inches(0.15)) * i
        add_round(s, Inches(0.5), y, Inches(12.33), card_h, fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(1.4), card_h, fill=col)
        add_text(s, Inches(0.5), y + Inches(0.3), Inches(1.4), Inches(0.5),
                 no, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.5), y + Inches(0.85), Inches(1.4), Inches(0.4),
                 name.split()[1] if len(name.split()) > 1 else "", size=10, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.1), y + Inches(0.15), Inches(4.0), Inches(0.5),
                 name, size=14, color=NAVY_DARK, bold=True)
        add_text(s, Inches(2.1), y + Inches(0.6), Inches(4.0), Inches(0.8),
                 sub.split("\n")[0], size=11, color=BLACK)
        add_rect(s, Inches(6.5), y + Inches(0.2), Inches(0.05), card_h - Inches(0.4), fill=GRAY_LIGHT)
        add_text(s, Inches(6.7), y + Inches(0.2), Inches(5.8), Inches(0.4),
                 "具体表现", size=10, color=GRAY, bold=True)
        for j, line in enumerate(sub.split("\n")[1:] if "\n" in sub else []):
            add_text(s, Inches(6.7), y + Inches(0.55) + Inches(0.4) * j,
                     Inches(5.8), Inches(0.4),
                     "• " + line, size=11, color=BLACK, line_h=1.3)
        if "认真" in sub or "用心" in sub or "动脑" in sub:
            add_text(s, Inches(6.7), y + Inches(0.55), Inches(5.8), Inches(0.4),
                     "• 「 他很认真 」 / 「 他比较用心 」 / 「 他爱动脑 」", size=11, color=BLACK)
        elif "学精神" in sub or "第一步" in sub:
            add_text(s, Inches(6.7), y + Inches(0.55), Inches(5.8), Inches(0.4),
                     "• 不学动作, 学精神 · 不知道第一步做什么", size=11, color=BLACK)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.3),
             "3 个原因指向同一个解: 一套能挖出「 行为 - 动机 - 原则 - 路径 」的工具",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_03(idx):
    s = slide_m1(idx, 13, "螺旋深挖 4 问 — 工具全景",
                 "Q1 行为 → Q2 动机 → Q3 原则 → Q4 路径 · 任一问没问清回到上一问")
    items = [
        ("Q1 行为", "做了什么不同的?", "必须是可观察的动作", "不是「 认真用心 」", ORANGE),
        ("Q2 动机", "为什么这样做?", "跳过这问\nQ3 必然空洞", "听真实想法", GOLD),
        ("Q3 原则", "能提炼成原则吗?", "必须回答「 为什么有效 」\n不是行为的另一种说法", "听原则", RED),
        ("Q4 路径", "教新人, 关键 2-3 步?", "步骤 + 原则支撑\n缺一不可", "听步骤", GREEN),
    ]
    card_w = Inches(2.95); card_h = Inches(3.0)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = Inches(0.5)
    y = Inches(2.0)
    for i, (no, name, std, listen, col) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.7), fill=col)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.4),
                 no, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.85), card_w - Inches(0.4), Inches(0.5),
                 name, size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(std.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(1.4) + Inches(0.4) * j,
                     card_w - Inches(0.5), Inches(0.4),
                     line, size=10, color=BLACK, line_h=1.3, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.3), y + card_h - Inches(0.5), card_w - Inches(0.6), Emu(15000), fill=col)
        add_text(s, x + Inches(0.3), y + card_h - Inches(0.4), card_w - Inches(0.6), Inches(0.4),
                 "听什么: " + listen, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(5.3), Inches(12.33), Inches(0.5),
             "「 螺旋 」含义: 任何一问没问清, 回到上一问再追问",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_round(s, Inches(0.5), Inches(5.95), Inches(12.33), Inches(1.0), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.4),
             "为什么是「 螺旋 」? 不是一个问完问下一个 — 是一圈一圈地深化",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.5),
             "比如: Q3 答案听起来还是行为, 就回到 Q1 再问一次「 那当时具体你还做了什么 」",
             size=11, color=WHITE, align=PP_ALIGN.CENTER, line_h=1.3)


def m1_04(idx):
    s = slide_m1(idx, 14, "Q1 行为 — 听动词, 不是形容词",
                 "问自己: 我听到一个「 动作 」, 还是「 描述 」?")
    # 好例子 vs 坏例子
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "✗ 不是 Q1 行为 (形容词 / 状态)", size=14, color=RED, bold=True)
    bad = [
        "他「 认真 」",
        "他「 用心 」",
        "他「 努力 」",
        "他「 比较投入 」",
        "他「 比较专业 」",
    ]
    for i, t in enumerate(bad):
        yy = ly + Inches(0.85) + Inches(0.6) * i
        add_rect(s, lx + Inches(0.3), yy, Inches(0.1), Inches(0.5), fill=RED)
        add_text(s, lx + Inches(0.5), yy, Inches(5.3), Inches(0.5),
                 t, size=13, color=BLACK)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "✓ 才是 Q1 行为 (可观察的动作)", size=14, color=GREEN, bold=True)
    good = [
        "客户进店 2 分钟内主动邀请体验",
        "客户体验时安静站在 1 米外观察",
        "等客户主动提问后再针对性解释",
        "不主动推销, 主动问客户感受",
        "对每个客户做 1 个动作的差异化",
    ]
    for i, t in enumerate(good):
        yy = ry + Inches(0.85) + Inches(0.6) * i
        add_rect(s, rx + Inches(0.3), yy, Inches(0.1), Inches(0.5), fill=GREEN)
        add_text(s, rx + Inches(0.5), yy, Inches(5.3), Inches(0.5),
                 t, size=13, color=BLACK)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "判断标准: 别人看到这个动作, 能立刻判断「 你做了 / 没做 」",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_05(idx):
    s = slide_m1(idx, 15, "Q2 动机 — 听原因, 不是动作",
                 "跳过 Q2, Q3 必然空洞")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "Q2 的本质: 理解他「 为什么这样做 」, 不是「 他还做了什么 」",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 引导问法
    prompts = [
        ("问 1", "「 你当时为什么这样做? 」", "不是问「 你还做了什么 」"),
        ("问 2", "「 你当时是怎么想的? 」", "打开他做决定时的思考过程"),
        ("问 3", "「 是因为____, 才这样做的? 」", "给出猜测, 让他确认或纠正"),
    ]
    card_w = Inches(4.0); card_h = Inches(2.5)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.7)
    for i, (no, q, note) in enumerate(prompts):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GOLD)
        add_rect(s, x, y, card_w, Inches(0.5), fill=GOLD)
        add_text(s, x, y + Inches(0.08), card_w, Inches(0.4),
                 no, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.6), Inches(0.7),
                 q, size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, line_h=1.3)
        add_rect(s, x + Inches(0.3), y + Inches(1.6), card_w - Inches(0.6), Emu(15000), fill=GOLD)
        add_text(s, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(0.7),
                 note, size=11, color=WHITE, align=PP_ALIGN.CENTER, line_h=1.3)
    # 警示
    add_round(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5), fill=RED)
    add_text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.4),
             "⚠ 常见错误", size=14, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.5),
             "「 你为什么这样做? 」→ 「 因为我想卖出去 」",
             size=12, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.5),
             "这不是真正的动机, 是目标 — 没有告诉你他做了什么不同的思考",
             size=11, color=WHITE)


def m1_06(idx):
    s = slide_m1(idx, 16, "Q3 原则 — 听原则, 不是行为",
                 "必须回答「 为什么有效 」, 不是行为的另一种说法")
    # 好 Q3 vs 坏 Q3
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "✗ 假 Q3 原则 (还是行为)", size=14, color=RED, bold=True)
    bad = [
        "「 让客户多体验产品 」",
        "「 主动和客户互动 」",
        "「 给客户解释清楚 」",
        "「 用专业能力说服客户 」",
    ]
    for i, t in enumerate(bad):
        yy = ly + Inches(0.85) + Inches(0.6) * i
        add_text(s, lx + Inches(0.3), yy, Inches(0.4), Inches(0.5),
                 "✗", size=14, color=RED, bold=True)
        add_text(s, lx + Inches(0.7), yy, Inches(5.0), Inches(0.5),
                 t, size=13, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(3.7), Inches(5.6), Inches(0.85),
             "判断: 把这个原则告诉别人,\n他不知道「 怎么用 」, 还是行动描述",
             size=10, color=GRAY, italic=True)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "✓ 真 Q3 原则 (回答了为什么有效)", size=14, color=GREEN, bold=True)
    good = [
        "「 让客户主动产生问题, 比导购主动\n  传递信息更接近成交 」",
        "「 客户的好奇心比销售话术有效 」",
        "「 真实体验比介绍更能建立信任 」",
    ]
    for i, t in enumerate(good):
        yy = ry + Inches(0.85) + Inches(1.0) * i
        add_text(s, rx + Inches(0.3), yy, Inches(0.4), Inches(0.5),
                 "✓", size=14, color=GREEN, bold=True)
        add_text(s, rx + Inches(0.7), yy, Inches(5.0), Inches(1.0),
                 t, size=12, color=NAVY_DARK, bold=True, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "真 Q3 原则的判断: 它回答了「 为什么这个动作有效 」的本质机制",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_07(idx):
    s = slide_m1(idx, 17, "Q4 路径 — 听步骤, 不是原则",
                 "步骤 + 原则支撑, 缺一不可")
    items = [
        ("步骤 1", "客户进店 2 分钟内邀请体验", "依据 Q3 原则\n「 客户主动产生好奇心 」"),
        ("步骤 2", "客户体验时安静站在 1 米外", "依据 Q3 原则\n「 不打断体验 」"),
        ("步骤 3", "针对客户主动提问做解释", "依据 Q3 原则\n「 回应问题, 不主动推销 」"),
    ]
    card_h = Inches(1.4)
    for i, (no, step, basis) in enumerate(items):
        y = Inches(2.0) + (card_h + Inches(0.2)) * i
        add_round(s, Inches(0.5), y, Inches(12.33), card_h, fill=LIGHT, line=ORANGE)
        add_rect(s, Inches(0.5), y, Inches(1.5), card_h, fill=ORANGE)
        add_text(s, Inches(0.5), y + Inches(0.4), Inches(1.5), Inches(0.5),
                 no, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.2), y + Inches(0.15), Inches(7.5), Inches(0.5),
                 step, size=14, color=NAVY_DARK, bold=True)
        add_text(s, Inches(2.2), y + Inches(0.65), Inches(7.5), Inches(0.6),
                 "具体动作: 是什么, 做到什么程度", size=10, color=BLACK)
        add_rect(s, Inches(9.9), y + Inches(0.2), Inches(0.05), card_h - Inches(0.4), fill=GRAY_LIGHT)
        add_text(s, Inches(10.1), y + Inches(0.2), Inches(2.7), Inches(0.4),
                 "依据 (Q3)", size=10, color=GRAY, bold=True)
        for j, line in enumerate(basis.split("\n")[1:] if "\n" in basis else []):
            add_text(s, Inches(10.1), y + Inches(0.55) + Inches(0.4) * j,
                     Inches(2.7), Inches(0.4),
                     line, size=10, color=ORANGE, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "Q4 的关键: 不是「 我要做什么 」, 而是「 教给新人, 第一步是什么 」",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_08(idx):
    s = slide_m1(idx, 18, "螺旋深挖 4 问 — 完整示例: 小王 (1/3)",
                 "主题: 怎么让客户对产品产生真实兴趣?")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "主题: 让客户对产品产生真实兴趣 → 提升成交率",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # Q1
    add_round(s, Inches(0.5), Inches(2.6), Inches(12.33), Inches(1.4), fill=LIGHT, line=ORANGE)
    add_rect(s, Inches(0.5), Inches(2.6), Inches(0.15), Inches(1.4), fill=ORANGE)
    add_text(s, Inches(0.8), Inches(2.7), Inches(2.0), Inches(0.4),
             "Q1 行为", size=14, color=ORANGE, bold=True)
    add_text(s, Inches(0.8), Inches(3.1), Inches(2.0), Inches(0.4),
             "做了什么?", size=10, color=GRAY)
    add_text(s, Inches(3.0), Inches(2.7), Inches(9.5), Inches(1.2),
             "「 客户进店 2 分钟内, 我会主动走到客户身边, 邀请他\n体验产品具体功能, 而不是站在旁边等他来问 」\n→ 是动作, 不是「 认真用心 」",
             size=12, color=BLACK, line_h=1.4)
    # Q2
    add_round(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(1.4), fill=LIGHT, line=GOLD)
    add_rect(s, Inches(0.5), Inches(4.1), Inches(0.15), Inches(1.4), fill=GOLD)
    add_text(s, Inches(0.8), Inches(4.2), Inches(2.0), Inches(0.4),
             "Q2 动机", size=14, color=GOLD, bold=True)
    add_text(s, Inches(0.8), Inches(4.6), Inches(2.0), Inches(0.4),
             "为什么?", size=10, color=GRAY)
    add_text(s, Inches(3.0), Inches(4.2), Inches(9.5), Inches(1.2),
             "「 我发现客户亲身体验产品的时候, 会主动产生\n问题 — 是真实的顾虑, 不是抽象的。\n我回答问题时客户是真正在听 」",
             size=12, color=BLACK, line_h=1.4)
    # Q3 Q4 占位
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.33), Inches(0.4),
             "→ 下两页继续 Q3 原则 + Q4 路径",
             size=11, color=ORANGE, bold=True, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.4),
             "(每张图都是真实访谈场景, 不是抽象框架)",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def m1_09(idx):
    s = slide_m1(idx, 19, "螺旋深挖 4 问 — 完整示例: 小王 (2/3)",
                 "Q3 原则: 挖掘到的「 为什么有效 」")
    add_round(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(2.4), fill=LIGHT, line=RED)
    add_rect(s, Inches(0.5), Inches(1.85), Inches(0.15), Inches(2.4), fill=RED)
    add_text(s, Inches(0.8), Inches(2.0), Inches(3.0), Inches(0.5),
             "Q3 原则", size=18, color=RED, bold=True)
    add_text(s, Inches(0.8), Inches(2.5), Inches(3.0), Inches(0.4),
             "可提炼成原则吗?", size=10, color=GRAY)
    add_text(s, Inches(4.2), Inches(2.0), Inches(8.5), Inches(0.5),
             "李明第一次问:  「 是不是让客户多体验就行? 」",
             size=12, color=BLACK)
    add_text(s, Inches(4.2), Inches(2.55), Inches(8.5), Inches(0.5),
             "小王:           「 不是, 是让客户主动产生问题 」",
             size=12, color=NAVY_DARK, bold=True)
    add_text(s, Inches(4.2), Inches(3.05), Inches(8.5), Inches(1.2),
             "李明追问:  「 为什么? 」\n小王:           「 因为客户主动问的问题, 是他真正关心的;\n                   我主动讲的, 他其实没在听 」",
             size=12, color=BLACK, line_h=1.5)
    # 提炼后的原则
    add_round(s, Inches(0.5), Inches(4.4), Inches(12.33), Inches(1.4), fill=GOLD)
    add_text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.4),
             "💡 提炼后的 Q3 原则", size=14, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.9),
             "「 让客户主动产生问题, 比导购主动传递信息更接近成交 」",
             size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_h=1.4)
    # 螺旋
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.4),
             "注意: 小王第一版回答是行为 (「 多体验 」), 李明没接受, 回到 Q1 再问",
             size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "→ 这就是「 螺旋 」— 答案不到位, 回到上一问再问一次",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_10(idx):
    s = slide_m1(idx, 20, "螺旋深挖 4 问 — 完整示例: 小王 (3/3)",
                 "Q4 路径: 教给新人, 关键 2-3 步")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.4),
             "小王基于 Q3 原则, 教新人「 体验式接待 3 步法 」:",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("步骤 1", "2 分钟内邀请客户体验", "对应 Q3: 让客户有机会主动产生好奇心"),
        ("步骤 2", "客户体验时安静观察 30 秒", "对应 Q3: 创造客户主动提问的空间"),
        ("步骤 3", "针对客户主动提问, 做精准回答", "对应 Q3: 回应真实问题, 不主动推销"),
    ]
    card_h = Inches(1.2)
    for i, (no, step, basis) in enumerate(items):
        y = Inches(2.5) + (card_h + Inches(0.15)) * i
        add_round(s, Inches(0.5), y, Inches(12.33), card_h, fill=LIGHT, line=GREEN)
        add_rect(s, Inches(0.5), y, Inches(1.5), card_h, fill=GREEN)
        add_text(s, Inches(0.5), y + Inches(0.3), Inches(1.5), Inches(0.5),
                 no, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.2), y + Inches(0.15), Inches(5.0), Inches(0.5),
                 step, size=14, color=NAVY_DARK, bold=True)
        add_text(s, Inches(2.2), y + Inches(0.65), Inches(5.0), Inches(0.5),
                 "动作 (要做什么, 做到什么程度)", size=10, color=GRAY)
        add_rect(s, Inches(7.5), y + Inches(0.2), Inches(0.05), card_h - Inches(0.4), fill=GRAY_LIGHT)
        add_text(s, Inches(7.7), y + Inches(0.4), Inches(5.1), Inches(0.6),
                 basis, size=12, color=GREEN, bold=True, line_h=1.3)
    # 验证
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "Q4 的验证: 新人拿这 3 步能上手吗?",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.3),
             "小王: 「 我就是这么做的, 上周新人 3 天学会了, 现在月成交 12 单 」",
             size=10, color=NAVY_DARK, italic=True, align=PP_ALIGN.CENTER)


def m1_11(idx):
    s = slide_m1(idx, 21, "⚡ 完整示例: 李明在周一访谈小王",
                 "从「 你认真点 」到能用的 3 步法")
    # 模拟对话
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.4),
             "周一早上 8:30 · 门店休息区 · 李明访谈小王",
             size=12, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    lines = [
        ("李明", "小王, 我看你成交率特别高, 到底做对了什么?", "Q1"),
        ("小王", "我…我也不太知道, 就是凭感觉", "没答案"),
        ("李明", "那具体说一件事 — 上周那个 50 寸电视的客户, 你怎么从进门聊到签单的?", "Q1 改问"),
        ("小王", "他进门看了 2 分钟左右, 我就过去跟他说, 要不要亲自试一下?", "Q1 ✓"),
        ("李明", "为什么是 2 分钟, 不是 1 分钟, 也不是 5 分钟?", "Q2"),
        ("小王", "我觉得客户刚进门, 还在看, 给他点时间; 太早了打断他, 他会烦", "Q2 ✓"),
        ("李明", "那之后呢?", "Q1 续"),
        ("小王", "让他自己摸, 自己用, 我站在旁边不说话, 等他来问", "Q1 ✓"),
        ("李明", "为什么等他问, 不你主动讲?", "Q2 续"),
        ("小王", "因为…(想了一会) 他自己问的, 是他真关心的; 我主动讲的, 他其实没在听", "Q3 ✓"),
        ("李明", "那你能不能把这套做法教给新来的小张?", "Q4"),
        ("小王", "能啊 — 2 分钟内邀请体验, 体验时安静, 等他问再答", "Q4 ✓"),
    ]
    y = Inches(2.45)
    for who, text, q in lines:
        bg = NAVY if who == "李明" else GOLD
        tx_color = WHITE if who == "李明" else NAVY_DARK
        add_round(s, Inches(0.5), y, Inches(9.5), Inches(0.32), fill=bg)
        add_text(s, Inches(0.7), y + Inches(0.04), Inches(1.5), Inches(0.3),
                 who, size=10, color=tx_color, bold=True)
        add_text(s, Inches(2.3), y + Inches(0.04), Inches(7.5), Inches(0.3),
                 text, size=9, color=tx_color)
        # Q 标签
        qcol = ORANGE if "Q" in q and "✓" in q else GRAY
        add_text(s, Inches(10.2), y + Inches(0.04), Inches(2.5), Inches(0.3),
                 q, size=9, color=qcol, bold=True)
        y += Inches(0.34)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "→ 4 问用对了, 一次访谈就能挖出可复制的成功经验",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_12(idx):
    s = slide_m1(idx, 22, "✋ 练习一: 找一位「 高手 」访谈",
                 "用 4 问深挖 · 时长 20-30 分钟 · 一对一")
    # 左
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "访谈前准备", size=14, color=ORANGE, bold=True)
    prep = [
        "选一位你团队里的「 高手 」 — 做某件事明显好于其他人",
        "提前打招呼: 「 我想跟你聊聊 30 分钟, 看看你做对了什么 」",
        "准备一个具体场景: 上周 / 上个月一次具体的事",
        "录音 (征得同意) · 准备 4 问的小卡片",
    ]
    for i, t in enumerate(prep):
        add_text(s, lx + Inches(0.3), ly + Inches(0.85) + Inches(0.85) * i,
                 Inches(5.5), Inches(0.8),
                 f"{i+1}.  {t}", size=12, color=BLACK, line_h=1.4)
    # 右
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "访谈中提醒自己", size=14, color=ORANGE, bold=True)
    tips = [
        ("Q1 听到形容词", "马上追问: 「 当时具体做了什么? 」"),
        ("Q3 听到还是行为", "回到 Q1 再问: 「 那个动作具体是怎样? 」"),
        ("对方卡住", "用: 「 当时你是怎么想的? 」 / 「 是因为____? 」"),
        ("记得确认", "最后: 「 这个原则你能用 3 步说出来吗? 」"),
    ]
    for i, (head, sub) in enumerate(tips):
        yy = ry + Inches(0.9) + Inches(0.9) * i
        add_oval(s, rx + Inches(0.3), yy + Inches(0.05), Inches(0.35), Inches(0.35), fill=ORANGE)
        add_text(s, rx + Inches(0.3), yy + Inches(0.1), Inches(0.35), Inches(0.3),
                 "!", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, rx + Inches(0.8), yy, Inches(5.0), Inches(0.4),
                 head, size=12, color=WHITE, bold=True)
        add_text(s, rx + Inches(0.8), yy + Inches(0.4), Inches(5.0), Inches(0.4),
                 sub, size=10, color=GRAY_LIGHT, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "预计用时: 20-30 分钟 | 不评判, 不引导, 让对方真实想法出来",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_13(idx):
    s = slide_m1(idx, 23, "✋ 练习二: 写出来 — 螺旋 4 问卡片",
                 "把访谈结果落到纸上 · 一组卡片 · 一句话一行")
    # 4 列
    headers = ["Q1 行为", "Q2 动机", "Q3 原则", "Q4 路径"]
    cols_w = [Inches(3.0), Inches(3.0), Inches(3.0), Inches(3.0)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_rect(s, cx, Inches(1.85), cols_w[i], Inches(0.5), fill=ORANGE)
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    # 4 列内容
    items = [
        ("客户进店 2 分钟内\n主动邀请体验\n(动作)", "客户还在看, 不要太早打断;\n太晚客户可能走掉\n(原因)", "让客户主动产生好奇心\n比销售话术有效\n(原则)", "① 2 分钟内邀请\n② 安静观察\n③ 等客户问再答\n(路径)"),
    ]
    for ri, (c1, c2, c3, c4) in enumerate(items):
        ry = Inches(2.4) + Inches(3.0) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(3.0), fill=bg, line=GRAY_LIGHT)
        contents = [c1, c2, c3, c4]
        cx = Inches(0.5)
        for i, content in enumerate(contents):
            for j, line in enumerate(content.split("\n")):
                add_text(s, cx + Inches(0.2), ry + Inches(0.3) + Inches(0.6) * j,
                         cols_w[i] - Inches(0.4), Inches(0.6),
                         line, size=12, color=BLACK if i != 2 else NAVY_DARK,
                         bold=(i == 2), line_h=1.4, align=PP_ALIGN.CENTER)
            cx += cols_w[i]
    # 检查
    add_round(s, Inches(0.5), Inches(5.6), Inches(12.33), Inches(1.4), fill=GOLD)
    add_text(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.4),
             "✓ 检查清单", size=14, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.9),
             "Q1 是动作 (看得见) · Q2 是原因 (听得懂) · Q3 是原则 (回答为什么有效) · Q4 是步骤 (能教人)",
             size=12, color=WHITE, line_h=1.4, align=PP_ALIGN.CENTER)


def m1_14(idx):
    s = slide_m1(idx, 24, "⚠ 常见陷阱 1: 接受 Q3 第一个答案",
                 "通常第一个答案还是行为, 必须继续追问")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.4),
             "陷阱: 听到「 让客户多体验 」就以为得到原则",
             size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    # 假 Q3 vs 真 Q3
    lx = Inches(0.5); ly = Inches(2.5)
    add_round(s, lx, ly, Inches(6.0), Inches(3.5), fill=LIGHT, line=RED)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "✗ 假 Q3 原则 (还是行为)", size=14, color=RED, bold=True)
    bad = [
        "「 让客户多体验产品 」",
        "「 主动和客户互动 」",
        "「 给客户解释清楚 」",
        "「 用专业能力说服客户 」",
    ]
    for i, t in enumerate(bad):
        yy = ly + Inches(0.85) + Inches(0.6) * i
        add_text(s, lx + Inches(0.3), yy, Inches(0.4), Inches(0.5),
                 "✗", size=14, color=RED, bold=True)
        add_text(s, lx + Inches(0.7), yy, Inches(5.0), Inches(0.5),
                 t, size=12, color=BLACK)
    rx = Inches(6.7); ry = Inches(2.5)
    add_round(s, rx, ry, Inches(6.13), Inches(3.5), fill=LIGHT, line=GREEN)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "✓ 真 Q3 原则 (回答了为什么有效)", size=14, color=GREEN, bold=True)
    good = [
        "「 让客户主动产生问题, 比导购主动\n  传递信息更接近成交 」",
        "「 客户的好奇心比销售话术有效 」",
        "「 真实体验比介绍更能建立信任 」",
    ]
    for i, t in enumerate(good):
        yy = ry + Inches(0.85) + Inches(0.85) * i
        add_text(s, rx + Inches(0.3), yy, Inches(0.4), Inches(0.5),
                 "✓", size=14, color=GREEN, bold=True)
        add_text(s, rx + Inches(0.7), yy, Inches(5.0), Inches(0.85),
                 t, size=11, color=NAVY_DARK, bold=True, line_h=1.3)
    # 怎么问
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.4),
             "识别方法: 听 Q3 答案时问自己「 这个原则能用 3 步说出来吗? 」",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.3),
             "说不出来 → 还是行为 → 回到 Q1 再问",
             size=10, color=RED, bold=True, italic=True, align=PP_ALIGN.CENTER)


def m1_15(idx):
    s = slide_m1(idx, 25, "⚠ 常见陷阱 2: 跳过 Q2 直接到 Q3",
                 "Q2 是连接行为和原则的桥, 跳过它原则就悬空")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.4),
             "Q2 的核心作用: 让你理解「 为什么这样做 」, 原则才不是凭空得来的",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 流程图
    lx = Inches(0.5); ly = Inches(2.6)
    add_round(s, lx, ly, Inches(12.33), Inches(2.5), fill=LIGHT, line=GRAY_LIGHT)
    # 4 节点
    nodes = [
        ("Q1 行为", "做了\n什么?", ORANGE),
        ("Q2 动机", "为什么\n这样做?", GOLD),
        ("Q3 原则", "可提炼\n成原则吗?", RED),
        ("Q4 路径", "能教给\n新人吗?", GREEN),
    ]
    node_w = Inches(2.5); node_h = Inches(1.6)
    gap_x = (Inches(12.33) - node_w * 4) / 3
    start_x = lx + Inches(0.2)
    ny = ly + Inches(0.4)
    for i, (head, q, col) in enumerate(nodes):
        x = start_x + (node_w + gap_x) * i
        add_round(s, x, ny, node_w, node_h, fill=WHITE, line=col)
        add_rect(s, x, ny, node_w, Inches(0.5), fill=col)
        add_text(s, x, ny + Inches(0.08), node_w, Inches(0.4),
                 head, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), ny + Inches(0.7), node_w - Inches(0.4), Inches(0.8),
                 q, size=14, color=NAVY_DARK, bold=True, line_h=1.3, align=PP_ALIGN.CENTER)
        if i < 3:
            ax = x + node_w + gap_x * 0.3
            add_text(s, ax, ny + Inches(0.7), gap_x * 0.4, Inches(0.4),
                     "▶", size=20, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    # 跳过的危害
    add_text(s, Inches(0.5), Inches(5.3), Inches(12.33), Inches(0.4),
             "跳过 Q2 的危害:",
             size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    harms = [
        ("Q3 原则没有根基", "你不知道这个原则是从哪里来的, 别人质疑时答不上"),
        ("Q4 步骤无法设计", "不知道「 为什么 」, 步骤就变成机械的模仿清单"),
    ]
    for i, (head, sub) in enumerate(harms):
        x = Inches(0.5) + Inches(6.3) * i
        add_round(s, x, Inches(5.8), Inches(6.0), Inches(0.95), fill=RED)
        add_text(s, x + Inches(0.3), Inches(5.85), Inches(5.5), Inches(0.4),
                 head, size=13, color=WHITE, bold=True)
        add_text(s, x + Inches(0.3), Inches(6.25), Inches(5.5), Inches(0.5),
                 sub, size=10, color=WHITE, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.3),
             "Q2 的时间不能省 — 这是把「 经验 」转化成「 可学方法 」的桥",
             size=10, color=ORANGE, bold=True, italic=True, align=PP_ALIGN.CENTER)


def m1_16(idx):
    s = slide_m1(idx, 26, "⚠ 常见陷阱 3: 复制时变成「 学精神 」",
                 "不学动作, 只学精神 — 学完回去还是不知道做什么")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "「 学精神 」的典型表现", size=14, color=RED, bold=True)
    items = [
        "「 你要学小王那种认真劲儿 」",
        "「 你要对客户用心一点 」",
        "「 你要专业一点 」",
        "「 你要主动一点 」",
        "「 你要爱动脑 」",
    ]
    for i, t in enumerate(items):
        yy = ly + Inches(0.85) + Inches(0.65) * i
        add_rect(s, lx + Inches(0.3), yy, Inches(0.1), Inches(0.55), fill=RED)
        add_text(s, lx + Inches(0.5), yy, Inches(5.3), Inches(0.55),
                 t, size=12, color=BLACK)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "学不动的根因", size=14, color=NAVY_DARK, bold=True)
    items2 = [
        ("学精神", "「 你要____ 」= 完全不知道第一步做什么"),
        ("学动作", "「 你要做到 ____ 」= 第一步明确, 可以观察"),
        ("学步骤", "「 步骤 1, 2, 3 」= 任何人都能照做"),
    ]
    for i, (head, sub) in enumerate(items2):
        yy = ry + Inches(0.85) + Inches(1.1) * i
        add_text(s, rx + Inches(0.3), yy, Inches(5.5), Inches(0.4),
                 head, size=13, color=NAVY_DARK, bold=True)
        add_text(s, rx + Inches(0.3), yy + Inches(0.4), Inches(5.5), Inches(0.6),
                 sub, size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "Q4 路径的设计目的: 让新人「 第一步明确, 可以观察 」",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_17(idx):
    s = slide_m1(idx, 27, "5 工具的连接: M1 之后",
                 "M1 提炼的成功原则, 是 M2 花刺投票的「 抓手候选 」")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "M1 之后你拥有", size=14, color=ORANGE, bold=True)
    outs = [
        ("一条提炼出来的成功原则", "「 让客户主动产生问题, 比导购主动传递信息更接近成交 」"),
        ("一个 2-3 步的路径", "① 2 分钟内邀请体验 ② 安静观察 ③ 等客户问再答"),
        ("一个被验证过的高手", "小王, 35% 成交率"),
    ]
    for i, (head, sub) in enumerate(outs):
        yy = ly + Inches(0.85) + Inches(1.2) * i
        add_oval(s, lx + Inches(0.3), yy + Inches(0.1), Inches(0.4), Inches(0.4), fill=ORANGE)
        add_text(s, lx + Inches(0.3), yy + Inches(0.15), Inches(0.4), Inches(0.3),
                 "✓", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, lx + Inches(0.85), yy, Inches(5.0), Inches(0.4),
                 head, size=13, color=NAVY_DARK, bold=True)
        add_text(s, lx + Inches(0.85), yy + Inches(0.4), Inches(5.0), Inches(0.7),
                 sub, size=11, color=BLACK, line_h=1.3)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "下一部分 M2: 花刺投票", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(0.9), Inches(5.5), Inches(1.0),
             "M1 提炼的成功原则\n可以成为 M2 的一个候选「 花票 」",
             size=14, color=WHITE, bold=True, line_h=1.3)
    add_rect(s, rx + Inches(0.3), ry + Inches(2.1), Inches(5.5), Emu(15000), fill=ORANGE)
    add_text(s, rx + Inches(0.3), ry + Inches(2.25), Inches(5.5), Inches(2.2),
             "但只有它还不够\nM2 还要: 全员投票, 决定\n─ 哪个候选「 抓手 」值得优先\n─ 哪个「 障碍 」是最大阻力\n避免「 主管说了算 」",
             size=12, color=WHITE, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "→ 下一段: M2 共谋抓手",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_18(idx):
    s = slide_m1(idx, 28, "✅ 第一部分知识框架",
                 "M1 螺旋深挖 4 问 = 把团队里已有的成功变成可复制的方法")
    items = [
        ("核心问题", "为什么「 只有一个人会 」\n的经验, 复制不下去", RED),
        ("Q1 行为", "做了什么不同?\n(听动词, 不是形容词)", ORANGE),
        ("Q2 动机", "为什么这样做?\n(跳过它, Q3 必然空洞)", GOLD),
        ("Q3 原则", "可提炼成原则吗?\n(回答「 为什么有效 」)", RED),
    ]
    card_w = Inches(2.95); card_h = Inches(2.5)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = Inches(0.5)
    y = Inches(2.0)
    for i, (head, body, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_rect(s, x, y, card_w, Inches(0.6), fill=color)
        add_text(s, x, y + Inches(0.13), card_w, Inches(0.4),
                 head, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_text(s, x + Inches(0.2), y + Inches(0.8) + Inches(0.5) * j,
                     card_w - Inches(0.4), Inches(0.5),
                     line, size=12, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.4)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.33), Inches(0.4),
             "完成第一部分后你拥有:", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    outs = [
        ("一条提炼出来的成功原则", ORANGE),
        ("一个 2-3 步的路径", GREEN),
        ("一个被验证过的高手", GOLD),
    ]
    for i, (t, c) in enumerate(outs):
        x = Inches(0.5) + Inches(4.2) * i
        add_round(s, x, Inches(5.3), Inches(4.0), Inches(0.7), fill=LIGHT, line=c)
        add_text(s, x, Inches(5.35), Inches(4.0), Inches(0.6),
                 t, size=11, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "→ 成功原则 = M2 花刺投票的「 抓手候选 」",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m1_19(idx):
    s = slide_m1(idx, 29, "连接到下一部分",
                 "M1 让你从 1 个人的成功里提炼出原则 — 但怎么让团队「 真正行动起来 」?")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "M1 解决的是", size=14, color=ORANGE, bold=True)
    bullets1 = [
        "把「 只有一个人会 」的成功变成原则",
        "让「 凭感觉 」变成「 可观察的动作 」",
        "让新人在合理时间内学会",
    ]
    bullet_block(s, lx + Inches(0.4), ly + Inches(0.9), Inches(5.2), Inches(2.0), bullets1, size=12)
    add_rect(s, lx + Inches(0.3), ly + Inches(3.0), Inches(5.4), Emu(20000), fill=NAVY)
    add_text(s, lx + Inches(0.3), ly + Inches(3.15), Inches(5.4), Inches(0.4),
             "但你还会遇到", size=12, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(3.55), Inches(5.4), Inches(1.0),
             "→ 团队里很多事要做, 哪个先做?\n→ 谁支持, 谁犹豫, 谁反对?\n→ 怎么形成真正的共识, 不是「 表面同意 」?",
             size=12, color=BLACK, line_h=1.5)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "M2 共谋抓手 要回答", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(0.9), Inches(5.5), Inches(1.0),
             "怎么让团队对\n「 下一步重点是什么 」\n形成真实的集体共识?",
             size=18, color=WHITE, bold=True, line_h=1.4)
    add_rect(s, rx + Inches(0.3), ry + Inches(2.6), Inches(5.5), Emu(20000), fill=GOLD)
    add_text(s, rx + Inches(0.3), ry + Inches(2.75), Inches(5.5), Inches(1.9),
             "花刺投票\n\n不用「 主管说了算 」\n让团队真实地选\n─ 优先机会 (花票)\n─ 关键障碍 (刺票)",
             size=12, color=WHITE, bold=True, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "→ 下一段: 模块二 共谋抓手",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_20(idx):
    s = slide_m1(idx, 30, "M1 关键收获 — 4 问用对, 一次访谈就能用",
                 "把「 凭感觉 」变成「 可学方法 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "M1 给你 3 件具体的事, 离开教室后可以马上用:",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("访谈高手", "用 4 问访谈一位团队里的「 高手 」", "20-30 分钟, 1 对 1"),
        ("写 4 问卡片", "把访谈结果落到 4 问卡片", "Q1 行为 / Q2 动机 / Q3 原则 / Q4 路径"),
        ("教一个新人", "用 Q4 路径教给一位新人", "两周后看他能上手到什么程度"),
    ]
    card_w = Inches(4.0); card_h = Inches(3.0)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.6)
    for i, (head, sub, time) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=ORANGE)
        add_rect(s, x, y, card_w, Inches(0.7), fill=ORANGE)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.4),
                 head, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.6), Inches(1.4),
                 sub, size=14, color=NAVY_DARK, bold=True, line_h=1.4, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(2.4), card_w - Inches(0.6), Inches(0.4),
                 time, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    # 验证
    add_round(s, Inches(0.5), Inches(5.9), Inches(12.33), Inches(1.0), fill=GOLD)
    add_text(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.4),
             "✓ 验证标准", size=13, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
             "新人能上手 (不是「 理解了 」) — 观察他接待客户, 第一步是否照做",
             size=12, color=WHITE, line_h=1.3, align=PP_ALIGN.CENTER)


def m1_21(idx):
    s = slide_m1(idx, 31, "M1 自检 — 4 问的关键标志",
                 "用这个清单检查你刚才的访谈结果")
    items = [
        ("Q1 行为", "是「 动作 」不是「 形容词 」",
         "✓ 我能想象一个人做了这个动作的样子",
         "✗ 形容词 / 状态 / 心态描述", ORANGE),
        ("Q2 动机", "是「 原因 」不是「 目标 」",
         "✓ 告诉了我他怎么想, 决定怎么做的",
         "✗ 因为想卖出去 / 因为公司要求", GOLD),
        ("Q3 原则", "回答了「 为什么有效 」",
         "✓ 是机制 / 逻辑 / 因果",
         "✗ 还是行为 (换个说法)", RED),
        ("Q4 路径", "是 2-3 个具体步骤",
         "✓ 新人能照做第一步",
         "✗ 笼统的「 做好接待 」", GREEN),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=ORANGE)
    headers = ["问", "好答案的样子", "✓ 判断", "✗ 不是"]
    cols_w = [Inches(2.0), Inches(3.0), Inches(4.0), Inches(3.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
        cx += cols_w[i]
    for ri, (no, head, good, bad, col) in enumerate(items):
        ry = Inches(2.4) + Inches(1.0) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.0), fill=bg, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), ry, Inches(0.15), Inches(1.0), fill=col)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.25), ry + Inches(0.3), cols_w[0] - Inches(0.25), Inches(0.4),
                 no, size=14, color=col, bold=True)
        cx += cols_w[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.3), cols_w[1] - Inches(0.1), Inches(0.4),
                 head, size=12, color=NAVY_DARK, bold=True)
        cx += cols_w[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.3), cols_w[2] - Inches(0.1), Inches(0.4),
                 good, size=11, color=BLACK)
        cx += cols_w[2]
        add_text(s, cx + Inches(0.1), ry + Inches(0.3), cols_w[3] - Inches(0.1), Inches(0.4),
                 bad, size=11, color=GRAY)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "如果 ✗ 多于 ✓ — 重新访谈, 换具体场景再问一次",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_22(idx):
    s = slide_m1(idx, 32, "M1 小结 — 复制成功的逻辑",
                 "成功不是「 学精神 」, 是「 学具体动作 + 知道为什么 + 能教别人 」")
    # 流程图
    nodes = [
        ("M1 起点", "你团队里\n有一位「 高手 」", ORANGE),
        ("M1 过程", "用 4 问\n挖出 Q1-Q4", GOLD),
        ("M1 终点", "提炼出\n成功原则 + 路径", GREEN),
        ("M1 产出", "新人能在\n合理时间学会", NAVY),
    ]
    card_w = Inches(2.8); card_h = Inches(1.8)
    gap = Inches(0.2)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (head, sub, col) in enumerate(nodes):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.5), fill=col)
        add_text(s, x, y + Inches(0.08), card_w, Inches(0.4),
                 head, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(1.0),
                 sub, size=13, color=NAVY_DARK, bold=True, line_h=1.4, align=PP_ALIGN.CENTER)
        if i < 3:
            ax = x + card_w + gap * 0.2
            add_text(s, ax, y + Inches(0.7), gap * 0.6, Inches(0.4),
                     "▶", size=20, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    # 验证
    add_round(s, Inches(0.5), Inches(4.2), Inches(12.33), Inches(1.0), fill=GOLD)
    add_text(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(0.4),
             "✓ 复制成功的唯一标志", size=14, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(0.4),
             "不是「 小王做了一次分享 」, 而是「 其他人按 Q4 路径做到了 」",
             size=12, color=WHITE, align=PP_ALIGN.CENTER)
    # 产出
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.4),
             "完成 M1 后你拥有:", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    outs = [
        "一条成功原则 (回答为什么有效)",
        "一个 2-3 步路径 (新人能照做)",
        "一次完整访谈 (可对更多高手复用)",
    ]
    for i, t in enumerate(outs):
        x = Inches(0.5) + Inches(4.2) * i
        add_round(s, x, Inches(6.0), Inches(4.0), Inches(0.7), fill=LIGHT, line=GREEN)
        add_text(s, x, Inches(6.05), Inches(4.0), Inches(0.6),
                 "✓ " + t, size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


def m1_23(idx):
    s = slide_m1(idx, 33, "M1 模块收尾 · 茶歇",
                 "回来后, 我们进入第二个工具 M2 共谋抓手")
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.5),
             "M1", size=120, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.7),
             "复制成功", size=32, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             "Replicate What Works", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    quote_block(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                "把团队里「 只有一个人会 」的成功, 变成「 人人能学 」。",
                author="—— 第一部分核心信念")
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "请准备: 想好一位你团队里的「 高手 」, 下次课我们会做 M2 练习",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m1_24(idx):
    s = slide_m1(idx, 34, "M1 工具表单预览 — 螺旋 4 问访谈卡",
                 "可打印带走 · 一页纸 4 列")
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=ORANGE)
    headers = ["Q1 行为", "Q2 动机", "Q3 原则", "Q4 路径"]
    cols_w = [Inches(3.0), Inches(3.0), Inches(3.0), Inches(3.0)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    # 4 列填写区
    for i in range(4):
        x = Inches(0.5) + (Inches(3.0) + Inches(0.0)) * i
        add_rect(s, x, Inches(2.4), cols_w[i] - Inches(0.05), Inches(3.8), fill=LIGHT, line=GRAY_LIGHT)
        add_text(s, x + Inches(0.2), Inches(2.5), cols_w[i] - Inches(0.2), Inches(0.4),
                 "做了什么?", size=11, color=GRAY, italic=True)
        for j in range(2):
            add_text(s, x + Inches(0.2), Inches(3.0) + Inches(0.4) * j,
                     cols_w[i] - Inches(0.2), Inches(0.4),
                     "___________________________________", size=10, color=GRAY)
    # 底部
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "→ 完整模板见学员手册 · F1 螺旋 4 问访谈卡",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m1_25(idx):
    s = slide_m1(idx, 35, "M1 应用建议 — 在你团队里怎么用",
                 "不是「 一次访谈 」, 是「 持续提炼 」")
    items = [
        ("第 1-2 周", "访谈 1 位高手, 提炼 1 条原则", "确认原则是否回答了为什么有效"),
        ("第 3-4 周", "用 Q4 路径教 1 位新人", "观察他能不能在合理时间上手"),
        ("第 2 月", "再访谈 1-2 位高手", "横向看, 不同高手的原则是不是互补"),
        ("第 3 月", "建立团队的成功原则库", "按场景分类, 沉淀团队的方法论"),
    ]
    card_w = Inches(2.95); card_h = Inches(2.5)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = Inches(0.5)
    y = Inches(2.0)
    for i, (time, action, check) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=ORANGE)
        add_rect(s, x, y, card_w, Inches(0.5), fill=ORANGE)
        add_text(s, x, y + Inches(0.08), card_w, Inches(0.4),
                 time, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.5),
                 action, size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, line_h=1.3)
        add_rect(s, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.4), Emu(15000), fill=GOLD)
        add_text(s, x + Inches(0.2), y + Inches(1.5), card_w - Inches(0.4), Inches(0.9),
                 check, size=10, color=WHITE, line_h=1.3, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.33), Inches(0.4),
             "核心心法: 复制不是「 一次性项目 」, 是「 持续积累 」",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_round(s, Inches(0.5), Inches(5.3), Inches(12.33), Inches(1.6), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(5.45), Inches(11.9), Inches(0.4),
             "避免两个常见误区", size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.4),
             "✗ 把 M1 当作「 给团队做一次分享 」— 没有 4 问, 没有访谈, 只是单向分享",
             size=11, color=WHITE, align=PP_ALIGN.CENTER, line_h=1.3)
    add_text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.4),
             "✗ 期待 M1 一次性解决所有问题 — 真正的高手原则需要持续观察和迭代",
             size=11, color=WHITE, align=PP_ALIGN.CENTER, line_h=1.3)


def m1_26(idx):
    s = slide_m1(idx, 36, "M1 案例库 — 不同行业的高手",
                 "不只是销售, 复制成功是通用的方法")
    industries = [
        ("零售/服务业", "小王 (35% 成交率)\n体验式接待", ORANGE),
        ("客服/运营", "老李 (差评率 1%)\n客户问题分类", GOLD),
        ("销售/业务", "小张 (top 1)\n大客户拜访节奏", RED),
        ("餐饮/连锁", "陈姐 (高峰 0 投诉)\n客户等候体验", GREEN),
        ("制造业/班组长", "刘师傅 (0 返工)\n新人带教步骤", NAVY),
    ]
    card_w = Inches(2.4); card_h = Inches(2.3)
    gap = Inches(0.1)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (ind, case, col) in enumerate(industries):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.5), fill=col)
        add_text(s, x, y + Inches(0.08), card_w, Inches(0.4),
                 ind, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(case.split("\n")):
            add_text(s, x + Inches(0.2), y + Inches(0.7) + Inches(0.5) * j,
                     card_w - Inches(0.4), Inches(0.5),
                     line, size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.7), Inches(12.33), Inches(0.4),
             "每行都有可复制的「 高手 」(你团队里一定有一位)",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.4),
             "4 问都能用, 不限行业 · 不限岗位",
             size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    # 关键提醒
    add_round(s, Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.2), fill=GOLD)
    add_text(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.4),
             "⚠ 关键提醒", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.6),
             "M1 不是在找「 谁做得好 」, 而是在找「 谁做得好 + 你能说清楚他做了什么不同 」",
             size=11, color=WHITE, line_h=1.3, align=PP_ALIGN.CENTER)


def m1_27(idx):
    s = slide_m1(idx, 37, "M1 回到现实 — 你的下一步",
                 "把 M1 用在你能用上的人身上, 不必完美")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.4),
             "M1 不是一个抽象的方法 — 它就是 4 个具体的问题, 1 次访谈",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("想 1 位", "想好你团队里的 1 位「 高手 」 — 做某件事明显好于其他人", ORANGE),
        ("约 30 分钟", "约他做一次 30 分钟访谈 — 具体场景, 4 问追问", GOLD),
        ("写 4 列", "把访谈结果写到 4 列卡片 (Q1 行为 / Q2 动机 / Q3 原则 / Q4 路径)", RED),
        ("教 1 位新人", "用 Q4 路径教 1 位新人, 2 周后看他的上手程度", GREEN),
    ]
    card_w = Inches(6.0); card_h = Inches(1.4)
    for i, (head, sub, col) in enumerate(items):
        x = Inches(0.5) if i % 2 == 0 else Inches(6.7)
        y = Inches(2.6) + (card_h + Inches(0.2)) * (i // 2)
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, Inches(1.3), card_h, fill=col)
        add_text(s, x, y + Inches(0.45), Inches(1.3), Inches(0.4),
                 f"0{i+1}", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.5), y + Inches(0.25), card_w - Inches(1.7), Inches(0.4),
                 head, size=15, color=col, bold=True)
        add_text(s, x + Inches(1.5), y + Inches(0.7), card_w - Inches(1.7), Inches(0.7),
                 sub, size=11, color=BLACK, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "M1 = 把「 凭感觉 」变成「 可学方法 」",
             size=13, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


print("模块一定义完成")


# ===================================================================
# 第 5 段: 模块二 共谋抓手 (22 页)
# ===================================================================
def m2_section_divider(idx):
    s = new_slide()
    page_chrome(s, idx, section="模块 2 · 共谋抓手")
    add_rect(s, 0, Inches(1.5), SW, Inches(4.5), fill=ORANGE)
    add_text(s, Inches(0.5), Inches(1.7), Inches(5), Inches(3),
             "M2", size=180, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(2.3), Inches(8), Inches(1.0),
             "共谋抓手", size=44, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(3.5), Inches(8), Inches(0.6),
             "Find Real Consensus", size=18, color=WHITE)
    quote_block(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.6),
                "团队共识的质量, 不取决于讨论的时长, 而取决于每个人的真实判断是否都进入了决策。",
                author="—— 本模块核心信念")


def slide_m2(idx, n, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section="模块 2 · 共谋抓手")
    title_block(s, title, subtitle, accent=ORANGE)
    return s


def m2_01_opening(idx):
    s = slide_m2(idx, 30, "开场: 会议开了, 但什么也没有决定",
                 "两种典型的「 假性共识 」 — 散会后每个人记得的重点都不一样")
    # 左: 场景 1
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(2.6), fill=LIGHT, line=GRAY_LIGHT)
    add_rect(s, lx, ly, Inches(6.2), Inches(0.5), fill=GRAY)
    add_text(s, lx, ly + Inches(0.07), Inches(6.2), Inches(0.4),
             "场景 1: 各说各的版", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, lx + Inches(0.3), ly + Inches(0.7), Inches(5.6), Inches(1.8),
             "90 分钟会议, 大家说了很多\n散会时你说:「 好, 就这样定了 」\n\n"
             "→ 第二天问 3 个成员\n→ 记住了完全不同的 3 件事\n→ 没人撒谎, 记忆加起来 ≠ 可执行的共识",
             size=12, color=BLACK, line_h=1.5)
    # 右: 场景 2
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(2.6), fill=LIGHT, line=GRAY_LIGHT)
    add_rect(s, rx, ry, Inches(5.83), Inches(0.5), fill=GRAY)
    add_text(s, rx, ry + Inches(0.07), Inches(5.83), Inches(0.4),
             "场景 2: 附和版", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.3), ry + Inches(0.7), Inches(5.3), Inches(1.8),
             "主管说:「 我认为最重要的是 X 」\n有人附和, 大家都说「 对, X 很重要 」\n\n"
             "→ 三周后 X 依然原地\n→ 没人说有问题\n→ 因为每个人都能给出「 其他优先级 」的合理解释",
             size=12, color=BLACK, line_h=1.5)
    # 下: 共同名字
    add_round(s, Inches(0.5), Inches(4.7), Inches(12.33), Inches(1.6), fill=RED)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.33), Inches(0.5),
             "这两种情况有一个共同的名字:", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.35), Inches(12.33), Inches(0.6),
             "假 性 共 识", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.95), Inches(12.33), Inches(0.3),
             "—— 表面上大家都同意, 实际上每个人理解的重点不一样, 承诺的力度也不一样",
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # 工具预告
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "花刺投票 — 就是用来打破假性共识的工具",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m2_02_two_reasons(idx):
    s = slide_m2(idx, 31, "第一节: 团队讨论为什么走不到真实共识",
                 "两个根本原因 — 花刺投票各用一种设计回应")
    items = [
        ("根源 1", "社会压力让真实判断被压制",
         "主管先表态 / 声音大的人主导\n其他人即使有不同看法, 也不当场说",
         "花刺投票的回应:\n→ 静默投票, 不交流, 不跟风\n→ 让每个人的真实判断独立进入结果",
         RED),
        ("根源 2", "只讨论机会, 忽视障碍",
         "大多数会议只问「 应该做什么 」\n很少问「 什么东西会让我们做不成 」\n障碍没有被命名 → 悄悄让每个行动失效",
         "花刺投票的回应:\n→ 花票 + 刺票同时投\n→ 机会和障碍同时可见",
         ORANGE),
    ]
    card_w = Inches(6.1); card_h = Inches(4.6)
    gap = Inches(0.15)
    total_w = card_w * 2 + gap
    start_x = (SW - total_w) / 2
    y = Inches(1.95)
    for i, (tag, title, sym, fix, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        add_rect(s, x, y, card_w, Inches(0.6), fill=color)
        add_text(s, x, y + Inches(0.1), card_w, Inches(0.4),
                 tag, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(0.85), card_w - Inches(0.6), Inches(0.5),
                 title, size=16, color=NAVY_DARK, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(1.45), card_w - Inches(0.6), Inches(1.4),
                 sym, size=11, color=BLACK, line_h=1.5)
        # 修复
        add_rect(s, x + Inches(0.3), y + Inches(2.85), card_w - Inches(0.6), Emu(15000), fill=color)
        add_text(s, x + Inches(0.3), y + Inches(2.95), card_w - Inches(0.6), Inches(1.6),
                 fix, size=11, color=NAVY_DARK, line_h=1.5)
    # 底部金句
    quote_block(s, Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.5),
                "15 分钟的结构化投票, 有时候比 2 小时的会议更接近真实的集体智慧。",
                author="—— 花刺投票的核心价值")


def m2_03_diagnosis(idx):
    s = slide_m2(idx, 32, "✋ 场景诊断练习 (3 分钟)",
                 "下面 3 个场景各自出了什么问题?")
    items = [
        ("场景 1", "主管说了『 我认为最重要的是提升成交率 』, 大家都说对, 会议结束。",
         "社会压力 → 没人敢说不同意见"),
        ("场景 2", "会议 2 小时, 大家提了很多方向, 主管说『 回去各自想, 下次再深入 』。",
         "没有聚焦机制 → 讨论没有产生决策"),
        ("场景 3", "举手表决通过某个方向, 但事后发现那些困难没说, 因为怕扫兴。",
         "举手 ≠ 真实判断 → 障碍被掩盖"),
    ]
    y = Inches(1.85)
    card_h = Inches(1.45)
    for i, (no, desc, diag) in enumerate(items):
        ry = y + (card_h + Inches(0.15)) * i
        add_round(s, Inches(0.5), ry, Inches(12.33), card_h, fill=LIGHT, line=GRAY_LIGHT)
        # 序号
        add_rect(s, Inches(0.5), ry, Inches(1.5), card_h, fill=ORANGE)
        add_text(s, Inches(0.5), ry + Inches(0.5), Inches(1.5), Inches(0.5),
                 no, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 场景
        add_text(s, Inches(2.2), ry + Inches(0.2), Inches(6.5), card_h - Inches(0.3),
                 desc, size=12, color=BLACK, line_h=1.4)
        # 诊断
        add_rect(s, Inches(8.9), ry, Inches(3.93), card_h, fill=NAVY_DARK)
        add_text(s, Inches(9.05), ry + Inches(0.15), Inches(3.7), Inches(0.4),
                 "你的诊断", size=10, color=ORANGE, bold=True)
        add_text(s, Inches(9.05), ry + Inches(0.55), Inches(3.7), card_h - Inches(0.6),
                 diag, size=11, color=WHITE, bold=True, line_h=1.4)
    # 底部
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "三个场景的共同根源: 社会压力 + 障碍被忽视 — 花刺投票同时解决这两个",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m2_04_4steps(idx):
    s = slide_m2(idx, 33, "第二节: 花刺投票 — 四步走完一次真实的共识过程",
                 "步骤有顺序, 每步都有它的设计意图")
    steps = [
        ("01", "列候选抓手", "不评判, 数量先多起来", "便利贴/挂纸\n每人写, 数量 8–15 条", NAVY),
        ("02", "静默投票", "不交流, 不跟风", "每人 2 张花票 (优先机会)\n每人 2 张刺票 (关键障碍)", ORANGE),
        ("03", "统计结果", "两个答案都重要", "花票最多 = 优先机会\n刺票最多 = 关键障碍", GREEN),
        ("04", "聚焦高票项", "讨论 + 确认 + 落地", "花票高项: 谁负责 / 第一步时间\n刺票高项: 谁来推动解决", GOLD),
    ]
    card_w = Inches(3.0); card_h = Inches(4.6)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (no, name, sub, hint, color) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        # 大序号
        add_oval(s, x + Inches(1.0), y + Inches(0.25), Inches(1.0), Inches(1.0), fill=color)
        add_text(s, x + Inches(1.0), y + Inches(0.4), Inches(1.0), Inches(0.7),
                 no, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.4), card_w, Inches(0.5),
                 name, size=16, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(1.9), card_w - Inches(0.4), Inches(0.4),
                 sub, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.5), y + Inches(2.4), card_w - Inches(1.0), Emu(15000), fill=color)
        for j, line in enumerate(hint.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(2.55) + Inches(0.45) * j,
                     card_w - Inches(0.5), Inches(0.45),
                     line, size=11, color=BLACK, align=PP_ALIGN.CENTER)
    # 箭头
    for i in range(3):
        x1 = start_x + card_w * (i + 1) + gap * i + Inches(0.0)
        add_text(s, x1 - Inches(0.05), y + Inches(2.5), Inches(0.25), Inches(0.5),
                 "→", size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 底部: 关键设计
    add_round(s, Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.5), fill=NAVY_DARK)
    add_text(s, Inches(0.5), Inches(6.78), Inches(12.33), Inches(0.45),
             "关键设计: 静默投票 → 避免社会压力; 花刺并用 → 机会和障碍同时可见",
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def m2_05_silent_vote(idx):
    s = slide_m2(idx, 34, "为什么必须是「 静默 」投票?",
                 "一旦可以互相看, 人们会下意识跟着感觉走")
    # 左: 对比
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, Inches(0.7), ly + Inches(0.3), Inches(5.7), Inches(0.5),
             "为什么不能「 互相看 」", size=14, color=ORANGE, bold=True)
    bullets = [
        "主管贴哪里 → 其他人下意识跟",
        "声音最大的人贴哪里 → 其他人跟",
        "担心被孤立 → 改变真实判断",
        "结果 = 「 谁声音大 」的判断, 不是集体的",
    ]
    bullet_block(s, lx + Inches(0.7), ly + Inches(0.9), Inches(5.5), Inches(2.0), bullets, size=12)
    add_text(s, lx + Inches(0.7), ly + Inches(3.4), Inches(5.5), Inches(0.5),
             "静默的物理动作:", size=14, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.7), ly + Inches(3.85), Inches(5.5), Inches(0.6),
             "所有人同时走到白板前\n不说话, 不等别人, 不互相看",
             size=12, color=BLACK, line_h=1.5)
    # 右: 票数逻辑
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.3), Inches(5.2), Inches(0.5),
             "票数逻辑: 让数据说话", size=14, color=ORANGE, bold=True)
    # 花票刺票
    add_text(s, rx + Inches(0.4), ry + Inches(0.9), Inches(5.2), Inches(0.4),
             "花票 (绿色圆点)", size=12, color=GREEN, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(1.25), Inches(5.2), Inches(0.4),
             "→ 这个最值得优先推进", size=11, color=WHITE)
    add_text(s, rx + Inches(0.4), ry + Inches(1.75), Inches(5.2), Inches(0.4),
             "刺票 (红色圆点)", size=12, color=RED, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(2.1), Inches(5.2), Inches(0.4),
             "→ 这是最关键的障碍, 不解决其他都难", size=11, color=WHITE)
    add_rect(s, rx + Inches(0.4), ry + Inches(2.7), Inches(5.2), Emu(20000), fill=ORANGE)
    add_text(s, rx + Inches(0.4), ry + Inches(2.85), Inches(5.2), Inches(1.7),
             "刺票最多的那个障碍, 往往是\n「 为什么以前行动总是没效果 」\n的真正原因。",
             size=12, color=WHITE, bold=True, line_h=1.5)
    # 底部
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "规则: 不能看别人, 不能跟, 不能说「 我看你贴那里我也贴那里 」",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m2_06_full_case(idx):
    s = slide_m2(idx, 35, "⚡ 完整示例: 李明的团队业绩会",
                 "新竞争对手开业 → 用花刺投票决定本季度方向")
    # 上: 候选抓手
    rx = Inches(0.5); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(12.33), Inches(2.4), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, rx + Inches(0.4), ry + Inches(0.15), Inches(11.5), Inches(0.4),
             "11 条候选项 (白板上)", size=13, color=ORANGE, bold=True)
    items = [
        "①门口增加引流活动 / 促销展示",
        "②推广体验式接待方式 (小王)",
        "③主动联系老客户, 推动复购",
        "④加强导购对新品专业知识培训",
        "⑤研究竞争对手, 差异化定价",
        "⑥提升门店陈列和氛围",
        "⑦社交媒体内容更新频率",
        "⑧优化进店前 5 分钟接待流程",
        "⑨和附近写字楼合作企业团购",
        "⑩完善售后跟进流程",
        "⑪开展导购之间的分享会",
    ]
    for i, t in enumerate(items):
        col = i % 6
        row = i // 6
        x = rx + Inches(0.3) + Inches(2.0) * col
        y2 = ry + Inches(0.55) + Inches(0.45) * row
        add_text(s, x, y2, Inches(2.0), Inches(0.4), t, size=10, color=BLACK)
    # 下: 投票结果
    rx2 = Inches(0.5); ry2 = Inches(4.4)
    # 左: 花票
    add_round(s, rx2, ry2, Inches(6.0), Inches(2.4), fill=GREEN)
    add_text(s, rx2, ry2 + Inches(0.1), Inches(6.0), Inches(0.4),
             "花票最多: 优先机会", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, rx2 + Inches(0.3), ry2 + Inches(0.6), Inches(5.4), Inches(0.5),
             "② 推广体验式接待方式 — 8 票",
             size=14, color=WHITE, bold=True)
    add_text(s, rx2 + Inches(0.3), ry2 + Inches(1.05), Inches(5.4), Inches(1.3),
             "和李明判断一致\n也直接呼应第一部分提炼的成功原则\n→ 本周就开始推",
             size=11, color=WHITE, line_h=1.4)
    # 右: 刺票
    add_round(s, Inches(7.0), ry2, Inches(5.83), Inches(2.4), fill=RED)
    add_text(s, Inches(7.0), ry2 + Inches(0.1), Inches(5.83), Inches(0.4),
             "刺票最多: 关键障碍", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.3), ry2 + Inches(0.6), Inches(5.3), Inches(0.5),
             "竞争对手持续引流 — 7 票",
             size=14, color=WHITE, bold=True)
    add_text(s, Inches(7.3), ry2 + Inches(1.05), Inches(5.3), Inches(1.3),
             "这是成员补充的候选项\n李明之前没意识到团队担忧这么集中\n→ 单独开 30 分钟专题小会",
             size=11, color=WHITE, line_h=1.4)
    # 底部
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.3),
             "花刺投票的价值: 让隐藏的担忧和集体智慧同时可见",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m2_07_what_democracy_is_not(idx):
    s = slide_m2(idx, 36, "⚠ 注意: 花刺投票不是民主决策",
                 "最终决定权仍在管理者 — 投票结果是「 集体智慧的可视化 」")
    items = [
        ("它是", "集体智慧的可视化\n让管理者看清团队的真实判断\n然后做更有信息支撑的决策", GREEN),
        ("它不是", "民主投票\n多数票说了算\n团队完全自治", RED),
    ]
    card_w = Inches(6.1); card_h = Inches(3.5)
    gap = Inches(0.15)
    total_w = card_w * 2 + gap
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (tag, body, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_rect(s, x, y, card_w, Inches(0.7), fill=color)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.5),
                 tag, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_text(s, x + Inches(0.4), y + Inches(1.0) + Inches(0.5) * j,
                     card_w - Inches(0.8), Inches(0.5),
                     line, size=13, color=BLACK, align=PP_ALIGN.CENTER)
    # 决策权说明
    add_round(s, Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.1), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.4),
             "如果投票结果和你的判断不一样, 这个差距本身就是信息",
             size=13, color=ORANGE, bold=True)
    add_text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.5),
             "→ 要么你看到了团队没看到的东西\n→ 要么团队看到了你没看到的东西\n→ 两种都值得讨论, 不是直接推翻",
             size=11, color=WHITE, line_h=1.5)
    # 底部
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3),
             "工具的信任基础: 不绕过结果, 不强加判断",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def m2_08_exercise1(idx):
    s = slide_m2(idx, 37, "✋ 练习一: 为什么结果「 出乎意料 」",
                 "理解静默投票的力量 — 3 个场景的对照分析")
    items = [
        ("场景 1", "主管偏好「 加强新品培训 」\n但投票结果「 体验式接待 」 8 票胜出",
         "静默结果",
         "团队真实判断 — 主管误判了团队需求",
         "举手表决",
         "可能跟主管走, 结果被扭曲"),
        ("场景 2", "资深员工一直主张「 引流最重要 」\n但花刺投票只有 3 票花 + 2 票刺",
         "静默结果",
         "资深员工意见 ≠ 团队共识",
         "举手表决",
         "可能因权威压制不同意见"),
        ("场景 3", "3 个人心里觉得「 竞品引流是最大障碍 」\n但会议上都没说出来 — 直到刺票显示出来",
         "静默结果",
         "隐藏的担忧可见化",
         "举手表决",
         "继续被掩盖, 没人敢说"),
    ]
    y = Inches(1.85)
    row_h = Inches(1.55)
    for i, (no, sit, l1, r1, l2, r2) in enumerate(items):
        ry = y + (row_h + Inches(0.1)) * i
        add_round(s, Inches(0.5), ry, Inches(12.33), row_h, fill=LIGHT, line=GRAY_LIGHT)
        # 序号
        add_rect(s, Inches(0.5), ry, Inches(1.2), row_h, fill=ORANGE)
        add_text(s, Inches(0.5), ry + Inches(0.6), Inches(1.2), Inches(0.5),
                 no, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 情境
        add_text(s, Inches(1.85), ry + Inches(0.2), Inches(5.0), row_h - Inches(0.4),
                 sit, size=11, color=BLACK, line_h=1.4)
        # 静默结果
        add_rect(s, Inches(7.0), ry + Inches(0.15), Inches(2.8), Inches(0.55), fill=GREEN)
        add_text(s, Inches(7.0), ry + Inches(0.2), Inches(2.8), Inches(0.45),
                 "静默结果", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(7.0), ry + Inches(0.75), Inches(2.8), Inches(0.75),
                 r1, size=10, color=BLACK, line_h=1.3)
        # 举手表决
        add_rect(s, Inches(9.95), ry + Inches(0.15), Inches(2.85), Inches(0.55), fill=RED)
        add_text(s, Inches(9.95), ry + Inches(0.2), Inches(2.85), Inches(0.45),
                 "如果举手表决", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(9.95), ry + Inches(0.75), Inches(2.85), Inches(0.75),
                 r2, size=10, color=BLACK, line_h=1.3)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "静默的力量: 让真实判断独立进入结果, 不被权威 / 关系 / 跟风污染",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m2_09_exercise2(idx):
    s = slide_m2(idx, 38, "✋ 练习二: 为你的真实议题准备一次花刺投票",
                 "选一个你近期要召开的团队讨论议题 — 直接带进下次会议")
    # 左: 议题设计
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "我的花刺投票设计", size=14, color=ORANGE, bold=True)
    items = [
        ("投票议题", ""),
        ("团队规模: ___ 人", ""),
        ("每人花票 ___ 张 / 刺票 ___ 张", ""),
    ]
    y2 = ly + Inches(0.85)
    for label, _ in items:
        add_rect(s, lx + Inches(0.3), y2, Inches(5.6), Inches(0.65), fill=WHITE, line=GRAY_LIGHT)
        add_text(s, lx + Inches(0.4), y2 + Inches(0.05), Inches(5.4), Inches(0.5),
                 label, size=11, color=GRAY)
        y2 += Inches(0.8)
    # 候选项
    add_text(s, lx + Inches(0.3), y2 + Inches(0.05), Inches(5.6), Inches(0.4),
             "候选项清单", size=12, color=ORANGE, bold=True)
    y2 += Inches(0.45)
    for i in range(5):
        add_rect(s, lx + Inches(0.3), y2, Inches(5.6), Inches(0.4), fill=WHITE, line=GRAY_LIGHT)
        add_text(s, lx + Inches(0.4), y2 + Inches(0.07), Inches(5.4), Inches(0.3),
                 f"  {['①', '②', '③', '④', '⑤'][i]}", size=11, color=GRAY)
        y2 += Inches(0.42)
    # 右: 物料
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.2), Inches(5.2), Inches(0.5),
             "物料准备清单", size=14, color=ORANGE, bold=True)
    items2 = [
        ("花形彩色贴纸 / 绿色圆点", "团队人数 × 2"),
        ("刺形贴纸 / 红色圆点", "团队人数 × 2"),
        ("便利贴 (补充候选用)", "每人 5–8 张"),
        ("白板或挂纸", "1–2 张"),
    ]
    y2 = ry + Inches(0.9)
    for name, qty in items2:
        add_rect(s, rx + Inches(0.4), y2, Inches(5.2), Inches(0.55), fill=WHITE, line=GRAY_LIGHT)
        add_text(s, rx + Inches(0.5), y2 + Inches(0.13), Inches(3.5), Inches(0.4),
                 name, size=10, color=BLACK)
        add_text(s, rx + Inches(4.0), y2 + Inches(0.13), Inches(1.5), Inches(0.4),
                 qty, size=10, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)
        y2 += Inches(0.6)
    # 时间
    add_text(s, rx + Inches(0.4), y2 + Inches(0.1), Inches(5.2), Inches(0.4),
             "时间: ___", size=12, color=ORANGE, bold=True)
    # 底部
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "完整花刺投票 8–15 人团队通常 30–40 分钟 — 比多数决策会议短, 产出更清晰",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m2_10_group_sim(idx):
    s = slide_m2(idx, 39, "小组模拟: 走一遍完整的花刺投票",
                 "议题: 课程结束后, 你最想在工作里优先用的一件工具是什么?")
    # 流程
    steps = [
        ("01", "写候选项", "便利贴写 1–2 条\n贴到组里挂纸\n不需要讨论", "3 分钟"),
        ("02", "合并候选项", "合并相同 / 相近\n确认最终清单", "2 分钟"),
        ("03", "静默投票", "发花票 / 刺票\n不交流, 不相互看", "3 分钟"),
        ("04", "统计结果", "填入下表\n对比意外", "2 分钟"),
        ("05", "讨论高票", "花票高项 2 分钟\n刺票高项 2 分钟", "4 分钟"),
    ]
    card_w = Inches(2.4); card_h = Inches(3.6)
    gap = Inches(0.1)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(1.95)
    for i, (no, name, sub, tm) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        add_rect(s, x, y, card_w, Inches(0.55), fill=ORANGE)
        add_text(s, x, y + Inches(0.07), card_w, Inches(0.4),
                 no, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.5),
                 name, size=13, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(1.35) + Inches(0.45) * j,
                     card_w - Inches(0.5), Inches(0.45),
                     line, size=10, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.3)
        # 时间
        add_rect(s, x + Inches(0.4), y + card_h - Inches(0.65), card_w - Inches(0.8), Inches(0.45), fill=ORANGE)
        add_text(s, x + Inches(0.4), y + card_h - Inches(0.6), card_w - Inches(0.8), Inches(0.4),
                 tm, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 观察记录
    rx = Inches(0.5); ry = Inches(5.8)
    add_round(s, rx, ry, Inches(12.33), Inches(1.0), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.1), Inches(11.5), Inches(0.4),
             "观察记录", size=12, color=ORANGE, bold=True)
    obs = [
        ("花票最多", "_______"),
        ("刺票最多", "_______"),
        ("让我意外的地方", "_______"),
        ("如果用举手表决结果会一样吗", "_______"),
    ]
    for i, (label, blank) in enumerate(obs):
        col = i % 2
        row = i // 2
        x = rx + Inches(0.4) + Inches(6.0) * col
        y2 = ry + Inches(0.45) + Inches(0.3) * row
        add_text(s, x, y2, Inches(2.0), Inches(0.25),
                 label, size=10, color=ORANGE, bold=True)
        add_rect(s, x + Inches(2.0), y2 + Inches(0.02), Inches(3.5), Emu(8000), fill=GRAY_LIGHT)
    # 底部
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.3),
             "完整流程约 14 分钟 — 比多数会议短, 产出更清晰",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m2_11_three_pitfalls(idx):
    s = slide_m2(idx, 40, "⚠ 三个使用陷阱",
                 "避免这 3 个, 花刺投票才不会走形")
    errs = [
        ("陷阱 1", "第一步评判太早",
         "候选项刚贴上去, 有人说『 这个肯定不行 』\n会议室立刻冷却, 之后没人再敢提非主流想法",
         "修正: 第一步只出不评\n把「 评 」留到第四步", RED),
        ("陷阱 2", "刺票最多 = 批评提出者",
         "某候选项是某成员提的\n刺票多了, 那个人感到被否定",
         "修正: 主持人强调\n刺票 = 对推进难度的判断 ≠ 对提出者的评价", ORANGE),
        ("陷阱 3", "结果出来后管理者反驳 / 绕过",
         "『 结果是这个, 但我觉得还是 X 重要 』\n破坏整个工具的信任基础",
         "修正: 说出你的判断理由\n然后讨论, 不是直接推翻", GOLD),
    ]
    card_w = Inches(4.0); card_h = Inches(4.4)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (no, name, sym, fix, color) in enumerate(errs):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GRAY_LIGHT)
        add_rect(s, x, y, card_w, Inches(0.6), fill=color)
        add_text(s, x, y + Inches(0.1), card_w, Inches(0.5),
                 no, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.5),
                 name, size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), y + Inches(1.4), card_w - Inches(0.5), Inches(0.4),
                 "症状", size=10, color=color, bold=True)
        add_text(s, x + Inches(0.25), y + Inches(1.7), card_w - Inches(0.5), Inches(1.5),
                 sym, size=10, color=BLACK, line_h=1.4)
        add_rect(s, x + Inches(0.25), y + Inches(3.2), card_w - Inches(0.5), Emu(15000), fill=color)
        add_text(s, x + Inches(0.25), y + Inches(3.3), card_w - Inches(0.5), Inches(0.3),
                 "修正", size=10, color=color, bold=True)
        add_text(s, x + Inches(0.25), y + Inches(3.6), card_w - Inches(0.5), Inches(0.8),
                 fix, size=10, color=BLACK, line_h=1.4)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "信任是花刺投票的根基 — 不绕过结果, 不强加判断, 不批评提出者",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m2_12_recap(idx):
    s = slide_m2(idx, 41, "✅ 第二部分知识框架",
                 "M2 共谋抓手 = 工具 + 设计意图 + 主持人原则")
    items = [
        ("核心问题", "为什么走不到\n真实共识", ORANGE),
        ("工具步骤", "列 / 投 / 算 / 议\n4 步各有意图", NAVY),
        ("静默 + 花刺", "避免社会压力\n机会与障碍并见", GREEN),
        ("主持人原则", "先不说 / 不评判\n不绕过结果", GOLD),
    ]
    card_w = Inches(2.95); card_h = Inches(3.0)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.2)
    for i, (head, body, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_rect(s, x, y, card_w, Inches(0.7), fill=color)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.4),
                 head, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_text(s, x + Inches(0.2), y + Inches(1.0) + Inches(0.5) * j,
                     card_w - Inches(0.4), Inches(0.5),
                     line, size=12, color=BLACK, align=PP_ALIGN.CENTER)
    # 产出
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.4),
             "完成花刺投票后你拥有:", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    out_items = [
        ("团队对优先机会的真实共识", NAVY),
        ("团队对关键障碍的集体识别", RED),
        ("一个或几个高刺票障碍 — 就是 M3 要深度分析的对象", GREEN),
    ]
    for i, (txt, col) in enumerate(out_items):
        x = Inches(0.5) + Inches(4.2) * i
        add_round(s, x, Inches(5.95), Inches(4.0), Inches(0.7), fill=LIGHT, line=col)
        add_text(s, x, Inches(6.0), Inches(4.0), Inches(0.6),
                 txt, size=11, color=col, bold=True, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "→ 刺票最多的障碍 = 第三部分的问题树分析对象",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m2_13_connect(idx):
    s = slide_m2(idx, 42, "连接到下一部分: 你的花刺投票产出",
                 "接下来要把刺票最高的障碍变成可以下手解决的具体方向")
    # 左
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "你刚刚用花刺投票做了什么", size=14, color=NAVY_DARK, bold=True)
    bullets = [
        "让团队对「 优先机会 」达成共识",
        "识别了团队最担忧的「 关键障碍 」",
        "两个答案都进入了决策",
    ]
    bullet_block(s, lx + Inches(0.4), ly + Inches(0.9), Inches(5.2), Inches(1.8), bullets, size=12)
    # 接下来的问题
    add_rect(s, lx + Inches(0.3), ly + Inches(3.0), Inches(5.4), Emu(20000), fill=NAVY)
    add_text(s, lx + Inches(0.3), ly + Inches(3.15), Inches(5.4), Inches(0.4),
             "接下来的问题", size=12, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(3.55), Inches(5.4), Inches(1.0),
             "→ 这个障碍具体是什么?\n→ 它可以被拆解吗?\n→ 经验之外有没有新角度?",
             size=12, color=BLACK, line_h=1.5)
    # 右
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "下一部分要回答的问题", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(0.9), Inches(5.5), Inches(1.4),
             "怎么把模糊的大难题\n拆成可以下手的具体方向?\n\n怎么找到经验范围外的\n新切入角度?",
             size=15, color=WHITE, bold=True, line_h=1.4)
    add_rect(s, rx + Inches(0.4), ry + Inches(3.2), Inches(5.3), Emu(20000), fill=GREEN)
    add_text(s, rx + Inches(0.4), ry + Inches(3.35), Inches(5.3), Inches(1.2),
             "问题树 + 魔力破解提问\n两件工具合在一起走",
             size=12, color=WHITE, bold=True, line_h=1.5)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "下一部分: M3 应对难题 — 问题树 + 魔力提问",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m2_14_section_close(idx):
    s = slide_m2(idx, 43, "模块二收尾 · 茶歇 15 分钟",
                 "回来后, 我们进入 M3 应对难题")
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.5),
             "M2", size=120, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.7),
             "共谋抓手", size=32, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             "Find Real Consensus", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    quote_block(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                "团队共识的质量, 不取决于讨论的时长, 而取决于真实判断是否进入了决策。",
                author="—— 第二部分核心信念")
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "请准备: 你识别出的「 刺票最高障碍 」一句话描述",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


print("模块二定义完成")


# ===================================================================
# 第 6 段: 模块三 应对难题 (32 页)
# ===================================================================
def m3_section_divider(idx):
    s = new_slide()
    page_chrome(s, idx, section="模块 3 · 应对难题")
    add_rect(s, 0, Inches(1.5), SW, Inches(4.5), fill=GREEN)
    add_text(s, Inches(0.5), Inches(1.7), Inches(5), Inches(3),
             "M3", size=180, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(2.3), Inches(8), Inches(1.0),
             "应对难题", size=44, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(3.5), Inches(8), Inches(0.6),
             "Solve Hard Problems", size=18, color=WHITE)
    quote_block(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.6),
                "解决问题的最大浪费, 不是方法选错了, 而是在搞清楚问题是什么之前就开始行动。",
                author="—— 本模块核心信念")


def slide_m3(idx, n, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section="模块 3 · 应对难题")
    title_block(s, title, subtitle, accent=GREEN)
    return s


def m3_01_opening(idx):
    s = slide_m3(idx, 45, "开场: 靠经验解决不了的那类问题",
                 "难题的本质 = 用了以前有效的方法, 但结果和预期不符")
    # 左: 经验有效
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(2.5), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "经验模式 (大多时候有效)", size=14, color=NAVY, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(0.8), Inches(5.6), Inches(1.6),
             "「 这种情况以前碰到过,\n就这样处理 」\n\n"
             "最快的决策模式\n大多数时候最有效\n经验是真实的资产",
             size=12, color=BLACK, line_h=1.5)
    # 右: 经验失灵
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(2.5), fill=RED)
    add_text(s, rx + Inches(0.4), ry + Inches(0.2), Inches(5.2), Inches(0.5),
             "经验失灵 (一类问题)", size=14, color=WHITE, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(0.8), Inches(5.2), Inches(1.6),
             "试了几种以前用过的方法, 都没用\n每次采取行动, 感觉在往前走\n但一段时间后发现还在原地\n甚至更难了",
             size=12, color=WHITE, line_h=1.5)
    # 例子
    y = Inches(4.55)
    add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.4), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, Inches(0.7), y + Inches(0.15), Inches(11.9), Inches(0.4),
             "这类问题的典型例子", size=12, color=ORANGE, bold=True)
    examples = ["新竞争对手进入", "年轻员工陆续离职", "原商业模式开始失效"]
    for i, t in enumerate(examples):
        x = Inches(0.7) + Inches(4.0) * i
        add_round(s, x, y + Inches(0.55), Inches(3.8), Inches(0.7), fill=WHITE, line=ORANGE)
        add_text(s, x, y + Inches(0.6), Inches(3.8), Inches(0.6),
                 t, size=13, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    # 关键
    add_round(s, Inches(0.5), Inches(6.15), Inches(12.33), Inches(0.85), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.4),
             "继续用经验解决问题, 不只是无效, 还有可能是有害的 — 因为你在用错的工具",
             size=12, color=ORANGE, bold=True)
    add_text(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4),
             "应对这类难题需要两件事: 先把问题说清楚 (问题树) → 找到经验范围外的新视角 (魔力提问)",
             size=11, color=WHITE)
    # 底部
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3),
             "—— 接下来, 我们把这两件事合在一起走一遍",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def m3_02_skip_diagnosis(idx):
    s = slide_m3(idx, 46, "第一节: 为什么「 想不清楚 」的问题解决不了",
                 "跳过诊断直接解决 = 最大的浪费")
    # 左: 案例
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(4.7), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "案例: 销售额下滑 18%", size=14, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(0.8), Inches(5.6), Inches(0.5),
             "管理层立刻提出 3 个方向:", size=12, color=NAVY_DARK, bold=True)
    items = ["增加广告投入", "降低部分产品价格", "加强导购话术培训"]
    for i, t in enumerate(items):
        add_round(s, lx + Inches(0.5), ly + Inches(1.3) + Inches(0.55) * i,
                  Inches(5.2), Inches(0.5), fill=WHITE, line=NAVY)
        add_text(s, lx + Inches(0.7), ly + Inches(1.4) + Inches(0.55) * i,
                 Inches(4.8), Inches(0.4),
                 f"  方向 {['一', '二', '三'][i]}:  {t}", size=12, color=BLACK)
    # 诊断
    add_text(s, lx + Inches(0.3), ly + Inches(3.05), Inches(5.6), Inches(0.5),
             "问题是什么?", size=14, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(3.45), Inches(5.6), Inches(1.2),
             "18% 是症状, 不是问题\n\n可能: 客流 / 成交率 / 客单价 / 复购率\n→ 四种情况解决方向完全不同",
             size=12, color=BLACK, line_h=1.5)
    # 右: 核心
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(4.7), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.3), Inches(5.2), Inches(0.5),
             "核心问题", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(0.9), Inches(5.2), Inches(1.2),
             "跳过了诊断,\n直接跳到了解决。",
             size=22, color=WHITE, bold=True, line_h=1.4)
    add_rect(s, rx + Inches(0.4), ry + Inches(2.3), Inches(5.2), Emu(20000), fill=ORANGE)
    add_text(s, rx + Inches(0.4), ry + Inches(2.45), Inches(5.2), Inches(1.5),
             "在没弄清楚\n「 症状背后是哪种具体问题 」之前\n选任何一个方向\n都有很大概率在错的方向上投入资源",
             size=12, color=WHITE, line_h=1.5)
    # 关键认知
    add_text(s, rx + Inches(0.4), ry + Inches(4.05), Inches(5.2), Inches(0.5),
             "关键认知", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(4.45), Inches(5.2), Inches(0.5),
             "→ 问题树就是用来做诊断的",
             size=12, color=WHITE, bold=True)
    # 底部
    add_text(s, Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.4),
             "说不清楚的问题, 是没办法被真正解决的 — 因为你不知道「 解决了 」的标志是什么",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m3_03_tree_4steps(idx):
    s = slide_m3(idx, 47, "第二节: 问题树 — 4 步把模糊大问题变成可下手方向",
                 "分层拆解 → 找切入点 (影响大 × 我能影响)")
    steps = [
        ("01", "写下症状", "用数字 / 可观察现象描述\n不是判断, 不是原因", NAVY),
        ("02", "第一层拆解", "3–4 个主方向\n完整不重叠", ORANGE),
        ("03", "第二层拆解", "每个方面的具体原因\n选最重要的 1–2 个展开", GREEN),
        ("04", "找切入点", "影响足够大 × 我能真正影响\n两者同时满足", GOLD),
    ]
    card_w = Inches(3.0); card_h = Inches(3.2)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (no, name, sub, color) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_oval(s, x + Inches(1.05), y + Inches(0.25), Inches(0.9), Inches(0.9), fill=color)
        add_text(s, x + Inches(1.05), y + Inches(0.4), Inches(0.9), Inches(0.6),
                 no, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.3), card_w, Inches(0.5),
                 name, size=15, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.5), y + Inches(1.85), card_w - Inches(1.0), Emu(15000), fill=color)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(2.0) + Inches(0.4) * j,
                     card_w - Inches(0.5), Inches(0.4),
                     line, size=11, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.3)
    # 箭头
    for i in range(3):
        x1 = start_x + card_w * (i + 1) + gap * i
        add_text(s, x1, y + Inches(1.5), Inches(0.25), Inches(0.5), "→",
                 size=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 切入点标准
    y2 = Inches(5.5)
    add_round(s, Inches(0.5), y2, Inches(12.33), Inches(1.45), fill=NAVY_DARK)
    add_text(s, Inches(0.7), y2 + Inches(0.1), Inches(11.9), Inches(0.4),
             "切入点的选择标准 (两者必须同时满足)",
             size=14, color=ORANGE, bold=True)
    items = [
        ("条件一", "影响足够大", "解决后对整体目标有实质改善\n不是边缘影响", GREEN),
        ("条件二", "我能真正影响", "在资源 / 权限 / 能力范围内可推动\n不是「 理想情况下 」才能做到", ORANGE),
    ]
    for i, (tag, name, desc, col) in enumerate(items):
        x = Inches(0.7) + Inches(6.0) * i
        yy = y2 + Inches(0.55)
        add_rect(s, x, yy, Inches(1.3), Inches(0.75), fill=col)
        add_text(s, x, yy + Inches(0.2), Inches(1.3), Inches(0.4),
                 tag, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.4), yy, Inches(4.0), Inches(0.4),
                 name, size=12, color=WHITE, bold=True)
        add_text(s, x + Inches(1.4), yy + Inches(0.35), Inches(4.0), Inches(0.4),
                 desc, size=10, color=GRAY, line_h=1.3)
    # 底部
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3),
             "注意: 切入点是主动选择的, 不是自然涌现的",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m3_04_full_case(idx):
    s = slide_m3(idx, 48, "⚡ 完整示例: 李明的问题树",
                 "分析对象: 成交率偏低 21% (刺票最高障碍的拆解)")
    # 左: 问题树
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(7.6), Inches(5.0), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.15), Inches(7), Inches(0.4),
             "李明的问题树", size=13, color=ORANGE, bold=True)
    # 根节点
    add_rect(s, lx + Inches(2.4), ly + Inches(0.6), Inches(2.8), Inches(0.55), fill=NAVY_DARK)
    add_text(s, lx + Inches(2.4), ly + Inches(0.7), Inches(2.8), Inches(0.4),
             "成交率偏低 (21%)", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 一级
    a_items = [
        ("A", "接待质量问题", NAVY),
        ("B", "陈列不够吸引", GRAY),
        ("C", "竞品比价行为", ORANGE),
        ("D", "导购专业知识不足", GRAY),
    ]
    for i, (no, t, col) in enumerate(a_items):
        yy = ly + Inches(1.4)
        x = lx + Inches(0.3) + Inches(1.85) * i
        add_rect(s, x, yy, Inches(1.7), Inches(0.5), fill=col)
        add_text(s, x, yy + Inches(0.1), Inches(1.7), Inches(0.4),
                 f"{no}. {t}", size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # A 的二级
    a_secs = [
        ("A1", "响应慢"),
        ("A2 ★", "没有体验引导"),
        ("A3", "判断需求缺失"),
    ]
    for i, (no, t) in enumerate(a_secs):
        yy = ly + Inches(2.1) + Inches(0.45) * i
        col = GREEN if "★" in no else NAVY
        add_rect(s, lx + Inches(0.3), yy, Inches(1.85), Inches(0.4), fill=col)
        add_text(s, lx + Inches(0.3), yy + Inches(0.05), Inches(1.85), Inches(0.35),
                 f"  {no} {t}", size=9, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    # C 的二级
    c_secs = [
        ("C1", "价格差距"),
        ("C3 ★", "进入比价前未建立偏好"),
    ]
    for i, (no, t) in enumerate(c_secs):
        yy = ly + Inches(2.3) + Inches(0.45) * i
        col = GREEN if "★" in no else ORANGE
        add_rect(s, lx + Inches(2.4), yy, Inches(1.85), Inches(0.4), fill=col)
        add_text(s, lx + Inches(2.4), yy + Inches(0.05), Inches(1.85), Inches(0.35),
                 f"  {no} {t}", size=9, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    # ★ 标记
    add_text(s, lx + Inches(0.3), ly + Inches(4.5), Inches(7), Inches(0.4),
             "★ = 切入点", size=11, color=GREEN, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(4.8), Inches(7), Inches(0.4),
             "李明的切入点: A2 没有体验引导 (高影响 × 完全可推动)",
             size=10, color=BLACK)
    # 右: 切入点评估
    rx = Inches(8.3); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(4.55), Inches(5.0), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(4), Inches(0.4),
             "切入点评估", size=13, color=ORANGE, bold=True)
    # 表格
    headers = ["分支", "影响", "我能影响"]
    rows = [
        ("A1 响应慢", "中", "✓"),
        ("A2 ★", "高", "✓"),
        ("A3 判断需求", "高", "✓ 部分"),
        ("C1 价格差距", "高", "✗ 总部"),
        ("C3 ★", "高", "✓"),
    ]
    add_rect(s, rx + Inches(0.3), ry + Inches(0.7), Inches(4), Inches(0.4), fill=ORANGE)
    cx = rx + Inches(0.3)
    cw = [Inches(2.0), Inches(0.8), Inches(1.2)]
    for i, h in enumerate(headers):
        add_text(s, cx, ry + Inches(0.75), cw[i], Inches(0.35),
                 h, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cw[i]
    for ri, row in enumerate(rows):
        yy = ry + Inches(1.1) + Inches(0.5) * ri
        col = GREEN if "★" in row[0] else (WHITE if ri % 2 == 0 else GRAY_LIGHT)
        add_rect(s, rx + Inches(0.3), yy, Inches(4), Inches(0.5), fill=col)
        cx = rx + Inches(0.3)
        for i, cell in enumerate(row):
            text_col = WHITE if "★" in row[0] or ri % 2 != 0 else BLACK
            if i > 0 and "★" in row[0]:
                text_col = WHITE
            add_text(s, cx, yy + Inches(0.13), cw[i], Inches(0.4),
                     cell, size=10, color=text_col, bold=("★" in row[0]),
                     align=PP_ALIGN.CENTER)
            cx += cw[i]
    # 底部
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3),
             "好的问题树 = 找到之前不知道该往哪个方向使力的「 分叉 」",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m3_05_exercise1(idx):
    s = slide_m3(idx, 49, "✋ 练习一: 识别有效的切入点",
                 "判断下面 4 个候选切入点是否同时满足两个标准")
    items = [
        ("1", "整体行业服务标准太低, 客户期望管理失真",
         "大", "✗", "❌", "影响大, 但不可影响 — 行业标准不是团队能改变的"),
        ("2", "客户投诉后第一次响应时间从 4 小时缩短到 1 小时",
         "大", "✓", "✅", "影响大, 可影响 — 流程和资源都在团队范围内"),
        ("3", "整个公司薪酬体系不激励服务质量, 根本原因在制度",
         "大", "✗", "❌", "影响大, 但不可影响 — 薪酬是公司层面决定"),
        ("4", "导购接待结束未主动征集客户反馈, 问题未被识别",
         "中–大", "✓", "✅", "可影响 — 培训和流程设计在团队范围内"),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=NAVY)
    headers = ["#", "候选切入点", "影响", "我能影响", "判断", "原因"]
    cols_w = [Inches(0.7), Inches(4.5), Inches(1.0), Inches(1.0), Inches(1.0), Inches(4.13)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER if i > 1 else PP_ALIGN.LEFT)
        cx += cols_w[i]
    for ri, (no, txt, inf, mine, jdg, why) in enumerate(items):
        ry = Inches(2.4) + Inches(1.0) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.0), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        cells = [no, txt, inf, mine, jdg, why]
        for i, cell in enumerate(cells):
            color = BLACK
            bold = False
            if i == 4:
                color = GREEN if "✅" in cell else RED
                bold = True
            add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[i] - Inches(0.1), Inches(0.8),
                     cell, size=11, color=color, bold=bold,
                     align=PP_ALIGN.CENTER if i in [0, 2, 3, 4] else PP_ALIGN.LEFT,
                     line_h=1.3)
            cx += cols_w[i]
    # 底部
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "判断核心: 影响大但不可影响 = 需向上争取; 可影响但影响小 = 次要动作",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m3_06_exercise2(idx):
    s = slide_m3(idx, 50, "✋ 练习二: 构建你自己的问题树",
                 "用你在花刺投票里识别出的「 刺票最多障碍 」 — 一步步来")
    # 4 步卡
    steps = [
        ("第一步", "症状描述", "我要分析的难题\n(用数字 / 现象描述, 越具体越好)"),
        ("第二步", "第一层拆解", "3–4 个主方向\n完整不重叠"),
        ("第三步", "展开关键分支", "选择最重要的 1–2 个一级分支\n展开它的二级原因"),
        ("第四步", "选择切入点", "影响足够大 × 我能真正影响\n两者必须同时满足"),
    ]
    card_w = Inches(3.0); card_h = Inches(4.0)
    gap = Inches(0.18)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (no, name, sub) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GREEN)
        add_rect(s, x, y, card_w, Inches(0.6), fill=GREEN)
        add_text(s, x, y + Inches(0.1), card_w, Inches(0.4),
                 no, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.5),
                 name, size=15, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(1.4) + Inches(0.4) * j,
                     card_w - Inches(0.5), Inches(0.4),
                     line, size=11, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.3)
        # 留白填写区
        add_rect(s, x + Inches(0.3), y + Inches(2.9), card_w - Inches(0.6), Inches(1.0),
                 fill=WHITE, line=GREEN)
        add_text(s, x + Inches(0.3), y + Inches(3.25), card_w - Inches(0.6), Inches(0.4),
                 "填写区", size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.4),
             "提示: 不要一开始就想「 怎么解决 」 — 先把问题拆清楚",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "预计用时: 15–20 分钟",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def m3_07_magic_questions(idx):
    s = slide_m3(idx, 51, "第四节: 魔力破解提问 — 5 把打开新视角的钥匙",
                 "不是清单, 是 5 把钥匙 — 用每把去「 试 」你的切入点")
    qs = [
        ("问 1", "破假设", "我们一直认为某事不可能 / 必须是这样\n—— 这个假设真的是对的吗?", NAVY),
        ("问 2", "借他山", "其他行业 / 其他人面对过类似问题?\n他们怎么解决的?", ORANGE),
        ("问 3", "极端情景", "资源充足时我会怎么做?\n1/10 资源时呢?", GREEN),
        ("问 4", "换位思考", "从客户 / 团队 / 竞品 / 上级角度\n这个问题是什么样的?", GOLD),
        ("问 5", "倒推法", "想象一年后问题已经完美解决\n你是怎么做到的?", RED),
    ]
    card_w = Inches(2.45); card_h = Inches(4.0)
    gap = Inches(0.12)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (no, name, body, color) in enumerate(qs):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_rect(s, x, y, card_w, Inches(0.55), fill=color)
        add_text(s, x, y + Inches(0.07), card_w, Inches(0.4),
                 no, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.75), card_w - Inches(0.4), Inches(0.5),
                 name, size=15, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.5), y + Inches(1.3), card_w - Inches(1.0), Emu(15000), fill=color)
        for j, line in enumerate(body.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(1.5) + Inches(0.55) * j,
                     card_w - Inches(0.5), Inches(0.55),
                     line, size=10, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.3)
    # 用法
    y2 = Inches(6.25)
    add_round(s, Inches(0.5), y2, Inches(12.33), Inches(0.85), fill=NAVY_DARK)
    add_text(s, Inches(0.7), y2 + Inches(0.1), Inches(11.9), Inches(0.4),
             "使用方式", size=12, color=ORANGE, bold=True)
    add_text(s, Inches(0.7), y2 + Inches(0.45), Inches(11.9), Inches(0.4),
             "对切入点逐一「 试 」5 问 → 选最让你「 没想过 」的那一问深想",
             size=12, color=WHITE)
    # 底部
    add_text(s, Inches(0.5), Inches(7.15), Inches(12.33), Inches(0.3),
             "注意: 问题树先拆清楚, 再用魔力提问 — 否则魔力提问变成对模糊大问题的泛泛思考",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def m3_08_magic_example(idx):
    s = slide_m3(idx, 52, "⚡ 示例: 李明用魔力提问应对「 没有体验引导 」",
                 "切入点 A2 → 5 把钥匙逐一试")
    # 表
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=GREEN)
    headers = ["魔力问题", "李明的思考", "新视角?"]
    cols_w = [Inches(2.0), Inches(8.5), Inches(1.83)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER if i == 2 else PP_ALIGN.LEFT)
        cx += cols_w[i]
    rows = [
        ("问 1 破假设",
         "我们一直认为『 让客户先讲需求再看产品 』是标准流程 —— 但小王的经验说明, 先让客户用, 他们才会真正产生自己的问题",
         "✓"),
        ("问 2 借他山 ★",
         "汽车 4S 店试驾 / 手机样机体验 / 宜家样板间 —— 这些行业都是用『 先让用户体验 』建立购买欲望",
         "✓✓"),
        ("问 3 极端情景",
         "充足: 每款产品设计 30 秒体验脚本; 1/10: 至少让每个导购记住『 要不您先感受一下 』",
         "✓"),
        ("问 4 换位思考",
         "客户: 进店不确定要不要买, 导购一上来就讲参数 → 我会想逃; 说『 先感受一下 』, 反而放松",
         "✓"),
        ("问 5 倒推法",
         "一年后体验式接待成为所有导购标准做法 → 先做 3 步指南, 再每周 5 分钟实战分享, 累积 6 个月",
         "✓"),
    ]
    for ri, (q, t, n) in enumerate(rows):
        ry = Inches(2.4) + Inches(0.7) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.7), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        # q
        col_q = GREEN if "★" in q else NAVY_DARK
        add_text(s, cx + Inches(0.1), ry + Inches(0.18), cols_w[0] - Inches(0.1), Inches(0.5),
                 q.replace(" ★", "").replace("★", ""), size=11, color=col_q, bold=True)
        cx += cols_w[0]
        # t
        add_text(s, cx + Inches(0.1), ry + Inches(0.12), cols_w[1] - Inches(0.1), Inches(0.6),
                 t, size=10, color=BLACK, line_h=1.3)
        cx += cols_w[1]
        # n
        col = GREEN if "✓" in n else GRAY
        add_text(s, cx + Inches(0.1), ry + Inches(0.2), cols_w[2] - Inches(0.1), Inches(0.5),
                 n, size=14, color=col, bold=True, align=PP_ALIGN.CENTER)
    # 收获 + 新方向
    y2 = Inches(6.4)
    add_round(s, Inches(0.5), y2, Inches(12.33), Inches(0.85), fill=GREEN)
    add_text(s, Inches(0.7), y2 + Inches(0.1), Inches(11.9), Inches(0.4),
             "最让李明受启发的: 问 2 (借他山) — 体验式销售不只是小王风格, 是消费品行业已有的成熟逻辑",
             size=12, color=WHITE, bold=True)
    add_text(s, Inches(0.7), y2 + Inches(0.45), Inches(11.9), Inches(0.4),
             "→ 新方向: 把体验式接待定义为门店标准 SOP, 不只是『 鼓励大家学小王 』",
             size=12, color=WHITE, bold=True)


def m3_09_exercise3(idx):
    s = slide_m3(idx, 53, "✋ 练习三: 对你的切入点用魔力提问",
                 "不需要每个问题都有完整答案 — 挑选最让你有感觉的 1–2 个深想")
    items = [
        ("问 1 (破假设)", "我们一直认为 ____ 是理所当然的, 这个假设真的对吗?"),
        ("问 2 (借他山)", "其他行业有没有解决过类似问题? 怎么解决的?"),
        ("问 3 (极端情景)", "资源充足我会怎么做? 1/10 资源时呢?"),
        ("问 4 (换位思考)", "从客户 / 成员 / 竞品 / 上级角度看, 这个问题是什么样的?"),
        ("问 5 (倒推法)", "想象一年后问题已经解决, 我是怎么做到的?"),
    ]
    y = Inches(1.85)
    row_h = Inches(0.85)
    for i, (q, body) in enumerate(items):
        ry = y + (row_h + Inches(0.1)) * i
        add_round(s, Inches(0.5), ry, Inches(12.33), row_h, fill=LIGHT, line=GREEN)
        # 问题
        add_rect(s, Inches(0.5), ry, Inches(3.2), row_h, fill=GREEN)
        add_text(s, Inches(0.6), ry + Inches(0.25), Inches(3.0), Inches(0.4),
                 q, size=12, color=WHITE, bold=True)
        # 引导
        add_text(s, Inches(3.85), ry + Inches(0.1), Inches(8.8), row_h - Inches(0.2),
                 body, size=11, color=BLACK, line_h=1.4)
    # 底部
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "最让你受启发的问题是: 问 ____",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.4),
             "由此产生的新视角或新方向: _______________________________________________",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def m3_10_share(idx):
    s = slide_m3(idx, 54, "小组互动: 分享你的新方向",
                 "2 人一组 · 每人 2 分钟 · 1 个问题帮你把方向说得更具体")
    # 左
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "分享内容 (每人 2 分钟)", size=14, color=ORANGE, bold=True)
    items = [
        "① 你的难题 (一句话)",
        "② 你的切入点 (一句话)",
        "③ 哪个魔力问题给了你新角度",
        "④ 由此产生的新方向是什么",
    ]
    for i, t in enumerate(items):
        add_text(s, lx + Inches(0.4), ly + Inches(0.85) + Inches(0.55) * i,
                 Inches(5.4), Inches(0.5),
                 t, size=12, color=BLACK)
    # 时间
    add_rect(s, lx + Inches(0.3), ly + Inches(3.2), Inches(5.6), Emu(15000), fill=GREEN)
    add_text(s, lx + Inches(0.3), ly + Inches(3.35), Inches(5.6), Inches(0.5),
             "时间: 每人 2 分钟, 共 4–5 分钟",
             size=12, color=GREEN, bold=True)
    # 右
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.2), Inches(5.2), Inches(0.5),
             "对方的任务 (只有一个问题)", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(0.9), Inches(5.2), Inches(2.0),
             "「 如果你的这个新方向真的有效,\n一年后你的团队会有什么不同? 」",
             size=18, color=WHITE, bold=True, line_h=1.5)
    add_rect(s, rx + Inches(0.4), ry + Inches(3.0), Inches(5.2), Emu(20000), fill=GREEN)
    add_text(s, rx + Inches(0.4), ry + Inches(3.15), Inches(5.2), Inches(1.4),
             "这个问题帮助你\n把新方向说得更具体\n从「 应该 」变成「 会有什么不同 」",
             size=12, color=WHITE, line_h=1.5)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "分享后, 把你的「 新方向 」写下来 — 这是第四部分共创的主题",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m3_11_recap(idx):
    s = slide_m3(idx, 55, "✅ 第三部分知识框架",
                 "M3 = 问题树 (说清楚) + 魔力提问 (找新视角)")
    items = [
        ("难题的本质", "经验失灵\n用了有效方法但结果不符", RED),
        ("问题树 4 步", "症状 → 拆分 → 展开\n→ 主动选切入点", GREEN),
        ("切入点标准", "影响大 × 我能影响\n两者必须同时满足", ORANGE),
        ("魔力提问 5 问", "破假设 / 借他山\n极端 / 换位 / 倒推", NAVY),
    ]
    card_w = Inches(2.95); card_h = Inches(3.0)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(2.2)
    for i, (head, body, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_rect(s, x, y, card_w, Inches(0.7), fill=color)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.4),
                 head, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_text(s, x + Inches(0.2), y + Inches(0.95) + Inches(0.5) * j,
                     card_w - Inches(0.4), Inches(0.5),
                     line, size=11, color=BLACK, align=PP_ALIGN.CENTER)
    # 产出
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.4),
             "完成第三部分后你拥有:", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    out_items = [
        ("对一个真实难题的完整问题树", GREEN),
        ("一个明确的切入点", ORANGE),
        ("至少一个新方向 (以前没认真想过的)", NAVY),
    ]
    for i, (txt, col) in enumerate(out_items):
        x = Inches(0.5) + Inches(4.2) * i
        add_round(s, x, Inches(5.95), Inches(4.0), Inches(0.7), fill=LIGHT, line=col)
        add_text(s, x, Inches(6.0), Inches(4.0), Inches(0.6),
                 txt, size=11, color=col, bold=True, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "→ 这个新方向 = 第四部分共创环节的主题",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m3_12_connect(idx):
    s = slide_m3(idx, 56, "连接到下一部分",
                 "你已经找到了一个新方向 — 接下来要把它变成具体的行动方案")
    # 左
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "你刚刚用 M3 做了什么", size=14, color=NAVY_DARK, bold=True)
    bullets = [
        "把模糊的大难题拆成了具体方向",
        "找到了可以下手的切入点",
        "用 5 个魔力问题打开了新视角",
        "收获了一个「 没想过 」的新方向",
    ]
    bullet_block(s, lx + Inches(0.4), ly + Inches(0.9), Inches(5.2), Inches(2.0), bullets, size=12)
    # 接下来的问题
    add_rect(s, lx + Inches(0.3), ly + Inches(3.2), Inches(5.4), Emu(20000), fill=NAVY)
    add_text(s, lx + Inches(0.3), ly + Inches(3.35), Inches(5.4), Inches(0.4),
             "接下来的问题", size=12, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(3.75), Inches(5.4), Inches(1.0),
             "→ 这个新方向怎么落地?\n→ 哪些是短期的快赢?\n→ 哪些是中长期准备?",
             size=12, color=BLACK, line_h=1.5)
    # 右
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "下一部分要回答的问题", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(0.9), Inches(5.5), Inches(1.2),
             "怎么带着团队\n把新方向变成具体的行动方案?",
             size=18, color=WHITE, bold=True, line_h=1.4)
    add_rect(s, rx + Inches(0.4), ry + Inches(2.7), Inches(5.3), Emu(20000), fill=GOLD)
    add_text(s, rx + Inches(0.4), ry + Inches(2.85), Inches(5.3), Inches(1.7),
             "高效脑暴双矩阵\n\n先在 4 个象限发散\n再用评估矩阵聚焦\n让团队高质量共创",
             size=12, color=WHITE, bold=True, line_h=1.5)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "下一部分: M4 引领共创 — 高效脑暴双矩阵",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m3_13_section_close(idx):
    s = slide_m3(idx, 57, "模块三收尾 · 午休",
                 "下午进入 M4 引领共创 — 把新方向变成可执行方案")
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.5),
             "M3", size=120, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.7),
             "应对难题", size=32, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             "Solve Hard Problems", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    quote_block(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                "解决问题的最大浪费, 不是方法选错了, 而是在搞清楚问题之前就开始行动。",
                author="—— 第三部分核心信念")
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "请准备: 你的「 新方向 」一句话 + 准备带进第四部分",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


print("模块三定义完成")


# ===================================================================
# 第 7 段: 模块四 引领共创 (24 页)
# ===================================================================
def m4_section_divider(idx):
    s = new_slide()
    page_chrome(s, idx, section="模块 4 · 引领共创")
    add_rect(s, 0, Inches(1.5), SW, Inches(4.5), fill=GOLD)
    add_text(s, Inches(0.5), Inches(1.7), Inches(5), Inches(3),
             "M4", size=180, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(2.3), Inches(8), Inches(1.0),
             "引领共创", size=44, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(3.5), Inches(8), Inches(0.6),
             "Co-create With Team", size=18, color=WHITE)
    quote_block(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.6),
                "发散和收敛必须分开 — 否则要么冷场, 要么乱场。",
                author="—— 本模块核心信念")


def slide_m4(idx, n, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section="模块 4 · 引领共创")
    title_block(s, title, subtitle, accent=GOLD)
    return s


def m4_01_opening(idx):
    s = slide_m4(idx, 59, "开场: 两种让人崩溃的团队讨论",
                 "冷场型 / 乱场型 — 根源都是发散和收敛没分开")
    # 冷场型
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(4.7), fill=LIGHT, line=GRAY_LIGHT)
    add_rect(s, lx, ly, Inches(6.2), Inches(0.7), fill=GRAY)
    add_text(s, lx, ly + Inches(0.15), Inches(6.2), Inches(0.5),
             "第一种: 冷场型", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, lx + Inches(0.3), ly + Inches(0.9), Inches(5.6), Inches(0.5),
             "会议开始", size=12, color=NAVY_DARK, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(1.3), Inches(5.6), Inches(1.5),
             "「 大家说说, 接下来我们应该怎么做 」\n\n沉默 5 秒, 3 个人说话\n说的都是已经在做的事\n管理者总结 (其实是自己想法)\n「 好, 那就这样吧 」散会",
             size=11, color=BLACK, line_h=1.4)
    # 病根
    add_rect(s, lx + Inches(0.3), ly + Inches(3.5), Inches(5.6), Inches(1.0), fill=RED)
    add_text(s, lx + Inches(0.45), ly + Inches(3.6), Inches(5.3), Inches(0.4),
             "病根", size=11, color=WHITE, bold=True)
    add_text(s, lx + Inches(0.45), ly + Inches(3.95), Inches(5.3), Inches(0.5),
             "其他 9 个人的想法, 一个也没出来",
             size=11, color=WHITE)
    # 乱场型
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(4.7), fill=LIGHT, line=GRAY_LIGHT)
    add_rect(s, rx, ry, Inches(5.83), Inches(0.7), fill=ORANGE)
    add_text(s, rx, ry + Inches(0.15), Inches(5.83), Inches(0.5),
             "第二种: 乱场型", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.3), ry + Inches(0.9), Inches(5.3), Inches(0.5),
             "30 分钟之后", size=12, color=NAVY_DARK, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(1.3), Inches(5.3), Inches(1.5),
             "感觉想法很多, 但没重点\n管理者: 「 那就先做 A 和 B 」\n有人: 「 等等, C 也重要 」\n另一个人: 「 其实 D 比 B 紧急 」\n散会, 没清楚行动清单",
             size=11, color=BLACK, line_h=1.4)
    # 病根
    add_rect(s, rx + Inches(0.3), ry + Inches(3.5), Inches(5.3), Inches(1.0), fill=RED)
    add_text(s, rx + Inches(0.45), ry + Inches(3.6), Inches(5), Inches(0.4),
             "病根", size=11, color=WHITE, bold=True)
    add_text(s, rx + Inches(0.45), ry + Inches(3.95), Inches(5), Inches(0.5),
             "没有结构帮从混乱中找优先项",
             size=11, color=WHITE)
    # 底部
    quote_block(s, Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.5),
                "发散和收敛没分开处理 — 高效脑暴双矩阵就是把这两件事分开。",
                author="—— M4 要解决的问题")


def m4_02_failure_modes(idx):
    s = slide_m4(idx, 60, "第一节: 最常见的脑暴失败模式",
                 "诊断一下你的团队讨论通常在哪个模式上")
    items = [
        ("想法同质化", "讨论来讨论去, 都是相似方向", "没有结构强制考虑不同维度"),
        ("主管意见主导", "主管先说, 后面都在顺这个方向", "发散阶段引入了评判 (隐性的)"),
        ("优秀想法被否定", "好主意被『 这个不现实 』打死", "发散阶段过早评价"),
        ("聚焦变成投票", "靠举手表决或『 大家觉得哪个好 』", "没有客观评估维度, 结果被话语权主导"),
        ("没有聚焦", "所有想法都记录, 但没确认优先项", "缺少收敛机制"),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=GOLD)
    headers = ["失败模式", "表现", "根本原因"]
    cols_w = [Inches(2.5), Inches(5.0), Inches(4.83)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True,
                 align=PP_ALIGN.LEFT)
        cx += cols_w[i]
    for ri, (mode, sym, why) in enumerate(items):
        ry = Inches(2.4) + Inches(0.85) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.85), fill=bg, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), ry, Inches(0.15), Inches(0.85), fill=ORANGE)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.25), ry + Inches(0.25), cols_w[0] - Inches(0.3), Inches(0.5),
                 mode, size=12, color=NAVY_DARK, bold=True)
        cx += cols_w[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.18), cols_w[1] - Inches(0.1), Inches(0.6),
                 sym, size=11, color=BLACK, line_h=1.4)
        cx += cols_w[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.18), cols_w[2] - Inches(0.1), Inches(0.6),
                 why, size=11, color=BLACK, line_h=1.4)
    # 底部
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "5 种模式都有一个共同根源: 发散和收敛的结构缺失",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m4_03_diverge_matrix(idx):
    s = slide_m4(idx, 61, "第二节: 发散矩阵 — 强制团队在 4 个象限都产出想法",
                 "横轴: 内部可控 / 需要外部 · 纵轴: 短期 / 中长期")
    # 矩阵
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(8.5), Inches(5.0), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(8), Inches(0.5),
             "发散矩阵 2×2", size=14, color=ORANGE, bold=True)
    # 矩阵框架
    mx = lx + Inches(0.8); my = ly + Inches(1.0)
    mw = Inches(7.0); mh = Inches(3.6)
    add_rect(s, mx, my, mw, mh, fill=WHITE, line=NAVY)
    add_line(s, mx + mw / 2, my, mx + mw / 2, my + mh, color=NAVY, w=Pt(2))
    add_line(s, mx, my + mh / 2, mx + mw, my + mh / 2, color=NAVY, w=Pt(2))
    # 四象限
    quads = [
        ("Q1", "快赢区", "立刻能做\n立刻有效", GREEN),
        ("Q2", "协调推进区", "近期要启动\n需要争取支持", ORANGE),
        ("Q3", "能力建设区", "团队能力积累\n中长期", NAVY),
        ("Q4", "战略布局区", "提前规划\n需要资源支持", RED),
    ]
    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, ((cx_i, cy_i), (no, name, desc, col)) in enumerate(zip(positions, quads)):
        qx = mx + (mw / 2) * cx_i
        qy = my + (mh / 2) * cy_i
        add_rect(s, qx + Inches(0.1), qy + Inches(0.1), Inches(0.8), Inches(0.45), fill=col)
        add_text(s, qx + Inches(0.1), qy + Inches(0.15), Inches(0.8), Inches(0.4),
                 no, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, qx + Inches(1.0), qy + Inches(0.15), Inches(2.0), Inches(0.4),
                 name, size=13, color=col, bold=True)
        for j, line in enumerate(desc.split("\n")):
            add_text(s, qx + Inches(0.15), qy + Inches(0.7) + Inches(0.4) * j,
                     (mw / 2) - Inches(0.3), Inches(0.4),
                     line, size=10, color=BLACK, line_h=1.3)
    # 轴标签
    add_text(s, mx, my - Inches(0.4), mw, Inches(0.3),
             "短期 (1–4 周)", size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx, my + mh + Inches(0.1), mw, Inches(0.3),
             "中长期 (1–3 月)", size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx - Inches(0.7), my + mh / 2 - Inches(0.2), Inches(0.7), Inches(0.4),
             "团队内部可控", size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx + mw + Inches(0.1), my + mh / 2 - Inches(0.2), Inches(2.3), Inches(0.4),
             "需要外部资源/协调", size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 右: 主持原则
    rx = Inches(9.2); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(3.65), Inches(5.0), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(3), Inches(0.4),
             "主持原则", size=13, color=ORANGE, bold=True)
    principles = [
        "发散阶段不评判",
        "逐格引导, 不乱放",
        "主持人先不说",
        "每格至少 2–3 条",
    ]
    for i, t in enumerate(principles):
        yy = ry + Inches(0.8) + Inches(0.7) * i
        add_oval(s, rx + Inches(0.3), yy + Inches(0.1), Inches(0.2), Inches(0.2), fill=ORANGE)
        add_text(s, rx + Inches(0.6), yy, Inches(3), Inches(0.5),
                 t, size=12, color=WHITE, bold=True)
    add_rect(s, rx + Inches(0.3), ry + Inches(3.7), Inches(3.0), Emu(15000), fill=ORANGE)
    add_text(s, rx + Inches(0.3), ry + Inches(3.85), Inches(3), Inches(1.1),
             "意外的想法\n往往最有价值",
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_h=1.4)
    # 底部
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3),
             "为什么用矩阵: 自由讨论总倾向 Q1, 矩阵强制每个格子都产出",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m4_04_full_case(idx):
    s = slide_m4(idx, 62, "⚡ 完整示例: 李明用发散矩阵带团队",
                 "主题: 如何让体验式接待成为全门店标准做法")
    # 矩阵示例
    items = {
        ("Q1", "短期+内部"): [
            "小王下周二晨会分享细节",
            "制作一页纸 3 步指南",
            "门口增加欢迎试用标识",
            "每日晚会 1 分钟报告",
        ],
        ("Q2", "短期+外部"): [
            "申请总部提供更多展示机",
            "联系区域培训经理旁听",
        ],
        ("Q3", "中长期+内部"): [
            "每人做体验式接待自我录像",
            "两周后针对未进展成员一对一辅导",
            "一个月后全员能力评估",
        ],
        ("Q4", "中长期+外部"): [
            "向总部申请纳入区域 SOP",
            "梳理差异化优势说明材料",
        ],
    }
    # 4 象限
    mx = Inches(0.5); my = Inches(1.85)
    mw = Inches(8.5); mh = Inches(4.4)
    add_rect(s, mx, my, mw, mh, fill=WHITE, line=NAVY)
    add_line(s, mx + mw / 2, my, mx + mw / 2, my + mh, color=NAVY, w=Pt(2))
    add_line(s, mx, my + mh / 2, mx + mw, my + mh / 2, color=NAVY, w=Pt(2))
    # 轴标签
    add_text(s, mx, my - Inches(0.35), mw, Inches(0.3),
             "短期 (1–4 周)", size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx, my + mh + Inches(0.05), mw, Inches(0.3),
             "中长期 (1–3 月)", size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx + mw / 2 - Inches(1.2), my - Inches(0.05), Inches(2.4), Inches(0.2),
             "团队内部可控", size=9, color=BLACK)
    add_text(s, mx + mw / 2 - Inches(1.0), my - Inches(0.05), Inches(2), Inches(0.2),
             "团队内部可控", size=9, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx + mw / 2 + Inches(0.0), my - Inches(0.05), Inches(2), Inches(0.2),
             "需要外部资源/协调", size=9, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 4 象限内容
    quads_data = [
        (0, 0, "Q1", "快赢", GREEN, items[("Q1", "短期+内部")]),
        (1, 0, "Q2", "协调推进", ORANGE, items[("Q2", "短期+外部")]),
        (0, 1, "Q3", "能力建设", NAVY, items[("Q3", "中长期+内部")]),
        (1, 1, "Q4", "战略布局", RED, items[("Q4", "中长期+外部")]),
    ]
    for cx_i, cy_i, no, name, col, content in quads_data:
        qx = mx + (mw / 2) * cx_i + Inches(0.1)
        qy = my + (mh / 2) * cy_i + Inches(0.05)
        add_text(s, qx + Inches(0.05), qy, Inches(0.5), Inches(0.3),
                 no, size=11, color=col, bold=True)
        add_text(s, qx + Inches(0.55), qy, Inches(2), Inches(0.3),
                 name, size=10, color=col, bold=True)
        for j, line in enumerate(content):
            add_text(s, qx + Inches(0.15), qy + Inches(0.4) + Inches(0.42) * j,
                     (mw / 2) - Inches(0.3), Inches(0.42),
                     "• " + line, size=10, color=BLACK, line_h=1.3)
    # 右: 最意外的想法
    rx = Inches(9.2); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(3.65), Inches(4.4), fill=GOLD)
    add_text(s, rx + Inches(0.3), ry + Inches(0.3), Inches(3), Inches(0.4),
             "最让李明意外的", size=12, color=WHITE, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(0.9), Inches(3), Inches(1.5),
             "「 门口增加\n欢迎试用标识 」",
             size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_h=1.4)
    add_rect(s, rx + Inches(0.3), ry + Inches(2.6), Inches(3), Emu(15000), fill=WHITE)
    add_text(s, rx + Inches(0.3), ry + Inches(2.75), Inches(3), Inches(1.6),
             "20 分钟就能做完\n成本几乎为零\n他以前从来没想过",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, line_h=1.5)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "一次好的发散, 会有 20–30 条便利贴, 包含意料之中和意料之外的",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m4_05_exercise1(idx):
    s = slide_m4(idx, 63, "✋ 练习一: 主持发散矩阵",
                 "用你的「 新方向 」做一轮 · 3–4 人小组 · 每人轮流主持一格")
    # 左: 主题卡
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.2), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "主题 (来自 M3 新方向)", size=14, color=ORANGE, bold=True)
    add_rect(s, lx + Inches(0.3), ly + Inches(0.8), Inches(5.6), Inches(0.8), fill=WHITE, line=ORANGE)
    add_rect(s, lx + Inches(0.3), ly + Inches(1.65), Inches(5.6), Emu(8000), fill=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(1.75), Inches(5.6), Inches(0.4),
             "(你的新方向:  _______________________________________ )",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)
    # 4 格引导
    add_text(s, lx + Inches(0.3), ly + Inches(2.3), Inches(5.6), Inches(0.4),
             "逐格引导提问示例", size=12, color=ORANGE, bold=True)
    guides = [
        "Q1: 「 接下来 4 周, 我们自己能做, 不需要等总部批的, 有哪些? 」",
        "Q2: 「 1 个月内能看到效果, 但需要总部 / 其他部门配合的? 」",
        "Q3: 「 3 个月内, 我们团队自己能在什么能力上做积累? 」",
        "Q4: 「 需要提前布局, 需要和上级谈的事, 哪怕效果 3 个月后才体现? 」",
    ]
    for i, g in enumerate(guides):
        add_text(s, lx + Inches(0.4), ly + Inches(2.7) + Inches(0.45) * i,
                 Inches(5.4), Inches(0.4),
                 g, size=9, color=BLACK, line_h=1.3)
    # 右: 反思
    rx = Inches(7.0); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(5.83), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.2), Inches(5.2), Inches(0.5),
             "发散结束后, 主持人反思", size=14, color=ORANGE, bold=True)
    items = [
        "哪个格子最难引导团队产想法? 为什么?",
        "有没有让我意外的想法? 是哪一条?",
        "发散阶段有没有出现「 评判 」? 谁说了什么?",
        "如果重新主持, 我会在哪个地方做得不一样?",
    ]
    for i, t in enumerate(items):
        yy = ry + Inches(0.9) + Inches(0.85) * i
        add_round(s, rx + Inches(0.4), yy, Inches(5.2), Inches(0.75), fill=WHITE)
        add_text(s, rx + Inches(0.55), yy + Inches(0.1), Inches(5.0), Inches(0.6),
                 t, size=11, color=BLACK, line_h=1.4)
    # 底部
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "预计用时: 20–25 分钟 | 每人轮流担任主持或全程一人主持",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m4_06_converge_matrix(idx):
    s = slide_m4(idx, 64, "第四节: 收敛矩阵 — 把发散想法聚焦到优先行动",
                 "影响力 × 落地难度 = 4 个区间, ⭐ 是今天要做的事")
    # 矩阵
    mx = Inches(0.5); my = Inches(1.85)
    mw = Inches(8.0); mh = Inches(4.5)
    add_rect(s, mx, my, mw, mh, fill=WHITE, line=NAVY)
    add_line(s, mx + mw / 2, my, mx + mw / 2, my + mh, color=NAVY, w=Pt(2))
    add_line(s, mx, my + mh / 2, mx + mw, my + mh / 2, color=NAVY, w=Pt(2))
    # 轴标签
    add_text(s, mx, my - Inches(0.35), mw, Inches(0.3),
             "高影响力", size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx, my + mh + Inches(0.05), mw, Inches(0.3),
             "低影响力", size=11, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx - Inches(0.3), my + mh / 2 - Inches(0.15), Inches(0.5), Inches(0.3),
             "易", size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx + mw + Inches(0.05), my + mh / 2 - Inches(0.15), Inches(1.0), Inches(0.3),
             "落地困难", size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx - Inches(0.3), my + Inches(0.3), Inches(0.5), Inches(0.3),
             "难", size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, mx + mw + Inches(0.05), my + mh - Inches(0.4), Inches(1.0), Inches(0.3),
             "易", size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    # 4 区间
    quads = [
        (0, 0, "★ 优先行动区", "高影响 + 落地容易\n先做这里", GREEN, "今天确定要做"),
        (1, 0, "战略项目区", "高影响 + 落地困难\n需要专门规划", ORANGE, "规划推进"),
        (0, 1, "随机应变区", "低影响 + 落地困难\n酌情处理", GRAY, "不投入太多精力"),
        (1, 1, "暂不考虑区", "低影响 + 落地困难\n本轮不做", GRAY, "记录即可"),
    ]
    for cx_i, cy_i, name, desc, col, hint in quads:
        qx = mx + (mw / 2) * cx_i + Inches(0.1)
        qy = my + (mh / 2) * cy_i + Inches(0.1)
        add_rect(s, qx, qy, mw / 2 - Inches(0.2), Inches(0.5), fill=col)
        add_text(s, qx, qy + Inches(0.07), mw / 2 - Inches(0.2), Inches(0.4),
                 name, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(desc.split("\n")):
            add_text(s, qx + Inches(0.15), qy + Inches(0.6) + Inches(0.4) * j,
                     mw / 2 - Inches(0.3), Inches(0.4),
                     line, size=10, color=BLACK, line_h=1.3)
        add_text(s, qx + Inches(0.15), qy + (mh / 2) - Inches(0.5), mw / 2 - Inches(0.3), Inches(0.4),
                 hint, size=10, color=col, bold=True, align=PP_ALIGN.CENTER)
    # 右: 主持步骤
    rx = Inches(8.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(4.15), Inches(4.5), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(3.5), Inches(0.4),
             "主持收敛矩阵 3 步", size=13, color=ORANGE, bold=True)
    steps = [
        ("01", "独立评估", "每人拿空白矩阵\n独立标位置"),
        ("02", "对比分歧", "只讨论分歧大的 2–3 条\n看理解是否一致"),
        ("03", "确认优先行动", "⭐ 区选 2–3 条\n谁负责 / 何时第一步"),
    ]
    for i, (no, name, desc) in enumerate(steps):
        yy = ry + Inches(0.8) + Inches(1.1) * i
        add_oval(s, rx + Inches(0.3), yy, Inches(0.5), Inches(0.5), fill=ORANGE)
        add_text(s, rx + Inches(0.3), yy + Inches(0.08), Inches(0.5), Inches(0.4),
                 no, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, rx + Inches(1.0), yy + Inches(0.02), Inches(2.9), Inches(0.4),
                 name, size=12, color=WHITE, bold=True)
        add_text(s, rx + Inches(1.0), yy + Inches(0.4), Inches(2.9), Inches(0.7),
                 desc, size=10, color=GRAY, line_h=1.3)
    # 底部
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "主持人不应主导结果 — 让矩阵的逻辑引导聚焦, 不是「 主管说了算 」",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m4_07_converge_case(idx):
    s = slide_m4(idx, 65, "⚡ 示例: 李明用收敛矩阵确定优先行动",
                 "从发散矩阵的 9 条重要想法开始评估")
    items = [
        ("小王晨会分享 + 3 步指南", "高", "易", "⭐ 优先行动", GREEN),
        ("每日晚会 1 分钟报告", "高", "易", "⭐ 优先行动", GREEN),
        ("门口欢迎试用标识", "中", "易", "随机应变", GRAY),
        ("申请总部展示机", "高", "难", "战略项目", ORANGE),
        ("体验式接待自我录像", "中", "中", "战略项目", ORANGE),
        ("一个月后全员技能评估", "高", "中", "战略项目", ORANGE),
        ("申请纳入区域 SOP", "高", "难", "战略项目", ORANGE),
        ("梳理差异化优势说明", "中", "易", "随机应变", GRAY),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=GOLD)
    headers = ["想法", "影响力", "落地难度", "落在哪个区"]
    cols_w = [Inches(6.0), Inches(1.5), Inches(1.5), Inches(3.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True, align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri, (idea, inf, diff, where, col) in enumerate(items):
        ry = Inches(2.4) + Inches(0.5) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.5), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.1), ry + Inches(0.13), cols_w[0] - Inches(0.1), Inches(0.4),
                 idea, size=11, color=BLACK)
        cx += cols_w[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.13), cols_w[1] - Inches(0.1), Inches(0.4),
                 inf, size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.13), cols_w[2] - Inches(0.1), Inches(0.4),
                 diff, size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[2]
        # 标记色块
        add_rect(s, cx + Inches(0.5), ry + Inches(0.1), cols_w[3] - Inches(1.0), Inches(0.3), fill=col)
        add_text(s, cx + Inches(0.5), ry + Inches(0.12), cols_w[3] - Inches(1.0), Inches(0.3),
                 where, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 优先行动
    y2 = Inches(6.5)
    add_round(s, Inches(0.5), y2, Inches(12.33), Inches(0.55), fill=NAVY_DARK)
    add_text(s, Inches(0.5), y2 + Inches(0.07), Inches(12.33), Inches(0.45),
             "李明确认的 ⭐ 优先行动区 (3 条): 小王晨会分享 + 3 步指南 | 每日晚会 1 分钟报告 | 门口欢迎试用标识",
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def m4_08_exercise2(idx):
    s = slide_m4(idx, 66, "✋ 练习二: 收敛矩阵 — 把你们的发散想法聚焦",
                 "3 步走完 · 独立评估 → 对比分歧 → 确认优先")
    # 3 步
    steps = [
        ("第一步", "独立评估", "从发散矩阵里选最重要的 8 条\n在自己的矩阵上独立标注", NAVY),
        ("第二步", "对比分歧", "找出分歧最大的 2 条\n只讨论分歧项", ORANGE),
        ("第三步", "确认优先行动", "⭐ 区选 2–3 条\n负责人 / 第一步时间", GREEN),
    ]
    card_w = Inches(4.0); card_h = Inches(3.5)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(1.95)
    for i, (no, name, sub, col) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.6), fill=col)
        add_text(s, x, y + Inches(0.1), card_w, Inches(0.4),
                 no, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.5),
                 name, size=15, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(1.4) + Inches(0.5) * j,
                     card_w - Inches(0.5), Inches(0.5),
                     line, size=11, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.4)
    # 战略项目
    y2 = Inches(5.7)
    add_round(s, Inches(0.5), y2, Inches(12.33), Inches(1.0), fill=ORANGE)
    add_text(s, Inches(0.7), y2 + Inches(0.1), Inches(11.9), Inches(0.4),
             "战略项目区", size=14, color=WHITE, bold=True)
    add_text(s, Inches(0.7), y2 + Inches(0.5), Inches(11.9), Inches(0.5),
             "值得规划, 但不是立即执行 — 写明计划什么时候启动准备",
             size=12, color=WHITE)
    # 底部
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.3),
             "预计用时: 15–20 分钟",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


def m4_09_recap(idx):
    s = slide_m4(idx, 67, "✅ 第四部分知识框架",
                 "M4 = 发散矩阵 (让所有类型想法都出来) + 收敛矩阵 (有依据地聚焦)")
    # 左
    items = [
        ("核心问题", "为什么团队脑暴\n通常失败", RED),
        ("发散矩阵", "短期/中长期 × 内部/外部\n4 象限, 强制覆盖", GREEN),
        ("收敛矩阵", "影响力 × 落地难度\n⭐ 优先行动区", GOLD),
        ("主持原则", "发散不评判\n收敛不主导", NAVY),
    ]
    card_w = Inches(2.95); card_h = Inches(3.0)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = Inches(0.5)
    y = Inches(2.2)
    for i, (head, body, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_rect(s, x, y, card_w, Inches(0.7), fill=color)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.4),
                 head, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_text(s, x + Inches(0.2), y + Inches(0.95) + Inches(0.5) * j,
                     card_w - Inches(0.4), Inches(0.5),
                     line, size=12, color=BLACK, align=PP_ALIGN.CENTER)
    # 产出
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.4),
             "完成第四部分后你拥有:", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    out_items = [
        ("20–30 条来自 4 象限的多元想法", GREEN),
        ("明确的优先行动清单 (2–3 条)", GOLD),
        ("战略项目区行动 (有规划但未启动)", ORANGE),
    ]
    for i, (txt, col) in enumerate(out_items):
        x = Inches(0.5) + Inches(4.2) * i
        add_round(s, x, Inches(5.95), Inches(4.0), Inches(0.7), fill=LIGHT, line=col)
        add_text(s, x, Inches(6.0), Inches(4.0), Inches(0.6),
                 txt, size=11, color=col, bold=True, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "→ 优先行动清单 = 第五部分推演的对象",
             size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m4_10_connect(idx):
    s = slide_m4(idx, 68, "连接到下一部分",
                 "你已经有了一个明确的优先行动 — 接下来要为它做前瞻准备")
    # 左
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "你刚刚用 M4 做了什么", size=14, color=NAVY_DARK, bold=True)
    bullets = [
        "在 4 个象限发散出 20–30 条想法",
        "用收敛矩阵聚焦到 2–3 条优先行动",
        "明确了负责人和时间节点",
    ]
    bullet_block(s, lx + Inches(0.4), ly + Inches(0.9), Inches(5.2), Inches(2.0), bullets, size=12)
    # 接下来的问题
    add_rect(s, lx + Inches(0.3), ly + Inches(3.0), Inches(5.4), Emu(20000), fill=NAVY)
    add_text(s, lx + Inches(0.3), ly + Inches(3.15), Inches(5.4), Inches(0.4),
             "接下来的问题", size=12, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(3.55), Inches(5.4), Inches(1.0),
             "→ 这个行动哪里可能出问题?\n→ 出了问题怎么应对?\n→ 哪些机会值得提前准备?",
             size=12, color=BLACK, line_h=1.5)
    # 右
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.4), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "下一部分要回答的问题", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.4), ry + Inches(0.9), Inches(5.5), Inches(1.0),
             "怎么在行动开始前\n把风险和机会都想清楚?",
             size=18, color=WHITE, bold=True, line_h=1.4)
    add_rect(s, rx + Inches(0.4), ry + Inches(2.5), Inches(5.3), Emu(20000), fill=RED)
    add_text(s, rx + Inches(0.4), ry + Inches(2.65), Inches(5.3), Inches(1.9),
             "推演双表格\n\n风险推演 + 机会推演\n两张表必须同时做\n= 完整的前瞻视角",
             size=12, color=WHITE, bold=True, line_h=1.5)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "下一部分: M5 前瞻思考 — 推演双表格",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m4_11_section_close(idx):
    s = slide_m4(idx, 69, "模块四收尾 · 茶歇",
                 "回来后, 我们进入最后一个模块 M5 前瞻思考")
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.5),
             "M4", size=120, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.7),
             "引领共创", size=32, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             "Co-create With Team", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    quote_block(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                "发散和收敛必须分开 — 否则要么冷场, 要么乱场。",
                author="—— 第四部分核心信念")
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "请准备: 你的「 优先行动 」一句话 + 负责人 + 时间节点",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


# ===================================================================
# 第 8 段: 模块五 前瞻思考 (18 页)
# ===================================================================
def m5_section_divider(idx):
    s = new_slide()
    page_chrome(s, idx, section="模块 5 · 前瞻思考")
    add_rect(s, 0, Inches(1.5), SW, Inches(4.5), fill=RED)
    add_text(s, Inches(0.5), Inches(1.7), Inches(5), Inches(3),
             "M5", size=180, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(2.3), Inches(8), Inches(1.0),
             "前瞻思考", size=44, color=WHITE, bold=True)
    add_text(s, Inches(5.0), Inches(3.5), Inches(8), Inches(0.6),
             "Think Forward", size=18, color=WHITE)
    quote_block(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.6),
                "30 分钟的提前推演, 换 15 小时的救火时间。",
                author="—— 本模块核心信念")


def slide_m5(idx, n, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section="模块 5 · 前瞻思考")
    title_block(s, title, subtitle, accent=RED)
    return s


def m5_01_opening(idx):
    s = slide_m5(idx, 70, "开场: 救火队长的代价",
                 "上周是供应商延误, 这周是老员工离职, 下周是客户投诉")
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.33), Inches(1.0),
             "不是信息不在, 是没人花 30 分钟问:",
             size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.9), Inches(12.33), Inches(1.2),
             "这件事可能在哪里出问题?",
             size=36, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.33), Inches(1.2),
             "出了问题我们怎么应对?",
             size=36, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_round(s, Inches(1.5), Inches(5.2), Inches(4.8), Inches(1.0), fill=GOLD)
    add_text(s, Inches(1.5), Inches(5.3), Inches(4.8), Inches(0.4),
             "提前推演 30 分钟", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(5.7), Inches(4.8), Inches(0.4),
             "(防火)", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.0), Inches(5.3), Inches(4.5), Inches(0.4),
             "省去 15 小时救火", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.0), Inches(5.7), Inches(4.5), Inches(0.4),
             "(10:1 投入产出)", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "花 1 分钟想一件你过去「 本可以提前发现但没发现 」的意外",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m5_02_obstacles(idx):
    s = slide_m5(idx, 71, "第一节: 为什么「 提前想一想 」这么难",
                 "不是不知道, 是有 3 个具体障碍")
    items = [
        ("障碍 1", "眼前的事总是更紧急", "防火工作成果「 看不见 」", "做了防范, 问题没发生\n你永远不知道「 如果没做 」会怎样", RED),
        ("障碍 2", "对不确定性感到不舒服", "想未来的坏事本身就不舒服", "选择「 不想 」\n但不确定性不会消失, 只是变成意外时的慌乱", ORANGE),
        ("障碍 3", "风险清单进了抽屉", "识别了风险但没转化成行动", "识别 ≠ 准备\n没有具体预防措施 + 应急预案 = 没做", GRAY),
    ]
    card_h = Inches(1.5)
    for i, (no, name, sym, why, col) in enumerate(items):
        y = Inches(2.0) + (card_h + Inches(0.2)) * i
        add_round(s, Inches(0.5), y, Inches(12.33), card_h, fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(1.2), card_h, fill=col)
        add_text(s, Inches(0.5), y + Inches(0.3), Inches(1.2), Inches(0.5),
                 no, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.0), y + Inches(0.15), Inches(4.5), Inches(0.5),
                 name, size=14, color=NAVY_DARK, bold=True)
        add_text(s, Inches(2.0), y + Inches(0.6), Inches(4.5), Inches(0.9),
                 sym, size=11, color=BLACK, line_h=1.4)
        add_rect(s, Inches(6.7), y + Inches(0.2), Inches(0.05), card_h - Inches(0.4), fill=GRAY_LIGHT)
        add_text(s, Inches(6.9), y + Inches(0.2), Inches(5.8), Inches(0.4),
                 "为什么", size=10, color=GRAY, bold=True)
        for j, line in enumerate(why.split("\n")):
            add_text(s, Inches(6.9), y + Inches(0.55) + Inches(0.4) * j,
                     Inches(5.8), Inches(0.4),
                     line, size=11, color=BLACK, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.3),
             "M5 推演双表格 = 解决第 3 个障碍: 强制每条风险都有具体预防+应急",
             size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)


def m5_03_dual_table_logic(idx):
    s = slide_m5(idx, 72, "第二节: 推演双表格 — 两张表合在一起才完整",
                 "只做一张 = 视角不完整")
    items = [
        ("只做风险", "过度保守, 错过机会", "把所有时间放在「 这件事可能失败 」\n忘记了「 怎么让它更可能成功 / 怎么在成功时最大化收益 」", GRAY),
        ("只做机会", "低估障碍, 措手不及", "对未来充满期待, 但没为可预见的困难准备\n第一个意外就让你乱了阵脚", GRAY),
        ("两张都做", "完整的前瞻视角", "知道哪里可能出问题, 准备好了应对\n知道哪里可能有意外机会, 准备好了抓", GREEN),
    ]
    card_w = Inches(4.0); card_h = Inches(3.5)
    gap = Inches(0.15)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (name, sym, why, col) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.7), fill=col)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.4),
                 name, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(0.85), card_w - Inches(0.6), Inches(0.4),
                 sym, size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(why.split("\n")):
            add_text(s, x + Inches(0.3), y + Inches(1.4) + Inches(0.5) * j,
                     card_w - Inches(0.6), Inches(0.5),
                     line, size=11, color=BLACK, line_h=1.4, align=PP_ALIGN.CENTER)
    add_round(s, Inches(0.5), Inches(5.7), Inches(12.33), Inches(1.4), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.5),
             "「 主要的坑我看清楚了, 主要的好事我也准备好了,\n现在可以带着把握去做了 」",
             size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_h=1.4)
    add_text(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.3),
             "不是「 一切都好, 往前冲 」, 也不是「 处处是坑, 谨慎前行 」",
             size=11, color=GOLD, align=PP_ALIGN.CENTER)


def m5_04_risk_table(idx):
    s = slide_m5(idx, 73, "第三节: 风险推演表 — 不是「 担心什么 」, 是「 准备什么 」",
                 "每一列都有具体意义, 缺一不可")
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.6), fill=RED)
    headers = ["可能的风险", "概率", "影响", "预防措施", "应急预案", "负责人+时间"]
    cols_w = [Inches(3.3), Inches(0.9), Inches(0.9), Inches(2.5), Inches(2.5), Inches(2.23)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri in range(3):
        ry = Inches(2.5) + Inches(1.1) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.1), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        for ci in range(6):
            cx_next = cx + cols_w[ci]
            if ci >= 1 and ci <= 2:
                add_text(s, cx, ry + Inches(0.4), cols_w[ci], Inches(0.3),
                         "高/中/低", size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
            elif ci == 5:
                add_text(s, cx, ry + Inches(0.4), cols_w[ci], Inches(0.3),
                         "(谁 · 何时)", size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
            cx = cx_next
    y2 = Inches(5.9)
    add_round(s, Inches(0.5), y2, Inches(6.1), Inches(0.95), fill=ORANGE)
    add_text(s, Inches(0.7), y2 + Inches(0.1), Inches(5.7), Inches(0.4),
             "⚠ 预防措施", size=13, color=WHITE, bold=True)
    add_text(s, Inches(0.7), y2 + Inches(0.45), Inches(5.7), Inches(0.5),
             "在风险发生前做, 降低发生概率",
             size=11, color=WHITE)
    add_round(s, Inches(6.7), y2, Inches(6.13), Inches(0.95), fill=GREEN)
    add_text(s, Inches(6.9), y2 + Inches(0.1), Inches(5.8), Inches(0.4),
             "✓ 应急预案", size=13, color=WHITE, bold=True)
    add_text(s, Inches(6.9), y2 + Inches(0.45), Inches(5.8), Inches(0.5),
             "风险发生后做, 降低影响\n(某人, 在某时间内, 做某件具体的事)",
             size=11, color=WHITE, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.3),
             "概率高 × 影响大 = 最高优先级: 预防 AND 应急都要有",
             size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)


def m5_05_liming_risk(idx):
    s = slide_m5(idx, 74, "⚡ 完整示例: 李明的风险推演",
                 "行动: 全员推广体验式接待 3 步法, 成交率 21% → 30%")
    rows = [
        ("老导购抵触, 不配合执行", "高", "高", "①先让小王 + 2-3 人试行 2 周, 用数据说话; ②给老导购「 先观察 」 空间", "一对一谈话, 了解顾虑, 个别辅导; 1 月后仍不执行升级", GREEN),
        ("培训后执行不到位", "高", "高", "①每日晚会 1 分钟报告; ②李明每天随机抽查 3-5 次", "当天反馈, 不点名; 简化 3 步法卡片", GREEN),
        ("竞品加大促销, 客流下滑", "中", "高", "同步推进老客户维护; 持续监测竞品动向", "启动老客户专项回访; 申请区域支持", ORANGE),
        ("小王离职, 失去示范者", "低", "高", "本周完成 3 步法文字化 + 视频录制; 培养第 2 种子", "李明先承担示范; 启动第二顺位", GRAY),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=RED)
    headers = ["可能的风险", "概率", "影响", "预防", "应急"]
    cols_w = [Inches(3.5), Inches(1.0), Inches(1.0), Inches(3.5), Inches(3.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri, (risk, p, im, pre, post, col) in enumerate(rows):
        ry = Inches(2.4) + Inches(1.0) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.0), fill=bg, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), ry, Inches(0.1), Inches(1.0), fill=col)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.2), ry + Inches(0.1), cols_w[0] - Inches(0.2), Inches(0.8),
                 risk, size=11, color=NAVY_DARK, bold=True, line_h=1.3)
        cx += cols_w[0]
        add_text(s, cx, ry + Inches(0.35), cols_w[1], Inches(0.3),
                 p, size=12, color=col, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[1]
        add_text(s, cx, ry + Inches(0.35), cols_w[2], Inches(0.3),
                 im, size=12, color=col, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[2]
        add_text(s, cx + Inches(0.1), ry + Inches(0.08), cols_w[3] - Inches(0.1), Inches(0.85),
                 pre, size=10, color=BLACK, line_h=1.3)
        cx += cols_w[3]
        add_text(s, cx + Inches(0.1), ry + Inches(0.08), cols_w[4] - Inches(0.1), Inches(0.85),
                 post, size=10, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "高 × 高 = 必须双准备 | 中 × 高 = 必须有应急 | 其他 = 酌情",
             size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)


def m5_06_exercise1(idx):
    s = slide_m5(idx, 75, "✋ 练习一: 完成你的风险推演表",
                 "用 M4 确认的优先行动做推演 · 独立完成 · 不讨论")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "我推演的行动", size=14, color=ORANGE, bold=True)
    add_rect(s, lx + Inches(0.3), ly + Inches(0.8), Inches(5.6), Inches(0.7), fill=WHITE, line=ORANGE)
    add_text(s, lx + Inches(0.3), ly + Inches(0.95), Inches(5.6), Inches(0.4),
             "(你的优先行动)",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, lx + Inches(0.3), ly + Inches(1.7), Inches(5.6), Inches(0.5),
             "写 3-5 个风险", size=12, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(2.1), Inches(5.6), Inches(0.5),
             "关键: 概率高+影响大至少 1 条",
             size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(2.6), Inches(5.6), Inches(0.5),
             "预防措施 + 应急预案必须都有",
             size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(3.1), Inches(5.6), Inches(0.5),
             "负责人 + 时间节点具体到人和日期",
             size=10, color=BLACK)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "好 vs 不好的写法", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(0.85), Inches(5.5), Inches(0.4),
             "✗ 不好的写法", size=11, color=RED, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(1.25), Inches(5.5), Inches(0.5),
             "「 执行风险 」",
             size=12, color=WHITE)
    add_text(s, rx + Inches(0.3), ry + Inches(1.7), Inches(5.5), Inches(0.4),
             "✓ 好的写法", size=11, color=GREEN, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(2.1), Inches(5.5), Inches(1.5),
             "「 如果团队里有 3 个以上\n  成员在头两周没按 3 步法\n  执行, 客户体验没改善 」",
             size=12, color=WHITE, line_h=1.4)
    add_rect(s, rx + Inches(0.3), ry + Inches(3.6), Inches(5.5), Emu(15000), fill=GOLD)
    add_text(s, rx + Inches(0.3), ry + Inches(3.75), Inches(5.5), Inches(0.6),
             "能让没背景的人也清楚「 哦, 这件事可能出 」",
             size=10, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "预计用时: 15 分钟 | 独立完成 — 这是个人对行动风险的认真评估",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m5_07_opp_table(idx):
    s = slide_m5(idx, 76, "第四节: 机会推演表 — 不是「 希望什么 」, 是「 准备怎么抓 」",
                 "和风险推演对称, 但方向相反")
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.6), fill=GREEN)
    headers = ["可能出现的机会", "概率", "价值", "捕捉行动", "准备工作", "负责人+时间"]
    cols_w = [Inches(3.3), Inches(0.9), Inches(0.9), Inches(2.5), Inches(2.5), Inches(2.23)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri in range(3):
        ry = Inches(2.5) + Inches(1.0) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.0), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        for ci in range(6):
            cx_next = cx + cols_w[ci]
            if ci >= 1 and ci <= 2:
                add_text(s, cx, ry + Inches(0.35), cols_w[ci], Inches(0.3),
                         "高/中/低", size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
            elif ci == 5:
                add_text(s, cx, ry + Inches(0.35), cols_w[ci], Inches(0.3),
                         "(谁 · 何时)", size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
            cx = cx_next
    y2 = Inches(5.6)
    add_round(s, Inches(0.5), y2, Inches(6.1), Inches(1.15), fill=ORANGE)
    add_text(s, Inches(0.7), y2 + Inches(0.1), Inches(5.7), Inches(0.4),
             "捕捉行动", size=13, color=WHITE, bold=True)
    add_text(s, Inches(0.7), y2 + Inches(0.5), Inches(5.7), Inches(0.65),
             "机会出现那一刻\n你做的第一件事",
             size=11, color=WHITE, line_h=1.3)
    add_round(s, Inches(6.7), y2, Inches(6.13), Inches(1.15), fill=GOLD)
    add_text(s, Inches(6.9), y2 + Inches(0.1), Inches(5.8), Inches(0.4),
             "准备工作", size=13, color=WHITE, bold=True)
    add_text(s, Inches(6.9), y2 + Inches(0.5), Inches(5.8), Inches(0.65),
             "现在就做, 让自己\n在机会出现时立刻行动",
             size=11, color=WHITE, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.3),
             "没有准备 = 机会出现时还没准备好 | 没有捕捉 = 不知道该做什么",
             size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


def m5_08_liming_opp(idx):
    s = slide_m5(idx, 77, "⚡ 完整示例: 李明的机会推演",
                 "同一行动: 全员推广体验式接待 3 步法")
    rows = [
        ("新品上市, 体验式销售优势凸显", "高", "高", "新品到货当天启动「 新品体验周 」, 老客户定向邀约", "①确认新品时间; ②设计体验脚本; ③通知老客户", GREEN),
        ("竞品服务质量下滑, 客户转入", "中", "高", "设计「 竞品客户专属接待流程 」", "整理差异化优势清单; 准备对比说明", ORANGE),
        ("成交率提升后, 区域参访", "中", "高", "邀请区域经理到访, 展示数据 + 现场", "系统记录每日数据; 整理成功案例", ORANGE),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=GREEN)
    headers = ["可能的机会", "概率", "价值", "捕捉行动", "准备工作"]
    cols_w = [Inches(3.5), Inches(1.0), Inches(1.0), Inches(3.5), Inches(3.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri, (opp, p, val, cap, prep, col) in enumerate(rows):
        ry = Inches(2.4) + Inches(1.05) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.05), fill=bg, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), ry, Inches(0.1), Inches(1.05), fill=col)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.2), ry + Inches(0.1), cols_w[0] - Inches(0.2), Inches(0.9),
                 opp, size=11, color=NAVY_DARK, bold=True, line_h=1.3)
        cx += cols_w[0]
        add_text(s, cx, ry + Inches(0.4), cols_w[1], Inches(0.3),
                 p, size=12, color=col, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[1]
        add_text(s, cx, ry + Inches(0.4), cols_w[2], Inches(0.3),
                 val, size=12, color=col, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[2]
        add_text(s, cx + Inches(0.1), ry + Inches(0.08), cols_w[3] - Inches(0.1), Inches(0.9),
                 cap, size=10, color=BLACK, line_h=1.3)
        cx += cols_w[3]
        add_text(s, cx + Inches(0.1), ry + Inches(0.08), cols_w[4] - Inches(0.1), Inches(0.9),
                 prep, size=10, color=BLACK, line_h=1.3)
    add_round(s, Inches(0.5), Inches(5.7), Inches(12.33), Inches(1.2), fill=GOLD)
    add_text(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.4),
             "💡 最让李明觉得「 以前没重视过 」的机会", size=13, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.65),
             "新品上市窗口 — 客户对新品本来就有好奇心, 体验式销售正好是让客户产生好奇和体验冲动",
             size=11, color=WHITE, line_h=1.3)


def m5_09_exercise2(idx):
    s = slide_m5(idx, 78, "✋ 练习二: 完成你的机会推演表",
                 "用同一行动 · 写 2-4 个机会 · 重点识别高概率+高价值")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.6), Inches(0.5),
             "机会写法的对照", size=14, color=ORANGE, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(0.85), Inches(5.6), Inches(0.4),
             "✗ 「 市场好转 」", size=12, color=RED)
    add_text(s, lx + Inches(0.3), ly + Inches(1.3), Inches(5.6), Inches(0.4),
             "✓ 「 如果本季度有新品上市, 体验式销售\n  在新品展示上的优势会特别突出 」",
             size=11, color=GREEN, bold=True, line_h=1.4)
    add_rect(s, lx + Inches(0.3), ly + Inches(2.6), Inches(5.6), Emu(20000), fill=ORANGE)
    add_text(s, lx + Inches(0.3), ly + Inches(2.75), Inches(5.6), Inches(0.4),
             "机会推演 4 步", size=12, color=ORANGE, bold=True)
    steps = [
        "找出 2-4 个可能的机会 (写具体场景)",
        "评估概率和价值",
        "对高价值机会写捕捉行动",
        "对所有机会都写准备工作",
    ]
    for i, t in enumerate(steps):
        add_text(s, lx + Inches(0.4), ly + Inches(3.2) + Inches(0.4) * i,
                 Inches(5.5), Inches(0.4),
                 f"  {i+1}.  {t}", size=11, color=BLACK, line_h=1.4)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "互评检验 (与同伴互换看)", size=14, color=ORANGE, bold=True)
    questions = [
        "风险描述够具体吗? (是泛描述, 还是可观察场景)",
        "预防和应急有区分吗? (没把应急当预防)",
        "机会的准备工作, 是现在就能做的吗? (具体到周)",
    ]
    for i, q in enumerate(questions):
        yy = ry + Inches(0.9) + Inches(1.2) * i
        add_oval(s, rx + Inches(0.3), yy + Inches(0.1), Inches(0.4), Inches(0.4), fill=ORANGE)
        add_text(s, rx + Inches(0.3), yy + Inches(0.15), Inches(0.4), Inches(0.3),
                 str(i+1), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, rx + Inches(0.8), yy, Inches(5.0), Inches(1.1),
                 q, size=12, color=WHITE, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.4),
             "预计用时: 12 分钟 | 写完后与同伴互换, 用 3 个问题互相检验",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def m5_10_connect(idx):
    s = slide_m5(idx, 79, "连接到收尾: 两天产出整合",
                 "你的 5 个工具产出, 现在整合成「 管理行动地图 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "管理行动地图 — 5 个工具的整合", size=18, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    rows_data = [
        ("螺旋深挖 4 问", "打算访谈 ____ · 案例 ____", "成功原则: ____", "访谈时间: ____"),
        ("花刺投票", "用在议题: ____", "花票最高: ____ · 刺票最高: ____", "会议时间: ____"),
        ("问题树 + 魔力提问", "分析的难题: ____", "切入点: ____ · 新方向: ____", "行动时间: ____"),
        ("高效脑暴双矩阵", "主持的议题: ____", "优先行动 1-2: ____", "开会时间: ____"),
        ("推演双表格", "推演的行动: ____", "最高风险 + 预防 / 最高机会 + 准备", "启动时间: ____"),
    ]
    add_rect(s, Inches(0.5), Inches(2.6), Inches(12.33), Inches(0.5), fill=NAVY_DARK)
    headers = ["工具", "用在什么场景", "产出 / 下一步", "时间"]
    cols_w = [Inches(2.5), Inches(3.5), Inches(4.0), Inches(2.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(2.7), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
        cx += cols_w[i]
    for ri, (tool, scn, out, when) in enumerate(rows_data):
        ry = Inches(3.2) + Inches(0.55) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.55), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[0] - Inches(0.1), Inches(0.4),
                 tool, size=11, color=ORANGE, bold=True)
        cx += cols_w[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[1] - Inches(0.1), Inches(0.4),
                 scn, size=10, color=BLACK)
        cx += cols_w[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[2] - Inches(0.1), Inches(0.4),
                 out, size=10, color=BLACK)
        cx += cols_w[2]
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[3] - Inches(0.1), Inches(0.4),
                 when, size=10, color=GRAY, italic=True)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "→ 拍照存在手机里, 每两周回看一次",
             size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.4),
             "→ 哪一步走了? 哪一步没开始? 为什么?",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m5_11_recap(idx):
    s = slide_m5(idx, 80, "✅ 第五部分知识框架",
                 "M5 = 风险推演 + 机会推演 = 完整的前瞻视角")
    items = [
        ("核心问题", "为什么一线管理者\n总是救火不防火", RED),
        ("风险推演", "预防 (降低概率)\n+ 应急 (降低影响)", ORANGE),
        ("机会推演", "捕捉 (立刻行动)\n+ 准备 (现在就做)", GREEN),
        ("并行原则", "必须两张表一起做\n只做一张 = 视角缺失", NAVY),
    ]
    card_w = Inches(2.95); card_h = Inches(2.5)
    gap = Inches(0.15)
    total_w = card_w * 4 + gap * 3
    start_x = Inches(0.5)
    y = Inches(2.0)
    for i, (head, body, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=color)
        add_rect(s, x, y, card_w, Inches(0.6), fill=color)
        add_text(s, x, y + Inches(0.13), card_w, Inches(0.4),
                 head, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_text(s, x + Inches(0.2), y + Inches(0.8) + Inches(0.5) * j,
                     card_w - Inches(0.4), Inches(0.5),
                     line, size=12, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.4)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.33), Inches(0.4),
             "完成第五部分后你拥有:", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    outs = [
        "一份风险预案 (含预防 + 应急)",
        "一份机会行动计划 (含捕捉 + 准备)",
        "一份管理行动地图 (5 工具整合)",
    ]
    for i, t in enumerate(outs):
        x = Inches(0.5) + Inches(4.2) * i
        add_round(s, x, Inches(5.3), Inches(4.0), Inches(0.7), fill=LIGHT, line=GREEN)
        add_text(s, x, Inches(5.35), Inches(4.0), Inches(0.6),
                 t, size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "→ 这就是你要带回去的东西",
             size=13, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


def m5_12_module_close(idx):
    s = slide_m5(idx, 81, "模块五收尾 · 转向两天整合",
                 "5 个工具的连接, 就是管理的完整链路")
    add_text(s, Inches(0.5), Inches(1.7), Inches(12.33), Inches(1.0),
             "M5", size=120, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.7),
             "前瞻思考", size=32, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             "Think Forward", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    quote_block(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                "30 分钟的提前推演, 换 15 小时的救火时间。",
                author="—— 第五部分核心信念")
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "下一段: 收尾 — 两天产出整合 + 个人行动承诺 + 30 天节奏",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


print("模块五定义完成")


# ===================================================================
# 第 9 段: 收尾 (6 页)
# ===================================================================
def slide_close(idx, n, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section="收尾 · Two-Day Integration")
    title_block(s, title, subtitle, accent=NAVY_DARK)
    return s


def close_01_integration(idx):
    s = slide_close(idx, 82, "两天产出的整合 — 5 件工具一张图",
                    "从「 我有一个问题 」到「 带着行动地图回去 」")
    # 大图: 5 工具串联
    flow = [
        ("M1 螺旋\n深挖 4 问", "提炼\n成功原则", ORANGE),
        ("M2 花刺\n投票", "确定\n优先机会", GOLD),
        ("M3 问题树\n+ 魔力提问", "找到\n新方向", RED),
        ("M4 双矩阵\n脑暴", "产出\n优先行动", NAVY),
        ("M5 推演\n双表格", "前瞻\n准备就绪", GREEN),
    ]
    card_w = Inches(2.3); card_h = Inches(2.3)
    gap = Inches(0.2)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (head, foot, col) in enumerate(flow):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(1.6), fill=col)
        add_text(s, x, y + Inches(0.2), card_w, Inches(1.2),
                 head, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_h=1.3)
        add_text(s, x, y + Inches(1.7), card_w, Inches(0.6),
                 foot, size=11, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, line_h=1.3)
        if i < 4:
            arr_x = x + card_w + gap * 0.1
            add_text(s, arr_x, y + Inches(0.9), Inches(0.4), Inches(0.4),
                     "▶", size=20, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    # 产出
    add_round(s, Inches(0.5), Inches(4.7), Inches(12.33), Inches(1.5), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.4),
             "产出 — 你的「 管理行动地图 」", size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.85),
             "M1 提炼的成功原则 → M2 确定的优先机会 → M3 找到的新方向\nM4 共创的优先行动 → M5 准备好的风险预案 + 机会行动",
             size=12, color=WHITE, line_h=1.4, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "这就是你带回去的东西 — 不是 5 个孤立工具, 而是一条完整链路",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def close_02_commit(idx):
    s = slide_close(idx, 83, "个人行动承诺 — 下周第一件事",
                    "不要等「 准备好了再做 」, 而是「 做了才有反馈 」")
    # 左: 写下承诺
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(7.5), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(7.0), Inches(0.5),
             "📝 我的下周第一件事", size=15, color=ORANGE, bold=True)
    add_rect(s, lx + Inches(0.3), ly + Inches(0.85), Inches(7.0), Inches(0.8), fill=WHITE, line=ORANGE)
    add_text(s, lx + Inches(0.3), ly + Inches(0.85), Inches(7.0), Inches(0.8),
             "(写一件具体的事 · 选一件工具 · 真实场景)",
             size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, lx + Inches(0.3), ly + Inches(1.85), Inches(7.0), Inches(0.5),
             "📅 时间节点", size=14, color=NAVY_DARK, bold=True)
    add_rect(s, lx + Inches(0.3), ly + Inches(2.3), Inches(7.0), Inches(0.6), fill=WHITE, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(2.3), Inches(7.0), Inches(0.6),
             "(日期: ____)",
             size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, lx + Inches(0.3), ly + Inches(3.05), Inches(7.0), Inches(0.5),
             "🤝 谁来见证", size=14, color=NAVY_DARK, bold=True)
    add_rect(s, lx + Inches(0.3), ly + Inches(3.5), Inches(7.0), Inches(0.6), fill=WHITE, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(3.5), Inches(7.0), Inches(0.6),
             "(同事 · 老板 · 团队 · 给自己)",
             size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 右: 提示
    rx = Inches(8.2); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(4.65), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(4.0), Inches(0.5),
             "好承诺的 3 个特征", size=14, color=ORANGE, bold=True)
    items = [
        ("具体", "不是「 用一下 M3 」, 而是「 用 M3 分析老客户流失这个难题 」"),
        ("可衡量", "「 一周内完成 3 个老员工的成功访谈 」"),
        ("有反馈", "你做完了能立刻知道效果, 不是「 准备一下 」"),
    ]
    for i, (head, sub) in enumerate(items):
        yy = ry + Inches(0.9) + Inches(1.15) * i
        add_oval(s, rx + Inches(0.3), yy, Inches(0.5), Inches(0.5), fill=ORANGE)
        add_text(s, rx + Inches(0.3), yy + Inches(0.1), Inches(0.5), Inches(0.4),
                 str(i+1), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, rx + Inches(0.95), yy, Inches(3.5), Inches(0.4),
                 head, size=13, color=WHITE, bold=True)
        add_text(s, rx + Inches(0.95), yy + Inches(0.4), Inches(3.5), Inches(0.7),
                 sub, size=10, color=GRAY_LIGHT, line_h=1.3)
    # 底部
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "→ 写在管理行动地图上 · 拍照 · 存在手机里 · 给自己一个提醒",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def close_03_30day(idx):
    s = slide_close(idx, 84, "30 天节奏 — 怎么用 5 件工具",
                    "不要一次用 5 件, 一周一件就够")
    weeks = [
        ("第 1 周", "M1 螺旋深挖 4 问", "访谈一位你团队里的「 高手 」\n把成功经验变成可学的方法", ORANGE),
        ("第 2 周", "M2 花刺投票", "在一个争议议题上用一次\n感受真实共识 vs 假性共识", GOLD),
        ("第 3 周", "M3 问题树 + 魔力提问", "选一个最近卡住你的难题\n用问题树 + 5 问找新方向", RED),
        ("第 4 周", "M4 + M5 串联", "用双矩阵共创方案\n用推演双表格做好前瞻", GREEN),
    ]
    card_w = Inches(3.0); card_h = Inches(3.6)
    gap = Inches(0.1)
    total_w = card_w * 4 + gap * 3
    start_x = (SW - total_w) / 2
    y = Inches(1.95)
    for i, (wk, tool, sub, col) in enumerate(weeks):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.8), fill=col)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.5),
                 wk, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.95), card_w - Inches(0.4), Inches(0.5),
                 tool, size=13, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(1.5) + Inches(0.5) * j,
                     card_w - Inches(0.5), Inches(0.5),
                     line, size=11, color=BLACK, align=PP_ALIGN.CENTER, line_h=1.3)
    # 底部
    add_round(s, Inches(0.5), Inches(5.85), Inches(12.33), Inches(1.0), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.4),
             "30 天后你回头看:", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.45),
             "5 件工具你都有过真实场景的体验 — 这就是开始熟练的标志",
             size=12, color=WHITE, align=PP_ALIGN.CENTER)


def close_04_qa(idx):
    s = slide_close(idx, 85, "Q & A",
                    "你今天最想搞清楚的一件事是什么?")
    # 大字
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.33), Inches(2.5),
             "Q & A", size=180, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # 引子
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.5),
             "把你的具体场景说出来 — 我们一起看 5 件工具能帮上什么",
             size=14, color=GRAY, align=PP_ALIGN.CENTER)
    # 装饰
    add_rect(s, Inches(0.5), Inches(6.3), Inches(2.0), Inches(0.05), fill=ORANGE)
    add_rect(s, Inches(5.5), Inches(6.3), Inches(2.0), Inches(0.05), fill=GREEN)
    add_rect(s, Inches(10.5), Inches(6.3), Inches(2.0), Inches(0.05), fill=RED)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "「 你今天最想搞清楚的一件事 」",
             size=12, color=NAVY_DARK, bold=True, italic=True, align=PP_ALIGN.CENTER)


def close_05_promise(idx):
    s = slide_close(idx, 86, "一个简单的承诺",
                    "下周, 选其中一件工具, 在你团队里用一次")
    # 大引言
    add_round(s, Inches(0.5), Inches(2.0), Inches(12.33), Inches(2.6), fill=LIGHT, line=NAVY)
    add_text(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(2.2),
             "「 选其中一件工具\n  在你团队里用一次\n  不需要完美 · 不需要五件一起用\n  先把一件工具在真实场景里跑一遍 」",
             size=18, color=NAVY_DARK, bold=True, line_h=1.6, align=PP_ALIGN.CENTER)
    # 三个要点
    items = [
        ("不追求完美", GRAY),
        ("不一次全用", ORANGE),
        ("先跑一遍", GREEN),
    ]
    for i, (t, c) in enumerate(items):
        x = Inches(0.5) + Inches(4.2) * i
        add_round(s, x, Inches(4.9), Inches(4.0), Inches(0.8), fill=LIGHT, line=NAVY)
        add_text(s, x, Inches(5.0), Inches(4.0), Inches(0.6),
                 "✓ " + t, size=15, color=c, bold=True, align=PP_ALIGN.CENTER)
    # 14 天后
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.4),
             "两周后, 回看你的管理行动地图",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "哪一步走了 · 哪一步没开始 · 为什么",
             size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "「 那个「 为什么没有开始 」, 往往比工具本身更值得你认真面对 」",
             size=11, color=RED, bold=True, italic=True, align=PP_ALIGN.CENTER)


def close_06_thanks(idx):
    s = new_slide()
    page_chrome(s, idx, section="收尾 · 感谢")
    add_rect(s, 0, Inches(2.0), SW, Inches(3.5), fill=NAVY_DARK)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.2),
             "谢谢", size=80, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5),
             "Thank You", size=18, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.5),
             "一线管理者的现代五项 · Two Days With You", size=12, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5),
             "用起来 · 不然白学",
             size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


print("收尾定义完成")


# ===================================================================
# 第 10 段: 附录 (10 页)
# ===================================================================
def slide_app(idx, n, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section="附录 · Appendix")
    title_block(s, title, subtitle, accent=GRAY)
    return s


def app_01_m1(idx):
    s = slide_app(idx, 88, "M1 螺旋深挖 4 问 速查卡", "")
    items = [
        ("Q1 行为", "具体做了什么不同的动作?\n必须是可观察的动作, 不是「 认真用心 」", "听动词"),
        ("Q2 动机", "为什么这样做? 当时怎么想的?\n跳过这问, Q3 必然空洞", "听原因"),
        ("Q3 原则", "提炼成一个核心原则?\n必须回答「 为什么有效 」, 不是行为的另一种说法", "听原则"),
        ("Q4 路径", "教给新人, 最关键的 2-3 步?\n步骤 + 原则支撑, 缺一不可", "听步骤"),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=ORANGE)
    headers = ["问", "挖的层 · 标准", "听什么"]
    cols_w = [Inches(2.0), Inches(7.5), Inches(2.83)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
        cx += cols_w[i]
    for ri, (no, content, listen) in enumerate(items):
        ry = Inches(2.4) + Inches(0.9) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.9), fill=bg, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), ry, Inches(0.1), Inches(0.9), fill=ORANGE)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.2), ry + Inches(0.3), cols_w[0] - Inches(0.2), Inches(0.4),
                 no, size=14, color=ORANGE, bold=True)
        cx += cols_w[0]
        for j, line in enumerate(content.split("\n")):
            add_text(s, cx + Inches(0.1), ry + Inches(0.1) + Inches(0.4) * j,
                     cols_w[1] - Inches(0.1), Inches(0.4),
                     line, size=11, color=BLACK, line_h=1.3)
        cx += cols_w[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.3), cols_w[2] - Inches(0.1), Inches(0.4),
                 listen, size=12, color=NAVY_DARK, bold=True)
    # 陷阱
    add_round(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.85), fill=RED)
    add_text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.4),
             "常见陷阱", size=12, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.4),
             "接受 Q3 第一个答案(通常还是行为) · 把形容词当行为 · 跳过 Q2 直接到 Q3",
             size=10, color=WHITE)


def app_02_m2(idx):
    s = slide_app(idx, 89, "M2 花刺投票 速查卡", "")
    items = [
        ("01 列候选", "便利贴 · 数量越多越好 · 不评判", "把团队里能想到的「 抓手 」全部列出来"),
        ("02 静默投票", "每人 2 票 = 1 花票 (机会) + 1 刺票 (障碍)", "站起来贴票 · 物理动作维持能量"),
        ("03 统计结果", "花票最多 = 优先机会 · 刺票最多 = 关键障碍", "让数据说话 · 不辩论"),
        ("04 讨论确认", "只讨论票数最高 · 明确负责人 + 时间", "前 2 步, 不反复"),
    ]
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=GOLD)
    headers = ["步骤", "关键动作", "提示"]
    cols_w = [Inches(2.0), Inches(5.5), Inches(4.83)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
        cx += cols_w[i]
    for ri, (no, content, hint) in enumerate(items):
        ry = Inches(2.4) + Inches(0.9) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.9), fill=bg, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), ry, Inches(0.1), Inches(0.9), fill=GOLD)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.2), ry + Inches(0.3), cols_w[0] - Inches(0.2), Inches(0.4),
                 no, size=14, color=GOLD, bold=True)
        cx += cols_w[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.3), cols_w[1] - Inches(0.1), Inches(0.4),
                 content, size=11, color=BLACK)
        cx += cols_w[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.3), cols_w[2] - Inches(0.1), Inches(0.4),
                 hint, size=11, color=NAVY_DARK, bold=True)
    # 关键设计
    add_round(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.85), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.4),
             "关键设计", size=12, color=GOLD, bold=True)
    add_text(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.4),
             "静默投票 (避免社会压力) + 站起来贴票 (物理动作维持能量) + 票数逻辑 (让数据说话)",
             size=10, color=WHITE)


def app_03_m3(idx):
    s = slide_app(idx, 90, "M3 问题树 + 魔力破解提问 速查卡", "")
    # 左: 问题树
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.8), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "问题树 4 步", size=14, color=RED, bold=True)
    steps = [
        ("1", "写下症状", "你观察到的「 不对劲 」"),
        ("2", "第一层拆解", "可能是哪几个方面的问题 (3-4 个方向)"),
        ("3", "第二层拆解", "每个方面的具体原因"),
        ("4", "找切入点", "影响大 × 我能影响 (两者同时满足)"),
    ]
    for i, (no, name, sub) in enumerate(steps):
        yy = ly + Inches(0.85) + Inches(0.9) * i
        add_oval(s, lx + Inches(0.3), yy, Inches(0.5), Inches(0.5), fill=RED)
        add_text(s, lx + Inches(0.3), yy + Inches(0.1), Inches(0.5), Inches(0.4),
                 no, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, lx + Inches(0.95), yy + Inches(0.02), Inches(4.7), Inches(0.4),
                 name, size=13, color=NAVY_DARK, bold=True)
        add_text(s, lx + Inches(0.95), yy + Inches(0.4), Inches(4.7), Inches(0.5),
                 sub, size=10, color=BLACK, line_h=1.3)
    # 右: 魔力提问
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.8), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "魔力破解提问 5 问", size=14, color=ORANGE, bold=True)
    qs = [
        ("Q1 破假设", "我们一直认为不可能/必须的 —— 真的吗?"),
        ("Q2 借他山", "其他行业/场景怎么解决类似问题?"),
        ("Q3 极端情景", "资源充足/只有 1/10 资源时怎么做?"),
        ("Q4 换位思考", "从客户/团队/对手/上级角度看?"),
        ("Q5 倒推法", "想象一年后问题已解决, 你是怎么做到的?"),
    ]
    for i, (no, q) in enumerate(qs):
        yy = ry + Inches(0.85) + Inches(0.7) * i
        add_text(s, rx + Inches(0.3), yy, Inches(1.6), Inches(0.4),
                 no, size=12, color=ORANGE, bold=True)
        add_text(s, rx + Inches(1.9), yy, Inches(4.0), Inches(0.65),
                 q, size=10, color=WHITE, line_h=1.3)
    # 底部
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "注意: 问题树先拆清楚, 再用魔力提问 — 否则魔力提问变成对模糊大问题的泛泛思考",
             size=10, color=RED, bold=True, align=PP_ALIGN.CENTER)


def app_04_m4(idx):
    s = slide_app(idx, 91, "M4 高效脑暴双矩阵 速查卡", "")
    # 发散
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(2.4), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.15), Inches(5.5), Inches(0.4),
             "发散矩阵 2×2", size=12, color=GOLD, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(0.55), Inches(5.5), Inches(0.4),
             "时间 × 资源 — 4 格都要填", size=10, color=BLACK)
    qs = ["Q1 快赢 (短期+内部)", "Q2 协调推进 (短期+外部)", "Q3 能力建设 (中长期+内部)", "Q4 战略布局 (中长期+外部)"]
    for i, t in enumerate(qs):
        yy = ly + Inches(0.95) + Inches(0.32) * i
        add_text(s, lx + Inches(0.3), yy, Inches(5.5), Inches(0.3),
                 "• " + t, size=10, color=BLACK)
    # 收敛
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(2.4), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, rx + Inches(0.3), ry + Inches(0.15), Inches(5.5), Inches(0.4),
             "收敛矩阵 2×2", size=12, color=GREEN, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(0.55), Inches(5.5), Inches(0.4),
             "影响力 × 落地难度", size=10, color=BLACK)
    qs2 = ["⭐ 优先行动 (高影响 + 易)", "战略项目 (高影响 + 难)", "随机应变 (低影响 + 易)", "暂不考虑 (低影响 + 难)"]
    for i, t in enumerate(qs2):
        yy = ry + Inches(0.95) + Inches(0.32) * i
        add_text(s, rx + Inches(0.3), yy, Inches(5.5), Inches(0.3),
                 "• " + t, size=10, color=BLACK)
    # 主持原则
    add_round(s, Inches(0.5), Inches(4.4), Inches(12.33), Inches(1.0), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.4),
             "主持原则", size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.4),
             "发散阶段: 不评判 · 逐格引导 · 主持人先不说自己的想法",
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.25), Inches(11.9), Inches(0.4),
             "收敛阶段: 不主导 · 让矩阵逻辑引导聚焦, 不是「 主管说了算 」",
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # 3 步
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.33), Inches(0.4),
             "收敛 3 步", size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.4),
             "① 独立评估 (每人标位置) → ② 对比分歧 (只讨论分歧大的 2-3 条) → ③ 确认优先 (⭐ 区 2-3 条)",
             size=11, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.3),
             "意外的想法往往最有价值 — 自由讨论总倾向 Q1, 矩阵强制 4 格都填",
             size=10, color=GOLD, bold=True, align=PP_ALIGN.CENTER)


def app_05_m5(idx):
    s = slide_app(idx, 92, "M5 推演双表格 速查卡", "")
    # 上: 风险
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=RED)
    add_text(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.4),
             "风险推演表", size=12, color=WHITE, bold=True)
    headers_r = ["可能的风险 (具体场景)", "概率", "影响", "预防 (事前降概率)", "应急 (事后降影响)", "负责人+时间"]
    cols_w = [Inches(3.3), Inches(0.9), Inches(0.9), Inches(2.5), Inches(2.5), Inches(2.23)]
    cx = Inches(0.5)
    for i, h in enumerate(headers_r):
        add_text(s, cx, Inches(2.4), cols_w[i], Inches(0.4),
                 h, size=9, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri in range(2):
        ry = Inches(2.85) + Inches(0.6) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.6), fill=bg, line=GRAY_LIGHT)
    # 下: 机会
    add_rect(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5), fill=GREEN)
    add_text(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.4),
             "机会推演表", size=12, color=WHITE, bold=True)
    headers_o = ["可能的机会 (具体场景)", "概率", "价值", "捕捉 (机会出现时)", "准备 (现在就做)", "负责人+时间"]
    cx = Inches(0.5)
    for i, h in enumerate(headers_o):
        add_text(s, cx, Inches(4.65), cols_w[i], Inches(0.4),
                 h, size=9, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri in range(2):
        ry = Inches(5.1) + Inches(0.6) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.6), fill=bg, line=GRAY_LIGHT)
    # 优先级
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "优先级逻辑", size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.3),
             "高 × 高 = 必须双准备 · 中 × 高 = 至少有应急 · 其他 = 酌情",
             size=10, color=RED, bold=True, align=PP_ALIGN.CENTER)


def app_06_connection(idx):
    s = slide_app(idx, 93, "5 工具的连接逻辑 — 整门课一张图", "")
    # 流程
    flow = [
        ("M1 螺旋深挖 4 问", "→ 成功原则", ORANGE),
        ("M2 花刺投票", "→ 优先机会 + 关键障碍", GOLD),
        ("M3 问题树 + 魔力提问", "→ 新方向", RED),
        ("M4 双矩阵", "→ 优先行动", NAVY),
        ("M5 推演双表格", "→ 准备就绪", GREEN),
    ]
    card_w = Inches(2.3); card_h = Inches(2.0)
    gap = Inches(0.2)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (head, foot, col) in enumerate(flow):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(1.4), fill=col)
        add_text(s, x, y + Inches(0.3), card_w, Inches(0.8),
                 head, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_h=1.2)
        add_text(s, x, y + Inches(1.5), card_w, Inches(0.4),
                 foot, size=10, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        if i < 4:
            arr_x = x + card_w + gap * 0.1
            add_text(s, arr_x, y + Inches(0.8), Inches(0.4), Inches(0.4),
                     "▶", size=20, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    # 产出
    add_round(s, Inches(0.5), Inches(4.5), Inches(12.33), Inches(1.6), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(4.65), Inches(11.9), Inches(0.4),
             "学员带走的产出", size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.0),
             "M1 提炼的成功原则\nM2 确定的优先机会\nM3 找到的新方向\nM4 共创的优先行动\nM5 准备的风险预案 + 机会行动",
             size=11, color=WHITE, line_h=1.4, align=PP_ALIGN.CENTER)
    # 底部
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "这 5 件不是 5 个独立工具, 是一条完整的链路",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


def app_07_mgmt_map(idx):
    s = slide_app(idx, 94, "管理行动地图 (学员带回模板)", "")
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=NAVY_DARK)
    headers = ["工具", "用在什么场景", "产出 / 下一步", "计划日期"]
    cols_w = [Inches(2.5), Inches(3.5), Inches(4.0), Inches(2.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4),
                 h, size=12, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
        cx += cols_w[i]
    tools = [
        ("M1 螺旋深挖 4 问", "访谈: ____ · 案例: ____", "成功原则: ____"),
        ("M2 花刺投票", "用在: ____", "花票: ____ · 刺票: ____"),
        ("M3 问题树 + 魔力", "分析: ____", "切入: ____ · 新方向: ____"),
        ("M4 双矩阵", "主持: ____", "优先行动 1: ____\n优先行动 2: ____"),
        ("M5 推演双表格", "推演: ____", "最高风险: ____\n最高机会: ____"),
    ]
    for ri, (tool, scn, out) in enumerate(tools):
        ry = Inches(2.4) + Inches(0.55) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.55), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[0] - Inches(0.1), Inches(0.4),
                 tool, size=11, color=ORANGE, bold=True)
        cx += cols_w[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[1] - Inches(0.1), Inches(0.4),
                 scn, size=10, color=BLACK)
        cx += cols_w[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[2] - Inches(0.1), Inches(0.4),
                 out, size=10, color=BLACK, line_h=1.3)
        cx += cols_w[2]
        add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[3] - Inches(0.1), Inches(0.4),
                 "日期: ____", size=10, color=GRAY, italic=True)
    # 下周第一件事
    add_round(s, Inches(0.5), Inches(5.4), Inches(12.33), Inches(1.7), fill=GOLD)
    add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.4),
             "📌 下周第一件事 (具体, 一件)", size=14, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.5),
             "____________________________________________________",
             size=12, color=WHITE)
    add_text(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
             "一个月后, 我希望团队有什么不同: _________________________________________________",
             size=11, color=WHITE)


def app_08_resources(idx):
    s = slide_app(idx, 95, "课程资源 · 课前 / 课后", "")
    items = [
        ("课前", "1. 一线管理者能力自评表 (9 题)\n2. 课程预习: 体验式接待 3 步法\n3. 思考 1 件本可提前发现的意外", NAVY),
        ("课中", "1. 学员手册 (含 5 工具完整模板)\n2. 工具表单 F1-F10 (10 个)\n3. 练习题库 G1-G7 (7 个)", ORANGE),
        ("课后", "1. 5 工具速查卡 (本附录)\n2. 管理行动地图 (带回模板)\n3. 30 天节奏表 (一周一件工具)", GREEN),
    ]
    card_w = Inches(4.0); card_h = Inches(3.5)
    gap = Inches(0.15)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.0)
    for i, (head, content, col) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.7), fill=col)
        add_text(s, x, y + Inches(0.15), card_w, Inches(0.4),
                 head, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(content.split("\n")):
            add_text(s, x + Inches(0.25), y + Inches(0.95) + Inches(0.5) * j,
                     card_w - Inches(0.5), Inches(0.5),
                     line, size=11, color=BLACK, line_h=1.3)
    # 评估
    add_round(s, Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.2), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.4),
             "评估 / 反馈", size=13, color=ORANGE, bold=True)
    add_text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.65),
             "课后 30 天评估 (1 次) · 课后 90 天追踪 (1 次) · 团队版工具使用情况反馈表",
             size=11, color=WHITE, line_h=1.3)


def app_09_faq(idx):
    s = slide_app(idx, 96, "常见问题 FAQ", "")
    faqs = [
        ("Q: 5 件工具必须一次全用吗?",
         "A: 不必。一周用一件就够。先把一件工具在真实场景里跑一遍。"),
        ("Q: M3 魔力提问是不是越多越好?",
         "A: 不是。问题树先拆清楚, 再对切入点逐一「 试 」5 问, 选最让你「 没想过 」的那一问深想。"),
        ("Q: 团队只有 3-4 人, M4 矩阵还管用吗?",
         "A: 管用。小团队反而更能聚焦 4 格, 不会冷场。"),
        ("Q: 推演表里的概率/影响, 需要很精确吗?",
         "A: 不需要精确, 只需要你认真想过。主客观判断都接受。"),
        ("Q: 5 件工具都用了, 但团队没变化, 为什么?",
         "A: 工具 ≠ 结果。看是不是「 主持原则 」没做到 — 是不是 M4 收敛时主管意见主导了。"),
    ]
    y = Inches(1.85)
    for q, a in faqs:
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(0.85), fill=LIGHT, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.85), fill=ORANGE)
        add_text(s, Inches(0.8), y + Inches(0.1), Inches(11.5), Inches(0.4),
                 q, size=12, color=ORANGE, bold=True)
        add_text(s, Inches(0.8), y + Inches(0.45), Inches(11.5), Inches(0.4),
                 a, size=11, color=BLACK, line_h=1.3)
        y += Inches(0.95)


def app_10_glossary(idx):
    s = slide_app(idx, 97, "术语表 / 致谢", "")
    # 左: 术语
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.7), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "术语表", size=14, color=NAVY_DARK, bold=True)
    terms = [
        ("螺旋深挖", "用 Q1→Q2→Q3→Q4 逐层挖深, 任何一问没问清回到上一问"),
        ("花刺投票", "用花票 (机会) + 刺票 (障碍) 形成真实共识"),
        ("问题树", "从症状到原因的结构化拆解"),
        ("魔力提问", "5 个破除思维惯性的提问视角"),
        ("发散矩阵", "时间 × 资源的 2×2 矩阵, 让所有类型想法都出来"),
        ("收敛矩阵", "影响力 × 落地难度的 2×2 矩阵, 聚焦优先行动"),
        ("推演双表格", "风险推演 + 机会推演, 两张表必须同时做"),
        ("管理行动地图", "5 工具产出的整合, 学员带回的真实产出"),
    ]
    for i, (t, sub) in enumerate(terms):
        yy = ly + Inches(0.7) + Inches(0.5) * i
        add_text(s, lx + Inches(0.3), yy, Inches(2.0), Inches(0.4),
                 t, size=11, color=ORANGE, bold=True)
        add_text(s, lx + Inches(2.3), yy, Inches(3.4), Inches(0.5),
                 sub, size=9, color=BLACK, line_h=1.3)
    # 右: 致谢
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.7), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "致谢", size=14, color=ORANGE, bold=True)
    credits = [
        ("课程设计", "竞越顾问公司 · 一线管理者发展项目组"),
        ("李明 / 小王", "基于真实连锁零售场景虚构人物 (12 人门店)"),
        ("参考框架", "DLP · Korn Ferry · CCL 国际版权课方法论"),
        ("视觉风格", "国际授权课标准 · 杂志风 · 信息图"),
        ("教学文档", "5 个模块完整内容 (35,000+ 字)"),
    ]
    for i, (head, sub) in enumerate(credits):
        yy = ry + Inches(0.85) + Inches(0.65) * i
        add_text(s, rx + Inches(0.3), yy, Inches(1.8), Inches(0.4),
                 head, size=12, color=ORANGE, bold=True)
        add_text(s, rx + Inches(2.1), yy, Inches(3.7), Inches(0.65),
                 sub, size=10, color=WHITE, line_h=1.3)
    # 底部
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.3),
             "© 2026 一线管理者的现代五项 · 仅供内部教学使用",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)


print("附录定义完成")


# ===================================================================
# 主入口: 顺序调用所有 slide 函数
# ===================================================================
def build():
    idx = 1
    # 1. 封面 (1)
    slide_cover(idx); idx += 1
    # 2. 课程导览 (4)
    slide_01_course_map(idx); idx += 1
    slide_02_three_pains(idx); idx += 1
    slide_03_yield(idx); idx += 1
    slide_04_learning_roadmap(idx); idx += 1
    # 3. 开场导入 (5)
    slide_05_self_check(idx); idx += 1
    slide_06_five_tools_big(idx); idx += 1
    slide_07_li_ming_intro(idx); idx += 1
    slide_08_promise(idx); idx += 1
    slide_09_ground_rules(idx); idx += 1
    # 4. 模块一 复制成功 (28)
    m1_section_divider(idx); idx += 1
    for fn in [m1_01, m1_02, m1_03, m1_04, m1_05, m1_06, m1_07, m1_08, m1_09,
               m1_10, m1_11, m1_12, m1_13, m1_14, m1_15, m1_16, m1_17, m1_18,
               m1_19, m1_20, m1_21, m1_22, m1_23, m1_24, m1_25, m1_26, m1_27]:
        try:
            fn(idx); idx += 1
        except NameError:
            pass
    # 5. 模块二 共谋抓手 (14 实际 + 8 扩展 = 22)
    m2_section_divider(idx); idx += 1
    for name in ["m2_01_opening", "m2_02_two_reasons", "m2_03_diagnosis", "m2_04_4steps",
                 "m2_05_silent_vote", "m2_06_full_case", "m2_07_what_democracy_is_not",
                 "m2_08_exercise1", "m2_09_exercise2", "m2_10_group_sim", "m2_11_three_pitfalls",
                 "m2_12_recap", "m2_13_connect", "m2_14_section_close",
                 "m2_15", "m2_16", "m2_17", "m2_18", "m2_19", "m2_20", "m2_21", "m2_22"]:
        fn = globals().get(name)
        if fn:
            fn(idx); idx += 1
    # 6. 模块三 应对难题 (13 实际 + 19 扩展 = 32)
    m3_section_divider(idx); idx += 1
    for name in ["m3_01_opening", "m3_02_skip_diagnosis", "m3_03_tree_4steps",
                 "m3_04_full_case", "m3_05_exercise1", "m3_06_exercise2", "m3_07_magic_questions",
                 "m3_08_magic_example", "m3_09_exercise3", "m3_10_share", "m3_11_recap",
                 "m3_12_connect", "m3_13_section_close",
                 "m3_14", "m3_15", "m3_16", "m3_17", "m3_18", "m3_19", "m3_20", "m3_21", "m3_22",
                 "m3_23", "m3_24", "m3_25", "m3_26", "m3_27", "m3_28", "m3_29", "m3_30", "m3_31", "m3_32"]:
        fn = globals().get(name)
        if fn:
            fn(idx); idx += 1
    # 7. 模块四 引领共创 (11 实际 + 13 扩展 = 24)
    m4_section_divider(idx); idx += 1
    for name in ["m4_01_opening", "m4_02_failure_modes", "m4_03_diverge_matrix",
                 "m4_04_full_case", "m4_05_exercise1", "m4_06_converge_matrix",
                 "m4_07_converge_case", "m4_08_exercise2", "m4_09_recap", "m4_10_connect",
                 "m4_11_section_close",
                 "m4_12", "m4_13", "m4_14", "m4_15", "m4_16", "m4_17", "m4_18", "m4_19",
                 "m4_20", "m4_21", "m4_22", "m4_23", "m4_24"]:
        fn = globals().get(name)
        if fn:
            fn(idx); idx += 1
    # 8. 模块五 前瞻思考 (12 实际 + 6 扩展 = 18)
    m5_section_divider(idx); idx += 1
    for name in ["m5_01_opening", "m5_02_obstacles", "m5_03_dual_table_logic",
                 "m5_04_risk_table", "m5_05_liming_risk", "m5_06_exercise1",
                 "m5_07_opp_table", "m5_08_liming_opp", "m5_09_exercise2",
                 "m5_10_connect", "m5_11_recap", "m5_12_module_close",
                 "m5_13", "m5_14", "m5_15", "m5_16", "m5_17", "m5_18"]:
        fn = globals().get(name)
        if fn:
            fn(idx); idx += 1
    # 9. 收尾 (6)
    for name in ["close_01_integration", "close_02_commit", "close_03_30day",
                 "close_04_qa", "close_05_promise", "close_06_thanks"]:
        fn = globals().get(name)
        if fn:
            fn(idx); idx += 1
    # 10. 附录 (10)
    for name in ["app_01_m1", "app_02_m2", "app_03_m3", "app_04_m4", "app_05_m5",
                 "app_06_connection", "app_07_mgmt_map", "app_08_resources",
                 "app_09_faq", "app_10_glossary"]:
        fn = globals().get(name)
        if fn:
            fn(idx); idx += 1
    return idx - 1


# ===================================================================
# 补全函数: M2(15-22) + M3(14-32) + M4(12-24) + M5(13-18) = 46 页
# ===================================================================
def _generic_slide(idx, section, accent, title, subtitle=""):
    s = new_slide()
    page_chrome(s, idx, section=section)
    title_block(s, title, subtitle, accent=accent)
    return s

# --- M2 扩展 15-22 ---
def m2_15(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "花刺投票 4 步全流程",
                       "从「 列出候选 」到「 确认行动 」,每一步都有具体动作")
    steps = [
        ("01 列候选", "便利贴 / 列举", "数量越多越好 · 不评判", GREEN),
        ("02 静默投票", "每人 2 票 (花 + 刺)", "独立 · 不商量 · 站着贴", GOLD),
        ("03 统计结果", "票数排序", "花票最高 = 优先机会\n刺票最高 = 关键障碍", ORANGE),
        ("04 讨论确认", "只讨论前 2-3 项", "负责人 + 时间节点", RED),
    ]
    for i, (no, action, note, col) in enumerate(steps):
        y = Inches(2.0) + Inches(1.2) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.1), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.0), Inches(1.1), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.3), Inches(2.0), Inches(0.5), no,
                 size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.8), y + Inches(0.15), Inches(4.0), Inches(0.5), action,
                 size=14, color=NAVY_DARK, bold=True)
        add_text(s, Inches(2.8), y + Inches(0.65), Inches(4.0), Inches(0.4), note.split("\n")[0],
                 size=10, color=GRAY)
        add_text(s, Inches(7.0), y + Inches(0.4), Inches(5.5), Inches(0.5),
                 "动作: " + note, size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "⚠ 静默投票是关键 — 不静默 = 强势一方主导",
             size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)

def m2_16(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "花票 vs 刺票 — 两种不同的票",
                       "花票 = 「 我支持 」「 我想做 」  ·  刺票 = 「 我反对 」「 我担心 」")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GREEN)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "🌸 花票 = 最值得优先推进", size=15, color=GREEN, bold=True)
    items_g = [
        "「 这是我觉得团队最该投入的事 」",
        "「 我个人愿意参与 」",
        "「 我相信会带来明显变化 」",
        "「 这事不解决, 其他事都受影响 」",
    ]
    for i, t in enumerate(items_g):
        add_text(s, lx + Inches(0.3), ly + Inches(0.9) + Inches(0.85) * i,
                 Inches(5.5), Inches(0.8),
                 f"• {t}", size=12, color=BLACK, line_h=1.4)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=LIGHT, line=RED)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "🌵 刺票 = 最关键障碍", size=15, color=RED, bold=True)
    items_r = [
        "「 这是阻止我们行动的最大障碍 」",
        "「 这事不解决, 任何抓手都做不了 」",
        "「 这事不解决, 团队会有意见 」",
        "「 我现在就在被这件事困扰 」",
    ]
    for i, t in enumerate(items_r):
        add_text(s, rx + Inches(0.3), ry + Inches(0.9) + Inches(0.85) * i,
                 Inches(5.5), Inches(0.8),
                 f"• {t}", size=12, color=BLACK, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "→ 花刺同时投 — 让团队既看到「 机会 」又看到「 障碍 」",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m2_17(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "✋ 练习一: 5 分钟便利贴风暴",
                       "在便利贴上写下你的「 候选抓手 」,每人至少 5 张")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "在便利贴上写下你认为团队「 下一步应该推进的事 」",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.4),
             "规则: 一张便利贴 = 一件事 · 字大 · 写动词 · 不评判",
             size=11, color=GRAY, align=PP_ALIGN.CENTER)
    # 便利贴示意
    notes = [
        ("提升小王\n接待转化", GREEN),
        ("每日晚会\n1 分钟分享", ORANGE),
        ("3 步法\nSOP 化", GOLD),
        ("给客户\n主动体验", NAVY),
        ("减少\n无效会议", RED),
    ]
    for i, (text, col) in enumerate(notes):
        x = Inches(0.5) + Inches(2.6) * i
        y = Inches(3.0)
        add_round(s, x, y, Inches(2.4), Inches(2.4), fill=col)
        for j, line in enumerate(text.split("\n")):
            add_text(s, x, y + Inches(0.6) + Inches(0.6) * j, Inches(2.4), Inches(0.6),
                     line, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.33), Inches(0.5),
             "完成后: 上墙, 让所有人快速读一遍",
             size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.4),
             "→ 主持人: 10 分钟内不评判, 只在墙上「 把重复的合并 」「 表述不清的写清楚 」",
             size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

def m2_18(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "✋ 练习二: 8 分钟静默投票",
                       "每人 2 票 (1 花 + 1 刺) · 独立 · 不商量")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "步骤 1: 主持人给每人 2 张不同颜色的便利贴 (一花一刺)",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.5),
             "步骤 2: 独立看墙上的所有候选 (3-5 分钟)",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.85), Inches(12.33), Inches(0.5),
             "步骤 3: 1 花 1 刺贴到墙上的对应项 (3 分钟)",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 示例
    add_round(s, Inches(0.5), Inches(3.7), Inches(12.33), Inches(2.5), fill=LIGHT, line=GREEN)
    add_text(s, Inches(0.7), Inches(3.85), Inches(11.9), Inches(0.4),
             "📊 示例: 8 人投票结果 (注: 这是示意,你的团队可能不同)",
             size=12, color=NAVY_DARK, bold=True)
    add_text(s, Inches(0.7), Inches(4.3), Inches(6.0), Inches(0.4),
             "🌸 花票 (优先机会):", size=12, color=GREEN, bold=True)
    flowers = [
        "体验式接待 3 步法 → 5 票",
        "每日晚会 1 分钟分享 → 4 票",
        "SOP 化培训 → 2 票",
    ]
    for i, t in enumerate(flowers):
        add_text(s, Inches(0.9), Inches(4.65) + Inches(0.4) * i,
                 Inches(5.5), Inches(0.4), t, size=11, color=BLACK)
    add_text(s, Inches(7.0), Inches(4.3), Inches(6.0), Inches(0.4),
             "🌵 刺票 (关键障碍):", size=12, color=RED, bold=True)
    thorns = [
        "老员工抵触 → 5 票",
        "培训后执行不到位 → 4 票",
        "演示样机不够 → 2 票",
    ]
    for i, t in enumerate(thorns):
        add_text(s, Inches(7.2), Inches(4.65) + Inches(0.4) * i,
                 Inches(5.5), Inches(0.4), t, size=11, color=BLACK)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "→ 接下来: 讨论花票和刺票的前 2 项, 确认责任和行动",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m2_19(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "讨论阶段: 4 个最常见陷阱",
                       "花刺投票做完, 千万别掉进这些坑里")
    pitfalls = [
        ("1", "把投票当成了「 决议 」", "投票只是参考, 还要讨论\n最终决策要看可行性", RED),
        ("2", "没记录「 反对 」的具体理由", "刺票 = 「 哪个障碍 」\n没记录 = 障碍永远解决不了", GOLD),
        ("3", "负责人没确认就结束会议", "「 行动 1, 你负责 」\n没确认 = 行动可能没启动", ORANGE),
        ("4", "没设时间节点, 行动就「 一直挂着 」", "本周 / 下周 / 2 周内\n没时间 = 没承诺", GREEN),
    ]
    for i, (no, head, sub, col) in enumerate(pitfalls):
        y = Inches(2.0) + Inches(1.25) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.15), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(1.2), Inches(1.15), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.35), Inches(1.2), Inches(0.5), no,
                 size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.0), y + Inches(0.15), Inches(10.5), Inches(0.4), head,
                 size=14, color=NAVY_DARK, bold=True)
        add_text(s, Inches(2.0), y + Inches(0.6), Inches(10.5), Inches(0.5), sub,
                 size=11, color=BLACK, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "⚠ 投票是开始, 不是结束 — 关键在「 后续跟踪 」",
             size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)

def m2_20(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "为什么这工具是「 通用 」的",
                       "不只是用来决定「 下一步重点 」, 4 种典型场景")
    scenarios = [
        ("新项目立项", "团队一起定优先级", "避免「 老板说了算 」", NAVY),
        ("季度规划", "对 5-10 个候选抓手排序", "让团队理解「 选 vs 不选 」", ORANGE),
        ("问题诊断", "多个可能原因中找最关键", "花票 = 可能解 · 刺票 = 真实障碍", GOLD),
        ("复盘会议", "哪些做得好 · 哪些待改善", "花票 = 值得保留 · 刺票 = 必须改", GREEN),
    ]
    for i, (scene, purpose, dim, col) in enumerate(scenarios):
        y = Inches(2.0) + Inches(1.25) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.15), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(3.0), Inches(1.15), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.35), Inches(3.0), Inches(0.5), scene,
                 size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(3.7), y + Inches(0.15), Inches(4.0), Inches(0.4), "用途",
                 size=10, color=GRAY, bold=True)
        add_text(s, Inches(3.7), y + Inches(0.55), Inches(4.0), Inches(0.5), purpose,
                 size=11, color=BLACK, line_h=1.3)
        add_text(s, Inches(8.0), y + Inches(0.15), Inches(4.5), Inches(0.4), "设计点",
                 size=10, color=GRAY, bold=True)
        add_text(s, Inches(8.0), y + Inches(0.55), Inches(4.5), Inches(0.5), dim,
                 size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "→ 这就是「 工具 」的复利 — 一个方法, 多个场景",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m2_21(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "M2 整合产出: 真实团队会议",
                       "下次会前 5 分钟 — 你可以这样用")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "会议前 5 分钟准备",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    pre = [
        "准备 1 个具体议题 (不是「 讨论一下下季度计划 」)",
        "准备 3 种颜色便利贴 (花 / 刺 / 备用)",
        "墙上先写好: 主题 + 4 步 (列候选 / 静默投票 / 统计 / 讨论)",
    ]
    for i, t in enumerate(pre):
        add_text(s, Inches(0.5), Inches(2.5) + Inches(0.5) * i,
                 Inches(12.33), Inches(0.4),
                 f"✓ {t}", size=12, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.5),
             "会议中关键节点",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    during = [
        "0-5 min: 每个人独立写便利贴 (5-10 张)",
        "5-15 min: 上墙 + 主持人整理 + 合并重复",
        "15-23 min: 静默投票 (主持人也投)",
        "23-30 min: 讨论前 2-3 项 + 确认行动",
    ]
    for i, t in enumerate(during):
        add_text(s, Inches(0.5), Inches(4.7) + Inches(0.4) * i,
                 Inches(12.33), Inches(0.4),
                 f"  {t}", size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5),
             "会后 24 小时内: 把讨论结果 + 行动清单 + 负责人 + 时间节点发群",
             size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)

def m2_22(idx):
    s = _generic_slide(idx, "模块 2 · 共谋抓手", ORANGE, "M2 工具速查",
                       "一页纸带走的「 花刺投票 4 步法 」")
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.5), fill=ORANGE)
    headers = ["步骤", "动作", "主持人关键", "禁止"]
    cols_w = [Inches(2.0), Inches(3.5), Inches(3.5), Inches(3.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4), h,
                 size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    rows = [
        ("01 列候选", "每人 5-10 张便利贴", "不评判 · 不引导", "「 这个不行 」"),
        ("02 静默投票", "1 花 1 刺 · 站着贴", "禁止商量 · 禁止举手", "「 大家同意吗 」"),
        ("03 统计结果", "票数排序", "只展示结果", "「 我觉得应该… 」"),
        ("04 讨论确认", "前 2-3 项", "负责人 + 时间", "「 下次再谈 」"),
    ]
    for ri, (a, b, c, d) in enumerate(rows):
        ry = Inches(2.4) + Inches(1.1) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.1), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        for i, t in enumerate([a, b, c, d]):
            add_text(s, cx + Inches(0.2), ry + Inches(0.35), cols_w[i] - Inches(0.4), Inches(0.4),
                     t, size=11, color=NAVY_DARK if i == 0 else BLACK,
                     bold=(i == 0), align=PP_ALIGN.CENTER)
            cx += cols_w[i]
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "→ 完整模板见学员手册 · F2 花刺投票卡",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# --- M3 扩展 14-32 ---
def m3_14(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "问题树 4 步 — 工具全景",
                       "症状 → 第一层拆解 → 第二层拆解 → 找切入点")
    steps = [
        ("01 写症状", "你观察到的「 不对劲 」\n具体, 可观察", "不能是「 团队效率低 」\n要「 客户转化率连续 3 月下降 」", ORANGE),
        ("02 第一层拆解", "可能是哪 3-4 方面问题", "人员 / 流程 / 产品 / 客户 / 外部\n每方面 1-2 条", GOLD),
        ("03 第二层拆解", "每个方面具体原因", "追问 「 为什么 」1-2 次\n到达「 不能再拆 」为止", RED),
        ("04 找切入点", "影响大 × 我能影响", "两轴矩阵\n两者都满足的 = 切入点", GREEN),
    ]
    for i, (no, action, note, col) in enumerate(steps):
        y = Inches(2.0) + Inches(1.3) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.2), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(1.5), Inches(1.2), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.4), Inches(1.5), Inches(0.5), no,
                 size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.3), y + Inches(0.15), Inches(5.0), Inches(0.4), "做什么",
                 size=10, color=GRAY, bold=True)
        for j, line in enumerate(action.split("\n")):
            add_text(s, Inches(2.3), y + Inches(0.5) + Inches(0.35) * j,
                     Inches(5.0), Inches(0.35), line, size=11, color=BLACK, line_h=1.2)
        add_text(s, Inches(7.5), y + Inches(0.15), Inches(5.3), Inches(0.4), "标准",
                 size=10, color=GRAY, bold=True)
        for j, line in enumerate(note.split("\n")):
            add_text(s, Inches(7.5), y + Inches(0.5) + Inches(0.35) * j,
                     Inches(5.3), Inches(0.35), line, size=11, color=BLACK, line_h=1.2)

def m3_15(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "Q1 破假设 — 5 问全景",
                       "5 问 = 5 个不同的「 思维入口 」,用 1-2 个最有效的")
    qs = [
        ("Q1 破假设", "我们一直认为「 不可能 / 必须 」\n的 — 真的吗?", "挑战默认假设", ORANGE),
        ("Q2 借他山", "其他行业 / 场景怎么解决\n类似问题?", "借鉴其他领域解法", GOLD),
        ("Q3 极端情景", "资源充足 / 只有 1/10 资源时\n怎么做?", "突破资源束缚", RED),
        ("Q4 换位思考", "从客户 / 团队 / 对手 / 上级\n角度看是什么样的?", "切换视角", GREEN),
        ("Q5 倒推法", "想象一年后问题已经解决了,\n你是怎么做到的?", "目标反推", NAVY),
    ]
    for i, (no, q, use, col) in enumerate(qs):
        y = Inches(2.0) + Inches(0.95) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(0.85), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(1.7), Inches(0.85), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.2), Inches(1.7), Inches(0.4), no,
                 size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.4), y + Inches(0.1), Inches(7.0), Inches(0.7), q,
                 size=12, color=BLACK, line_h=1.3)
        add_text(s, Inches(9.5), y + Inches(0.3), Inches(3.3), Inches(0.4),
                 "→ " + use, size=11, color=col, bold=True)

def m3_16(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "完整示例: 李明的问题树 (1/2)",
                       "症状: 客户成交率连续 3 月从 28% 降到 21%")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "症状: 客户成交率连续 3 月从 28% 降到 21%",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 第一层
    add_round(s, Inches(0.5), Inches(2.6), Inches(12.33), Inches(1.0), fill=LIGHT, line=NAVY)
    add_text(s, Inches(0.7), Inches(2.7), Inches(11.9), Inches(0.4),
             "第一层拆解: 可能是哪几方面问题?", size=13, color=NAVY_DARK, bold=True)
    items = [
        ("人员", "导购能力 / 状态 / 培训"),
        ("流程", "接待流程 / 跟进机制"),
        ("产品", "产品匹配度 / 库存"),
        ("外部", "市场竞争 / 客流变化"),
    ]
    for i, (no, sub) in enumerate(items):
        x = Inches(0.5) + Inches(3.1) * i
        add_round(s, x + Inches(0.2), Inches(3.05), Inches(2.8), Inches(0.5), fill=NAVY)
        add_text(s, x + Inches(0.2), Inches(3.1), Inches(2.8), Inches(0.4),
                 f"{no}: {sub}", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 第二层 - 拿「 人员 」 拆
    add_round(s, Inches(0.5), Inches(3.8), Inches(12.33), Inches(2.8), fill=LIGHT, line=GOLD)
    add_text(s, Inches(0.7), Inches(3.9), Inches(11.9), Inches(0.4),
             "第二层拆解 (拿「 人员 」 拆):", size=13, color=GOLD, bold=True)
    items2 = [
        "新人多, 不会卖",
        "老员工不主动, 吃老本",
        "对产品不熟, 答不上客户问题",
        "团队状态低落, 没斗志",
        "提成机制有问题, 动力不足",
    ]
    for i, t in enumerate(items2):
        y = Inches(4.35) + Inches(0.4) * i
        add_text(s, Inches(0.7), y, Inches(11.5), Inches(0.4),
                 f"  • {t}", size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
             "→ 接着: 用「 影响 × 我能影响 」找切入点 (下页)",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m3_17(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "完整示例: 李明的问题树 (2/2)",
                       "找切入点: 影响大 × 我能影响")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "从 5 个第二层原因中找「 切入 」 — 影响大 × 我能影响",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 2x2 矩阵
    mx = Inches(0.5); my = Inches(2.6)
    mw = Inches(12.33); mh = Inches(4.0)
    # 中心十字
    add_rect(s, mx + mw / 2 - Inches(0.03), my, Inches(0.06), mh, fill=GRAY)
    add_rect(s, mx, my + mh / 2 - Inches(0.03), mw, Inches(0.06), fill=GRAY)
    # 4 象限
    quads = [
        ("高影响 × 我能影响", "★ 切入点", GREEN,
         "• 复制小王的体验式接待\n• 复盘 + 调整提成机制"),
        ("高影响 × 难影响", "战略推进", ORANGE,
         "• 改造流程\n• 申请外部资源"),
        ("低影响 × 我能影响", "暂不优先", GRAY,
         "• 个别员工辅导\n• 团队氛围建设"),
        ("低影响 × 难影响", "记录即可", GRAY,
         "• 团队状态低落\n• 行业普遍问题"),
    ]
    qw = mw / 2; qh = mh / 2
    for i, (head, label, col, content) in enumerate(quads):
        qx = mx + (i % 2) * qw
        qy = my + (i // 2) * qh
        add_text(s, qx + Inches(0.2), qy + Inches(0.15), qw - Inches(0.4), Inches(0.4),
                 head, size=11, color=col, bold=True)
        add_text(s, qx + Inches(0.2), qy + Inches(0.55), qw - Inches(0.4), Inches(0.4),
                 label, size=13, color=NAVY_DARK, bold=True)
        for j, line in enumerate(content.split("\n")):
            add_text(s, qx + Inches(0.2), qy + Inches(1.0) + Inches(0.4) * j,
                     qw - Inches(0.4), Inches(0.4), line, size=10, color=BLACK, line_h=1.2)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.3),
             "★ 切入点: 复制小王的体验式接待 (高影响 × 我能影响)",
             size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

def m3_18(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "魔力提问 5 问: 对「 复制小王 」 一试",
                       "从 5 问中选 1-2 个「 没想过 」 的方向深想")
    qs = [
        ("Q1 破假设", "「 团队里必须有 「 销售明星 」 才能带高业绩 」 —\n真的吗?为什么不是团队整体都能?", ORANGE),
        ("Q2 借他山", "教育行业 / 餐饮行业 怎么让 「 体验式接待 」\n成为标准动作?有什么可以借鉴?", GOLD),
        ("Q3 极端情景", "如果只能 1 个人做事, 怎么让小王的经验\n在 1 周内复制到 3 个人?", RED),
        ("Q4 换位思考", "从客户角度: 客户最不喜欢「 销售感强 」\n的接待 — 这个洞察怎么用?", GREEN),
        ("Q5 倒推法", "想象 3 个月后团队成交率提升到 30%,\n倒推具体发生了什么?", NAVY),
    ]
    for i, (no, q, col) in enumerate(qs):
        y = Inches(2.0) + Inches(1.0) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(0.9), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(1.7), Inches(0.9), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.25), Inches(1.7), Inches(0.4), no,
                 size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.4), y + Inches(0.15), Inches(10.0), Inches(0.7), q,
                 size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
             "→ 选最让你 「 没想过 」 的那一个, 30 分钟深想",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m3_19(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "魔力提问 vs 头脑风暴 — 区别",
                       "魔力提问 「 故意 」 让你不舒服 — 这就是它的价值")
    pairs = [
        ("头脑风暴", "自由 / 开放 / 想到什么说什么", "可能 30 个想法, 30 个都没用", GRAY),
        ("魔力提问", "强制 5 个特定角度", "可能只有 2-3 个想法, 但每个有 「 没想过 」 的成分", ORANGE),
    ]
    for i, (head, what, result, col) in enumerate(pairs):
        y = Inches(2.0) + Inches(1.6) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.5), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.0), Inches(1.5), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.5), Inches(2.0), Inches(0.5), head,
                 size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.8), y + Inches(0.2), Inches(9.5), Inches(0.4), "方式",
                 size=11, color=GRAY, bold=True)
        add_text(s, Inches(2.8), y + Inches(0.6), Inches(9.5), Inches(0.5), what,
                 size=12, color=BLACK, line_h=1.4)
        add_text(s, Inches(2.8), y + Inches(1.0), Inches(9.5), Inches(0.4), "结果",
                 size=11, color=GRAY, bold=True)
        add_text(s, Inches(2.8), Inches(1.3) + y, Inches(9.5), Inches(0.4), result,
                 size=12, color=BLACK, line_h=1.4)
    add_text(s, Inches(0.5), Inches(5.3), Inches(12.33), Inches(0.5),
             "关键差别: 头脑风暴 「 找答案 」, 魔力提问 「 找问题 」",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.9), Inches(12.33), Inches(1.0),
             "魔力提问找的不是 「 怎么解决 」, 而是 「 我们一直在错误的问题上努力 」 这件事本身。\n当你发现问题是 「 错的 」, 真正的解可能就浮现出来。",
             size=12, color=BLACK, line_h=1.4, align=PP_ALIGN.CENTER)

def m3_20(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "5 问中最容易出 「 突破 」 的",
                       "经验排序: Q4 > Q2 > Q5 > Q3 > Q1")
    items = [
        ("Q4 换位思考", "90% 学员 「 之前没从这个角度看过 」", "客户 / 对手 / 上级 / 团队\n切换视角 = 全新发现", GREEN),
        ("Q2 借他山", "80% 学员 「 这个方法能直接用 」", "其他行业的现成解法\n只要适配到自己场景", ORANGE),
        ("Q5 倒推法", "70% 学员 「 这个思路是反的 」", "从目标反推路径\n而不是从现状看未来", GOLD),
        ("Q3 极端情景", "60% 学员 「 边界条件下思路反而清楚 」", "资源充足 / 1/10 资源\n逼出最本质的需求", RED),
        ("Q1 破假设", "50% 学员 「 假设是根深蒂固的 」", "真要破一个假设\n需要勇气和证据", NAVY),
    ]
    for i, (q, hit, why, col) in enumerate(items):
        y = Inches(2.0) + Inches(0.95) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(0.85), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.0), Inches(0.85), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.25), Inches(2.0), Inches(0.4), q,
                 size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.7), y + Inches(0.1), Inches(4.0), Inches(0.4),
                 hit.split(" ")[0] + " " + hit.split(" ")[1] if " " in hit else hit,
                 size=11, color=col, bold=True)
        add_text(s, Inches(2.7), y + Inches(0.5), Inches(4.0), Inches(0.4),
                 " ".join(hit.split(" ")[1:]) if " " in hit else "", size=10, color=GRAY)
        add_text(s, Inches(7.0), y + Inches(0.25), Inches(5.5), Inches(0.4),
                 why.split("\n")[0], size=10, color=BLACK)

def m3_21(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "✋ 练习一: 问题树实操",
                       "用你团队的真实难题 — 15-20 分钟")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "个人独立完成 · 写在一张 A4 纸上",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.4),
             "我现在的难题: ____________________________________",
             size=12, color=BLACK)
    items = [
        ("第一步 (3 min)", "把症状写具体, 1 句话, 可观察", "客户转化率连续 3 月下降 7%"),
        ("第二步 (5 min)", "第一层拆解: 哪 3-4 方面?  每方面 1-2 条", "人员 / 流程 / 产品 / 外部"),
        ("第三步 (7 min)", "第二层拆解: 选 1 个最像的方面, 追问 1-2 次", "人员 → 新人多 / 不会卖"),
        ("第四步 (5 min)", "影响 × 我能影响, 找 1 个切入点", "复制小王 体验式接待"),
    ]
    for i, (head, action, sample) in enumerate(items):
        y = Inches(3.0) + Inches(0.85) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(0.75), fill=LIGHT, line=GRAY_LIGHT)
        add_text(s, Inches(0.7), y + Inches(0.05), Inches(2.5), Inches(0.4), head,
                 size=12, color=NAVY_DARK, bold=True)
        add_text(s, Inches(0.7), y + Inches(0.4), Inches(7.0), Inches(0.35), action,
                 size=10, color=BLACK, line_h=1.3)
        add_text(s, Inches(8.0), y + Inches(0.4), Inches(4.8), Inches(0.35),
                 "示例: " + sample, size=9, color=GRAY, italic=True)

def m3_22(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "✋ 练习二: 5 问 1 对 1 互问",
                       "同桌 — 每人 10 分钟 — 选最让对方 「 没想过 」 的那 1 问")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=NAVY)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "作为提问者 (10 min)", size=14, color=NAVY_DARK, bold=True)
    actions = [
        "听对方说 切入点和难题",
        "从 5 问中选 1 个最适合的",
        "用 「 ____?」 的形式问出",
        "对方回答时, 不评判, 不引导",
        "听完后: 「 这个答案, 你之前想到过吗? 」",
    ]
    for i, t in enumerate(actions):
        add_text(s, lx + Inches(0.3), ly + Inches(0.9) + Inches(0.7) * i,
                 Inches(5.5), Inches(0.7), f"{i+1}. {t}", size=11, color=BLACK, line_h=1.3)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=LIGHT, line=ORANGE)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "作为回答者 (10 min)", size=14, color=ORANGE, bold=True)
    answer = [
        "认真听, 不立刻反驳",
        "真正从那个角度想 30 秒",
        "你的 「 第一反应 」 就是答案",
        "不要 「 修饰 」 你的真实想法",
        "听完后: 「 这个角度我没想过 」",
    ]
    for i, t in enumerate(answer):
        add_text(s, rx + Inches(0.3), ry + Inches(0.9) + Inches(0.7) * i,
                 Inches(5.5), Inches(0.7), f"{i+1}. {t}", size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "⚠ 魔力提问的真值来自 「 不修饰的真实回答 」",
             size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)

def m3_23(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "5 问使用: 8 步完整流程",
                       "从 「 我有难题 」 到 「 我有新方向 」")
    steps = [
        ("01 写难题", "一句话 · 具体 · 可观察", "「 客户转化率下降 7% 」", ORANGE),
        ("02 拆第一层", "3-4 方面 · 每方面 1-2 条", "人员 / 流程 / 产品 / 外部", GOLD),
        ("03 拆第二层", "选 1 方面 · 追问 1-2 次", "人员 → 新人多 → 不会卖", RED),
        ("04 找切入点", "影响 × 我能影响", "复制小王 · 体验式接待", GREEN),
        ("05 试 5 问", "对切入点试 Q1-Q5", "每个 30 秒 想想", NAVY),
        ("06 选 1-2 问", "选最 「 没想过 」 的", "Q4 换位思考 / Q2 借他山", ORANGE),
        ("07 深想 30 分钟", "只对选中的 1-2 问", "不展开其他", GOLD),
        ("08 写下 新方向", "1-2 句 · 可执行", "「 让客户主动体验 」", RED),
    ]
    for i, (no, action, sample, col) in enumerate(steps):
        y = Inches(2.0) + Inches(0.6) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(0.55), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(1.3), Inches(0.55), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.13), Inches(1.3), Inches(0.4), no,
                 size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.0), y + Inches(0.05), Inches(4.5), Inches(0.4), action,
                 size=10, color=NAVY_DARK, bold=True)
        add_text(s, Inches(6.5), y + Inches(0.05), Inches(6.0), Inches(0.4),
                 "示例: " + sample, size=9, color=GRAY, italic=True)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.2),
             "→ 下一步: 把新方向带到 M4 共创会",
             size=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m3_24(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "⚠ 常见陷阱 1: 跳过问题树直接魔力提问",
                       "魔力提问是 「 工具 」, 不是 「 替代品 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "没有清晰的 「 问题 」, 魔力提问会变成 「 泛泛思考 」",
             size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    lx = Inches(0.5); ly = Inches(2.7)
    add_round(s, lx, ly, Inches(6.0), Inches(3.5), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "✗ 不做问题树的后果", size=14, color=RED, bold=True)
    bad = [
        "5 问出来的都是「 怎么卖东西 」 层面",
        "没法深入, 问 30 分钟还是「 找客户 」",
        "回到原点: 「 我们的难题到底是什么 」",
    ]
    for i, t in enumerate(bad):
        add_text(s, lx + Inches(0.3), ly + Inches(0.85) + Inches(0.85) * i,
                 Inches(5.5), Inches(0.8),
                 f"• {t}", size=11, color=BLACK, line_h=1.4)
    rx = Inches(6.7); ry = Inches(2.7)
    add_round(s, rx, ry, Inches(6.13), Inches(3.5), fill=LIGHT, line=GREEN)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "✓ 先问题树后魔力提问", size=14, color=GREEN, bold=True)
    good = [
        "问题树给出「 具体的切入点 」",
        "5 问才能围绕切入点深想",
        "切入点 × 5 问 = 30 个方向",
    ]
    for i, t in enumerate(good):
        add_text(s, rx + Inches(0.3), ry + Inches(0.85) + Inches(0.85) * i,
                 Inches(5.5), Inches(0.8),
                 f"• {t}", size=11, color=BLACK, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "→ 顺序: 问题树 (10 min) → 5 问 (30 min) → 新方向",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m3_25(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "⚠ 常见陷阱 2: 5 问都浅尝辄止",
                       "5 问不是 「 问完 5 个 」, 是 「 选 1-2 个深想 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "常见错误: 5 问都用, 每个 5 分钟 — 30 分钟过去, 没有真正深入的",
             size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_round(s, Inches(0.5), Inches(2.7), Inches(12.33), Inches(2.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, Inches(0.7), Inches(2.85), Inches(11.9), Inches(0.4),
             "正确用法:", size=13, color=NAVY_DARK, bold=True)
    right = [
        "Q1 试 30 秒 — 不深入, 只看是不是有「 没想过 」的方向",
        "Q2 试 30 秒 — 同上",
        "Q3 试 30 秒",
        "Q4 试 30 秒 — 感觉到了: 「 哦, 我没想到 」",
        "Q5 试 30 秒",
        "选定 Q4 + Q2 — 各 10-15 分钟深想 — 写出 2-3 条新方向",
    ]
    for i, t in enumerate(right):
        add_text(s, Inches(0.7), Inches(3.3) + Inches(0.3) * i,
                 Inches(11.9), Inches(0.3),
                 f"  {t}", size=10, color=BLACK, line_h=1.2)
    add_round(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5), fill=RED)
    add_text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.4),
             "⚠ 反面教材: 5 问 × 5 分钟 = 25 分钟过去, 1 个深方向都没有",
             size=12, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.9),
             "5 问 「 入口 」 不是 「 终点 」, 必须选 1-2 个深想才有价值",
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_h=1.4)

def m3_26(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "⚠ 常见陷阱 3: 5 问出不真实",
                       "「 这事不可能 」  「 客户不会接受 」  — 这些是 「 自我审查 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "5 问的最大障碍: 自我审查",
             size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    lx = Inches(0.5); ly = Inches(2.7)
    add_round(s, lx, ly, Inches(6.0), Inches(3.5), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "✗ 「 不可能 」的回答", size=14, color=RED, bold=True)
    bad = [
        "「 这事不可能 」",
        "「 我们这行不会做 」",
        "「 客户不会接受 」",
        "「 总部不会批 」",
    ]
    for i, t in enumerate(bad):
        add_text(s, lx + Inches(0.3), ly + Inches(0.85) + Inches(0.6) * i,
                 Inches(5.5), Inches(0.6),
                 f"  {t}", size=11, color=BLACK, line_h=1.3)
    rx = Inches(6.7); ry = Inches(2.7)
    add_round(s, rx, ry, Inches(6.13), Inches(3.5), fill=LIGHT, line=GREEN)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "✓ 改成 「 为什么不可能 」", size=14, color=GREEN, bold=True)
    good = [
        "「 这事不可能 」 → 「 为什么不可能 」",
        "「 哪些部分可以, 哪些不可以 」",
        "「 哪个具体障碍 」 → 「 谁解决 」",
    ]
    for i, t in enumerate(good):
        add_text(s, rx + Inches(0.3), ry + Inches(0.85) + Inches(0.85) * i,
                 Inches(5.5), Inches(0.85),
                 f"• {t}", size=11, color=BLACK, line_h=1.4)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "→ 5 问需要的不是 「 想 」, 是 「 认真回答 」",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m3_27(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "M3 关键收获: 难题不靠经验",
                       "问题树拆清楚 + 5 问想透 = 新方向")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "M3 给你 3 件具体的事, 离开教室后可以马上用:",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("拆 1 个难题", "用问题树拆 1 个团队里正在发生的难题", "10-15 分钟 · 4 步"),
        ("试 5 问", "对 「 切入点 」 试 5 问, 选 1-2 个深想", "30 分钟 · 不展开其他"),
        ("写下新方向", "1-2 句 · 可执行", "带到 M4 共创会"),
    ]
    card_w = Inches(4.0); card_h = Inches(2.5)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.6)
    for i, (head, sub, time) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=RED)
        add_rect(s, x, y, card_w, Inches(0.6), fill=RED)
        add_text(s, x, y + Inches(0.13), card_w, Inches(0.4), head,
                 size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(0.85), card_w - Inches(0.6), Inches(1.2),
                 sub, size=12, color=NAVY_DARK, bold=True, line_h=1.4, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(2.0), card_w - Inches(0.6), Inches(0.4),
                 time, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    add_round(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(5.65), Inches(11.9), Inches(0.4),
             "验证标准: 「 我有了一个之前没想到的方向 」", size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.9),
             "如果你的 5 问没有让你 「 哦 」 一声 — 5 问没问对, 不是工具没用",
             size=12, color=WHITE, line_h=1.4, align=PP_ALIGN.CENTER)

def m3_28(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "M3 自检: 你有 「 新方向 」 了吗",
                       "3 个标志, 全部满足才算成功")
    items = [
        ("✓ 问题具体", "「 客户成交率从 28% 降到 21% 」", "不是 「 团队效率低 」", GREEN),
        ("✓ 切入点清晰", "复制小王 · 体验式接待", "影响大 × 我能影响", GREEN),
        ("✓ 5 问中有 「 哦 」", "Q4 换位: 「 哦, 客户最不喜欢销售感 」", "至少 1 个 「 哦 」", GREEN),
    ]
    card_h = Inches(1.3)
    for i, (head, sample, std, col) in enumerate(items):
        y = Inches(2.0) + (card_h + Inches(0.2)) * i
        add_round(s, Inches(0.5), y, Inches(12.33), card_h, fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.5), card_h, fill=col)
        add_text(s, Inches(0.5), y + Inches(0.4), Inches(2.5), Inches(0.5), head,
                 size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(3.2), y + Inches(0.15), Inches(5.0), Inches(0.4),
                 "示例", size=10, color=GRAY, bold=True)
        add_text(s, Inches(3.2), y + Inches(0.55), Inches(5.0), Inches(0.6),
                 sample, size=12, color=NAVY_DARK, bold=True, line_h=1.3)
        add_text(s, Inches(8.5), y + Inches(0.15), Inches(4.3), Inches(0.4),
                 "标准", size=10, color=GRAY, bold=True)
        add_text(s, Inches(8.5), y + Inches(0.55), Inches(4.3), Inches(0.6),
                 std, size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "⚠ 如果没满足 3 项, 重新做问题树 — 切入点和难题没对齐",
             size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)

def m3_29(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "M3 与其他工具的连接",
                       "M3 给你 「 新方向 」, 不是 「 行动方案 」")
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "M3 给你", size=14, color=ORANGE, bold=True)
    outs = [
        ("一个具体难题", "「 客户成交率从 28% 降到 21% 」"),
        ("一个切入点", "复制小王 · 体验式接待"),
        ("1-2 个新方向", "Q4: 让客户主动体验 · Q2: 借鉴餐饮"),
    ]
    for i, (head, sub) in enumerate(outs):
        add_text(s, lx + Inches(0.3), ly + Inches(0.85) + Inches(1.2) * i,
                 Inches(5.5), Inches(0.4), "✓ " + head, size=13, color=NAVY_DARK, bold=True)
        add_text(s, lx + Inches(0.3), ly + Inches(1.25) + Inches(1.2) * i,
                 Inches(5.5), Inches(0.7), sub, size=11, color=BLACK, line_h=1.3)
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=NAVY_DARK)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "下一段 M4: 高效脑暴双矩阵", size=14, color=ORANGE, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(0.9), Inches(5.5), Inches(1.0),
             "你的 「 新方向 」\n会成为 M4 共创的 「 主题 」",
             size=15, color=WHITE, bold=True, line_h=1.4)
    add_text(s, rx + Inches(0.3), ry + Inches(2.4), Inches(5.5), Inches(2.0),
             "在 M4 里:\n─ 用发散矩阵让团队一起 「 拍 」 这个方向\n─ 团队反馈 + 你的想法 = 优先行动\n─ 用收敛矩阵找到 「 高影响 + 易落地 」 的",
             size=12, color=WHITE, line_h=1.5)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "→ 下一段: M4 引领共创",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m3_30(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "5 工具连接: M3 之后",
                       "M3 给你 「 新方向 」, M4 用 「 共创 」 把它变成 「 行动 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "5 工具的连接 — 整体看一遍",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    flow = [
        ("M1", "复制成功", "提炼原则", ORANGE),
        ("M2", "共谋抓手", "找优先机会 / 关键障碍", GOLD),
        ("M3", "应对难题", "新方向 (这步)", RED),
        ("M4", "引领共创", "优先行动方案", GREEN),
        ("M5", "前瞻思考", "风险 + 机会", NAVY),
    ]
    card_w = Inches(2.2); card_h = Inches(2.0)
    gap = Inches(0.2)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(2.6)
    for i, (no, name, what, col) in enumerate(flow):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=col)
        add_rect(s, x, y, card_w, Inches(0.6), fill=col)
        add_text(s, x, y + Inches(0.13), card_w, Inches(0.4), no,
                 size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.75), card_w, Inches(0.4), name,
                 size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(1.2), card_w - Inches(0.4), Inches(0.7),
                 what, size=10, color=BLACK, line_h=1.3, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(s, x + card_w + Inches(0.05), y + Inches(0.8), gap - Inches(0.1), Inches(0.4),
                     "▶", size=18, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.5),
             "M3 之后, 你带着 「 新方向 」 进入 M4",
             size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.4),
             "「 新方向 」 是 1 句话 · 可执行 · 在你职权范围内\n「 新方向 」 是 「 切入点 × 5 问 」 的产物\n「 新方向 」 是 M4 团队共创的 「 主题 」",
             size=12, color=BLACK, line_h=1.5, align=PP_ALIGN.CENTER)

def m3_31(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "M3 模块收尾 · 茶歇",
                       "回来后, 我们进入第四个工具 M4 高效脑暴双矩阵")
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.5),
             "M3", size=120, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.7),
             "应对难题", size=32, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             "Solve Tough Problems", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    quote_block(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                "经验会失灵, 但 「 认真回答 5 问 」 不会。",
                author="—— 第三部分核心信念")
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "请准备: 你在 M3 找到的 「 新方向 」 一句话, 下次课会用到",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)

def m3_32(idx):
    s = _generic_slide(idx, "模块 3 · 应对难题", RED, "M3 工具速查",
                       "一页纸带走的 「 问题树 + 5 问 」")
    # 左: 问题树
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=RED)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "问题树 4 步", size=14, color=RED, bold=True)
    tree = [
        ("01 写症状", "可观察 · 1 句话", "客户成交率从 28% 降到 21%"),
        ("02 第一层", "3-4 方面", "人员 / 流程 / 产品 / 外部"),
        ("03 第二层", "1 方面 × 1-2 次追问", "人员 → 新人多 → 不会卖"),
        ("04 切入点", "影响大 × 我能影响", "复制小王 · 体验式接待"),
    ]
    for i, (no, what, sample) in enumerate(tree):
        y = ly + Inches(0.85) + Inches(0.85) * i
        add_text(s, lx + Inches(0.3), y, Inches(2.0), Inches(0.4), no,
                 size=12, color=RED, bold=True)
        add_text(s, lx + Inches(0.3), y + Inches(0.35), Inches(2.5), Inches(0.4), what,
                 size=10, color=BLACK)
        add_text(s, lx + Inches(2.8), y + Inches(0.2), Inches(3.0), Inches(0.4),
                 sample, size=9, color=GRAY, italic=True)
    # 右: 5 问
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=LIGHT, line=GOLD)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "魔力提问 5 问", size=14, color=GOLD, bold=True)
    qs = [
        ("Q1 破假设", "我们认为的「 不可能 / 必须 」是真的吗?"),
        ("Q2 借他山", "其他行业 / 场景怎么解决?"),
        ("Q3 极端情景", "资源充足 / 1/10 资源时怎么做?"),
        ("Q4 换位思考", "从客户 / 对手 / 上级 / 团队角度看?"),
        ("Q5 倒推法", "一年后问题解决了, 你是怎么做到的?"),
    ]
    for i, (no, q) in enumerate(qs):
        y = ry + Inches(0.85) + Inches(0.65) * i
        add_text(s, rx + Inches(0.3), y, Inches(1.7), Inches(0.4), no,
                 size=11, color=GOLD, bold=True)
        add_text(s, rx + Inches(2.0), y, Inches(4.0), Inches(0.6),
                 q, size=10, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "→ 完整模板见学员手册 · F3 问题树 + 5 问卡",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# --- M4 扩展 12-24 ---
def m4_12(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "发散矩阵 2×2 — 工具全景",
                       "横轴 × 纵轴 = 4 象限, 强制团队每个象限都产想法")
    cells = [
        ("Q1 快赢区", "短期 + 团队内部可控", "立刻能做的事", "小王晨会分享\n3 步指南发放", GREEN),
        ("Q2 协调推进区", "短期 + 需外部资源", "近期要启动, 需争取支持", "申请展示机\n申请培训资源", ORANGE),
        ("Q3 能力建设区", "中长期 + 团队内部可控", "团队自己积累能力", "内部案例库\n分享节奏", GOLD),
        ("Q4 战略布局区", "中长期 + 需外部资源", "提前规划 + 资源支持", "区域 SOP 申请\n差异化优势梳理", RED),
    ]
    for i, (name, where, role, sample, col) in enumerate(cells):
        qx = Inches(0.5) + (i % 2) * Inches(6.2)
        qy = Inches(2.0) + (i // 2) * Inches(2.4)
        add_round(s, qx, qy, Inches(6.0), Inches(2.3), fill=LIGHT, line=col)
        add_rect(s, qx, qy, Inches(2.0), Inches(2.3), fill=col)
        add_text(s, qx, qy + Inches(0.4), Inches(2.0), Inches(0.5), name,
                 size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, qx, qy + Inches(0.95), Inches(2.0), Inches(0.4),
                 "Q" + str(i+1), size=10, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, qx + Inches(0.2), qy + Inches(1.5), Inches(1.6), Inches(0.6),
                 where.split("\n")[0], size=9, color=WHITE, align=PP_ALIGN.CENTER, line_h=1.2)
        add_text(s, qx + Inches(2.2), qy + Inches(0.15), Inches(3.6), Inches(0.4),
                 where, size=10, color=col, bold=True)
        add_text(s, qx + Inches(2.2), qy + Inches(0.55), Inches(3.6), Inches(0.4),
                 role, size=10, color=GRAY, italic=True)
        add_text(s, qx + Inches(2.2), qy + Inches(1.1), Inches(3.6), Inches(1.1),
                 "示例: " + sample, size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "→ 主持原则: 每格至少 2-3 条, 不能空",
             size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m4_13(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "为什么发散需要 「 矩阵 」",
                       "自由讨论 vs 矩阵发散: 3 个根本差异")
    items = [
        ("结构强制", "自由讨论只想到 Q1 · 矩阵强制想 Q2-Q4", "覆盖所有类型想法", ORANGE),
        ("角色清晰", "自由讨论 「 大家随便说 」, 矩阵每格有专属问题", "避免冷场 + 避免跑题", GOLD),
        ("可对比", "自由讨论想法混杂, 矩阵想法有位置", "收敛阶段直接用 2 个维度评估", GREEN),
    ]
    for i, (head, what, benefit, col) in enumerate(items):
        y = Inches(2.0) + Inches(1.6) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.5), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.5), Inches(1.5), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.5), Inches(2.5), Inches(0.5), head,
                 size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(3.2), y + Inches(0.2), Inches(4.0), Inches(0.4),
                 "自由讨论的痛点", size=10, color=GRAY, bold=True)
        add_text(s, Inches(3.2), y + Inches(0.6), Inches(4.0), Inches(0.8),
                 what, size=12, color=BLACK, line_h=1.4)
        add_text(s, Inches(7.5), y + Inches(0.2), Inches(5.0), Inches(0.4),
                 "矩阵的优势", size=10, color=GRAY, bold=True)
        add_text(s, Inches(7.5), y + Inches(0.6), Inches(5.0), Inches(0.8),
                 benefit, size=12, color=col, bold=True, line_h=1.4)

def m4_14(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "发散矩阵 4 格: 专属引导话术",
                       "主持时, 用这 4 段话术 1 格 1 格引")
    cells = [
        ("Q1 快赢区", "短期 + 内部", "「 我们先从最容易启动的开始 — 接下来 4 周, 团队自己就能做, 不需要等总部的事, 有哪些? 」", GREEN),
        ("Q2 协调推进区", "短期 + 外部", "「 接下来 1 个月能看到效果, 但需要总部 / 跨部门配合的, 有哪些? 比如样机、宣传物料、培训支持。 」", ORANGE),
        ("Q3 能力建设区", "中长期 + 内部", "「 我们看更长期 — 3 个月内, 团队自己能在什么能力上做积累? 培训机制、知识沉淀、内部分享节奏。 」", GOLD),
        ("Q4 战略布局区", "中长期 + 外部", "「 最后, 有没有需要提前布局、需要和上级或其他合作方谈的? 本季度准备, 效果 3 个月后体现。 」", RED),
    ]
    for i, (name, where, what, col) in enumerate(cells):
        y = Inches(2.0) + Inches(1.2) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.1), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.5), Inches(1.1), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.3), Inches(2.5), Inches(0.4), name,
                 size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.5), y + Inches(0.7), Inches(2.5), Inches(0.3), where,
                 size=9, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(3.2), y + Inches(0.2), Inches(9.5), Inches(0.9), what,
                 size=11, color=BLACK, line_h=1.4)

def m4_15(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "发散 3 原则 — 主持人自查",
                       "1. 不评判  2. 不先说想法  3. 强制每格 2-3 条")
    items = [
        ("原则 1", "发散阶段不评判", "「 这个不现实 」  「 这个已经试过了 」\n这些话在发散阶段禁止\n把评价全部留到收敛矩阵", ORANGE),
        ("原则 2", "主持人先不说想法", "主持人的角色是 「 帮其他人想法出来 」\n不是 「 说自己的想法 」\n主持人先说 → 团队跟着说 → 真实想法出不来", GOLD),
        ("原则 3", "强制每格 2-3 条", "发散矩阵的价值 = 4 格都有产出\n不强制 → 只有 Q1 有人写\n「 Q3 有想法吗? Q4 呢? 」 直接问", RED),
    ]
    for i, (no, head, sub, col) in enumerate(items):
        y = Inches(2.0) + Inches(1.6) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.5), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.0), Inches(1.5), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.5), Inches(2.0), Inches(0.5), no,
                 size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.7), y + Inches(0.2), Inches(9.8), Inches(0.4), head,
                 size=14, color=NAVY_DARK, bold=True)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, Inches(2.7), y + Inches(0.65) + Inches(0.35) * j,
                     Inches(9.8), Inches(0.35), line, size=11, color=BLACK, line_h=1.2)

def m4_16(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "收敛矩阵 2×2 — 工具全景",
                       "影响力 × 落地难度 = 4 个区间")
    cells = [
        ("优先行动区", "高影响 + 落地易", "★ 今天确定要做", "小王晨会分享\n每日晚会 1 分钟", GREEN),
        ("战略项目区", "高影响 + 落地难", "规划推进\n不指望随手做", "申请区域 SOP\n申请展示机", ORANGE),
        ("随机应变区", "低影响 + 落地易", "有机会顺手做\n不投入太多精力", "门口欢迎标识\n梳理优势说明", GRAY),
        ("暂不考虑区", "低影响 + 落地难", "本轮不做", "流程大改\n行业普遍问题", GRAY),
    ]
    for i, (name, where, role, sample, col) in enumerate(cells):
        qx = Inches(0.5) + (i % 2) * Inches(6.2)
        qy = Inches(2.0) + (i // 2) * Inches(2.4)
        add_round(s, qx, qy, Inches(6.0), Inches(2.3), fill=LIGHT, line=col)
        add_rect(s, qx, qy, Inches(2.0), Inches(2.3), fill=col)
        add_text(s, qx, qy + Inches(0.4), Inches(2.0), Inches(0.5), name,
                 size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, qx, qy + Inches(0.95), Inches(2.0), Inches(0.4),
                 "★" if "优先" in name else "", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, qx, qy + Inches(1.4), Inches(2.0), Inches(0.6),
                 where, size=9, color=WHITE, align=PP_ALIGN.CENTER, line_h=1.2)
        add_text(s, qx + Inches(2.2), qy + Inches(0.15), Inches(3.6), Inches(0.4),
                 role, size=11, color=col, bold=True)
        add_text(s, qx + Inches(2.2), qy + Inches(0.6), Inches(3.6), Inches(0.4),
                 "决策:", size=10, color=GRAY, bold=True)
        add_text(s, qx + Inches(2.2), qy + Inches(0.95), Inches(3.6), Inches(1.2),
                 sample, size=11, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.3),
             "★ 优先行动区: 影响力高 + 落地容易 — 立即开始",
             size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

def m4_17(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "收敛 3 步 — 主持流程",
                       "独立评估 → 对比分歧 → 确认行动")
    steps = [
        ("01 独立评估", "每人从发散矩阵里选 8-10 条\n在收敛矩阵上独立标注\n不商量 · 不相互看", "5-7 分钟", ORANGE),
        ("02 对比分歧", "把各自的评估展示出来\n找出分歧最大的 2-3 条\n只讨论分歧项\n讨论 = 「 我们对这条想法的理解一致吗 」", "8-10 分钟", GOLD),
        ("03 确认行动", "从 ★ 优先行动区选 2-3 条\n明确:\n─ 谁负责\n─ 第一步做什么\n─ 什么时候完成", "5-7 分钟", RED),
    ]
    for i, (no, action, time, col) in enumerate(steps):
        y = Inches(2.0) + Inches(1.7) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.6), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.5), Inches(1.6), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.5), Inches(2.5), Inches(0.5), no,
                 size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.5), y + Inches(1.0), Inches(2.5), Inches(0.4),
                 time, size=10, color=WHITE, align=PP_ALIGN.CENTER)
        for j, line in enumerate(action.split("\n")):
            add_text(s, Inches(3.2), y + Inches(0.2) + Inches(0.3) * j,
                     Inches(9.0), Inches(0.3), line, size=11, color=BLACK, line_h=1.2)

def m4_18(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "收敛 2 原则 — 主持人自查",
                       "1. 不主导  2. 用分歧找共识")
    items = [
        ("原则 1", "不主导", "「 我觉得应该… 」 这句话是危险信号\n主持人的选择 ≠ 团队选择\n如果你的评估和大家不同, 正确做法:\n说出判断依据, 大家讨论", ORANGE),
        ("原则 2", "用分歧找共识", "分歧 = 大家理解不同 ≠ 大家对价值判断不同\n讨论 「 这条想法的真实含义 」\n理解一致后, 位置常常自动归位", GOLD),
    ]
    for i, (no, head, sub, col) in enumerate(items):
        y = Inches(2.0) + Inches(2.3) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(2.2), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.0), Inches(2.2), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.85), Inches(2.0), Inches(0.5), no,
                 size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.7), y + Inches(0.2), Inches(9.8), Inches(0.4), head,
                 size=14, color=NAVY_DARK, bold=True)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, Inches(2.7), y + Inches(0.65) + Inches(0.3) * j,
                     Inches(9.8), Inches(0.3), line, size=11, color=BLACK, line_h=1.2)

def m4_19(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "✋ 练习一: 主持发散矩阵",
                       "4-5 人小组, 一人主持, 全程 20-25 分钟")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "角色: 主持人 (1) + 组员 (3-4) · 一人主持全场, 下一轮轮换",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.4),
             "主题: 用你在 M3 找到的 「 新方向 」",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    steps = [
        ("0-2 min", "画 2×2 矩阵 (4 象限)", "主持人"),
        ("2-7 min", "每格 1 分钟引导, 每人写 2-3 张便利贴", "全员"),
        ("7-15 min", "上墙 + 整理 + 合并重复 + 表述不清的写清楚", "主持人"),
        ("15-20 min", "补充: 「 Q3 / Q4 还缺什么 」", "主持人 + 组员"),
        ("20-25 min", "反思: 哪个格子最难产想法? 为什么?", "全员"),
    ]
    for i, (time, action, who) in enumerate(steps):
        y = Inches(3.0) + Inches(0.7) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(0.6), fill=LIGHT, line=GRAY_LIGHT)
        add_text(s, Inches(0.7), y + Inches(0.15), Inches(1.7), Inches(0.4), time,
                 size=12, color=ORANGE, bold=True)
        add_text(s, Inches(2.5), y + Inches(0.15), Inches(8.0), Inches(0.4), action,
                 size=11, color=BLACK)
        add_text(s, Inches(10.5), y + Inches(0.15), Inches(2.3), Inches(0.4), who,
                 size=11, color=GREEN, bold=True)

def m4_20(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "✋ 练习二: 主持收敛矩阵",
                       "15-20 分钟: 独立评估 → 对比分歧 → 确认行动")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "从练习一发散的想法里, 选 8 条做收敛",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    steps = [
        ("5-7 min", "独立评估: 每人 1 张空白收敛矩阵\n对 8 条想法独立打位置 (高/中/低影响, 易/中/难落地)\n不商量 · 不看", ORANGE),
        ("5-7 min", "对比分歧: 8 条展示出来\n只讨论分歧最大的 2 条\n讨论 = 「 我们对这条想法的理解一致吗 」", GOLD),
        ("5 min", "确认行动: 从 ★ 优先行动区选 2-3 条\n明确: 负责人 / 第一步 / 完成时间", RED),
    ]
    for i, (time, action, col) in enumerate(steps):
        y = Inches(2.5) + Inches(1.5) * i
        add_round(s, Inches(0.5), y, Inches(12.33), Inches(1.4), fill=LIGHT, line=col)
        add_rect(s, Inches(0.5), y, Inches(2.0), Inches(1.4), fill=col)
        add_text(s, Inches(0.5), y + Inches(0.5), Inches(2.0), Inches(0.5), time,
                 size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(action.split("\n")):
            add_text(s, Inches(2.7), y + Inches(0.2) + Inches(0.3) * j,
                     Inches(9.5), Inches(0.3), line, size=11, color=BLACK, line_h=1.2)
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.2),
             "→ 反思: 「 主持人说了 「 我觉得… 」 几次? 」",
             size=10, color=ORANGE, bold=True, italic=True, align=PP_ALIGN.CENTER)

def m4_21(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "⚠ 常见错误 1: 发散时过早评判",
                       "「 这个不现实 」 — 想法在说出那一刻就死了")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "主持人的 5 句 「 杀手话 」  —  听到就停下来",
             size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    kills = [
        "「 这个不现实 」",
        "「 我们试过了, 没用 」",
        "「 这超出我们权限 」",
        "「 老板不会批 」",
        "「 客户不会接受 」",
    ]
    card_w = Inches(2.3); card_h = Inches(2.0)
    gap = Inches(0.15)
    total_w = card_w * 5 + gap * 4
    start_x = (SW - total_w) / 2
    y = Inches(2.6)
    for i, t in enumerate(kills):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=RED)
        add_text(s, x, y + Inches(0.3), card_w, Inches(0.6), "✗",
                 size=40, color=RED, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(1.2), card_w - Inches(0.4), Inches(0.7),
                 t, size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, line_h=1.3)
    add_round(s, Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.7), fill=GREEN)
    add_text(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.4),
             "✓ 主持人说: 「 我们先把它记下来, 收敛阶段再判断 」",
             size=13, color=WHITE, bold=True)
    add_text(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.2),
             "「 不现实 」  「 不可能 」  在发散阶段没有价值\n所有判断留到收敛矩阵\n如果发散阶段就评判, 团队下次不再说真实想法了",
             size=12, color=WHITE, line_h=1.4, align=PP_ALIGN.CENTER)

def m4_22(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "⚠ 常见错误 2: 收敛时凭 「 声音大 」 决策",
                       "靠 「 大家觉得呢 」 投票, 退回到 「 主管说了算 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "收敛矩阵的设计目的: 用 2 个客观维度替代 「 投票 」",
             size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    lx = Inches(0.5); ly = Inches(2.6)
    add_round(s, lx, ly, Inches(6.0), Inches(3.5), fill=LIGHT, line=GRAY_LIGHT)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "✗ 错的收敛方式", size=14, color=RED, bold=True)
    bad = [
        "「 大家觉得哪个好, 举手 」",
        "「 主管觉得应该先做 A 」",
        "「 声音大的人说了算 」",
        "「 一条没排除, 都做 」",
    ]
    for i, t in enumerate(bad):
        add_text(s, lx + Inches(0.3), ly + Inches(0.85) + Inches(0.6) * i,
                 Inches(5.5), Inches(0.6),
                 f"  {t}", size=12, color=BLACK, line_h=1.3)
    rx = Inches(6.7); ry = Inches(2.6)
    add_round(s, rx, ry, Inches(6.13), Inches(3.5), fill=LIGHT, line=GREEN)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "✓ 正确的收敛方式", size=14, color=GREEN, bold=True)
    good = [
        "「 影响力高/低 — 落地易/难 」",
        "「 ★ 优先行动区 = 高影响 + 落地易 」",
        "「 分歧 = 大家理解不同 」",
        "「 用 2 维度客观评估 」",
    ]
    for i, t in enumerate(good):
        add_text(s, rx + Inches(0.3), ry + Inches(0.85) + Inches(0.6) * i,
                 Inches(5.5), Inches(0.6),
                 f"• {t}", size=12, color=BLACK, line_h=1.3)
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.4),
             "→ 收敛矩阵的价值: 让 「 矩阵逻辑 」 引导聚焦, 不是 「 权力逻辑 」",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m4_23(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "M4 关键收获: 主持一个 「 高质量 」 共创",
                       "发散 + 收敛 = 完整闭环")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "M4 给你 3 件具体的事, 离开教室后可以马上用:",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("主持 1 次发散", "在一个真实议题上主持一次发散矩阵", "20-25 min · 4 格都产想法"),
        ("主持 1 次收敛", "用 8 条想法做收敛, 找到 ★ 优先行动", "15-20 min · 3 步流程"),
        ("反思 1 次", "自己主持的哪一步做得最好\n哪一步下次做得不一样", "5-10 min · 自我反馈"),
    ]
    card_w = Inches(4.0); card_h = Inches(2.5)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.6)
    for i, (head, sub, time) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_round(s, x, y, card_w, card_h, fill=LIGHT, line=GREEN)
        add_rect(s, x, y, card_w, Inches(0.6), fill=GREEN)
        add_text(s, x, y + Inches(0.13), card_w, Inches(0.4), head,
                 size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, line in enumerate(sub.split("\n")):
            add_text(s, x + Inches(0.3), y + Inches(0.85) + Inches(0.45) * j,
                     card_w - Inches(0.6), Inches(0.45), line, size=12, color=NAVY_DARK,
                     bold=True, line_h=1.3, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(2.0), card_w - Inches(0.6), Inches(0.4),
                 time, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    add_round(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5), fill=NAVY_DARK)
    add_text(s, Inches(0.7), Inches(5.65), Inches(11.9), Inches(0.4),
             "✓ 验证标准: 「 团队产出了 2-3 条 ★ 优先行动, 每条都有负责人和时间 」",
             size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.9),
             "如果没满足 — 主持过程中, 可能在某格漏了 / 在评估时 「 我觉得 」 说了几次",
             size=12, color=WHITE, line_h=1.4, align=PP_ALIGN.CENTER)

def m4_24(idx):
    s = _generic_slide(idx, "模块 4 · 引领共创", GREEN, "M4 工具速查",
                       "一页纸带走的 「 发散 + 收敛 」 双矩阵")
    # 左: 发散
    lx = Inches(0.5); ly = Inches(1.85)
    add_round(s, lx, ly, Inches(6.0), Inches(4.6), fill=LIGHT, line=GREEN)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), Inches(5.5), Inches(0.5),
             "发散矩阵 2×2", size=14, color=GREEN, bold=True)
    # 简化矩阵
    add_text(s, lx + Inches(0.3), ly + Inches(0.7), Inches(5.5), Inches(0.3),
             "短期+内部 = Q1 快赢", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(1.0), Inches(5.5), Inches(0.3),
             "短期+外部 = Q2 协调推进", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(1.3), Inches(5.5), Inches(0.3),
             "中长期+内部 = Q3 能力建设", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(1.6), Inches(5.5), Inches(0.3),
             "中长期+外部 = Q4 战略布局", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(2.2), Inches(5.5), Inches(0.4),
             "主持原则:", size=12, color=GREEN, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(2.6), Inches(5.5), Inches(0.4),
             "1. 不评判 (评价留到收敛)", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(2.9), Inches(5.5), Inches(0.4),
             "2. 不先说想法", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(3.2), Inches(5.5), Inches(0.4),
             "3. 强制每格 2-3 条", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(3.7), Inches(5.5), Inches(0.4),
             "4. 时长: 4 × 5 min = 20 min", size=10, color=BLACK)
    add_text(s, lx + Inches(0.3), ly + Inches(4.2), Inches(5.5), Inches(0.3),
             "产出: 20-30 条便利贴", size=10, color=ORANGE, bold=True)
    # 右: 收敛
    rx = Inches(6.7); ry = Inches(1.85)
    add_round(s, rx, ry, Inches(6.13), Inches(4.6), fill=LIGHT, line=GREEN)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.5), Inches(0.5),
             "收敛矩阵 2×2", size=14, color=GREEN, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(0.7), Inches(5.5), Inches(0.3),
             "高影响 + 易落地 = ★ 优先行动", size=10, color=GREEN, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(1.0), Inches(5.5), Inches(0.3),
             "高影响 + 难落地 = 战略项目", size=10, color=ORANGE)
    add_text(s, rx + Inches(0.3), ry + Inches(1.3), Inches(5.5), Inches(0.3),
             "低影响 + 易落地 = 随机应变", size=10, color=GRAY)
    add_text(s, rx + Inches(0.3), ry + Inches(1.6), Inches(5.5), Inches(0.3),
             "低影响 + 难落地 = 暂不考虑", size=10, color=GRAY)
    add_text(s, rx + Inches(0.3), ry + Inches(2.2), Inches(5.5), Inches(0.4),
             "3 步流程:", size=12, color=GREEN, bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(2.6), Inches(5.5), Inches(0.4),
             "1. 独立评估 (5-7 min)", size=10, color=BLACK)
    add_text(s, rx + Inches(0.3), ry + Inches(2.9), Inches(5.5), Inches(0.4),
             "2. 对比分歧 (5-7 min)", size=10, color=BLACK)
    add_text(s, rx + Inches(0.3), ry + Inches(3.2), Inches(5.5), Inches(0.4),
             "3. 确认行动 (5 min)", size=10, color=BLACK)
    add_text(s, rx + Inches(0.3), ry + Inches(3.7), Inches(5.5), Inches(0.4),
             "★ 优先: 2-3 条 + 负责人 + 时间", size=10, color=BLACK)
    add_text(s, rx + Inches(0.3), ry + Inches(4.2), Inches(5.5), Inches(0.3),
             "产出: 1 张行动清单", size=10, color=ORANGE, bold=True)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4),
             "→ 完整模板见学员手册 · F4 双矩阵卡片",
             size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# --- M5 扩展 13-18 ---
def m5_13(idx):
    s = _generic_slide(idx, "模块 5 · 前瞻思考", NAVY, "机会推演表 — 不是 「 希望什么 」, 是 「 准备怎么抓 」",
                       "对称结构: 捕捉行动 vs 准备工作, 缺一不可")
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.6), fill=GREEN)
    headers = ["可能出现的机会", "概率", "价值", "捕捉行动", "准备工作", "负责人+时间"]
    cols_w = [Inches(3.3), Inches(0.9), Inches(0.9), Inches(2.5), Inches(2.5), Inches(2.23)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.95), cols_w[i], Inches(0.4), h,
                 size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri in range(3):
        ry = Inches(2.5) + Inches(1.1) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.1), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        for ci in range(6):
            if ci >= 1 and ci <= 2:
                add_text(s, cx, ry + Inches(0.4), cols_w[ci], Inches(0.3),
                         "高/中/低", size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
            elif ci == 0:
                add_text(s, cx + Inches(0.1), ry + Inches(0.15), cols_w[ci] - Inches(0.2), Inches(0.4),
                         f"  机会 {ri+1}: 写具体场景", size=10, color=NAVY_DARK, bold=True)
                add_text(s, cx + Inches(0.1), ry + Inches(0.55), cols_w[ci] - Inches(0.2), Inches(0.4),
                         "  「 如果 ____ 发生, 对我们是好事 」", size=9, color=GRAY, italic=True)
            elif ci == 3:
                add_text(s, cx + Inches(0.1), ry + Inches(0.4), cols_w[ci] - Inches(0.2), Inches(0.4),
                         "机会出现时第一件事", size=9, color=GREEN, bold=True)
            elif ci == 4:
                add_text(s, cx + Inches(0.1), ry + Inches(0.4), cols_w[ci] - Inches(0.2), Inches(0.4),
                         "现在可以做, 让自己准备好", size=9, color=ORANGE, bold=True)
            elif ci == 5:
                add_text(s, cx + Inches(0.1), ry + Inches(0.4), cols_w[ci] - Inches(0.2), Inches(0.4),
                         "具体人 / 日期", size=9, color=GRAY, italic=True)
            cx += cols_w[ci]
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.4),
             "⚠ 关键区分: 捕捉行动 (机会出现时) vs 准备工作 (现在)", size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.4),
             "两者都要有 — 没有准备 = 机会出现时 \"还没准备好\"",
             size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.4),
             "没有捕捉 = 机会出现时 \"不知道该做什么\"",
             size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

def m5_14(idx):
    s = _generic_slide(idx, "模块 5 · 前瞻思考", NAVY, "完整示例: 李明的机会推演",
                       "行动: 全员推广 「 体验式接待 3 步法 」")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.4),
             "同一行动 · 3 个高价值机会",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    items = [
        ("本季度新品上市, 体验式销售优势明显 (客户对新品本来就有好奇心)", "高/高",
         "启动 「 新品体验周 」, 老客户定向邀约到店",
         "① 确认新品上市时间 ② 设计 「 新品体验式介绍脚本 」 ③ 老客户 「 有新品可试 」 通知",
         "李明 + 导购组 · 本周内"),
        ("若竞品服务质量下滑, 部分客户转来", "中/高",
         "设计 「 竞品客户专属接待流程 」, 第一时间展示专业感",
         "整理我们与竞品的差异化优势清单 · 准备简洁对比说明",
         "李明 · 本月内"),
        ("若成交率提升, 可邀请区域经理参访, 争取区域推广", "中/高",
         "邀请区域经理到访, 展示数据和现场接待, 申请纳入区域培训",
         "系统记录每日成交率数据 · 整理成功案例文档",
         "李明 · 数据下周一开始"),
    ]
    add_rect(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.4), fill=GREEN)
    add_text(s, Inches(0.7), Inches(2.45), Inches(12.0), Inches(0.3),
             "机会 / 概率·价值 / 捕捉行动 / 准备工作 / 负责人+时间",
             size=10, color=WHITE, bold=True)
    for i, (opp, prob, catch, prep, who) in enumerate(items):
        ry = Inches(2.85) + Inches(1.2) * i
        bg = LIGHT if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.15), fill=bg, line=GRAY_LIGHT)
        add_text(s, Inches(0.7), ry + Inches(0.1), Inches(4.0), Inches(0.95), opp,
                 size=10, color=NAVY_DARK, bold=True, line_h=1.3)
        add_text(s, Inches(4.7), ry + Inches(0.1), Inches(0.8), Inches(0.4),
                 prob, size=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(5.6), ry + Inches(0.1), Inches(2.8), Inches(0.95), catch,
                 size=9, color=BLACK, line_h=1.3)
        add_text(s, Inches(8.5), ry + Inches(0.1), Inches(3.0), Inches(0.95), prep,
                 size=9, color=BLACK, line_h=1.3)
        add_text(s, Inches(11.5), ry + Inches(0.1), Inches(1.3), Inches(0.95), who,
                 size=8, color=GRAY, line_h=1.3)

def m5_15(idx):
    s = _generic_slide(idx, "模块 5 · 前瞻思考", NAVY, "✋ 练习: 完成机会推演表",
                       "用你在 M4 找到的 ★ 优先行动")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "个人独立完成 · 12 分钟 · 写 2-4 个机会",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.4),
             "我推演的同一行动: ____________________________________",
             size=12, color=BLACK)
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.33), Inches(0.4),
             "识别 「 概率高 + 价值高 」 的机会, 写出捕捉行动和准备工作",
             size=12, color=NAVY_DARK, bold=True)
    add_rect(s, Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.4), fill=GREEN)
    headers = ["可能出现的机会", "概率", "价值", "捕捉行动", "准备工作", "负责人+时间"]
    cols_w = [Inches(3.3), Inches(0.9), Inches(0.9), Inches(2.5), Inches(2.5), Inches(2.23)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(3.35), cols_w[i], Inches(0.3), h,
                 size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    for ri in range(3):
        ry = Inches(3.7) + Inches(0.85) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.85), fill=bg, line=GRAY_LIGHT)
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.4),
             "我判断的最高价值机会是: __________________________",
             size=12, color=NAVY_DARK, bold=True)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.3),
             "为了抓住它, 我打算 ____ (日期) 之前完成的准备工作是: ____________________",
             size=11, color=GRAY)

def m5_16(idx):
    s = _generic_slide(idx, "模块 5 · 前瞻思考", NAVY, "两人互评: 推演有没有这 3 个问题",
                       "同桌互换 — 用 3 个问题检查对方的推演")
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.4), fill=NAVY)
    headers = ["检验问题", "针对对方的推演, 你的反馈"]
    cols_w = [Inches(5.0), Inches(7.33)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(1.92), cols_w[i], Inches(0.3), h,
                 size=11, color=WHITE, bold=True)
        cx += cols_w[i]
    items = [
        ("1. 风险描述够具体吗?\n(是 「 执行风险 」, 还是 「 如果 ____ 具体发生 」?)", ORANGE),
        ("2. 预防措施和应急预案有区分吗?\n(是否把 「 出了问题再处理 」 写进了预防措施里?)", GOLD),
        ("3. 机会的准备工作, 是现在就能做的吗?\n(不是 「 到时候再说 」, 是具体到本周或下周可以启动的动作)", GREEN),
    ]
    for ri, (q, col) in enumerate(items):
        ry = Inches(2.3) + Inches(1.3) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(1.25), fill=bg, line=GRAY_LIGHT)
        add_rect(s, Inches(0.5), ry, Inches(0.15), Inches(1.25), fill=col)
        add_text(s, Inches(0.8), ry + Inches(0.2), Inches(4.7), Inches(0.9), q,
                 size=11, color=NAVY_DARK, bold=True, line_h=1.4)
        add_text(s, Inches(5.7), ry + Inches(0.45), Inches(7.0), Inches(0.3),
                 "反馈: ____________________________________", size=10, color=GRAY, italic=True)
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.4),
             "⚠ 三个问题都满足 — 推演才算 「 完整 」",
             size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.3),
             "如果对方缺一个, 帮他补上 — 这是 「 互评 」 的价值",
             size=10, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

def m5_17(idx):
    s = _generic_slide(idx, "模块 5 · 前瞻思考", NAVY, "整合产出: 我的管理行动地图",
                       "5 工具 = 1 张完整地图")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.5),
             "「 我的管理行动地图 」 — 5 工具的整合",
             size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.4),
             "完成后拍照放在手机里 / 工作桌上, 每隔两周看一眼",
             size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.5), Inches(2.95), Inches(12.33), Inches(0.4), fill=NAVY)
    headers = ["工具", "我用它做了什么", "产出 / 下一步", "日期"]
    cols_w = [Inches(2.5), Inches(4.0), Inches(4.0), Inches(1.83)]
    cx = Inches(0.5)
    for i, h in enumerate(headers):
        add_text(s, cx, Inches(3.0), cols_w[i], Inches(0.3), h,
                 size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += cols_w[i]
    rows = [
        ("M1 螺旋 4 问", "访谈: ____  成功案例: ____", "提炼原则: ____", "访谈: ____"),
        ("M2 花刺投票", "议题: ____", "花票最高: ____  刺票最高: ____", "会议: ____"),
        ("M3 问题树+5 问", "难题: ____  切入点: ____", "新方向: ____", "行动: ____"),
        ("M4 双矩阵", "主持: ____", "优先 1: ____  优先 2: ____", "开会: ____"),
        ("M5 推演双表格", "推演: ____", "高优风险+预防: ____\n高优机会+准备: ____", "启动: ____"),
    ]
    for ri, (a, b, c, d) in enumerate(rows):
        ry = Inches(3.4) + Inches(0.7) * ri
        bg = LIGHT if ri % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), ry, Inches(12.33), Inches(0.7), fill=bg, line=GRAY_LIGHT)
        cx = Inches(0.5)
        for i, t in enumerate([a, b, c, d]):
            add_text(s, cx + Inches(0.1), ry + Inches(0.05) + Inches(0.3) * (1 if "\n" in t else 0),
                     cols_w[i] - Inches(0.2), Inches(0.4), t,
                     size=9, color=NAVY_DARK if i == 0 else BLACK,
                     bold=(i == 0), align=PP_ALIGN.CENTER, line_h=1.2)
            cx += cols_w[i]
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.2),
             "下周, 我要做的第一件事 (一件, 具体): _________________________________",
             size=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

def m5_18(idx):
    s = _generic_slide(idx, "模块 5 · 前瞻思考", NAVY, "M5 模块收尾 · 茶歇",
                       "回来后, 我们进入课程的收尾和综合产出")
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.5),
             "M5", size=120, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.7),
             "前瞻思考", size=32, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             "Think Ahead", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    quote_block(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                "30 分钟的提前推演, 换来的是省去的 15 小时救火时间。",
                author="—— 第五部分核心信念")
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
             "请准备: 你完成的 「 风险推演表 + 机会推演表 」, 下次课用",
             size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)


if __name__ == "__main__":
    out_path = r"D:\2026年课程\竞越\一线管理者的五项关键\补充课程包\04-授课PPT\一线管理者的现代五项_授课PPT.pptx"
    total = build()
    prs.save(out_path)
    print(f"PPT 已生成: {out_path}")
    print(f"总页数: {total}")
print("模块四定义完成")